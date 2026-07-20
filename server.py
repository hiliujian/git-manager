#!/usr/bin/env python3
"""
Git Manager Backend — 完整修复版 + 完善日志系统 + WebSocket实时通信
python3 server.py  →  http://localhost:7842
"""

import json
import time
import re
from collections import deque, OrderedDict
import subprocess, os, sys, logging, logging.handlers, threading, hashlib, difflib, shutil, mimetypes, base64, tempfile, shlex, atexit
import queue
from pathlib import Path
import urllib.request
import urllib.error
import urllib.parse
import uuid
import py_compile
import ssl
from http.server import ThreadingHTTPServer, BaseHTTPRequestHandler
from urllib.parse import urlparse, parse_qs, unquote
from datetime import datetime, timedelta, timezone
try:
    from websockets.sync.server import serve as ws_serve, ServerConnection
    WEBSOCKET_AVAILABLE = True
except ImportError:
    WEBSOCKET_AVAILABLE = False
    try:
        logger.warning("websockets 库未安装，WebSocket 功能将不可用 (pip install websockets)")
    except Exception:
        pass

try:
    from websockets.exceptions import ConnectionClosed  # type: ignore
except Exception:
    ConnectionClosed = None

try:
    import chardet
    CHARDET_AVAILABLE = True
except ImportError:
    CHARDET_AVAILABLE = False
    try:
        logger.warning("chardet 库未安装，中文编码检测可能不够准确 (pip install chardet)")
    except Exception:
        pass

try:
    from playwright.sync_api import sync_playwright  # type: ignore
except Exception:
    sync_playwright = None

# ════════════════════════════════════════════════════════
#  日志系统配置（带自动轮转和过期清理）
# ════════════════════════════════════════════════════════

# 日志轮转参数
_LOG_MAX_BYTES = 10 * 1024 * 1024   # 单个日志文件最大 10MB
_LOG_BACKUP_COUNT = 5               # 每种日志最多保留 5 个备份
_LOG_RETENTION_DAYS = 30            # 超过 30 天的日志文件自动清理

# 创建日志目录
LOG_DIR = Path(__file__).parent / "logs"
LOG_DIR.mkdir(exist_ok=True)


def _cleanup_old_logs():
    """清理超过 _LOG_RETENTION_DAYS 天的旧日志文件"""
    try:
        now = time.time()
        cutoff = now - _LOG_RETENTION_DAYS * 86400
        removed = 0
        for f in LOG_DIR.glob("*.log*"):
            try:
                if f.stat().st_mtime < cutoff:
                    f.unlink()
                    removed += 1
            except Exception:
                pass
        if removed > 0:
            print(f"[日志清理] 已清除 {removed} 个过期日志文件（超过 {_LOG_RETENTION_DAYS} 天）")
    except Exception:
        pass


# 启动时清理旧日志
_cleanup_old_logs()

# 配置日志格式
log_format = logging.Formatter(
    '%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    datefmt='%Y-%m-%d %H:%M:%S'
)

# 创建主日志记录器
logger = logging.getLogger('GitManager')
logger.setLevel(logging.DEBUG)

# 控制台处理器 - 显示 INFO 及以上级别
console_handler = logging.StreamHandler(sys.stdout)
console_handler.setLevel(logging.INFO)
console_handler.setFormatter(log_format)
logger.addHandler(console_handler)

# 文件处理器 - 按大小轮转，记录所有级别的日志
log_file = LOG_DIR / f"git_manager_{datetime.now().strftime('%Y%m%d')}.log"
file_handler = logging.handlers.RotatingFileHandler(
    log_file, encoding='utf-8', maxBytes=_LOG_MAX_BYTES, backupCount=_LOG_BACKUP_COUNT
)
file_handler.setLevel(logging.DEBUG)
file_handler.setFormatter(log_format)
logger.addHandler(file_handler)

# 错误日志文件处理器 - 按大小轮转，只记录 ERROR 及以上级别
error_log_file = LOG_DIR / f"errors_{datetime.now().strftime('%Y%m%d')}.log"
error_handler = logging.handlers.RotatingFileHandler(
    error_log_file, encoding='utf-8', maxBytes=_LOG_MAX_BYTES, backupCount=_LOG_BACKUP_COUNT
)
error_handler.setLevel(logging.ERROR)
error_handler.setFormatter(log_format)
logger.addHandler(error_handler)

logger.info("=" * 60)
logger.info("Git Manager Backend 启动初始化...")
logger.info(f"日志文件: {log_file}")
logger.info(f"错误日志文件: {error_log_file}")
logger.info(f"日志轮转: 单文件最大 {_LOG_MAX_BYTES // (1024*1024)}MB, 保留 {_LOG_BACKUP_COUNT} 个备份, {_LOG_RETENTION_DAYS} 天后自动清理")
logger.info("=" * 60)

# ════════════════════════════════════════════════════════
#  全局变量
# ════════════════════════════════════════════════════════

REPO_PATH = None
WS_PORT = 7843  # WebSocket端口号
ws_clients = set()  # WebSocket客户端集合
ws_clients_lock = threading.Lock()
ws_client_queues = {}  # client -> Queue[str]


def get_repo_status():
    if not REPO_PATH:
        return False, "未打开仓库", {}
    try:
        origin_url = ""
        try:
            out, _, code = run_git(["config", "--get", "remote.origin.url"])
            if code == 0:
                origin_url = (out or "").strip()
        except Exception:
            origin_url = ""

        has_staged = False
        staged_count = 0
        try:
            sf, err_sf = get_staged_files()
            if (not err_sf) and isinstance(sf, list):
                staged_count = len(sf)
                has_staged = staged_count > 0
        except Exception:
            has_staged = False
            staged_count = 0

        has_stash = False
        stash_count = 0
        stash_top = ""
        try:
            list_out, _, list_code = run_git(["stash", "list"], timeout=30)
            if list_code == 0:
                raw = (list_out or "").strip()
                if raw:
                    lines = [x for x in raw.splitlines() if x.strip()]
                    stash_count = len(lines)
                    has_stash = stash_count > 0
                    stash_top = lines[0].strip() if lines else ""
        except Exception:
            has_stash = False
            stash_count = 0
            stash_top = ""

        st0 = {
            "repo": REPO_PATH,
            "valid": bool(REPO_PATH and os.path.isdir(os.path.join(REPO_PATH, ".git"))),
            "origin_url": origin_url,
            "has_staged_changes": has_staged,
            "staged_count": staged_count,
            "has_stash": has_stash,
            "stash_count": stash_count,
            "stash_top": stash_top,
        }
        return True, "", st0
    except Exception as e:
        return False, str(e), {}


def set_origin_url(url: str):
    if not REPO_PATH:
        return False, "未打开仓库", ""
    raw = str(url or "").strip()
    if not raw:
        return False, "远程地址为空", ""
    try:
        out, err, code = run_git(["remote", "set-url", "origin", raw])
        if code != 0:
            msg = (err or out or "").strip()
            return False, msg or "设置远程地址失败", ""
    except Exception as e:
        return False, str(e), ""

    origin_url = ""
    try:
        out2, _, code2 = run_git(["config", "--get", "remote.origin.url"])
        if code2 == 0:
            origin_url = (out2 or "").strip()
    except Exception:
        origin_url = raw
    return True, "", origin_url


def get_current_branch():
    if not REPO_PATH:
        return ""
    try:
        out, _, code = run_git(["rev-parse", "--abbrev-ref", "HEAD"])
        if code == 0:
            return (out or "").strip()
    except Exception:
        return ""
    return ""


def open_repo(path: str):
    global REPO_PATH
    raw = str(path or "").strip()
    if not raw:
        return False, "路径为空"
    raw = os.path.expanduser(raw)
    if not os.path.isdir(raw):
        return False, f"目录不存在: {raw}"
    check = raw
    root = None
    for _ in range(15):
        if os.path.isdir(os.path.join(check, ".git")):
            root = check
            break
        parent = os.path.dirname(check)
        if parent == check:
            break
        check = parent
    if not root:
        return False, "不是 git 仓库（未找到 .git 目录）"
    REPO_PATH = root
    try:
        _undo_load_state()
    except Exception as e:
        logger.debug(f"Exception ignored: {e}")
    return True, ""


def workspace_context_text():
    return get_workspace_context()

ai_config_lock = threading.Lock()
ai_history_lock = threading.Lock()


_pw_lock = threading.Lock()
_pw_runtime = {
    "pw": None,
    "browser": None,
    "sessions": {},
}


def _pw_cleanup():
    """Clean up all Playwright resources (browser + playwright). Safe to call multiple times."""
    with _pw_lock:
        pw = _pw_runtime.pop("pw", None)
        browser = _pw_runtime.pop("browser", None)
    for obj in (browser, pw):
        try:
            if obj is not None:
                obj.close()
        except Exception:
            pass
    try:
        _pw_runtime["sessions"].clear()
    except Exception:
        pass


atexit.register(_pw_cleanup)


def _pw_is_available():
    return sync_playwright is not None


def _pw_get_session(session_id: str):
    sid = str(session_id or "").strip() or "global"
    with _pw_lock:
        sess = _pw_runtime.get("sessions", {}).get(sid)
        if isinstance(sess, dict) and sess.get("page") is not None:
            return sid, sess

        if not _pw_is_available():
            return sid, None

        if _pw_runtime.get("pw") is None:
            _pw_runtime["pw"] = sync_playwright().start()
        if _pw_runtime.get("browser") is None:
            _pw_runtime["browser"] = _pw_runtime["pw"].chromium.launch(headless=True)

        ctx = _pw_runtime["browser"].new_context()
        page = ctx.new_page()
        sess = {"ctx": ctx, "page": page}
        _pw_runtime.setdefault("sessions", {})[sid] = sess
        return sid, sess


def _pw_require_session(session_id: str):
    if not _pw_is_available():
        return False, "Playwright 未安装（可选能力）。如需启用请安装 playwright 并执行 playwright install。", None
    try:
        sid, sess = _pw_get_session(session_id)
        if sess is None:
            return False, "Playwright 未安装（可选能力）。如需启用请安装 playwright 并执行 playwright install。", None
        return True, "", sess
    except Exception as e:
        return False, str(e), None


def _pw_open(session_id: str, url: str, wait_ms: int = 0):
    ok, msg, sess = _pw_require_session(session_id)
    if not ok:
        return False, msg, {}
    u = str(url or "").strip()
    if not u:
        return False, "缺少 url", {}
    if not (u.startswith("http://") or u.startswith("https://")):
        return False, "仅支持 http/https", {}
    page = sess.get("page")
    page.goto(u)
    if int(wait_ms or 0) > 0:
        try:
            page.wait_for_timeout(int(wait_ms))
        except Exception:
            pass
    return True, "", {"session_id": str(session_id or "").strip() or "global", "obs": _pw_obs(sess)}


def _pw_click(session_id: str, selector: str, wait_ms: int = 0):
    ok, msg, sess = _pw_require_session(session_id)
    if not ok:
        return False, msg, {}
    sel = str(selector or "").strip()
    if not sel:
        return False, "缺少 selector", {}
    page = sess.get("page")
    page.click(sel)
    if int(wait_ms or 0) > 0:
        try:
            page.wait_for_timeout(int(wait_ms))
        except Exception:
            pass
    return True, "", {"session_id": str(session_id or "").strip() or "global", "obs": _pw_obs(sess)}


def _pw_type(session_id: str, selector: str, text: str, clear: bool = True, wait_ms: int = 0):
    ok, msg, sess = _pw_require_session(session_id)
    if not ok:
        return False, msg, {}
    sel = str(selector or "").strip()
    if not sel:
        return False, "缺少 selector", {}
    page = sess.get("page")
    if clear:
        try:
            page.fill(sel, "")
        except Exception:
            pass
    page.type(sel, str(text or ""))
    if int(wait_ms or 0) > 0:
        try:
            page.wait_for_timeout(int(wait_ms))
        except Exception:
            pass
    return True, "", {"session_id": str(session_id or "").strip() or "global", "obs": _pw_obs(sess)}


def _pw_eval(session_id: str, script: str):
    ok, msg, sess = _pw_require_session(session_id)
    if not ok:
        return False, msg, {}
    js = str(script or "")
    if not js.strip():
        return False, "缺少 script", {}
    page = sess.get("page")
    out = page.evaluate(js)
    return True, "", {"session_id": str(session_id or "").strip() or "global", "result": out, "obs": _pw_obs(sess)}


def _pw_text(session_id: str, selector: str = ""):
    ok, msg, sess = _pw_require_session(session_id)
    if not ok:
        return False, msg, {}
    page = sess.get("page")
    sel = str(selector or "").strip()
    if sel:
        out = page.text_content(sel)
        return True, "", {"session_id": str(session_id or "").strip() or "global", "text": str(out or ""), "obs": _pw_obs(sess)}
    out = page.content()
    out_s = str(out or "")
    if len(out_s) > 12000:
        out_s = out_s[:12000]
    return True, "", {"session_id": str(session_id or "").strip() or "global", "html": out_s, "obs": _pw_obs(sess)}


def _pw_screenshot(session_id: str, full_page: bool = True):
    ok, msg, sess = _pw_require_session(session_id)
    if not ok:
        return False, msg, {}
    page = sess.get("page")
    b = page.screenshot(full_page=bool(full_page))
    b64 = base64.b64encode(b).decode("ascii")
    return True, "", {"session_id": str(session_id or "").strip() or "global", "image_b64": b64, "obs": _pw_obs(sess)}


def _pw_wait(session_id: str, ms: int = 0):
    ok, msg, sess = _pw_require_session(session_id)
    if not ok:
        return False, msg, {}
    page = sess.get("page")
    try:
        page.wait_for_timeout(int(ms or 0))
    except Exception:
        pass
    return True, "", {"session_id": str(session_id or "").strip() or "global", "obs": _pw_obs(sess)}


def _pw_close_session(session_id: str):
    sid = str(session_id or "").strip() or "global"
    with _pw_lock:
        sess = _pw_runtime.get("sessions", {}).pop(sid, None)
    if isinstance(sess, dict):
        try:
            ctx = sess.get("ctx")
            if ctx is not None:
                ctx.close()
        except Exception:
            pass
    return True


def _pw_obs(sess: dict):
    try:
        page = sess.get("page")
        if page is None:
            return {}
        try:
            url = str(page.url or "")
        except Exception:
            url = ""
        try:
            title = str(page.title() or "")
        except Exception:
            title = ""
        return {"url": url, "title": title}
    except Exception:
        return {}

_hivo_mem_lock = threading.Lock()
_hivo_session_mem: dict = {}  # session_id -> {summary:str, chat:list[dict], tool_log:list[dict], tool_cache:dict}

_hivo_cfg_lock = threading.Lock()
_hivo_cfg_cache: dict = {}

# Hivo agent run control (cancel / concurrency)
_hivo_agent_run_lock = threading.Lock()
_hivo_agent_run_state: dict = {}  # run_id -> {session_id:str, started_at:float, cancel:bool, done:bool}
_hivo_agent_session_active: dict = {}  # session_id -> run_id
_hivo_agent_run_proc: dict = {}  # run_id -> subprocess.Popen
_hivo_agent_proc_lock = threading.Lock()
undo_lock = threading.Lock()
# Serialize revert/undo apply operations to prevent race conditions on rapid consecutive actions
_apply_revert_lock = threading.Lock()
_undo_groups = {}  # group_id -> list[entry]
_undo_group_order = []  # stack of group_id in commit order
last_file_state = None  # 上次的文件状态（用于检测变化）
_changed_files_cache = {
    "ts": 0.0,
    "files": None,
}
_changed_files_lock = threading.Lock()


def _undo_repo_key():
    try:
        rp = str(REPO_PATH or "").strip()
        if not rp:
            return "global"
        return hashlib.sha1(rp.encode("utf-8", errors="ignore")).hexdigest()[:12]
    except Exception:
        return "global"


def _undo_state_path():
    try:
        base = Path.home() / ".git-manager"
        base.mkdir(parents=True, exist_ok=True)
        return base / ("undo_state_" + _undo_repo_key() + ".json")
    except Exception:
        return Path.home() / ".git-manager" / "undo_state_global.json"


def _undo_save_state():
    try:
        p = _undo_state_path()
        with undo_lock:
            payload = {
                "version": 1,
                "repo": str(REPO_PATH or ""),
                "saved_at": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                "order": list(_undo_group_order),
                "groups": dict(_undo_groups),
            }
        _atomic_write_json(p, payload, indent=None)
    except Exception:
        return


def _undo_load_state():
    try:
        p = _undo_state_path()
        if not p.exists():
            return

        
        raw = p.read_text(encoding="utf-8")
        data = json.loads(raw or "{}")
        order = data.get("order")
        groups = data.get("groups")
        if not isinstance(order, list) or not isinstance(groups, dict):
            return
        with undo_lock:
            norm_groups = {}
            for k, v in groups.items():
                kk = str(k).strip()
                if not kk:
                    continue
                if not isinstance(v, list):
                    continue
                norm_groups[kk] = v

            _undo_groups.clear()
            _undo_groups.update(norm_groups)
            _undo_group_order[:] = [str(x) for x in order if str(x).strip() and (str(x).strip() in norm_groups)]
    except Exception:
        return

# Idempotency cache for write APIs (prevents repeated tool calls)
idempotency_lock = threading.Lock()
_idempotency_cache = {}  # key -> {"ts": float, "code": int, "payload": dict}
_IDEMPOTENCY_TTL_SEC = 300
_IDEMPOTENCY_MAX = 500


def _idempotency_get(key: str):
    k = (key or "").strip()
    if not k:
        return None
    now = time.time()
    with idempotency_lock:
        ent = _idempotency_cache.get(k)
        if not ent:
            return None
        ts = float(ent.get("ts") or 0.0)
        if (now - ts) > _IDEMPOTENCY_TTL_SEC:
            try:
                del _idempotency_cache[k]
            except Exception as e:
                logger.debug(f"Exception ignored: {e}")
            return None
        return ent


def _idempotency_set(key: str, payload: dict, code: int = 200):
    k = (key or "").strip()
    if not k:
        return
    now = time.time()
    with idempotency_lock:
        _idempotency_cache[k] = {"ts": now, "code": int(code), "payload": payload}
        if len(_idempotency_cache) > _IDEMPOTENCY_MAX:
            # purge expired first
            expired = []
            for kk, vv in list(_idempotency_cache.items()):
                ts = float((vv or {}).get("ts") or 0.0)
                if (now - ts) > _IDEMPOTENCY_TTL_SEC:
                    expired.append(kk)
            for kk in expired:
                try:
                    del _idempotency_cache[kk]
                except Exception as e:
                    logger.debug(f"Exception ignored: {e}")
            # still too large: drop oldest
            if len(_idempotency_cache) > _IDEMPOTENCY_MAX:
                items = sorted(_idempotency_cache.items(), key=lambda x: float((x[1] or {}).get("ts") or 0.0))
                for kk, _vv in items[: max(0, len(_idempotency_cache) - _IDEMPOTENCY_MAX)]:
                    try:
                        del _idempotency_cache[kk]
                    except Exception as e:
                        logger.debug(f"Exception ignored: {e}")


def _undo_capture_file_snapshot(rel_path: str):
    rp = (rel_path or "").replace("\\", "/").lstrip("/")
    if not rp:
        return None
    full = _safe_repo_abspath(rp)
    if not full:
        return None
    if os.path.isdir(full):
        return {"path": rp, "exists": False, "content": ""}
    if not os.path.exists(full):
        return {"path": rp, "exists": False, "content": ""}
    content, _enc = get_file_content(rp, return_encoding=True)
    if content is None:
        content = ""
    return {"path": rp, "exists": True, "content": str(content)}


def _undo_apply_file_snapshot(snapshot: dict):
    if not isinstance(snapshot, dict):
        return False, "snapshot 非法"
    rp = (snapshot.get("path") or "").replace("\\", "/").lstrip("/")
    if not rp:
        return False, "snapshot 缺少 path"
    if snapshot.get("exists"):
        return save_file_content(rp, str(snapshot.get("content") or ""))
    full = _safe_repo_abspath(rp)
    if not full:
        return False, "非法路径"
    try:
        if os.path.exists(full) and (not os.path.isdir(full)):
            os.remove(full)
        return True, ""
    except Exception as e:
        return False, str(e)


def _undo_git_status_map():
    """Return {path: info} based on git status porcelain (-z).

    info schema:
    - {"xy": "XY"} for normal entries
    - For renames/copies (R/C), both src and dst paths are included:
      src: {"xy": "R ", "kind": "R", "role": "src", "dst": "new"}
      dst: {"xy": "R ", "kind": "R", "role": "dst", "src": "old"}
    """
    try:
        out, _err, code = run_git(["status", "--porcelain=v1", "-u", "-z"], timeout=30)
        if code != 0:
            return {}
        m = {}
        entries = (out or "").split("\x00")
        i = 0
        while i < len(entries):
            ent = entries[i]
            if not ent or len(ent) < 4:
                i += 1
                continue
            xy = ent[:2]
            name0 = ent[3:]
            kind = xy[0] if xy else ""
            if name0:
                p0 = str(name0).replace("\\", "/")
                if kind in ("R", "C"):
                    name1 = entries[i + 1] if (i + 1) < len(entries) else ""
                    p1 = str(name1).replace("\\", "/") if name1 else ""
                    m[p0] = {"xy": xy, "kind": kind, "role": "src", "dst": p1}
                    if p1:
                        m[p1] = {"xy": xy, "kind": kind, "role": "dst", "src": p0}
                else:
                    m[p0] = {"xy": xy}
            if kind in ("R", "C"):
                i += 1
            i += 1
        return m
    except Exception:
        return {}


def _undo_capture_head_file_snapshot(rel_path: str):
    rp = (rel_path or "").replace("\\", "/").lstrip("/")
    if not rp:
        return None
    try:
        content = get_head_file_content(rp)
    except Exception:
        content = None
    if content is None:
        return None
    return {"path": rp, "exists": True, "content": str(content)}


def _undo_prepare_cmd_snapshots(undo_gid: str):
    gid = str(undo_gid or "").strip()
    if not gid:
        return {}, {}
    pre_map = _undo_git_status_map()
    pre_snaps = {}
    try:
        for rp in pre_map.keys():
            try:
                snap0 = _undo_capture_file_snapshot(rp)
            except Exception:
                snap0 = None
            if snap0 is not None:
                pre_snaps[str(rp)] = snap0
    except Exception:
        pre_snaps = {}
    return pre_map, pre_snaps


def _undo_finalize_cmd_snapshots(undo_gid: str, pre_map: dict, pre_snaps: dict):
    gid = str(undo_gid or "").strip()
    if not gid:
        return
    post_map = _undo_git_status_map()

    # If cmd performed rename/copy on a previously clean tracked file, pre_map/pre_snaps may not include it.
    # For such sources, capture snapshot from HEAD so undo can restore it.
    try:
        for rp, info in (post_map or {}).items():
            if not isinstance(info, dict):
                continue
            kind = str(info.get("kind") or "").strip()
            role = str(info.get("role") or "").strip()
            if kind != "R" or role != "src":
                continue
            if (pre_snaps or {}).get(rp) is not None:
                continue
            snap_h = _undo_capture_head_file_snapshot(rp)
            if snap_h is not None:
                _undo_record(gid, {"type": "file_snapshot", "op": "delete", "snapshot": snap_h})
    except Exception as e:
        logger.debug(f"Exception ignored: {e}")

    keys = set((pre_map or {}).keys()) | set((post_map or {}).keys())
    for rp in keys:
        pre_xy = (pre_map or {}).get(rp)
        post_xy = (post_map or {}).get(rp)
        pre_xy_s = str(pre_xy.get("xy") if isinstance(pre_xy, dict) else (pre_xy or ""))
        post_xy_s = str(post_xy.get("xy") if isinstance(post_xy, dict) else (post_xy or ""))
        if pre_xy is None and post_xy is not None:
            _undo_record(gid, {"type": "file_snapshot", "op": "add", "snapshot": {"path": rp, "exists": False, "content": ""}})
            continue
        if pre_xy is not None and post_xy is None:
            snap = (pre_snaps or {}).get(rp)
            if snap is not None:
                _undo_record(gid, {"type": "file_snapshot", "op": "delete", "snapshot": snap})
            continue
        if pre_xy is not None and post_xy is not None and pre_xy_s != post_xy_s:
            snap = (pre_snaps or {}).get(rp)
            if snap is not None:
                _undo_record(gid, {"type": "file_snapshot", "op": "modify", "snapshot": snap})


def _undo_record(group_id: str, entry: dict):
    gid = str(group_id or "").strip()
    if not gid:
        return
    if not isinstance(entry, dict):
        return
    try:
        if "ts" not in entry:
            entry["ts"] = time.time()
    except Exception as e:
        logger.debug(f"Exception ignored: {e}")
    with undo_lock:
        _undo_groups.setdefault(gid, []).append(entry)
        if gid not in _undo_group_order:
            _undo_group_order.append(gid)
            try:
                tp = str(entry.get("type") or "")
                logger.warning(f"undo_steps+1 (group_id={gid}, first_type={tp}, undo_steps={len(_undo_group_order)})")
            except Exception as e:
                logger.debug(f"Exception ignored: {e}")
    try:
        _undo_save_state()
    except Exception as e:
        logger.debug(f"Exception ignored: {e}")


def _undo_drop_groups(group_ids: list[str]):
    gids_in = [str(x).strip() for x in (group_ids or []) if x is not None and str(x).strip()]
    if not gids_in:
        return 0
    removed = 0
    with undo_lock:
        for gid in gids_in:
            if gid in _undo_groups:
                try:
                    del _undo_groups[gid]
                except Exception as e:
                    logger.debug(f"Exception ignored: {e}")
            try:
                if gid in _undo_group_order:
                    _undo_group_order[:] = [g for g in _undo_group_order if str(g) != gid]
            except Exception as e:
                logger.debug(f"Exception ignored: {e}")
            removed += 1
    try:
        _undo_save_state()
    except Exception:
        pass
    return removed


def _hivo_agent_mark_started(run_id: str, session_id: str):
    try:
        rid = str(run_id or "").strip()
        sid = str(session_id or "").strip()
        if not rid:
            return
        with _hivo_agent_run_lock:
            _hivo_agent_run_state[rid] = {"session_id": sid, "started_at": time.time(), "cancel": False, "done": False}
            if sid:
                _hivo_agent_session_active[sid] = rid
    except Exception as e:
        logger.debug(f"Exception ignored: {e}")


def _hivo_agent_mark_done(run_id: str):
    try:
        rid = str(run_id or "").strip()
        if not rid:
            return
        with _hivo_agent_run_lock:
            st = _hivo_agent_run_state.get(rid)
            if isinstance(st, dict):
                st["done"] = True
            try:
                sid = str(st.get("session_id") or "") if isinstance(st, dict) else ""
            except Exception:
                sid = ""
            if sid and _hivo_agent_session_active.get(sid) == rid:
                _hivo_agent_session_active.pop(sid, None)
            # 清理已完成的旧条目，防止内存泄漏（保留最近 50 条）
            if len(_hivo_agent_run_state) > 50:
                done_keys = [k for k, v in _hivo_agent_run_state.items() if isinstance(v, dict) and v.get("done")]
                # 保留最近 20 条已完成的
                done_keys.sort(key=lambda k: _hivo_agent_run_state[k].get("started_at", 0), reverse=True)
                for k in done_keys[20:]:
                    _hivo_agent_run_state.pop(k, None)
    except Exception as e:
        logger.debug(f"Exception ignored: {e}")


def _hivo_agent_is_cancelled(run_id: str) -> bool:
    try:
        rid = str(run_id or "").strip()
        if not rid:
            return False
        with _hivo_agent_run_lock:
            st = _hivo_agent_run_state.get(rid)
            if not isinstance(st, dict):
                return False
            return bool(st.get("cancel"))
    except Exception:
        return False


def _hivo_agent_request_cancel(run_id: str) -> bool:
    try:
        rid = str(run_id or "").strip()
        if not rid:
            return False
        with _hivo_agent_run_lock:
            st = _hivo_agent_run_state.get(rid)
            if not isinstance(st, dict):
                return False
            st["cancel"] = True
        try:
            proc = _hivo_agent_run_proc.get(rid)
            if proc is not None:
                try:
                    proc.terminate()
                except Exception as e:
                    logger.debug(f"Exception ignored: {e}")
        except Exception as e:
            logger.debug(f"Exception ignored: {e}")
        return True
    except Exception:
        return False


def _hivo_agent_clear_proc(run_id: str):
    try:
        rid = str(run_id or "").strip()
        if not rid:
            return
        _hivo_agent_run_proc.pop(rid, None)
    except Exception as e:
        logger.debug(f"Exception ignored: {e}")


def _undo_pop_latest_group():
    with undo_lock:
        while _undo_group_order:
            gid = _undo_group_order.pop()
            actions = _undo_groups.pop(gid, None)
            if actions:
                break
        else:
            gid, actions = "", []
    try:
        _undo_save_state()
    except Exception as e:
        logger.debug(f"Exception ignored: {e}")
    return gid, actions


def _undo_pop_latest_group_for_session(session_id: str):
    sid = str(session_id or "").strip()
    if not sid:
        return "", []
    prefix = "ai-" + sid + "-"
    with undo_lock:
        gid = ""
        actions = []
        for i in range(len(_undo_group_order) - 1, -1, -1):
            g = str(_undo_group_order[i] or "").strip()
            if not g.startswith(prefix):
                continue
            gid = g
            try:
                _undo_group_order.pop(i)
            except Exception:
                try:
                    _undo_group_order[:] = [x for x in _undo_group_order if str(x) != gid]
                except Exception as e:
                    logger.debug(f"Exception ignored: {e}")
            actions = _undo_groups.pop(gid, None)
            break
    try:
        _undo_save_state()
    except Exception as e:
        logger.debug(f"Exception ignored: {e}")
    if not isinstance(actions, list):
        actions = []
    return gid, actions


def _undo_pop_group_by_id(group_id: str):
    gid = str(group_id or "").strip()
    if not gid:
        return "", []
    with undo_lock:
        actions = _undo_groups.pop(gid, None)
        try:
            _undo_group_order[:] = [x for x in _undo_group_order if str(x) != gid]
        except Exception as e:
            logger.debug(f"Exception ignored: {e}")
    try:
        _undo_save_state()
    except Exception as e:
        logger.debug(f"Exception ignored: {e}")
    if not isinstance(actions, list):
        actions = []
    return gid, actions


def _undo_get_steps():
    with undo_lock:
        return len(_undo_group_order)


def _undo_get_steps_for_session(session_id: str):
    sid = str(session_id or "").strip()
    if not sid:
        return 0
    prefix = "ai-" + sid + "-"
    with undo_lock:
        n = 0
        for gid in _undo_group_order:
            if str(gid).startswith(prefix):
                n += 1
        return n


def _undo_clear_for_session(session_id: str):
    sid = str(session_id or "").strip()
    if not sid:
        return 0
    prefix = "ai-" + sid + "-"
    removed = 0
    with undo_lock:
        gids = [str(g) for g in list(_undo_group_order) if str(g).startswith(prefix)]
        if not gids:
            return 0
        try:
            _undo_group_order[:] = [g for g in _undo_group_order if not str(g).startswith(prefix)]
        except Exception as e:
            logger.debug(f"Exception ignored: {e}")
        for gid in gids:
            if gid in _undo_groups:
                try:
                    del _undo_groups[gid]
                except Exception as e:
                    logger.debug(f"Exception ignored: {e}")
            removed += 1
    try:
        _undo_save_state()
    except Exception as e:
        logger.debug(f"Exception ignored: {e}")
    return removed


def _undo_apply_actions(actions: list):
    if not isinstance(actions, list) or not actions:
        return True, "无可撤销操作"
    errs = []
    for a in reversed(actions):
        try:
            tp = str(a.get("type") or "")
            if tp == "file_snapshot":
                ok, msg = _undo_apply_file_snapshot(a.get("snapshot"))
                if not ok:
                    errs.append(msg or "恢复文件失败")
            elif tp == "rename":
                oldp = (a.get("old_path") or "").replace("\\", "/").lstrip("/")
                newp = (a.get("new_path") or "").replace("\\", "/").lstrip("/")
                if oldp and newp:
                    rename_file(newp, oldp)
                for snap in (a.get("snapshots") or []):
                    _undo_apply_file_snapshot(snap)
            elif tp == "ai_config":
                prev = a.get("prev")
                ok, msg = save_hivo_ai_config(prev)
                if not ok:
                    errs.append(msg or "恢复 AI 配置失败")
            elif tp == "hivo_ai_config":
                prev = a.get("prev")
                ok, msg = _hivo_save_cfg(prev if isinstance(prev, dict) else {})
                if not ok:
                    errs.append(msg or "恢复 Hivo AI 配置失败")
        except Exception as e:
            errs.append(str(e))
    if errs:
        return False, "\n".join([x for x in errs if x])
    try:
        invalidate_changed_files_cache()
        notify_files_updated()
    except Exception as e:
        logger.debug(f"Exception ignored: {e}")
    return True, ""


def invalidate_changed_files_cache():
    try:
        with _changed_files_lock:
            _changed_files_cache["ts"] = 0.0
            _changed_files_cache["files"] = None
    except Exception as e:
        logger.debug(f"Exception ignored: {e}")


def notify_files_updated():
    try:
        if not REPO_PATH:
            return
        with ws_clients_lock:
            if not ws_clients:
                return
        files = get_changed_files_cached(max_age_sec=0)
        broadcast_to_clients({
            'type': 'files_updated',
            'files': files
        })
    except Exception as e:
        logger.debug(f"Exception ignored: {e}")

# 最近一次读取文件的编码/是否发生 lossy(replace) 解码。
# 说明：为保持前端显示体验，读取时可能用 errors="replace"。
# 但若发生替换，文本已无法无损还原原始字节，此时必须阻止保存以免写坏文件。
_file_last_encoding = {}
_file_decode_lossy = {}
_file_encoding_lock = threading.Lock()

# 文件内容内存缓存：(path, mtime) → (content, encoding)。LRU 淘汰，单次运行内有效。
_file_content_cache = OrderedDict()  # type: OrderedDict[tuple, tuple]
_file_content_cache_lock = threading.Lock()
_FILE_CONTENT_CACHE_MAX = 120

# ════════════════════════════════════════════════════════
#  WebSocket 实时通信
# ════════════════════════════════════════════════════════

def broadcast_to_clients(message):
    """向所有WebSocket客户端广播消息"""
    if not WEBSOCKET_AVAILABLE:
        return

    try:
        payload = json.dumps(message)
    except Exception:
        return

    with ws_clients_lock:
        clients = list(ws_clients)
        for client in clients:
            q = ws_client_queues.get(client)
            if q is None:
                q = queue.Queue(maxsize=200)
                ws_client_queues[client] = q
            try:
                q.put_nowait(payload)
            except Exception:
                # queue full or other error -> drop client
                try:
                    ws_clients.discard(client)
                except Exception:
                    pass
                try:
                    ws_client_queues.pop(client, None)
                except Exception:
                    pass


def _hivo_cfg_path():
    try:
        return _hivo_ai_data_dir() / "hivo_ai_config.json"
    except Exception:
        return Path("hivo_ai_config.json")


def _hivo_ai_data_dir() -> Path:
    try:
        return Path(__file__).resolve().parent / "hivo_ai_data"
    except Exception:
        return Path("hivo_ai_data")


def _atomic_write_json(path, data, indent=2):
    """原子写入 JSON 文件：先写临时文件再替换，避免写入中途崩溃导致文件损坏。

    Args:
        path: 目标文件路径 (Path or str)
        data: 可序列化的 Python 对象
        indent: JSON 缩进，None 表示不缩进
    """
    p = Path(path)
    p.parent.mkdir(parents=True, exist_ok=True)
    tmp = p.with_suffix(p.suffix + ".tmp")
    text = json.dumps(data, ensure_ascii=False) if indent is None else json.dumps(data, ensure_ascii=False, indent=indent)
    tmp.write_text(text, encoding="utf-8")
    os.replace(tmp, p)


def _hivo_load_cfg():
    with _hivo_cfg_lock:
        try:
            p = _hivo_cfg_path()
            if p.exists():
                data = json.loads(p.read_text(encoding="utf-8") or "{}")
                if isinstance(data, dict):
                    _hivo_cfg_cache.clear()
                    _hivo_cfg_cache.update(data)
        except Exception as e:
            logger.debug(f"Exception ignored: {e}")
        if not _hivo_cfg_cache:
            _hivo_cfg_cache.update({
                "version": 1,
                "max_rounds": 12,
                "max_tool_calls_per_round": 3,
                "max_visible_history": 80,
                "repeat_block": {"window": 4, "max_same": 3, "signature": "tool_types"},
                "status_events": {"enabled": True},
                "status_messages": {
                    "sending": "发送中...",
                    "thinking": "思考中...",
                    "executing": "执行中...",
                    "done": "完成",
                },
                "features": {
                    "memory_enabled": True,
                    "tool_memory_enabled": True,
                    "tool_cache_enabled": True,
                    "context_refs_enabled": True,
                    "web_search_enabled": False,
                    "topic_isolation_enabled": True,
                    "browser_auto_screenshot": False,
                },
                "active_profile_id": "",
                "profiles": [],
                "web_search": {
                    "active_profile_id": "default",
                    "profiles": [
                        {
                            "id": "default",
                            "name": "",
                            "provider": "tavily",
                            "api_key": "",
                            "base_url": "https://api.tavily.com/search",
                            "timeout": 20,
                            "top_k": 5,
                        }
                    ],
                },
                "memory": {
                    "short_term_turns": 6,
                    "long_term_summary_chars": 3500,
                    "tool_log_items": 80,
                    "tool_cache_items": 120,
                },
                "agent": {
                    "mode": "fallback_chat",
                    "system_prompt": "",
                    "tools": [],
                },
            })
        else:
            # Fill required keys for unified config without overwriting existing values.
            try:
                if "features" not in _hivo_cfg_cache or not isinstance(_hivo_cfg_cache.get("features"), dict):
                    _hivo_cfg_cache["features"] = {}
                if "web_search_enabled" not in _hivo_cfg_cache["features"]:
                    _hivo_cfg_cache["features"]["web_search_enabled"] = False
                if "browser_auto_screenshot" not in _hivo_cfg_cache["features"]:
                    _hivo_cfg_cache["features"]["browser_auto_screenshot"] = False
            except Exception as e:
                logger.debug(f"Exception ignored: {e}")
            try:
                if "profiles" not in _hivo_cfg_cache or not isinstance(_hivo_cfg_cache.get("profiles"), list):
                    _hivo_cfg_cache["profiles"] = []
                if "active_profile_id" not in _hivo_cfg_cache:
                    _hivo_cfg_cache["active_profile_id"] = ""
            except Exception as e:
                logger.debug(f"Exception ignored: {e}")
            try:
                ws = _hivo_cfg_cache.get("web_search")
                if not isinstance(ws, dict):
                    ws = {}
                # If old scalar format exists, lift it to list-based format.
                if "profiles" not in ws or not isinstance(ws.get("profiles"), list):
                    prov = str(ws.get("provider") or "").strip()
                    ws_profiles = [{
                        "id": "default",
                        "name": "",
                        "provider": prov,
                        "api_key": str(ws.get("api_key") or ""),
                        "base_url": str(ws.get("base_url") or "").strip(),
                        "timeout": int(ws.get("timeout") or 20),
                        "top_k": int(ws.get("top_k") or 5),
                    }]
                    ws = {"active_profile_id": "default", "profiles": ws_profiles}
                if "active_profile_id" not in ws:
                    ws["active_profile_id"] = str((ws.get("profiles") or [{}])[0].get("id") or "").strip() if isinstance(ws.get("profiles"), list) and ws.get("profiles") else ""
                _hivo_cfg_cache["web_search"] = ws
            except Exception as e:
                logger.debug(f"Exception ignored: {e}")
        return dict(_hivo_cfg_cache)


def _hivo_save_cfg(cfg: dict):
    if not isinstance(cfg, dict):
        return False, "config 必须是对象"
    with _hivo_cfg_lock:
        try:
            p = _hivo_cfg_path()
            _atomic_write_json(p, cfg, indent=2)
            _hivo_cfg_cache.clear()
            _hivo_cfg_cache.update(cfg)
            return True, "保存成功"
        except Exception as e:
            return False, str(e)


def _hivo_web_search(query: str, top_k: int = 5, timeout: int = 20):
    enabled, sc0 = _ai_load_web_search_cfg()
    if not enabled:
        return False, "web_search 未启用", {"result": {"query": str(query or ""), "items": []}}

    sc0 = sc0 if isinstance(sc0, dict) else {}

    sc = sc0
    profiles = sc0.get("profiles")
    if isinstance(profiles, list) and profiles:
        active_id = str(sc0.get("active_profile_id") or "").strip()
        pick = None
        for p in profiles:
            if not isinstance(p, dict):
                continue
            if active_id and str(p.get("id") or "").strip() == active_id:
                pick = p
                break
        if pick is None:
            for p in profiles:
                if isinstance(p, dict):
                    pick = p
                    break
        if isinstance(pick, dict):
            sc = pick

    provider = str(sc.get("provider") or "").strip().lower()
    api_key = str(sc.get("api_key") or "").strip()
    base_url = str(sc.get("base_url") or "").strip().rstrip("/")
    try:
        k = int(top_k or sc.get("top_k") or 5)
    except Exception:
        k = 5
    k = max(1, min(10, k))
    try:
        to = int(timeout or sc.get("timeout") or 20)
    except Exception:
        to = 20
    to = max(3, min(60, to))

    q = str(query or "").strip()
    if not q:
        return False, "缺少 query", {"result": {"query": "", "items": []}}

    if provider in ("serper", "google_serper", "serper_dev"):
        if not api_key:
            return False, "Serper 缺少 api_key", {"result": {"query": q, "provider": provider, "items": []}}
        url = "https://google.serper.dev/search"
        body = {"q": q, "num": k}
        headers = {"Content-Type": "application/json", "X-API-KEY": api_key, "Connection": "close"}
        req = urllib.request.Request(url, data=json.dumps(body).encode("utf-8"), headers=headers, method="POST")
        try:
            with urllib.request.urlopen(req, timeout=to) as resp:
                raw = resp.read().decode("utf-8", errors="replace")
                data = json.loads(raw or "{}")
        except Exception as e:
            return False, str(e), {"result": {"query": q, "provider": provider, "items": []}}
        items = []
        for it in (data.get("organic") or [])[:k]:
            if not isinstance(it, dict):
                continue
            items.append({
                "title": str(it.get("title") or "").strip(),
                "url": str(it.get("link") or "").strip(),
                "snippet": str(it.get("snippet") or "").strip(),
                "source": "serper",
            })
        return True, "", {"result": {"query": q, "provider": provider, "items": items}}

    if provider in ("tavily",):
        if not api_key:
            return False, "Tavily 缺少 api_key", {"result": {"query": q, "provider": provider, "items": []}}
        url = "https://api.tavily.com/search"
        body = {"api_key": api_key, "query": q, "max_results": k, "search_depth": "basic"}
        headers = {"Content-Type": "application/json", "Connection": "close"}
        req = urllib.request.Request(url, data=json.dumps(body).encode("utf-8"), headers=headers, method="POST")
        try:
            with urllib.request.urlopen(req, timeout=to) as resp:
                raw = resp.read().decode("utf-8", errors="replace")
                data = json.loads(raw or "{}")
        except Exception as e:
            return False, str(e), {"result": {"query": q, "provider": provider, "items": []}}
        items = []
        for it in (data.get("results") or [])[:k]:
            if not isinstance(it, dict):
                continue
            items.append({
                "title": str(it.get("title") or "").strip(),
                "url": str(it.get("url") or "").strip(),
                "snippet": str(it.get("content") or "").strip(),
                "source": "tavily",
            })
        return True, "", {"result": {"query": q, "provider": provider, "items": items}}

    if provider in ("searxng", "searx"):
        if not base_url:
            return False, "SearxNG 缺少 base_url", {"result": {"query": q, "provider": provider, "items": []}}
        url = base_url + "/search?" + urllib.parse.urlencode({"q": q, "format": "json"})
        req = urllib.request.Request(url, headers={"Connection": "close"}, method="GET")
        try:
            with urllib.request.urlopen(req, timeout=to) as resp:
                raw = resp.read().decode("utf-8", errors="replace")
                data = json.loads(raw or "{}")
        except Exception as e:
            return False, str(e), {"result": {"query": q, "provider": provider, "items": []}}
        items = []
        for it in (data.get("results") or [])[:k]:
            if not isinstance(it, dict):
                continue
            items.append({
                "title": str(it.get("title") or "").strip(),
                "url": str(it.get("url") or "").strip(),
                "snippet": str(it.get("content") or it.get("snippet") or "").strip(),
                "source": "searxng",
            })
        return True, "", {"result": {"query": q, "provider": provider, "items": items}}

    return False, "未配置可用 provider（支持：searxng/serper/tavily）", {"result": {"query": q, "provider": provider, "items": []}}


def _hivo_web_fetch(url: str, timeout: int = 20):
    enabled, _sc0 = _ai_load_web_search_cfg()
    if not enabled:
        return False, "web_fetch 未启用", {"result": {"url": str(url or ""), "text": ""}}

    u = str(url or "").strip()
    if not u:
        return False, "缺少 url", {"result": {"url": "", "text": ""}}
    if not (u.startswith("http://") or u.startswith("https://")):
        return False, "仅支持 http/https", {"result": {"url": u, "text": ""}}
    try:
        to = int(timeout or 20)
    except Exception:
        to = 20
    to = max(3, min(60, to))
    req = urllib.request.Request(u, headers={"User-Agent": "Mozilla/5.0", "Connection": "close"}, method="GET")
    try:
        with urllib.request.urlopen(req, timeout=to) as resp:
            raw = resp.read(1024 * 512)
            txt = raw.decode("utf-8", errors="replace")
    except Exception as e:
        return False, str(e), {"result": {"url": u, "text": ""}}
    try:
        txt = re.sub(r"<script[\s\S]*?</script>", " ", txt, flags=re.IGNORECASE)
        txt = re.sub(r"<style[\s\S]*?</style>", " ", txt, flags=re.IGNORECASE)
        txt = re.sub(r"<[^>]+>", " ", txt)
        txt = re.sub(r"[ \t\x0b\x0c\r]+", " ", txt)
        txt = re.sub(r"\n{3,}", "\n\n", txt)
        txt = txt.strip()
    except Exception:
        txt = (txt or "").strip()
    if len(txt) > 12000:
        txt = txt[:12000]
    return True, "", {"result": {"url": u, "text": txt}}


def _hivo_status_message(cfg: dict, stage: str, **kwargs):
    try:
        st = str(stage or "").strip()
        sm = cfg.get("status_messages") if isinstance(cfg, dict) else None
        sm = sm if isinstance(sm, dict) else {}
        base = str(sm.get(st) or "").strip()
        if not base:
            base = st

        if st == "executing":
            n = kwargs.get("tool_count")
            if isinstance(n, int) and n >= 0:
                return f"{base} ({n} 个工具)"
            return base

        return base
    except Exception:
        return str(stage or "")


def _hivo_get_session_state(session_id: str):
    sid = str(session_id or "").strip()
    if not sid:
        sid = "global"
    with _hivo_mem_lock:
        st = _hivo_session_mem.get(sid)
        if not isinstance(st, dict):
            st = {"summary": "", "chat": [], "tool_log": [], "tool_cache": {}, "topic_id": 1, "topic_archive": []}
            _hivo_session_mem[sid] = st
        if not isinstance(st.get("summary"), str):
            st["summary"] = ""
        if not isinstance(st.get("chat"), list):
            st["chat"] = []
        if not isinstance(st.get("tool_log"), list):
            st["tool_log"] = []
        if not isinstance(st.get("tool_cache"), dict):
            st["tool_cache"] = {}
        try:
            if not isinstance(st.get("topic_id"), int) or int(st.get("topic_id") or 0) <= 0:
                st["topic_id"] = 1
        except Exception:
            st["topic_id"] = 1
        if not isinstance(st.get("topic_archive"), list):
            st["topic_archive"] = []
        return st


def _hivo_detect_topic_switch(user_text: str):
    try:
        t = str(user_text or "").strip()
    except Exception:
        t = ""
    if not t:
        return False
    tl = t.lower()
    if tl.startswith("/new_topic") or tl.startswith("/topic") or tl.startswith("/reset_topic"):
        return True
    # explicit intent only (avoid accidental resets)
    kws = [
        "换个话题",
        "切换话题",
        "新话题",
        "另一个话题",
        "另一个问题",
        "换个问题",
        "忽略以上",
        "忽略之前",
        "不相关了",
        "重新开始",
        "从头开始",
        "reset context",
        "new topic",
    ]
    for k in kws:
        if k and k in t:
            return True
    return False


def _hivo_start_new_topic(st: dict, reason: str = ""):
    if not isinstance(st, dict):
        return
    try:
        arc = st.get("topic_archive")
        if not isinstance(arc, list):
            arc = []
            st["topic_archive"] = arc
        entry = {
            "ts": time.time(),
            "topic_id": int(st.get("topic_id") or 1),
            "reason": str(reason or ""),
            "summary": str(st.get("summary") or ""),
        }
        # keep only a small archive to avoid unbounded growth
        arc.append(entry)
        while len(arc) > 12:
            arc.pop(0)
    except Exception as e:
        logger.debug(f"Exception ignored: {e}")
    try:
        st["topic_id"] = int(st.get("topic_id") or 1) + 1
    except Exception:
        st["topic_id"] = 2
    # hard isolation: drop old topic memory/caches
    st["summary"] = ""
    st["chat"] = []
    st["tool_log"] = []
    st["tool_cache"] = OrderedDict()


def _hivo_mem_conf(cfg: dict):
    m = cfg.get("memory") if isinstance(cfg, dict) else None
    if not isinstance(m, dict):
        m = {}
    try:
        short_turns = int(m.get("short_term_turns") or 6)
    except Exception:
        short_turns = 6
    short_turns = max(3, min(12, short_turns))
    try:
        long_summary_chars = int(m.get("long_term_summary_chars") or 3500)
    except Exception:
        long_summary_chars = 3500
    long_summary_chars = max(800, min(12000, long_summary_chars))
    try:
        tool_log_items = int(m.get("tool_log_items") or 80)
    except Exception:
        tool_log_items = 80
    tool_log_items = max(10, min(300, tool_log_items))
    try:
        tool_cache_items = int(m.get("tool_cache_items") or 120)
    except Exception:
        tool_cache_items = 120
    tool_cache_items = max(20, min(500, tool_cache_items))
    return {
        "short_term_turns": short_turns,
        "long_term_summary_chars": long_summary_chars,
        "tool_log_items": tool_log_items,
        "tool_cache_items": tool_cache_items,
    }


def _hivo_summarize_for_long_term(messages: list, max_chars: int = 3500, profile_id: str = ""):
    """Generate a long-term summary from older messages.
    Uses LLM when profile_id is provided; falls back to heuristic otherwise.
    """
    try:
        arr = messages if isinstance(messages, list) else []
        if not arr:
            return ""
        
        # Format messages for the prompt (truncate each to 300 chars to fit budget)
        formatted = []
        for m in arr:
            if not isinstance(m, dict):
                continue
            role = str(m.get("role") or "").strip()
            if role not in ("user", "assistant"):
                continue
            c = str(m.get("content") or "").strip()
            if not c:
                continue
            c = c[:300].replace("\n", " ").strip()
            formatted.append(f"[{role}] {c}")
        
        if not formatted:
            return ""
        
        history_text = "\n".join(formatted[-40:])  # Last 40 messages max
        
        # Try LLM summarization if profile_id is available
        if profile_id:
            try:
                prompt = (
                    "以下是一段对话历史，请按中文简要总结（控制在" + str(max(200, int(max_chars))) + "个字符以内）：\n"
                    "1) 用户的主要目标和需求\n"
                    "2) AI 已完成的操作和结论\n"
                    "3) 当前进展状态\n\n"
                    "对话历史：\n"
                    + history_text
                )
                msgs = [{"role": "user", "content": prompt}]
                ok, msg, result = ai_chat(msgs, profile_id=profile_id, timeout_s=15, stream=False)
                if ok and result and isinstance(result.get("content"), str):
                    summary = result["content"].strip()
                    if len(summary) > 50:
                        summary = "【长期摘要记忆】\n" + summary
                        if len(summary) > int(max_chars):
                            summary = summary[:int(max_chars)]
                        return summary
            except Exception:
                pass  # Fall through to heuristic
        
        # Heuristic fallback: first line of each message, keyword classification
        items = []
        for m in arr:
            if not isinstance(m, dict):
                continue
            role = str(m.get("role") or "").strip()
            if role not in ("user", "assistant"):
                continue
            c = str(m.get("content") or "").strip()
            if not c:
                continue
            head = c.splitlines()[0].strip()
            if len(head) > 180:
                head = head[:180]
            items.append((role, head))
        if not items:
            return ""

        goal = []
        status = []
        for role, head in items:
            if role == "user":
                if re.search(r"(目标|需求|希望|请|实现|修复|优化)", head):
                    goal.append(head)
                else:
                    status.append(head)
            else:
                if re.search(r"(已完成|完成|结论|原因|根因|修复|结果|验证)", head):
                    status.append(head)
                else:
                    status.append(head)

        lines = []
        lines.append("【长期摘要记忆】")
        if goal:
            lines.append("【目标】")
            for x in goal[-8:]:
                lines.append(f"- {x}")
        if status:
            lines.append("【进展/结论】")
            for x in status[-12:]:
                lines.append(f"- {x}")
        out = "\n".join(lines).strip()
        if len(out) > int(max_chars):
            out = out[: int(max_chars)]
        return out
    except Exception:
        return ""


def _hivo_is_timeout_error(msg: str) -> bool:
    try:
        s = str(msg or "").strip().lower()
    except Exception:
        return False
    if not s:
        return False
    if "timed out" in s:
        return True
    if "timeout" in s:
        return True
    if "read operation" in s and "timed" in s:
        return True
    return False


def _hivo_resolve_path_if_needed(rp: str) -> str:
    """Resolve a possibly-ambiguous file name to a repo-relative path.

    Root goal: even if the model provides only a basename (e.g. get_ip.py / 123.png),
    tools should still work by locating the file in workspace first.
    """
    try:
        p0 = str(rp or "").strip()
        if not p0:
            return ""
        p1 = p0.replace("\\", "/")
        # Normalize chat attachment shorthand to repo-relative path
        try:
            if p1.startswith("attachments/"):
                return f"hivo_ai_data/{p1}"
        except Exception:
            pass
        if "/" in p1:
            return p1
        hits = find_files_by_name(p1, max_results=5)
        hits = hits if isinstance(hits, list) else []
        if len(hits) == 1:
            return str(hits[0] or "").strip()
        return p1
    except Exception:
        return str(rp or "").strip()


def _hivo_resolve_full_path(rp: str):
    """将相对路径解析为绝对路径，支持附件路径和仓库路径。
    返回 (relative_path, absolute_path) 元组。
    """
    try:
        p0 = str(rp or "").strip()
        p0 = _hivo_resolve_path_if_needed(p0)
        base_data = str(_hivo_ai_data_dir())
        p_norm = str(p0 or "").replace("\\", "/")
        if p_norm.startswith("hivo_ai_data/attachments/"):
            suffix = p_norm.split("hivo_ai_data/", 1)[1]
            rp_out = suffix
            full = os.path.abspath(os.path.join(base_data, suffix))
        elif p_norm.startswith("attachments/"):
            rp_out = p_norm
            full = os.path.abspath(os.path.join(base_data, p_norm))
        else:
            rp_out = p_norm
            full = _safe_repo_abspath(p_norm)
        return rp_out, full
    except Exception:
        p_out = str(rp or "").strip()
        return p_out, _safe_repo_abspath(p_out)


def _hivo_tool_call_sig(call: dict):
    try:
        if not isinstance(call, dict):
            return ""
        try:
            nm0 = str((call or {}).get("type") or "").strip()
        except Exception:
            nm0 = ""
        if nm0 in ("web_search", "web_fetch"):
            try:
                enabled0, _sc0 = _ai_load_web_search_cfg()
            except Exception:
                enabled0 = False
            payload = {"call": call, "web_search_enabled": bool(enabled0)}
            return json.dumps(payload, ensure_ascii=False, sort_keys=True)
        return json.dumps(call, ensure_ascii=False, sort_keys=True)
    except Exception:
        return ""


def _hivo_tool_log_record(tool_type: str, call: dict, ok: bool, msg: str, data: dict | None = None):
    try:
        rec = {
            "ts": time.time(),
            "type": str(tool_type or ""),
            "args": call if isinstance(call, dict) else {},
            "ok": bool(ok),
            "msg": str(msg or "")[:400],
        }
        if isinstance(data, dict):
            slim = {}
            for k in ("path", "paths", "count", "output", "content", "file", "files", "hits", "entries"):
                if k in data:
                    v = data.get(k)
                    if isinstance(v, str) and len(v) > 800:
                        v = v[:800]
                    slim[k] = v
            if slim:
                rec["data"] = slim
        return rec
    except Exception:
        return {"ts": time.time(), "type": str(tool_type or ""), "args": {}, "ok": bool(ok), "msg": str(msg or "")[:400]}


def _hivo_format_tool_memory_block(tool_log: list, limit: int = 60):
    try:
        arr = tool_log if isinstance(tool_log, list) else []
        picked = arr[-max(0, int(limit)) :]
        if not picked:
            return ""
        # Compact entries to avoid token blowup: keep type/ok/path, truncate long payloads
        compact = []
        for entry in picked:
            if not isinstance(entry, dict):
                compact.append(entry)
                continue
            c = {"type": entry.get("type"), "ok": entry.get("ok")}
            path = entry.get("path")
            if path:
                c["path"] = str(path)
            msg = entry.get("msg")
            if msg:
                c["msg"] = str(msg)[:200]
            compact.append(c)
        payload = {"recent_tool_calls": compact}
        s = json.dumps(payload, ensure_ascii=False, indent=2)
        return "【工具执行记忆（用于避免重复调用/支持结果复用）】\n" + s
    except Exception:
        return ""


def _hivo_parse_context_refs_structured(context_refs: list, per_file_chars: int = 6000, total_chars: int = 24000):
    if not isinstance(context_refs, list) or not context_refs:
        return "", []
    out = []
    used = 0
    for ref0 in context_refs[:8]:
        if used >= int(total_chars):
            break
        raw = str(ref0 or "").strip().lstrip("@")
        if not raw:
            continue

        rp = raw.replace("\\", "/").lstrip("/")

        def _resolve_file_ref(rp_in: str):
            try:
                rp2 = str(rp_in or "").replace("\\", "/").lstrip("/")
            except Exception:
                rp2 = ""
            candidates2 = []
            resolved2 = ""
            parse_way2 = ""
            conf2 = 0.0

            def _score_candidate2(rel: str):
                try:
                    p = str(rel or "")
                    if not p:
                        return 0
                    bn_in = rp2.split("/")[-1].lower()
                    bn = p.split("/")[-1].lower()
                    score = 10
                    if bn == bn_in:
                        score += 50
                    depth = p.count("/")
                    score += max(0, 12 - depth)
                    if "/" in rp2:
                        pref = "/".join(rp2.split("/")[:-1]).lower()
                        if pref and p.lower().startswith(pref + "/"):
                            score += 8
                    return score
                except Exception:
                    return 0

            if "/" in rp2:
                safe2 = _safe_repo_abspath(rp2)
                if safe2 and os.path.isfile(safe2):
                    resolved2 = rp2
                    candidates2 = [rp2]
                    parse_way2 = "精确匹配"
                    conf2 = 0.99

            if not resolved2:
                bn2 = rp2.split("/")[-1] if rp2 else ""
                cands2 = find_files_by_name(bn2, max_results=20) or []
                candidates2 = [str(x) for x in cands2 if x is not None and str(x).strip()]
                if candidates2:
                    try:
                        candidates2.sort(key=lambda x: _score_candidate2(x), reverse=True)
                    except Exception as e:
                        logger.debug(f"Exception ignored: {e}")
                if len(candidates2) == 1:
                    resolved2 = candidates2[0]
                    parse_way2 = "搜索匹配"
                    conf2 = 0.92
                elif len(candidates2) > 1:
                    resolved2 = candidates2[0]
                    parse_way2 = "搜索匹配(多候选)"
                    conf2 = 0.72

            return resolved2, candidates2, parse_way2, conf2

        try:
            if re.match(r"^https?://", str(raw or "").strip(), flags=re.I):
                url0 = str(raw or "").strip()
                item = {
                    "解析结果": "resolved",
                    "资源类型": "链接",
                    "用户输入名称": raw,
                    "真实路径": url0,
                    "文件内容": "",
                    "文件内容分段": [],
                    "分段信息": {"strategy": "none", "chunk_type": "", "target_lines": 0, "max_lines": 0, "total_parts": 0, "provided_parts": 0, "complete": True},
                    "候选列表": [],
                    "解析方式": "URL 识别",
                    "置信度": 0.99,
                    "是否需要工具调用": True,
                    "提示": "这是一个 URL 链接；如需读取网页正文，请使用 web_fetch。",
                    "建议工具调用": [{"type": "web_fetch", "url": url0}],
                }
                out.append(item)
                continue
        except Exception as e:
            logger.debug(f"Exception ignored: {e}")

        try:
            rp_dir = rp
            if rp_dir.endswith("/"):
                rp_dir = rp_dir.rstrip("/")
            safe_dir = _safe_repo_abspath(rp_dir) if rp_dir else None
            if safe_dir and os.path.isdir(safe_dir):
                tree, err = list_dir_tree(rp_dir, max_depth=3, max_entries=260)
                room = int(total_chars) - used
                if room < 0:
                    room = 0
                tree_txt = str(tree or "")
                complete0 = True
                if tree_txt and len(tree_txt) > room:
                    tree_txt = tree_txt[:room]
                    complete0 = False
                if tree_txt:
                    used += len(tree_txt)
                item = {
                    "解析结果": "resolved" if tree else "unreadable",
                    "资源类型": "目录",
                    "用户输入名称": raw,
                    "真实路径": rp_dir,
                    "文件内容": tree_txt,
                    "文件内容分段": [],
                    "分段信息": {"strategy": "dir_tree", "chunk_type": "", "target_lines": 0, "max_lines": 0, "total_parts": 1 if tree else 0, "provided_parts": 1 if tree_txt else 0, "complete": bool(complete0)},
                    "候选列表": [],
                    "解析方式": "目录匹配",
                    "置信度": 0.95,
                    "是否需要工具调用": True,
                }
                if not tree:
                    item["提示"] = str(err or "目录内容不可读取")
                else:
                    item["提示"] = "已注入目录树摘要；如需进一步确认文件位置/文件名，可用 find_files 或 search_code。"
                    item["建议工具调用"] = [{"type": "list_dir_tree", "path": rp_dir, "max_depth": 4, "max_entries": 600}]
                    if tree_txt and (not complete0):
                        item["提示"] = "目录树摘要过长，受上下文长度限制已截断；如需完整目录树，请使用 list_dir_tree 继续查看。"
                out.append(item)
                continue
        except Exception as e:
            logger.debug(f"Exception ignored: {e}")

        try:
            is_drive_path = bool(re.match(r"^[A-Za-z]:/", rp))
        except Exception:
            is_drive_path = False

        # File range reference: path:line[-end] or path#Lx-Ly
        try:
            path_part = ""
            s0 = 0
            e0 = 0
            m1 = re.match(r"^(.+?)#L(\d+)(?:-L?(\d+))?$", raw)
            if m1:
                path_part = str(m1.group(1) or "").strip().replace("\\", "/").lstrip("/")
                s0 = int(m1.group(2) or 0)
                e0 = int(m1.group(3) or 0) if m1.group(3) else s0
            else:
                m2 = re.match(r"^(.+?):(\d+)(?:-(\d+))?$", raw)
                if m2 and (not is_drive_path):
                    path_part = str(m2.group(1) or "").strip().replace("\\", "/").lstrip("/")
                    s0 = int(m2.group(2) or 0)
                    e0 = int(m2.group(3) or 0) if m2.group(3) else s0

            if path_part and s0 > 0:
                if e0 <= 0:
                    e0 = s0
                if e0 < s0:
                    s0, e0 = e0, s0
                if (e0 - s0) > 500:
                    e0 = s0 + 500

                resolved2, candidates2, parse_way2, conf2 = _resolve_file_ref(path_part)

                status2 = "not_found"
                body2 = ""
                if resolved2:
                    data2, _err2 = read_file_range(resolved2, start=int(s0), end=int(e0))
                    if data2 is not None and isinstance(data2, dict):
                        lines2 = data2.get("lines") if isinstance(data2.get("lines"), list) else []
                        body2 = "\n".join([str(x or "") for x in lines2])
                        status2 = "resolved" if body2 else "unreadable"
                    else:
                        status2 = "unreadable"
                else:
                    status2 = "ambiguous" if candidates2 else "not_found"

                room = int(total_chars) - used
                if room < 0:
                    room = 0
                complete2 = True
                if body2 and len(body2) > room:
                    body2 = body2[:room]
                    complete2 = False
                if body2:
                    used += len(body2)

                item = {
                    "解析结果": status2,
                    "资源类型": "文件片段",
                    "用户输入名称": raw,
                    "真实路径": resolved2,
                    "文件内容": body2,
                    "文件内容分段": ([{
                        "chunk_index": 1,
                        "chunk_total": 1,
                        "chunk_type": "unstructured",
                        "source": {"type": "file", "path": resolved2, "ext": resolved2.split(".")[-1].lower() if resolved2 and ("." in resolved2) else ""},
                        "range": {"start_line": int(s0), "end_line": int(e0)},
                        "content": body2,
                    }] if body2 else []),
                    "分段信息": {"strategy": "read_file_range", "chunk_type": "", "target_lines": 0, "max_lines": 0, "total_parts": 1 if body2 else 0, "provided_parts": 1 if body2 else 0, "complete": bool(complete2)},
                    "候选列表": candidates2[:8],
                    "解析方式": parse_way2,
                    "置信度": conf2,
                    "是否需要工具调用": True,
                }
                if status2 == "ambiguous":
                    item["提示"] = "存在多个候选路径，当前真实路径为最佳猜测，可能不确定；建议先确认真实路径再读取行范围。"
                    item["建议工具调用"] = [{"type": "find_files", "name": path_part.split("/")[-1]}]
                elif status2 == "not_found":
                    item["提示"] = "未找到对应文件；请确认路径或文件名。"
                    item["建议工具调用"] = [{"type": "find_files", "name": path_part.split("/")[-1]}]
                elif resolved2 and body2 and (not complete2):
                    item["提示"] = "文件片段过长，受上下文长度限制已截断；如需完整范围，请继续分批读取。"
                    item["建议工具调用"] = [{"type": "read_file_range", "path": resolved2, "start": int(s0), "end": int(e0)}]
                elif resolved2 and status2 == "unreadable":
                    item["提示"] = "文件片段读取失败或为空；请检查路径，或改用 file_content/read_file_range 进一步读取。"
                    item["建议工具调用"] = [{"type": "read_file_range", "path": resolved2, "start": int(s0), "end": int(e0)}]

                out.append(item)
                continue
        except Exception as e:
            logger.debug(f"Exception ignored: {e}")

        # Git reference: HEAD/branch/tag/commit-ish
        try:
            git_ref = str(raw or "").strip()
        except Exception:
            git_ref = ""
        if git_ref and (not is_drive_path) and ("/" not in rp):
            try:
                ext0 = rp.split("/")[-1].split(".")[-1].lower() if ("." in rp) else ""
            except Exception:
                ext0 = ""
            try:
                is_common_file = bool(ext0 and ext0 in ("py", "js", "ts", "jsx", "tsx", "html", "htm", "css", "md", "markdown", "json", "yml", "yaml", "txt", "ini", "toml", "xml", "csv"))
            except Exception:
                is_common_file = False
            if is_common_file:
                git_ref = ""

        if git_ref and (not is_drive_path) and ("/" not in rp):
            try:
                out0, err0, code0 = run_git(["rev-parse", "--verify", f"{git_ref}^{{commit}}"], timeout=15)
            except Exception:
                out0, err0, code0 = "", "", 1
            if code0 == 0 and out0:
                full_hash = str(out0 or "").strip().splitlines()[0].strip()
                short_hash = full_hash[:8] if full_hash else ""
                detail = {}
                try:
                    detail = get_commit_detail(full_hash) if full_hash else {}
                except Exception:
                    detail = {}

                lines = []
                try:
                    if isinstance(detail, dict) and ("error" not in detail):
                        h1 = str(detail.get("full_hash") or detail.get("hash") or full_hash)
                        subj = str(detail.get("subject") or detail.get("message") or "").strip()
                        author = str(detail.get("author_name") or detail.get("author") or "").strip()
                        date1 = str(detail.get("author_date") or detail.get("date") or "").strip()
                        if h1:
                            lines.append(f"commit {h1}")
                        if author or date1:
                            lines.append(f"Author: {author}  Date: {date1}".strip())
                        if subj:
                            lines.append(f"Subject: {subj}")
                        files0 = detail.get("files")
                        if isinstance(files0, list) and files0:
                            lines.append("Files:")
                            for f0 in files0[:40]:
                                if not isinstance(f0, dict):
                                    continue
                                p0 = str(f0.get("path") or "")
                                st0 = str(f0.get("status") or "")
                                if p0:
                                    lines.append(f"  {st0}  {p0}".rstrip())
                except Exception:
                    lines = []

                body0 = "\n".join(lines).strip() if lines else f"commit {full_hash}".strip()
                room = int(total_chars) - used
                if room < 0:
                    room = 0
                complete0 = True
                if body0 and len(body0) > room:
                    body0 = body0[:room]
                    complete0 = False
                if body0:
                    used += len(body0)

                item = {
                    "解析结果": "resolved",
                    "资源类型": "Git引用",
                    "用户输入名称": raw,
                    "真实路径": full_hash,
                    "文件内容": body0,
                    "文件内容分段": ([{
                        "chunk_index": 1,
                        "chunk_total": 1,
                        "chunk_type": "unstructured",
                        "source": {"type": "git", "ref": git_ref, "hash": full_hash},
                        "range": {"start_line": 1, "end_line": max(1, len(body0.splitlines())) if body0 else 1},
                        "content": body0,
                    }] if body0 else []),
                    "分段信息": {"strategy": "git_resolve", "chunk_type": "", "target_lines": 0, "max_lines": 0, "total_parts": 1 if body0 else 0, "provided_parts": 1 if body0 else 0, "complete": bool(complete0)},
                    "候选列表": [],
                    "解析方式": "git rev-parse",
                    "置信度": 0.95,
                    "是否需要工具调用": True,
                    "提示": "这是 Git 引用解析结果；如需更完整信息，请使用 commit_detail/commit_file_diff。" + ("（内容过长，已截断）" if (body0 and (not complete0)) else ""),
                    "建议工具调用": ([{"type": "commit_detail", "hash": short_hash or full_hash}] if (short_hash or full_hash) else []),
                }
                out.append(item)
                continue

        candidates = []
        resolved = ""
        parse_way = ""
        conf = 0.0

        resolved, candidates, parse_way, conf = _resolve_file_ref(rp)

        status = "not_found"
        content = ""
        parts = []
        total_parts = 0
        provided_parts = 0
        complete = False
        media_meta = {}
        if resolved:
            c0 = get_file_content(resolved)
            if c0 is not None:
                content = str(c0 or "")
            if content:
                status = "resolved"
            else:
                status = "unreadable"
        else:
            status = "ambiguous" if candidates else "not_found"

        if status == "not_found":
            try:
                q0 = str(raw or "").strip()
            except Exception:
                q0 = ""
            if q0 and len(q0) >= 2:
                hits, err = search_code(q0, case_sensitive=False, max_results=18)
                if isinstance(hits, list) and hits:
                    room = int(total_chars) - used
                    if room < 0:
                        room = 0
                    lines = []
                    for h in hits[:18]:
                        if not isinstance(h, dict):
                            continue
                        p1 = str(h.get("path") or "")
                        ln1 = int(h.get("line") or 0) if str(h.get("line") or "").strip() else 0
                        tx1 = str(h.get("text") or "")
                        if ln1 > 0:
                            lines.append(f"{p1}:{ln1}  {tx1}")
                        else:
                            lines.append(f"{p1}  {tx1}")
                    body0 = "\n".join(lines).strip()
                    complete0 = True
                    if body0 and len(body0) > room:
                        body0 = body0[:room]
                        complete0 = False
                    if body0:
                        used += len(body0)
                    item = {
                        "解析结果": "resolved",
                        "资源类型": "关键词",
                        "用户输入名称": raw,
                        "真实路径": "",
                        "文件内容": body0,
                        "文件内容分段": ([{
                            "chunk_index": 1,
                            "chunk_total": 1,
                            "chunk_type": "unstructured",
                            "source": {"type": "search", "query": q0},
                            "range": {"start_line": 1, "end_line": max(1, len(body0.splitlines())) if body0 else 1},
                            "content": body0,
                        }] if body0 else []),
                        "分段信息": {"strategy": "search_code", "chunk_type": "", "target_lines": 0, "max_lines": 0, "total_parts": 1 if body0 else 0, "provided_parts": 1 if body0 else 0, "complete": bool(complete0)},
                        "候选列表": [],
                        "解析方式": "仓库搜索",
                        "置信度": 0.70,
                        "是否需要工具调用": True,
                        "搜索命中": hits[:18],
                        "提示": "这是关键词搜索命中（非文件精确引用）。如需进一步读取上下文，请对命中项使用 file_content/read_file_range。" + ("（搜索结果过长，已截断）" if (body0 and (not complete0)) else ""),
                        "建议工具调用": [{"type": "search_code", "query": q0, "case_sensitive": False, "max_results": 50}],
                    }
                    out.append(item)
                    continue
                if err:
                    item = {
                        "解析结果": "not_found",
                        "资源类型": "关键词",
                        "用户输入名称": raw,
                        "真实路径": "",
                        "文件内容": "",
                        "文件内容分段": [],
                        "分段信息": {"strategy": "search_code", "chunk_type": "", "target_lines": 0, "max_lines": 0, "total_parts": 0, "provided_parts": 0, "complete": True},
                        "候选列表": [],
                        "解析方式": "仓库搜索(失败)",
                        "置信度": 0.40,
                        "是否需要工具调用": True,
                        "提示": str(err),
                        "建议工具调用": [{"type": "search_code", "query": q0, "case_sensitive": False, "max_results": 50}],
                    }
                    out.append(item)
                    continue

        item = {
            "解析结果": status,
            "资源类型": "文件",
            "用户输入名称": raw,
            "真实路径": resolved,
            "文件内容": "",
            "文件内容分段": [],
            "分段信息": {"strategy": "chunk", "chunk_type": "", "target_lines": 260, "max_lines": 400, "total_parts": 0, "provided_parts": 0, "complete": False},
            "候选列表": candidates[:8],
            "解析方式": parse_way,
            "置信度": conf,
            "是否需要工具调用": True,
        }

        if status == "ambiguous":
            try:
                bn = rp.split("/")[-1]
            except Exception:
                bn = raw
            item["提示"] = "存在多个候选路径，当前真实路径为最佳猜测，可能不确定；请优先通过 find_files/search_code 进一步确认后再做修改。"
            item["建议工具调用"] = [{"type": "find_files", "name": bn}]

        if resolved:
            try:
                ext = resolved.replace("\\", "/").split("/")[-1].split(".")[-1].lower() if "." in resolved else ""
            except Exception:
                ext = ""
            try:
                full = _safe_repo_abspath(resolved)
            except Exception:
                full = ""
            try:
                if full:
                    mime0, _enc0 = mimetypes.guess_type(full)
                else:
                    mime0, _enc0 = (None, None)
            except Exception:
                mime0 = None

            is_media = False
            try:
                if mime0 and (mime0.startswith("image/") or mime0.startswith("audio/") or mime0.startswith("video/")):
                    is_media = True
            except Exception:
                is_media = False
            try:
                if ext in ("png", "jpg", "jpeg", "gif", "webp", "bmp", "ico", "tiff", "mp3", "wav", "flac", "mp4", "mkv", "mov", "avi", "pdf"):
                    is_media = True
            except Exception as e:
                logger.debug(f"Exception ignored: {e}")
            try:
                if "\x00" in content:
                    is_media = True
            except Exception as e:
                logger.debug(f"Exception ignored: {e}")

            # PDF 结构化解析
        if is_pdf:
            page_count = 0
            pages = []
            ok = False
            try:
                try:
                    import PyPDF2  # type: ignore
                    with open(full, "rb") as fp:
                        r = PyPDF2.PdfReader(fp)
                        page_count = len(r.pages)
                        for i, p in enumerate(r.pages):
                            try:
                                txt = p.extract_text() or ""
                            except Exception:
                                txt = ""
                            pages.append({"page": i+1, "content": txt})
                    ok = True
                except Exception:
                    ok = False
            except Exception:
                ok = False
            return True, "", {"file": {"path": resolved, "file_type": "application/pdf", "file_size": int(os.path.getsize(full) if os.path.exists(full) else 0), "page_count": int(page_count), "pages": pages, **({"note": "pdf_parser_unavailable"} if (not ok and not pages) else {})}}

        # Office 文档结构化解析（尽力而为）
        if is_docx:
            sections = []
            ok = False
            try:
                try:
                    import docx  # type: ignore
                    d = docx.Document(full)
                    cur = {"title": "", "content": []}
                    for para in d.paragraphs:
                        t = (para.text or "").strip()
                        if not t:
                            continue
                        if para.style and ("Heading" in str(para.style.name or "")) and cur["content"]:
                            sections.append({"title": cur["title"], "content": "\n".join(cur["content"])})
                            cur = {"title": t, "content": []}
                        else:
                            if not cur["title"]:
                                cur["title"] = t[:30]
                            cur["content"].append(t)
                    if cur["content"]:
                        sections.append({"title": cur["title"], "content": "\n".join(cur["content"])})
                    ok = True
                except Exception:
                    ok = False
            except Exception:
                ok = False
            if not ok:
                # 退化为简单正文（不可用时留空）
                return True, "", {"file": {"path": resolved, "file_type": "application/vnd.openxmlformats-officedocument.wordprocessingml.document", "content": "", "note": "docx_parser_unavailable"}}
            return True, "", {"file": {"path": resolved, "file_type": "application/vnd.openxmlformats-officedocument.wordprocessingml.document", "sections": sections}}

        if is_xlsx:
            sheets = []
            ok = False
            try:
                try:
                    import openpyxl  # type: ignore
                    wb = openpyxl.load_workbook(full, read_only=True, data_only=True)
                    for ws in wb.worksheets:
                        rows = []
                        for row in ws.iter_rows(values_only=True):
                            try:
                                rows.append(["" if (v is None) else (str(v)) for v in row])
                            except Exception:
                                rows.append([""])
                        sheets.append({"name": ws.title, "rows": rows})
                    ok = True
                except Exception:
                    ok = False
            except Exception:
                ok = False
            if not ok:
                return True, "", {"file": {"path": resolved, "file_type": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", "sheets": [], "note": "xlsx_parser_unavailable"}}
            return True, "", {"file": {"path": resolved, "file_type": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", "sheets": sheets}}

        # Legacy Office: DOC (binary)
        if is_doc:
            try:
                size0 = int(os.path.getsize(full) if os.path.exists(full) else 0)
            except Exception:
                size0 = 0
            return True, "", {"file": {"path": resolved, "file_type": "application/msword", "file_size": size0, "note": "legacy_office_doc_unparsed"}}

        # Legacy Office: XLS (binary)
        if is_xls:
            try:
                size0 = int(os.path.getsize(full) if os.path.exists(full) else 0)
            except Exception:
                size0 = 0
            return True, "", {"file": {"path": resolved, "file_type": "application/vnd.ms-excel", "file_size": size0, "sheets": [], "note": "legacy_office_xls_unparsed"}}

        # Legacy Office: PPT (binary)
        if is_ppt:
            try:
                size0 = int(os.path.getsize(full) if os.path.exists(full) else 0)
            except Exception:
                size0 = 0
            return True, "", {"file": {"path": resolved, "file_type": "application/vnd.ms-powerpoint", "file_size": size0, "slides": [], "note": "legacy_office_ppt_unparsed"}}

        if is_pptx:
            slides = []
            ok = False
            try:
                try:
                    from pptx import Presentation  # type: ignore
                    prs = Presentation(full)
                    for i, s in enumerate(prs.slides):
                        title = ""
                        content = []
                        try:
                            title = s.shapes.title.text if s.shapes.title else ""
                        except Exception:
                            title = ""
                        for shp in s.shapes:
                            try:
                                if hasattr(shp, "text") and shp.text:
                                    content.append(str(shp.text))
                            except Exception:
                                continue
                        slides.append({"slide": i+1, "title": title or ("Slide " + str(i+1)), "content": "\n".join(content)})
                    ok = True
                except Exception:
                    ok = False
            except Exception:
                ok = False
            if not ok:
                return True, "", {"file": {"path": resolved, "file_type": "application/vnd.openxmlformats-officedocument.presentationml.presentation", "slides": [], "note": "pptx_parser_unavailable"}}
            return True, "", {"file": {"path": resolved, "file_type": "application/vnd.openxmlformats-officedocument.presentationml.presentation", "slides": slides}}

        if is_zip:
            entries = []
            try:
                import zipfile
                with zipfile.ZipFile(full, 'r') as zf:
                    for n in zf.namelist():
                        t, _ = mimetypes.guess_type(n)
                        entries.append({"name": n, "type": t or "application/octet-stream"})
            except Exception:
                entries = []
            return True, "", {"file": {"path": resolved, "file_type": "application/zip", "entries": entries}}

        if is_media:
            item["资源类型"] = "媒体/二进制"
            try:
                size0 = int(os.path.getsize(full)) if full and os.path.exists(full) else 0
            except Exception:
                size0 = 0
            media_meta = {
                "path": resolved,
                "mime": mime0 or "application/octet-stream",
                "size": size0,
            }
            item["媒体元信息"] = media_meta
            item["提示"] = "媒体/二进制内容不进行文本分块传输；请使用 get_file 工具读取获取实际内容（文本或 data_url），并在需要理解时提供可解析文本（如 OCR/转写/关键帧描述等）。"
        elif content and used < int(total_chars):
                # Chunking (prefer semantic boundaries; fallback to fixed lines)
                try:
                    lines = str(content or "").splitlines(True)
                except Exception:
                    lines = []

                target_lines = 260
                max_lines = 400
                try:
                    if per_file_chars:
                        # keep legacy knob but interpret as approx char budget -> line target
                        pass
                except Exception as e:
                    logger.debug(f"Exception ignored: {e}")

                def _is_boundary(li: str):
                    try:
                        s = str(li or "")
                    except Exception:
                        return False
                    if not s:
                        return False
                    if ext == "py":
                        return bool(re.match(r"^(class|def)\\s+\\w+", s))
                    if ext in ("js", "ts", "jsx", "tsx"):
                        if re.match(r"^\s*(export\s+)?class\s+\w+", s):
                            return True
                        if re.match(r"^\s*(export\s+)?(async\s+)?function\s+\w+", s):
                            return True
                        if re.match(r"^\s*(export\s+)?(const|let|var)\s+\w+\s*=\s*\(.*=>", s):
                            return True
                        return False
                    if ext in ("md", "markdown"):
                        return bool(re.match(r"^#{1,6}\\s+", s))
                    if ext in ("txt", "log"):
                        return bool(re.match(r"^\s*\d+(\\.|\\))\s+", s))
                    if ext in ("html", "htm"):
                        return bool(re.match(r"^\s*<h[1-6][\s>]", s, flags=re.I))
                    return False

                # Build segments by boundary lines
                segs = []
                if lines:
                    starts = [0]
                    for i in range(1, len(lines)):
                        if _is_boundary(lines[i]):
                            starts.append(i)
                    starts = sorted(list(dict.fromkeys(starts)))
                    for si, st0 in enumerate(starts):
                        ed0 = (starts[si + 1] if (si + 1) < len(starts) else len(lines))
                        if ed0 > st0:
                            segs.append((st0, ed0))

                # Fallback: treat whole file as one segment
                if not segs and lines:
                    segs = [(0, len(lines))]

                # Normalize segments into chunk-sized pieces
                chunks_full = []
                chunk_kind = "structured" if (len(segs) > 1 and any((b - a) > 0 for a, b in segs)) else "unstructured"
                for (a, b) in segs:
                    n = b - a
                    if n <= 0:
                        continue
                    if n <= max_lines:
                        chunks_full.append((a, b, chunk_kind))
                    else:
                        # split oversized structure by fixed lines
                        i0 = a
                        while i0 < b:
                            j0 = min(b, i0 + max_lines)
                            chunks_full.append((i0, j0, "split"))
                            i0 = j0

                total_parts = len(chunks_full)
                room = int(total_chars) - used
                if room < 0:
                    room = 0

                provided = []
                room_left = room
                for ci, (a, b, ck) in enumerate(chunks_full, start=1):
                    txt = "".join(lines[a:b])
                    if not txt:
                        continue
                    # do not cut a chunk in the middle; if doesn't fit, stop.
                    if len(txt) > room_left:
                        if (b - a) > 20:
                            # try further split by smaller blocks to fit
                            step = max(20, min(120, max_lines // 3))
                            i0 = a
                            while i0 < b:
                                j0 = min(b, i0 + step)
                                txt2 = "".join(lines[i0:j0])
                                if txt2 and len(txt2) <= room_left:
                                    provided.append({
                                        "chunk_index": len(provided) + 1,
                                        "chunk_total": 0,
                                        "chunk_type": "unstructured" if ck == "split" else "structured",
                                        "source": {"type": "file", "path": resolved, "ext": ext},
                                        "range": {"start_line": int(i0 + 1), "end_line": int(j0)},
                                        "content": txt2,
                                    })
                                    used += len(txt2)
                                    room_left -= len(txt2)
                                    i0 = j0
                                    continue
                                break
                            # stop after attempting micro-split
                        break
                    provided.append({
                        "chunk_index": len(provided) + 1,
                        "chunk_total": 0,
                        "chunk_type": "structured" if ck != "split" else "unstructured",
                        "source": {"type": "file", "path": resolved, "ext": ext},
                        "range": {"start_line": int(a + 1), "end_line": int(b)},
                        "content": txt,
                    })
                    used += len(txt)
                    room_left -= len(txt)

                provided_parts = len(provided)
                complete = bool(provided_parts >= total_parts and total_parts > 0)
                for ch in provided:
                    ch["chunk_total"] = int(total_parts)
                item["文件内容分段"] = provided
                item["分段信息"] = {
                    "strategy": "semantic_chunk_first",
                    "chunk_type": chunk_kind,
                    "target_lines": target_lines,
                    "max_lines": max_lines,
                    "total_parts": int(total_parts),
                    "provided_parts": int(provided_parts),
                    "complete": bool(complete),
                }
                if not complete and total_parts > 0:
                    item["提示"] = "内容过长，已按语义优先分块注入；但受单次上下文长度限制，本次仅提供前若干 Chunk。后续如需完整内容，请继续分批读取（建议用 read_file_range/file_content）。"

        out.append(item)
        if used >= int(total_chars):
            break

    if not out:
        return "", []
    block = "【@引用解析结果（高优先级上下文，禁止忽略）】\n【CHUNKS BEGIN】\n" + json.dumps({"refs": out}, ensure_ascii=False, indent=2) + "\n[END OF CHUNKS]"
    return block, out


def _hivo_repeat_signature(calls: list, mode: str = "tool_types"):
    try:
        m = str(mode or "tool_types").strip().lower()
        if not isinstance(calls, list):
            calls = []
        if m == "full":
            return json.dumps(calls, ensure_ascii=False, sort_keys=True)
        # default: only tool type sequence; much less sensitive to params ordering/details
        seq = [str((c or {}).get("type") or "").strip() for c in calls if isinstance(c, dict)]
        seq = [x for x in seq if x]
        return "|".join(seq)
    except Exception:
        return ""


def _hivo_repeat_escalation_prompt(level: int, last_sig: str = ""):
    lv = int(level or 1)
    sig = str(last_sig or "")
    base = (
        "【重复调用自我修正指引】\n"
        "你刚才的工具调用方案出现重复趋势。禁止输出类似 <|tool_call...|> 的文本模拟；若需要调用工具，必须输出严格合法的工具 JSON。\n"
        "你必须改变策略，而不是重复上一轮同样的工具序列。\n"
    )
    if sig:
        base += f"- 最近重复的工具序列签名: {sig}\n"

    if lv <= 1:
        return base + (
            "一级策略（先补齐信息/澄清）：\n"
            "- 如果 path/文件名不明确：先用 find_files 或 search_code 确认真实路径\n"
            "- 如果需要文件内容：用 file_content/read_file_range 获取证据后再决定 save/update\n"
            "- 如果参数缺失：先向用户提 1-2 个关键澄清问题\n"
            "输出要求：要么给出最终结论；要么输出新的（不同于上一轮策略的）工具 JSON。"
        )
    if lv == 2:
        return base + (
            "二级策略（换工具/换路径）：\n"
            "- 不要再次调用同一组工具；尝试替代路径：例如从 search_code -> list_dir_tree -> file_content\n"
            "- 将任务拆分为更小的子目标，每轮只做一个明确动作\n"
            "输出要求：必须输出与上一轮不同的工具 JSON（或明确提出你需要的关键信息）。"
        )
    return base + (
        "三级策略（重新理解目标/请求补充）：\n"
        "- 重新总结用户目标与当前已知事实\n"
        "- 明确指出卡点，并请求用户补充最小必要信息（例如：目标文件名、期望修改点、报错文本）\n"
        "输出要求：不要再尝试重复工具；以澄清问题为主。"
    )


def _hivo_agent_conf(cfg: dict):
    try:
        a = cfg.get("agent") if isinstance(cfg, dict) else None
        if not isinstance(a, dict):
            a = {}
        mode = str(a.get("mode") or "hivo_agent").strip() or "hivo_agent"
        sys_prompt = str(a.get("system_prompt") or "")
        tools = a.get("tools")
        if not isinstance(tools, list):
            tools = []
        return {
            "mode": mode,
            "system_prompt": sys_prompt,
            "tools": tools,
        }
    except Exception:
        return {"mode": "fallback_chat", "system_prompt": "", "tools": []}


def _hivo_ws_emit(run_id: str, session_id: str, stage: str, message: str = "", extra: dict | None = None):
    try:
        payload = {
            "type": "ai_agent_status",
            "run_id": str(run_id or ""),
            "session_id": str(session_id or ""),
            "stage": str(stage or ""),
            "message": str(message or ""),
            "ts": time.time(),
        }
        if isinstance(extra, dict):
            payload.update(extra)
        broadcast_to_clients(payload)
    except Exception as e:
        logger.debug(f"Exception ignored: {e}")


def _hivo_ws_emit_final(run_id: str, session_id: str, content: str, ok: bool = True, extra: dict | None = None):
    try:
        payload = {
            "type": "ai_agent_final",
            "run_id": str(run_id or ""),
            "session_id": str(session_id or ""),
            "ok": bool(ok),
            "content": str(content or ""),
            "ts": time.time(),
        }
        if isinstance(extra, dict):
            payload.update(extra)
        broadcast_to_clients(payload)
    except Exception as e:
        logger.debug(f"Exception ignored: {e}")


def _hivo_ws_emit_delta(run_id: str, session_id: str, delta: str, delta_type: str = "content"):
    try:
        payload = {
            "type": "ai_agent_delta",
            "run_id": str(run_id or ""),
            "session_id": str(session_id or ""),
            "delta": str(delta or ""),
            "delta_type": str(delta_type or "content"),
            "ts": time.time(),
        }
        broadcast_to_clients(payload)
    except Exception as e:
        logger.debug(f"Exception ignored: {e}")


def get_file_state_hash():
    """获取当前文件状态的哈希值"""
    if not REPO_PATH:
        return None
    try:
        # 轻量级状态：基于 git status porcelain 输出 + 相关文件 mtime
        # 仅用 status 输出会导致"持续编辑但状态不变（一直是 M）"时无法触发推送。
        # 这里额外叠加 mtime，避免触发昂贵的 git diff 统计。
        out, err, code = run_git(["status", "--porcelain=v1", "-u", "-z"], timeout=30)
        if code != 0:
            logger.debug(f"获取文件状态失败: {err}")
            return None

        entries = (out or "").split("\x00")
        state_items = []
        i = 0
        while i < len(entries):
            entry = entries[i]
            if not entry:
                i += 1
                continue
            if len(entry) < 4:
                i += 1
                continue

            xy = entry[:2]
            path = entry[3:]
            idx_s = xy[0]

            paths = [(path, xy)]
            # Rename/Copy 在 -z 下会额外带一个 path
            if idx_s in ("R", "C") and (i + 1) < len(entries):
                new_path = entries[i + 1]
                if new_path:
                    paths.append((new_path, xy))
                i += 1

            for p, flag in paths:
                p_key = (p or "").replace("\\", "/")
                if not p_key:
                    continue
                full = os.path.join(REPO_PATH, p_key)
                if os.path.exists(full):
                    try:
                        mtime = os.path.getmtime(full)
                    except Exception:
                        mtime = 0
                else:
                    # 文件不存在（删除/重命名旧路径等）：必须使用稳定值
                    mtime = 0
                state_items.append((p_key, flag, mtime))

            i += 1

        state_str = json.dumps(sorted(state_items), sort_keys=True)
        return hashlib.md5(state_str.encode("utf-8", errors="replace")).hexdigest()
    except Exception as e:
        logger.error(f"获取文件状态哈希失败: {e}")
        return None


def get_changed_files_cached(max_age_sec=1.0):
    """带短缓存的变更文件列表。

    典型场景：
    - 前端 /api/files 轮询备份
    - WebSocket 初始化/按需请求
    """
    now = time.time()
    with _changed_files_lock:
        try:
            ts = float(_changed_files_cache.get("ts") or 0.0)
        except Exception:
            ts = 0.0

        cached = _changed_files_cache.get("files")
        if cached is not None and (now - ts) <= float(max_age_sec):
            return cached

    files = get_changed_files()
    with _changed_files_lock:
        _changed_files_cache["ts"] = now
        _changed_files_cache["files"] = files
    return files


def watch_repository():
    """监控仓库变化并通知客户端"""
    global last_file_state
    
    while True:
        try:
            has_clients = False
            with ws_clients_lock:
                has_clients = bool(ws_clients)

            if REPO_PATH and has_clients:
                current_state = get_file_state_hash()
                if current_state and current_state != last_file_state:
                    logger.debug(f"检测到文件变化: {last_file_state} -> {current_state}")
                    last_file_state = current_state
                    
                    # 获取最新的文件列表
                    files = get_changed_files_cached(max_age_sec=0)
                    
                    # 广播更新消息
                    broadcast_to_clients({
                        'type': 'files_updated',
                        'files': files
                    })
            
            time.sleep(1)  # 每秒检查一次
        except Exception as e:
            logger.error(f"仓库监控异常: {e}", exc_info=True)
            time.sleep(5)


def handle_websocket(websocket: ServerConnection):
    """处理WebSocket连接"""
    logger.info(f"新的WebSocket连接: {websocket.remote_address}")
    with ws_clients_lock:
        ws_clients.add(websocket)
        if websocket not in ws_client_queues:
            ws_client_queues[websocket] = queue.Queue(maxsize=200)
    
    try:
        # 发送欢迎消息
        websocket.send(json.dumps({
            'type': 'connected',
            'message': 'WebSocket连接成功'
        }))
        
        # 如果已经打开了仓库，立即发送当前状态
        if REPO_PATH:
            files = get_changed_files_cached()
            websocket.send(json.dumps({
                'type': 'files_updated',
                'files': files
            }))
        
        # 主循环：仅在此线程调用 websocket.send（线程安全）
        while True:
            # drain outgoing queue
            try:
                q = None
                with ws_clients_lock:
                    q = ws_client_queues.get(websocket)
                if q is not None:
                    while True:
                        try:
                            payload = q.get_nowait()
                        except Exception:
                            break
                        websocket.send(payload)
            except Exception as e:
                logger.warning(f"WebSocket发送失败: {e}")
                break

            # receive client message with short timeout so we can keep draining queue
            try:
                message = websocket.recv(timeout=0.5)
            except TimeoutError:
                continue
            except Exception as e:
                try:
                    if ConnectionClosed is not None and isinstance(e, ConnectionClosed):
                        code = getattr(e, "code", None)
                        reason = getattr(e, "reason", "")
                        logger.info(f"WebSocket接收关闭: code={code}, reason={reason}")
                    else:
                        logger.error(f"WebSocket接收异常: {e}")
                except Exception:
                    logger.error(f"WebSocket接收异常: {e}")
                break

            if message is None:
                continue
            try:
                data = json.loads(message)
                msg_type = data.get('type')

                if msg_type == 'ping':
                    websocket.send(json.dumps({'type': 'pong'}))
                elif msg_type == 'request_files':
                    if REPO_PATH:
                        files = get_changed_files_cached()
                        websocket.send(json.dumps({
                            'type': 'files_updated',
                            'files': files
                        }))
                else:
                    logger.warning(f"未知的WebSocket消息类型: {msg_type}")

            except json.JSONDecodeError:
                logger.warning(f"无效的JSON消息: {message}")
            except Exception as e:
                logger.error(f"处理WebSocket消息异常: {e}", exc_info=True)
                
    except Exception as e:
        logger.error(f"WebSocket连接异常: {e}")
    finally:
        with ws_clients_lock:
            ws_clients.discard(websocket)
            try:
                ws_client_queues.pop(websocket, None)
            except Exception:
                pass
        ra = None
        try:
            ra = websocket.remote_address
        except Exception:
            ra = None
        logger.info(f"WebSocket连接关闭: {ra}")


def start_websocket_server():
    """启动WebSocket服务器"""
    if not WEBSOCKET_AVAILABLE:
        logger.warning("WebSocket库未安装，跳过WebSocket服务器启动")
        return
    
    try:
        ws_port = 7843
        max_attempts = 10
        
        for attempt in range(max_attempts):
            try:
                with ws_serve(handle_websocket, "127.0.0.1", ws_port) as server:
                    logger.info(f"WebSocket服务器启动成功: ws://localhost:{ws_port}")
                    
                    # 更新全局变量，让前端知道正确的WebSocket端口
                    global WS_PORT
                    WS_PORT = ws_port
                    
                    # 启动文件监控线程
                    watch_thread = threading.Thread(target=watch_repository, daemon=True)
                    watch_thread.start()
                    logger.info("文件监控线程已启动")
                    
                    # 保持服务器运行
                    server.serve_forever()
            except OSError as e:
                if e.errno == 10048:  # 端口被占用
                    logger.warning(f"端口 {ws_port} 被占用，尝试下一个端口...")
                    ws_port += 1
                    continue
                else:
                    raise
    except Exception as e:
        logger.error(f"WebSocket服务器启动失败: {e}", exc_info=True)


# ════════════════════════════════════════════════════════
#  Git 执行核心 — 处理 Windows GBK / 中文路径 / 空格
# ════════════════════════════════════════════════════════

def run_git(args, cwd=None, input_data=None, timeout=60):
    """执行 Git 命令并记录日志"""
    cwd = cwd or REPO_PATH
    env = os.environ.copy()
    env["GIT_TERMINAL_PROMPT"] = "0"
    env["LANG"]                = "en_US.UTF-8"
    env["LC_ALL"]              = "en_US.UTF-8"
    env["PYTHONIOENCODING"]    = "utf-8"

    # 记录命令执行
    cmd_str = " ".join(["git"] + args)
    logger.debug(f"执行 Git 命令: {cmd_str} (工作目录: {cwd})")

    def decode_bytes(b):
        if not b:
            return ""
        for enc in ("utf-8", "utf-8-sig", "gbk", "gb2312", "cp936", "latin-1"):
            try:
                return b.decode(enc)
            except Exception:
                continue
        return b.decode("utf-8", errors="replace")

    try:
        inp = input_data.encode("utf-8") if isinstance(input_data, str) else input_data
        # -c core.quotepath=false  → 不转义中文/特殊字符路径
        r = subprocess.run(
            ["git", "-c", "core.quotepath=false"] + args,
            cwd=cwd,
            capture_output=True,
            text=False,
            input=inp,
            timeout=timeout,
            env=env,
        )
        stdout_str = decode_bytes(r.stdout)
        stderr_str = decode_bytes(r.stderr)
        
        if r.returncode == 0:
            logger.debug(f"Git 命令执行成功: {cmd_str}")
            if stdout_str:
                logger.debug(f"标准输出: {stdout_str[:200]}...")  # 只记录前200个字符
        else:
            logger.warning(f"Git 命令执行失败 (返回码: {r.returncode}): {cmd_str}")
            if stderr_str:
                logger.warning(f"错误输出: {stderr_str}")
        
        return stdout_str, stderr_str, r.returncode
    except FileNotFoundError:
        error_msg = "git 未安装或不在 PATH 中"
        logger.error(f"Git 命令执行失败: {error_msg}")
        return "", error_msg, 1
    except subprocess.TimeoutExpired:
        error_msg = f"命令超时(超过{timeout}秒)"
        logger.error(f"Git 命令超时: {cmd_str} - {error_msg}")
        return "", error_msg, 1
    except Exception as e:
        logger.error(f"Git 命令执行异常: {cmd_str} - {str(e)}", exc_info=True)
        return "", str(e), 1

def unstage_file(rel_path: str):
    if not REPO_PATH:
        return False, "未打开仓库"
    rp = str(rel_path or "").strip()
    if not rp:
        return False, "缺少 path"
    _, err, code = run_git(["restore", "--staged", "--", rp], timeout=120)
    if code != 0:
        return False, err or "取消暂存失败"
    return True, ""

def discard_staged_file(rel_path: str):
    if not REPO_PATH:
        return False, "未打开仓库"
    rp = str(rel_path or "").strip()
    if not rp:
        return False, "缺少 path"
    _, err, code = run_git(["restore", "--staged", "--worktree", "--source=HEAD", "--", rp], timeout=120)
    if code != 0:
        return False, err or "丢弃暂存区失败"
    return True, ""

def unstage_all_staged():
    if not REPO_PATH:
        return False, "未打开仓库"
    _, err, code = run_git(["restore", "--staged", "."], timeout=180)
    if code != 0:
        return False, err or "取消全部暂存失败"
    return True, ""

def discard_all_staged():
    if not REPO_PATH:
        return False, "未打开仓库"
    _, err, code = run_git(["restore", "--staged", "--worktree", "--source=HEAD", "."], timeout=180)
    if code != 0:
        return False, err or "丢弃全部暂存失败"
    return True, ""


def _extract_conflict_hunks(content):
    """从文件内容中提取所有 <<<<<<< / ======= / >>>>>>> 冲突块。
    返回: [{start: 行号(1-based), end: 行号, ours: str, theirs: str, base_section: str}, ...]
    """
    hunks = []
    if not content:
        return hunks
    lines = content.split("\n")
    i = 0
    n = len(lines)
    while i < n:
        line = lines[i]
        if line.startswith("<<<<<<<"):
            start = i + 1  # 1-based
            ours_lines = []
            theirs_lines = []
            base_lines = []
            sep_seen = False
            end_seen = False
            marker_end = None
            j = i + 1
            while j < n:
                lj = lines[j]
                if lj.startswith("======="):
                    sep_seen = True
                    j += 1
                    continue
                if lj.startswith(">>>>>>>"):
                    end_seen = True
                    marker_end = j
                    break
                # 处理 ||||||| ours base（三方合并时可能存在）
                if lj.startswith("|||||||"):
                    # 旧版本块开始，跳到下一个 =======
                    j += 1
                    base_section = []
                    while j < n and not lines[j].startswith("======="):
                        base_section.append(lines[j])
                        j += 1
                    base_lines = base_section
                    if j < n and lines[j].startswith("======="):
                        sep_seen = True
                        j += 1
                    continue
                if not sep_seen:
                    ours_lines.append(lj)
                else:
                    theirs_lines.append(lj)
                j += 1
            if end_seen and marker_end is not None:
                hunks.append({
                    "start": start,
                    "end": marker_end + 1,  # 1-based, 含 >>>>>>> 行
                    "ours": "\n".join(ours_lines),
                    "theirs": "\n".join(theirs_lines),
                    "base": "\n".join(base_lines)
                })
            i = marker_end + 1 if marker_end is not None else i + 1
        else:
            i += 1
    return hunks


def get_unmerged_files():
    """Return a list of unmerged (conflicted) files in the working tree.

    Used by /api/pull to detect merge conflicts.
    Returns: (files: List[str], raw_output: str)
    """
    if not REPO_PATH:
        return [], ""

    # Prefer diff-filter=U which directly lists unmerged paths
    out, err, code = run_git(["diff", "--name-only", "--diff-filter=U", "-z"], timeout=60)
    raw = (out or "")
    if code != 0:
        # Fallback to ls-files -u when diff fails
        out2, err2, code2 = run_git(["ls-files", "-u", "-z"], timeout=60)
        raw = (out2 or "")
        if code2 != 0:
            logger.debug(f"获取冲突文件失败: {err or err2}")
            return [], raw

        # ls-files -u -z format: <mode> <sha> <stage>\t<path>\0 ...
        files = []
        for item in raw.split("\x00"):
            if not item:
                continue
            if "\t" in item:
                path = item.split("\t", 1)[1]
                if path:
                    files.append(path)
        uniq = sorted({p.replace("\\", "/") for p in files if p})
        return uniq, raw

    items = [p for p in raw.split("\x00") if p]
    uniq = sorted({p.replace("\\", "/") for p in items if p})
    return uniq, raw

def has_any_staged_changes():
    if not REPO_PATH:
        return False, "未打开仓库"
    out, err, code = run_git(["diff", "--cached", "--name-only", "-z"], timeout=60)
    if code != 0:
        return False, (err or "检查暂存区失败")
    raw = out or ""
    files = [p for p in raw.split("\x00") if p]
    return (len(files) > 0), None

def get_staged_files():
    if not REPO_PATH:
        return [], "未打开仓库"
    out, err, code = run_git(["diff", "--cached", "--name-status", "-z"], timeout=60)
    if code != 0:
        return [], (err or "读取暂存区失败")
    raw = out or ""
    parts = [x for x in raw.split("\x00") if x]
    files = []
    i = 0
    while i < len(parts):
        item = parts[i]
        i += 1
        if not item:
            continue
        cols = item.split("\t")
        st = (cols[0] or "").strip()
        path = ""
        if len(cols) >= 2:
            path = cols[1]
        else:
            if i < len(parts):
                path = parts[i]
                i += 1
        st2 = (st[:1] or "M").upper()
        if st2 not in ("A", "M", "D", "R", "C"):
            st2 = "M"
        p = (path or "").replace("\\", "/")
        if not p:
            continue
        files.append({"path": p, "status": st2})
    return files, None

def _safe_repo_abspath(rel_path: str):
    """Resolve a repo-relative path to an absolute path, preventing path traversal."""
    if not REPO_PATH:
        return None
    repo_root = os.path.abspath(REPO_PATH)
    rel_path = (rel_path or "")
    try:
        rel_path = rel_path.replace("\x00", "")
    except Exception as e:
        logger.debug(f"Exception ignored: {e}")
    rel_path = str(rel_path).strip()
    rel_path = rel_path.replace("\\", "/")
    if rel_path.lower().startswith("file://"):
        rel_path = rel_path[7:]

    try:
        if os.path.isabs(rel_path) or re.match(r"^[A-Za-z]:[\\/]", rel_path or ""):
            abs_in = os.path.abspath(rel_path)
            try:
                if os.path.commonpath([repo_root, abs_in]) != repo_root:
                    return None
            except Exception:
                return None
            rel_path = os.path.relpath(abs_in, repo_root).replace("\\", "/")
    except Exception as e:
        logger.debug(f"Exception ignored: {e}")

    rel_path = rel_path.lstrip("/")
    rel_path = rel_path.replace("\r", "").replace("\n", "")
    if not rel_path:
        return None

    bad_chars = '<>:"|?*'
    parts = [p for p in rel_path.split("/") if p not in ("", ".")]
    for p in parts:
        if p in ("..",):
            return None
        if any((c in p) for c in bad_chars):
            return None
        if any((ord(c) < 32) for c in p):
            return None
        if p.endswith(" ") or p.endswith("."):
            return None

    rel_path = "/".join(parts)
    abs_path = os.path.abspath(os.path.join(repo_root, rel_path.replace("/", os.sep)))
    try:
        if os.path.commonpath([repo_root, abs_path]) != repo_root:
            return None
    except Exception:
        return None
    return abs_path


def get_file_content(filepath: str, return_encoding=False):
    """Read working tree file content as text.
    
    Args:
        filepath: 文件路径
        return_encoding: 是否返回检测到的编码，默认False（向后兼容）
    
    Returns:
        如果return_encoding=False: 返回文件内容字符串或None
        如果return_encoding=True: 返回(内容字符串, 编码名称)或(None, None)
    """
    try:
        full = _safe_repo_abspath(filepath)
        if not full or (not os.path.exists(full)) or os.path.isdir(full):
            return (None, None) if return_encoding else None
        
        # Check file content cache (keyed by path + mtime)
        try:
            mtime = int(os.path.getmtime(full))
            cache_key = (full, mtime)
            with _file_content_cache_lock:
                if cache_key in _file_content_cache:
                    cached = _file_content_cache[cache_key]
                    _file_content_cache.move_to_end(cache_key)  # LRU: touch
                    return (cached[0], cached[1]) if return_encoding else cached[0]
        except Exception:
            cache_key = None

        # Helper: store decoded result into in-memory cache then return
        def _store_and_return(content, enc):
            try:
                if cache_key and content is not None:
                    with _file_content_cache_lock:
                        _file_content_cache[cache_key] = (content, enc)
                        if len(_file_content_cache) > _FILE_CONTENT_CACHE_MAX:
                            try: _file_content_cache.popitem(last=False)
                            except Exception: pass
            except Exception:
                pass
            return (content, enc) if return_encoding else content
        
        with open(full, "rb") as f:
            data = f.read()
        
        # 如果数据为空,直接返回空字符串
        if not data:
            return ("", "utf-8") if return_encoding else ""
        
        detected_encoding = None

        # 首先尝试 UTF-8 / UTF-8-BOM（最快路径）
        try:
            result = data.decode("utf-8")
            detected_encoding = "utf-8"
            logger.debug(f"文件 {filepath} 使用 UTF-8 编码")
            try:
                with _file_encoding_lock:
                    _file_last_encoding[filepath] = detected_encoding
                    _file_decode_lossy[filepath] = False
            except Exception as e:
                logger.debug(f"Exception ignored: {e}")
            return _store_and_return(result, detected_encoding)
        except UnicodeDecodeError:
            pass

        try:
            result = data.decode("utf-8-sig")
            detected_encoding = "utf-8-sig"
            logger.debug(f"文件 {filepath} 使用 UTF-8-BOM 编码")
            try:
                with _file_encoding_lock:
                    _file_last_encoding[filepath] = detected_encoding
                    _file_decode_lossy[filepath] = False
            except Exception as e:
                logger.debug(f"Exception ignored: {e}")
            return _store_and_return(result, detected_encoding)
        except UnicodeDecodeError:
            pass

        # UTF-8 严格解码失败：优先尝试常见中文编码（严格模式），避免探测误判为单字节编码导致乱码。
        for enc0 in ("gb18030", "gbk", "cp936", "gb2312"):
            try:
                result = data.decode(enc0)
                detected_encoding = enc0
                logger.debug(f"文件 {filepath} 使用常见中文编码读取: {enc0}")
                try:
                    with _file_encoding_lock:
                        _file_last_encoding[filepath] = detected_encoding
                        _file_decode_lossy[filepath] = False
                except Exception as e:
                    logger.debug(f"Exception ignored: {e}")
                return _store_and_return(result, detected_encoding)
            except UnicodeDecodeError:
                continue

        # UTF-8 严格解码失败：探测编码（gbk/gb2312/gb18030 等），用于前端显示。
        # 注意：这里仅用于"读取展示"，保存仍由 save_file_content() 决定编码，避免扩大损害。
        try:
            enc = _detect_text_encoding_from_bytes(data)
            detected_encoding = enc
            lossy = False
            try:
                # strict 解码成功：无损
                result = data.decode(enc)
            except UnicodeDecodeError:
                # 为保证前端可显示，回退 replace，但标记为 lossy
                result = data.decode(enc, errors="replace")
                lossy = True
            logger.debug(f"文件 {filepath} 使用检测编码读取: {enc}{' (replace)' if lossy else ''}")
            try:
                with _file_encoding_lock:
                    _file_last_encoding[filepath] = detected_encoding
                    _file_decode_lossy[filepath] = bool(lossy)
            except Exception as e:
                logger.debug(f"Exception ignored: {e}")
            return _store_and_return(result, detected_encoding)
        except Exception as e:
            logger.warning(f"文件 {filepath} 编码探测/解码失败: {e}，回退 UTF-8+replace")

        result = data.decode("utf-8", errors="replace")
        detected_encoding = "utf-8"
        try:
            with _file_encoding_lock:
                _file_last_encoding[filepath] = detected_encoding
                _file_decode_lossy[filepath] = True
        except Exception as e:
            logger.debug(f"Exception ignored: {e}")
        return _store_and_return(result, detected_encoding)
    except Exception as e:
        logger.error(f"读取文件内容失败: {filepath} - {e}")
        return (None, None) if return_encoding else None


def _detect_text_encoding_from_bytes(data: bytes):
    """检测文本数据的编码格式"""
    if data is None or len(data) == 0:
        return "utf-8"
    
    # 检查 UTF-8 BOM
    try:
        if data.startswith(b"\xef\xbb\xbf"):
            return "utf-8-sig"
    except Exception as e:
        logger.debug(f"Exception ignored: {e}")
    
    # 首先尝试UTF-8（最常见）
    try:
        data.decode('utf-8')
        return "utf-8"
    except UnicodeDecodeError:
        pass
    
    # 使用chardet检测（如果可用）
    if CHARDET_AVAILABLE:
        try:
            detected = chardet.detect(data)
            if detected and detected.get('encoding'):
                detected_enc = detected['encoding']
                confidence = detected.get('confidence', 0)
                
                if confidence > 0.5:
                    detected_enc_lower = detected_enc.lower()
                    # 避免单字节编码误判：这类编码几乎总能"解码成功"，但中文会变成乱码
                    if (
                        detected_enc_lower in ('latin-1', 'iso-8859-1') or
                        detected_enc_lower.startswith('windows-125') or
                        detected_enc_lower.startswith('iso-8859-')
                    ):
                        detected_enc = None
                        detected_enc_lower = ''

                    # 避免误判为罕见的 cpXXXX 编码（Windows 默认代码页/阿拉伯语等），
                    # 这些编码对中文项目极易造成保存失败或乱码。
                    if detected_enc_lower.startswith('cp') and detected_enc_lower not in ('cp936',):
                        detected_enc = None
                        detected_enc_lower = ''

                    # 标准化编码
                    if detected_enc_lower in ('gb2312',):
                        return 'gb2312'
                    if detected_enc_lower in ('gbk', 'windows-1252'):
                        return 'gbk'
                    if detected_enc_lower in ('gb18030',):
                        return 'gb18030'
                    elif detected_enc_lower.startswith('utf'):
                        return 'utf-8'
                    elif detected_enc_lower in ('cp936',):
                        return 'cp936'
                    else:
                        # 验证检测到的编码是否有效
                        try:
                            data.decode(detected_enc)
                            return detected_enc
                        except Exception as e:
                            logger.debug(f"Exception ignored: {e}")
        except Exception as e:
            logger.debug(f"Exception ignored: {e}")
    
    # Fallback: 按优先级尝试常见编码
    for enc in ("gb18030", "gbk", "gb2312", "cp936"):
        try:
            data.decode(enc)
            return enc
        except Exception:
            continue
    
    return "utf-8"


def get_head_file_content(filepath: str):
    """Read HEAD version file content as text via git show."""
    filepath = (filepath or "").replace("\\", "/").lstrip("/")
    if not filepath:
        return None
    out, err, code = run_git(["show", f"HEAD:{filepath}"])
    if code != 0:
        logger.debug(f"读取 HEAD 文件内容失败: {filepath} - {err}")
        return None
    return out


def list_dir_tree(rel_path: str = "", max_depth: int = 3, max_entries: int = 500):
    if not REPO_PATH:
        return None, "未打开仓库"
    try:
        depth = int(max_depth)
    except Exception:
        depth = 3
    if depth < 0:
        depth = 0
    if depth > 6:
        depth = 6
    try:
        cap = int(max_entries)
    except Exception:
        cap = 500
    if cap <= 0:
        cap = 200
    if cap > 2000:
        cap = 2000

    rel = (rel_path or "").replace("\\", "/").lstrip("/")
    root_abs = _safe_repo_abspath(rel) if rel else os.path.abspath(REPO_PATH)
    if not root_abs or (not os.path.isdir(root_abs)):
        return None, "目录不存在"
    repo_root = os.path.abspath(REPO_PATH)

    lines = []
    count = 0
    for cur_root, dirs, files in os.walk(root_abs):
        try:
            rel_cur = os.path.relpath(cur_root, repo_root).replace("\\", "/")
        except Exception:
            rel_cur = rel
        cur_depth = 0
        if rel_cur and rel_cur != ".":
            cur_depth = rel_cur.count("/") + 1
        base_depth = 0
        if rel:
            base_depth = rel.count("/") + 1
        if (cur_depth - base_depth) >= depth:
            dirs[:] = []

        dirs[:] = [d for d in sorted(dirs) if d != ".git" and not d.startswith(".")]
        files = [f for f in sorted(files) if not f.startswith(".")]

        indent_level = max(0, cur_depth - base_depth)
        if rel_cur == ".":
            show_dir = "."
        else:
            show_dir = rel_cur
        if show_dir:
            lines.append("  " * indent_level + show_dir + "/")
            count += 1
            if count >= cap:
                break

        for d in dirs:
            if count >= cap:
                break
            lines.append("  " * (indent_level + 1) + d + "/")
            count += 1
        for f in files:
            if count >= cap:
                break
            lines.append("  " * (indent_level + 1) + f)
            count += 1
        if count >= cap:
            break

    if count >= cap:
        lines.append("…")
    return "\n".join(lines), ""


def read_file_range(rel_path: str, start: int = 1, end: int = 200):
    if not REPO_PATH:
        return None, "未打开仓库"
    p = (rel_path or "").replace("\\", "/").lstrip("/")
    if not p:
        return None, "缺少 path"
    full = _safe_repo_abspath(p)
    if not full:
        return None, "非法路径"
    if not os.path.exists(full) or os.path.isdir(full):
        return None, "文件不存在"

    try:
        s = int(start)
    except Exception:
        s = 1
    try:
        e = int(end)
    except Exception:
        e = s + 200
    if s < 1:
        s = 1
    if e < s:
        e = s
    if (e - s) > 500:
        e = s + 500

    try:
        with open(full, "r", encoding="utf-8", errors="replace") as f:
            out = []
            idx = 0
            for line in f:
                idx += 1
                if idx < s:
                    continue
                if idx > e:
                    break
                out.append(line.rstrip("\n"))
        return {"path": p, "start": s, "end": e, "lines": out}, ""
    except Exception as ex:
        return None, str(ex)


def search_code(query: str, case_sensitive: bool = False, max_results: int = 50, max_file_size: int = 512 * 1024):
    if not REPO_PATH:
        return None, "未打开仓库"
    q = str(query or "")
    if not q.strip():
        return None, "缺少 query"
    try:
        cap = int(max_results)
    except Exception:
        cap = 50
    if cap <= 0:
        cap = 20
    if cap > 200:
        cap = 200

    repo_root = os.path.abspath(REPO_PATH)
    flags = 0 if case_sensitive else re.IGNORECASE
    try:
        pat = re.compile(q, flags)
    except Exception:
        pat = None

    out = []
    for root, dirs, files in os.walk(repo_root):
        if ".git" in dirs:
            dirs.remove(".git")
        dirs[:] = [d for d in dirs if not d.startswith(".")]
        for fn in files:
            if fn.startswith("."):
                continue
            abs_path = os.path.join(root, fn)
            try:
                if os.path.getsize(abs_path) > int(max_file_size):
                    continue
            except Exception:
                continue
            rel = os.path.relpath(abs_path, repo_root).replace("\\", "/")
            try:
                with open(abs_path, "r", encoding="utf-8", errors="replace") as f:
                    ln = 0
                    for line in f:
                        ln += 1
                        hit = False
                        if pat is not None:
                            hit = bool(pat.search(line))
                        else:
                            if case_sensitive:
                                hit = (q in line)
                            else:
                                hit = (q.lower() in line.lower())
                        if not hit:
                            continue
                        out.append({"path": rel, "line": ln, "text": line.rstrip("\n")})
                        if len(out) >= cap:
                            return out, ""
            except Exception:
                continue

    return out, ""


def save_file_content(filepath: str, content: str, force_encoding: str = None):
    """Save content to working tree file (text).
    
    Args:
        filepath: 文件路径
        content: 文件内容
        force_encoding: 强制使用的编码（如果提供，将覆盖自动检测）
    """
    try:
        full = _safe_repo_abspath(filepath)
        if not full:
            logger.error(f"保存文件失败: 非法路径 {filepath}")
            return False, "非法路径"
        
        logger.info(f"开始保存文件: {filepath}, 内容长度: {len(content) if content else 0}, 强制编码: {force_encoding}")
        
        parent = os.path.dirname(full)
        if parent and (not os.path.exists(parent)):
            os.makedirs(parent, exist_ok=True)
            logger.debug(f"创建目录: {parent}")

        target_enc = "utf-8"
        file_exists = os.path.exists(full) and (not os.path.isdir(full))
        
        # 如果指定了编码，优先使用
        if force_encoding:
            target_enc = force_encoding
            logger.debug(f"使用强制指定的编码: {target_enc}")
        else:
            try:
                if file_exists:
                    with open(full, "rb") as rf:
                        raw0 = rf.read()
                    target_enc = _detect_text_encoding_from_bytes(raw0)
                    logger.debug(f"检测到已存在文件编码: {target_enc}")
                else:
                    logger.debug(f"新文件，使用默认编码: {target_enc}")
            except Exception as e:
                logger.warning(f"检测文件编码失败: {e}, 使用默认UTF-8")
                target_enc = "utf-8"

        rel_path = (filepath or "").replace("\\", "/").lstrip("/")

        # 如果该文件在最近一次读取时发生过 replace（lossy 解码），说明原始字节无法从文本无损还原。
        # 这种情况下禁止保存，避免写坏文件（VS2019/编译器会报错）。
        try:
            with _file_encoding_lock:
                lossy = _file_decode_lossy.get(filepath)
                enc_hint = _file_last_encoding.get(filepath) or target_enc
            if lossy:
                return False, f"该文件读取时发生编码替换（{enc_hint} + replace），无法保证无损保存。请用正确编码打开/修复原始字节后再保存。"
        except Exception as e:
            logger.debug(f"Exception ignored: {e}")

        # Decide target EOL:
        # Priority: 1) detect from original file (most reliable)
        #           2) git attributes (eol)
        #           3) core.autocrlf (working tree convention)
        #           4) fallback: LF
        target_eol = None

        # Step 1: Detect from existing file (highest priority - preserves original)
        try:
            if file_exists:
                with open(full, "rb") as rf:
                    raw_eol = rf.read()
                if b"\r\n" in raw_eol:
                    target_eol = "\r\n"
                elif b"\n" in raw_eol:
                    target_eol = "\n"
        except Exception as e:
            logger.warning(f"检测文件换行符失败: {e}")

        # Step 2: Check git attributes if we couldn't detect from file
        if target_eol is None:
            try:
                out, _, code = run_git(["check-attr", "-z", "eol", "--", rel_path], timeout=30)
                if code == 0 and out:
                    parts = out.split("\x00")
                    # format: path\0attr\0value\0
                    if len(parts) >= 3:
                        val = (parts[2] or "").strip().lower()
                        if val == "crlf":
                            target_eol = "\r\n"
                            logger.debug(f"从 git 属性检测到 {filepath} 使用 CRLF")
                        elif val == "lf":
                            target_eol = "\n"
                            logger.debug(f"从 git 属性检测到 {filepath} 使用 LF")
            except Exception as e:
                logger.warning(f"检查 git 属性失败: {e}")

        # Step 3: Check core.autocrlf if still unknown
        if target_eol is None:
            try:
                out2, _, code2 = run_git(["config", "--get", "core.autocrlf"], timeout=10)
                if code2 == 0 and (out2 or "").strip().lower() == "true":
                    target_eol = "\r\n"
                    logger.debug("从 core.autocrlf 检测到使用 CRLF")
            except Exception as e:
                logger.warning(f"检查 core.autocrlf 失败: {e}")

        # Step 4: Default to LF if still unknown
        if target_eol is None:
            target_eol = "\n"
            logger.debug("使用默认换行符 LF")

        txt = content if content is not None else ""
        
        # Check if content already has the correct line endings
        has_crlf = "\r\n" in txt
        has_lf_only = "\n" in txt and not has_crlf
        
        if target_eol == "\r\n" and has_crlf:
            # Content already has CRLF and we want CRLF - no conversion needed
            logger.debug("内容已是 CRLF 格式,无需转换")
            pass
        elif target_eol == "\n" and has_lf_only:
            # Content already has LF only and we want LF - no conversion needed
            logger.debug("内容已是 LF 格式,无需转换")
            pass
        else:
            # Need to normalize and convert
            logger.debug(f"转换换行符: 当前={'CRLF' if has_crlf else 'LF'} -> 目标={repr(target_eol)}")
            txt = str(txt).replace("\r\n", "\n").replace("\r", "\n")
            if target_eol != "\n":
                txt = txt.replace("\n", target_eol)

        # Validate encoding. Prefer preserving original encoding.
        # If the original encoding cannot represent the edited content, try gb18030 first
        # (compatible superset for common Chinese encodings). Avoid silently transcoding
        # to UTF-8 here, because it may lead to "乱码" when toolchains expect ANSI/GBK.
        enc_used = target_enc
        transcoded = False
        try:
            _ = txt.encode(enc_used)
        except UnicodeEncodeError as e:
            try:
                # Prefer gb18030 as a safer fallback within GB-family.
                _ = txt.encode('gb18030')
                logger.warning(f"内容无法用 {enc_used} 编码，将回退 gb18030 保存: {filepath} - {e}")
                transcoded = True
                enc_used = 'gb18030'
            except UnicodeEncodeError as e2:
                logger.error(f"内容包含无法用 {enc_used} 编码的字符: {filepath} - {e}")
                return False, f"内容包含无法用 {enc_used} 编码的字符: {e}"

        logger.debug(f"准备写入文件: {full}, 编码: {enc_used}, 换行符: {repr(target_eol)}, 内容长度: {len(txt)}")

        tmp_path = None
        try:
            base_dir = os.path.dirname(full)
            tmp_name = f".{os.path.basename(full)}.tmp.{os.getpid()}.{int(time.time() * 1000)}"
            tmp_path = os.path.join(base_dir, tmp_name)

            data_bytes = txt.encode(enc_used)
            with open(tmp_path, "wb") as f:
                f.write(data_bytes)
                try:
                    f.flush()
                except Exception as e:
                    logger.debug(f"Exception ignored: {e}")
                try:
                    os.fsync(f.fileno())
                except Exception as e:
                    logger.debug(f"Exception ignored: {e}")

            os.replace(tmp_path, full)
            tmp_path = None
        finally:
            if tmp_path:
                try:
                    if os.path.exists(tmp_path):
                        os.remove(tmp_path)
                except Exception as e:
                    logger.debug(f"Exception ignored: {e}")
        
        if transcoded and (enc_used != target_enc):
            logger.info(f"✓ 文件保存成功: {filepath} (编码: {target_enc} -> {enc_used}, 换行符: {repr(target_eol)}, {len(txt)}字符)")
            return True, f"保存成功（已从 {target_enc} 转为 {enc_used}）"
        logger.info(f"✓ 文件保存成功: {filepath} (编码: {enc_used}, 换行符: {repr(target_eol)}, {len(txt)}字符)")
        return True, "保存成功"
    except Exception as e:
        logger.error(f"保存文件失败: {filepath} - {e}", exc_info=True)
        return False, str(e)


# ════════════════════════════════════════════════════════
#  工作区变更文件
# ════════════════════════════════════════════════════════

def get_changed_files():
    """获取工作区变更文件列表"""
    logger.debug("开始获取变更文件列表")
    # -z 用 NUL 分隔,彻底避免中文/空格路径问题
    stdout, _, code = run_git(["status", "--porcelain=v1", "-u", "-z"])
    if code != 0:
        logger.warning("获取变更文件列表失败")
        return []

    def parse_numstat_z(out_text):
        m = {}
        if not out_text:
            return m
        parts = out_text.split("\x00")
        for item in parts:
            if not item:
                continue
            cols = item.split("\t")
            if len(cols) < 3:
                continue
            a, r, p = cols[0], cols[1], cols[2]
            try:
                added = int(a) if a != "-" else 0
                removed = int(r) if r != "-" else 0
            except Exception:
                added = removed = 0
            p = (p or "").replace("\\", "/")
            m[p] = (added, removed)
        return m

    # 批量统计增删行：最多执行 3 次 diff（而不是每个文件 3 次）
    ns_out, _, _ = run_git(["diff", "HEAD",     "--numstat", "-z"])
    ns_map = parse_numstat_z(ns_out)
    if not ns_map:
        ns_out2, _, _ = run_git(["diff",             "--numstat", "-z"])
        ns_map = parse_numstat_z(ns_out2)
    ns_cached_out, _, _ = run_git(["diff", "--cached", "--numstat", "-z"])
    ns_cached_map = parse_numstat_z(ns_cached_out)

    entries = stdout.split("\x00")
    files   = []
    i       = 0
    while i < len(entries):
        entry = entries[i]
        if len(entry) < 4:
            i += 1
            continue
        xy       = entry[:2]
        name     = entry[3:]
        old_name = None
        idx_s    = xy[0]
        work_s   = xy[1]

        if idx_s in ("R", "C"):
            i += 1
            old_name = entries[i] if i < len(entries) else ""

        if idx_s == "?" and work_s == "?":
            status = "U"
        elif idx_s == "D" or work_s == "D":
            status = "D"
        elif idx_s == "A":
            status = "A"
        elif idx_s in ("R", "C"):
            status = "R"
        else:
            status = "M"

        p_key = (name or "").replace("\\", "/")
        added, removed = (0, 0)
        if status == "U":
            try:
                txt_u = get_file_content(p_key) or ""
                added = len(txt_u.splitlines())
                removed = 0
            except Exception as e:
                logger.debug(f"读取新文件行数失败: {p_key} - {e}")
                added = removed = 0
        else:
            if p_key in ns_map:
                added, removed = ns_map.get(p_key, (0, 0))
            elif p_key in ns_cached_map:
                added, removed = ns_cached_map.get(p_key, (0, 0))
        files.append({
            "path":     name,
            "status":   status,
            "old_path": old_name,
            "added":    added,
            "removed":  removed,
        })
        i += 1
    
    logger.debug(f"成功获取变更文件列表，共 {len(files)} 个文件")
    return files

# ════════════════════════════════════════════════════════
#  Diff 解析
# ════════════════════════════════════════════════════════

def get_file_diff(filepath, status, ctx_lines=5):
    """获取文件的 diff 内容"""
    try:
        ctx_lines = int(ctx_lines)
    except Exception:
        ctx_lines = 5
    if ctx_lines < 0:
        ctx_lines = 0
    if ctx_lines > 200:
        ctx_lines = 200
    logger.debug(f"获取文件 diff: {filepath} (状态: {status}, ctx: {ctx_lines})")
    if status == "U":
        try:
            txt_u = get_file_content(filepath) or ""
            content = txt_u.splitlines(keepends=True)
            # Empty new file: still render one editable empty line so user can type.
            if not content:
                content = ["\n"]
            # Preserve trailing empty line for files ending with EOL.
            # Python splitlines(keepends=True) does not include the final empty line.
            if txt_u.endswith("\n") or txt_u.endswith("\r"):
                content.append("\n")
            logger.debug(f"读取新文件内容: {filepath} - {len(content)} 行")
            return [{
                "header":    f"@@ -0,0 +1,{len(content)} @@ 新文件",
                "old_start": 0,
                "new_start": 1,
                "lines": [{"type":"added","old":None,"new":i+1,"text":l.rstrip("\n").rstrip("\r")}
                          for i, l in enumerate(content)]
            }], None
        except Exception as e:
            logger.error(f"读取新文件 diff 失败: {filepath} - {e}")
            return [], str(e)

    raw = None
    unified_arg = f"--unified={ctx_lines}"
    for args in (
        ["diff", "HEAD",     unified_arg, "--", filepath],
        ["diff",             unified_arg, "--", filepath],
        ["diff", "--cached", unified_arg, "--", filepath],
    ):
        out, _, code = run_git(args)
        if code == 0 and (out or "").strip():
            raw = out
            break
    if not raw:
        logger.debug(f"未找到文件 diff: {filepath}")
        return [], None
    
    logger.debug(f"成功获取文件 diff: {filepath}")
    return parse_diff(raw), None


def parse_diff(text):
    """解析 diff 文本"""
    logger.debug("开始解析 diff 文本")
    hunks = []
    cur   = None
    ol = nl = 0
    removed_block = []
    added_block = []

    def _norm_line(s):
        if s is None:
            return ""
        return s[:-1] if s.endswith("\r") else s

    def _lcs_pairs(a, b):
        n = len(a)
        m = len(b)
        if n == 0 or m == 0:
            return []
        dp = [[0] * (m + 1) for _ in range(n + 1)]
        for i in range(n - 1, -1, -1):
            ai = a[i]
            row = dp[i]
            row_next = dp[i + 1]
            for j in range(m - 1, -1, -1):
                if ai == b[j]:
                    row[j] = row_next[j + 1] + 1
                else:
                    v1 = row_next[j]
                    v2 = row[j + 1]
                    row[j] = v1 if v1 >= v2 else v2

        pairs = []
        i = j = 0
        while i < n and j < m:
            if a[i] == b[j]:
                pairs.append((i, j))
                i += 1
                j += 1
            else:
                if dp[i + 1][j] >= dp[i][j + 1]:
                    i += 1
                else:
                    j += 1
        return pairs

    def _line_similarity(s1, s2):
        """计算两行文本的相似度 (0.0 - 1.0)"""
        if s1 == s2:
            return 1.0
        if not s1 or not s2:
            return 0.0

        s1_stripped = s1.strip()
        s2_stripped = s2.strip()
        if not s1_stripped or not s2_stripped:
            return 0.0

        def _collapse_ws(s: str):
            try:
                return " ".join((s or "").split())
            except Exception:
                return (s or "").strip()

        a = _collapse_ws(s1_stripped)
        b = _collapse_ws(s2_stripped)
        if not a or not b:
            return 0.0

        len1, len2 = len(a), len(b)
        max_len = max(len1, len2)
        if max_len > 0:
            len_ratio = min(len1, len2) / max_len
            if len_ratio < 0.15:
                return 0.0

        try:
            if a.startswith(b) or b.startswith(a):
                return max(len_ratio, 0.6)
        except Exception as e:
            logger.debug(f"Exception ignored: {e}")

        try:
            seq = difflib.SequenceMatcher(None, a, b, autojunk=False).ratio()
        except Exception:
            seq = 0.0

        # Token similarity: more tolerant for identifier/number tweaks
        try:
            ta = re.findall(r"[A-Za-z_][A-Za-z_0-9]*|\d+|\S", a)
            tb = re.findall(r"[A-Za-z_][A-Za-z_0-9]*|\d+|\S", b)
            if not ta or not tb:
                return seq
            sa = set(ta)
            sb = set(tb)
            inter = len(sa & sb)
            if inter <= 0:
                return seq
            token_sim = (2.0 * inter) / (len(sa) + len(sb))
            return max(seq, token_sim)
        except Exception:
            return seq

    def flush_change_block():
        nonlocal removed_block, added_block, ol, nl, cur
        if not cur:
            removed_block = []
            added_block = []
            return
        if not removed_block and not added_block:
            return

        pairs = _lcs_pairs(removed_block, added_block)
        ri = ai = 0

        def emit_unmatched(r_end, a_end):
            nonlocal ri, ai, ol, nl
            r_len = r_end - ri
            a_len = a_end - ai
            
            # 纯新增场景
            if r_len == 0:
                for k in range(a_len):
                    cur["lines"].append({"type": "added", "old": None, "new": nl, "text": added_block[ai + k]})
                    nl += 1
                ri = r_end
                ai = a_end
                return
            
            # 纯删除场景
            if a_len == 0:
                for k in range(r_len):
                    cur["lines"].append({"type": "removed", "old": ol, "new": None, "text": removed_block[ri + k]})
                    ol += 1
                ri = r_end
                ai = a_end
                return
            
            # 同时存在删除和新增: 使用智能匹配
            # 构建相似度矩阵，找出最佳匹配
            removed_lines = removed_block[ri:r_end]
            added_lines = added_block[ai:a_end]

            sm = difflib.SequenceMatcher(a=removed_lines, b=added_lines, autojunk=False)
            opcodes = list(sm.get_opcodes())
            op_i = 0
            while op_i < len(opcodes):
                tag, i1, i2, j1, j2 = opcodes[op_i]
                # Merge adjacent delete+insert into replace (common for single-line modifications)
                if tag == "delete" and (op_i + 1) < len(opcodes):
                    t2, ii1, ii2, jj1, jj2 = opcodes[op_i + 1]
                    if t2 == "insert" and ii1 == i2 and jj1 == j1:
                        tag = "replace"
                        j2 = jj2
                        op_i += 1

                if tag == "equal":
                    for k in range(i2 - i1):
                        cur["lines"].append({"type": "context", "old": ol, "new": nl, "text": removed_lines[i1 + k]})
                        ol += 1
                        nl += 1

                elif tag == "delete":
                    for k in range(i1, i2):
                        cur["lines"].append({"type": "removed", "old": ol, "new": None, "text": removed_lines[k]})
                        ol += 1

                elif tag == "insert":
                    for k in range(j1, j2):
                        cur["lines"].append({"type": "added", "old": None, "new": nl, "text": added_lines[k]})
                        nl += 1

                else:  # replace
                    r_seg = removed_lines[i1:i2]
                    a_seg = added_lines[j1:j2]

                    if not r_seg:
                        for txt in a_seg:
                            cur["lines"].append({"type": "added", "old": None, "new": nl, "text": txt})
                            nl += 1
                        op_i += 1
                        continue
                    if not a_seg:
                        for txt in r_seg:
                            cur["lines"].append({"type": "removed", "old": ol, "new": None, "text": txt})
                            ol += 1
                        op_i += 1
                        continue


                    pairs = []
                    for rr, r_txt in enumerate(r_seg):
                        for aa, a_txt in enumerate(a_seg):
                            sim = _line_similarity(r_txt, a_txt)
                            if sim >= 0.20:
                                pairs.append((sim, rr, aa))

                    pairs.sort(reverse=True)
                    used_r = set()
                    used_a = set()
                    chosen = []
                    for sim, rr, aa in pairs:
                        if rr in used_r or aa in used_a:
                            continue
                        used_r.add(rr)
                        used_a.add(aa)
                        chosen.append((rr, aa))

                    chosen.sort()
                    r_i = 0
                    a_i = 0
                    c_i = 0
                    while r_i < len(r_seg) or a_i < len(a_seg):
                        if c_i < len(chosen):
                            mr, ma = chosen[c_i]
                            if r_i == mr and a_i == ma:
                                cur["lines"].append({
                                    "type": "modified",
                                    "old": ol,
                                    "new": nl,
                                    "text": a_seg[a_i],
                                    "old_text": r_seg[r_i]
                                })
                                ol += 1
                                nl += 1
                                r_i += 1
                                a_i += 1
                                c_i += 1
                                continue
                            if r_i < mr:
                                cur["lines"].append({"type": "removed", "old": ol, "new": None, "text": r_seg[r_i]})
                                ol += 1
                                r_i += 1
                                continue
                            if a_i < ma:
                                cur["lines"].append({"type": "added", "old": None, "new": nl, "text": a_seg[a_i]})
                                nl += 1
                                a_i += 1
                                continue

                        if r_i < len(r_seg):
                            cur["lines"].append({"type": "removed", "old": ol, "new": None, "text": r_seg[r_i]})
                            ol += 1
                            r_i += 1
                        elif a_i < len(a_seg):
                            cur["lines"].append({"type": "added", "old": None, "new": nl, "text": a_seg[a_i]})
                            nl += 1
                            a_i += 1

                op_i += 1
            
            ri = r_end
            ai = a_end

        for rj, aj in pairs:
            if rj > ri or aj > ai:
                emit_unmatched(rj, aj)

            cur["lines"].append({"type": "context", "old": ol, "new": nl, "text": removed_block[rj]})
            ol += 1
            nl += 1
            ri = rj + 1
            ai = aj + 1

        if ri < len(removed_block) or ai < len(added_block):
            emit_unmatched(len(removed_block), len(added_block))

        removed_block = []
        added_block = []

    
    for raw_line in (text or "").splitlines():
        if raw_line.startswith("\\ No newline at end of file"):
            continue
        if raw_line.startswith("+++") or raw_line.startswith("---"):
            continue
            
        m = re.match(r'^@@ -(\d+)(?:,\d+)? \+(\d+)(?:,\d+)? @@(.*)', raw_line)
        if m:
            if cur:
                flush_change_block()
                hunks.append(cur)
            ol  = int(m.group(1))
            nl  = int(m.group(2))
            ctx = m.group(3).strip()
            cur = {
                "header":    raw_line,
                "old_start": ol,
                "new_start": nl,
                "context":   ctx,
                "lines": []
            }
            continue

        if not cur:
            continue

        if raw_line.startswith("+"):
            added_block.append(_norm_line(raw_line[1:]))

        elif raw_line.startswith("-"):
            removed_block.append(_norm_line(raw_line[1:]))

        else:
            flush_change_block()
            text = raw_line[1:] if raw_line.startswith(" ") else raw_line
            text = _norm_line(text)
            cur["lines"].append({"type":"context", "old":ol, "new":nl, "text":text})
            ol += 1
            nl += 1

    if cur:
        flush_change_block()
        hunks.append(cur)
    
    logger.debug(f"解析完成，共 {len(hunks)} 个 hunk")
    return hunks


def _get_raw_file_diff_patch(filepath, status, ctx_lines=5):
    try:
        ctx_lines = int(ctx_lines)
    except Exception:
        ctx_lines = 5
    if ctx_lines < 0:
        ctx_lines = 0
    if ctx_lines > 200:
        ctx_lines = 200

    filepath = (filepath or "").replace("\\", "/").lstrip("/")
    if not filepath:
        return None, "缺少 path"

    st = (status or "").strip().upper() or "M"
    # U 类型文件在提交时会通过 git add -N 处理，可以支持按行/按块操作
    # 只有在撤回操作时才真正不支持

    unified_arg = f"--unified={ctx_lines}"
    for args in (
        ["diff", "HEAD", unified_arg, "--", filepath],
        ["diff", unified_arg, "--", filepath],
        ["diff", "--cached", unified_arg, "--", filepath],
    ):
        out, err, code = run_git(args)
        if code == 0 and (out or "").strip():
            return out, None

    return "", None


def get_raw_file_diff_patch(filepath, status, ctx_lines=5):
    """Compatibility wrapper: return unified diff patch text for file."""
    return _get_raw_file_diff_patch(filepath, status, ctx_lines=ctx_lines)


def _parse_unified_patch_with_mapping(patch_text: str):
    lines = (patch_text or "").splitlines(True)  # keep line endings
    file_header = []
    hunks = []

    cur = None
    removed_buf = []  # [raw_idx]

    def flush_removed():
        nonlocal removed_buf, cur
        while removed_buf:
            raw_idx = removed_buf.pop(0)
            pi = cur["parsed_idx"]
            cur["map"][pi] = [raw_idx]
            cur["parsed_idx"] += 1

    for ln in lines:
        m = re.match(r'^@@ -(\d+)(?:,(\d+))? \+(\d+)(?:,(\d+))? @@', ln)
        if m:
            if cur:
                flush_removed()
                hunks.append(cur)
            cur = {
                "header": ln.rstrip("\n"),
                "old_start": int(m.group(1)),
                "new_start": int(m.group(3)),
                "raw_lines": [],
                "map": {},
                "parsed_idx": 0,
            }
            continue

        if not cur:
            file_header.append(ln)
            continue

        cur["raw_lines"].append(ln)
        raw_idx = len(cur["raw_lines"]) - 1

        if ln.startswith("+"):
            if removed_buf:
                old_raw_idx = removed_buf.pop(0)
                pi = cur["parsed_idx"]
                cur["map"][pi] = [old_raw_idx, raw_idx]
                cur["parsed_idx"] += 1
            else:
                pi = cur["parsed_idx"]
                cur["map"][pi] = [raw_idx]
                cur["parsed_idx"] += 1
        elif ln.startswith("-"):
            removed_buf.append(raw_idx)
        else:
            flush_removed()
            pi = cur["parsed_idx"]
            cur["map"][pi] = [raw_idx]
            cur["parsed_idx"] += 1

    if cur:
        flush_removed()
        hunks.append(cur)

    for h in hunks:
        h.pop("parsed_idx", None)
    return file_header, hunks


def _build_partial_hunk_patch(file_header_lines, hunk, include_raw_idx_set):
    raw_lines = hunk.get("raw_lines") or []
    include = set(int(x) for x in (include_raw_idx_set or set()) if x is not None)
    if not raw_lines or not include:
        return None

    old_ln = int(hunk.get("old_start") or 0)
    new_ln = int(hunk.get("new_start") or 0)

    hunks_out = []
    cur_start_old = None
    cur_start_new = None
    cur_old_len = 0
    cur_new_len = 0
    cur_lines = []

    def flush_current():
        nonlocal cur_start_old, cur_start_new, cur_old_len, cur_new_len, cur_lines
        if cur_start_old is None or not cur_lines:
            cur_start_old = None
            cur_start_new = None
            cur_old_len = 0
            cur_new_len = 0
            cur_lines = []
            return
        hdr = f"@@ -{cur_start_old},{cur_old_len} +{cur_start_new},{cur_new_len} @@\n"
        hunks_out.append(hdr + "".join(cur_lines))
        cur_start_old = None
        cur_start_new = None
        cur_old_len = 0
        cur_new_len = 0
        cur_lines = []

    for i, ln in enumerate(raw_lines):
        inc = (i in include)

        # If we are in a running hunk but current line is excluded, close the hunk.
        if (not inc) and cur_start_old is not None:
            flush_current()

        if inc:
            if cur_start_old is None:
                cur_start_old = old_ln
                cur_start_new = new_ln
            cur_lines.append(ln)
            if ln.startswith(" "):
                cur_old_len += 1
                cur_new_len += 1
            elif ln.startswith("-"):
                cur_old_len += 1
            elif ln.startswith("+"):
                cur_new_len += 1

        # advance counters regardless of inclusion
        if ln.startswith(" "):
            old_ln += 1
            new_ln += 1
        elif ln.startswith("-"):
            old_ln += 1
        elif ln.startswith("+"):
            new_ln += 1

    if cur_start_old is not None:
        flush_current()

    if not hunks_out:
        return None

    patch = "".join(file_header_lines) + "".join(hunks_out)
    if not patch.endswith("\n"):
        patch += "\n"
    return patch


def extract_selected_hunks_from_patch(patch_text: str, hunk_indices):
    """Extract selected hunks (by index) from a unified diff patch."""
    file_header, hunks = _parse_unified_patch_with_mapping(patch_text or "")
    if not hunks:
        return ""

    picked = []
    for hi in (hunk_indices or []):
        try:
            picked.append(int(hi))
        except Exception:
            continue

    picked = sorted(set(picked))

    hunks_out = []
    for hi in picked:
        if hi < 0 or hi >= len(hunks):
            continue
        h = hunks[hi] or {}
        include = set(range(len(h.get("raw_lines") or [])))
        p = _build_partial_hunk_patch([], h, include)
        if p:
            hunks_out.append(p)

    if not hunks_out:
        return ""
    return "".join(file_header) + "".join(hunks_out)


def extract_selected_lines_from_patch(patch_text: str, line_keys):
    """Extract selected changed lines from a unified diff patch.

    line_keys: list like ["hunkIndex:lineIndex", ...] where lineIndex matches the
    parsed line index within that hunk (aligned with unified patch parsing order).
    """
    file_header, hunks = _parse_unified_patch_with_mapping(patch_text or "")
    if not hunks:
        return ""

    include_map = {}  # hunk_idx -> set(raw_idx)
    for k in (line_keys or []):
        if k is None:
            continue
        s = str(k)
        if ":" not in s:
            continue
        a, b = s.split(":", 1)
        try:
            hi = int(a)
            li = int(b)
        except Exception:
            continue
        if hi < 0 or hi >= len(hunks):
            continue
        h = hunks[hi] or {}
        m = h.get("map") or {}
        raw_idxs = m.get(li)
        if not raw_idxs:
            continue
        inc = include_map.get(hi)
        if inc is None:
            inc = set()
            include_map[hi] = inc
        for ridx in raw_idxs:
            try:
                inc.add(int(ridx))
            except Exception as e:
                logger.debug(f"Exception ignored: {e}")

    hunks_out = []
    for hi, inc in include_map.items():
        h = hunks[hi] or {}
        p = _build_partial_hunk_patch([], h, inc)
        if p:
            hunks_out.append(p)

    if not hunks_out:
        return ""
    return "".join(file_header) + "".join(hunks_out)


def apply_patch_to_index(patch_text: str):
    """Apply a unified diff patch to Git index (staging area) only."""
    if not patch_text or not (patch_text or "").strip():
        return False, "空 patch"

    attempts = [
        ["apply", "--cached", "--whitespace=nowarn", "--recount"],
        ["apply", "--cached", "--whitespace=nowarn", "--recount", "-C0"],
        ["apply", "--cached", "--whitespace=nowarn", "--recount", "--unidiff-zero"],
        ["apply", "--cached", "--whitespace=nowarn", "--recount", "--unidiff-zero", "-C0"],
    ]

    last_msg = ""
    for args in attempts:
        out, err, code = run_git(args, input_data=patch_text, timeout=120)
        if code == 0:
            return True, ""
        last_msg = (err or out or "git apply --cached 失败")

    return False, last_msg


def restore_file_from_commit(commit: str, filepath: str):
    if not REPO_PATH:
        return False, "未打开仓库"

    commit = (commit or "").strip()
    filepath = (filepath or "").replace("\\", "/").lstrip("/")
    if not commit:
        return False, "缺少 hash"
    if not filepath:
        return False, "缺少 path"

    full = _safe_repo_abspath(filepath)
    if not full:
        return False, "非法路径"

    _, err, code = run_git(["checkout", commit, "--", filepath], timeout=120)
    if code != 0:
        return False, (err or "恢复文件失败")
    return True, ""


def restore_workspace_to_commit(commit: str, force=False):
    if not REPO_PATH:
        return False, "未打开仓库"

    commit = (commit or "").strip()
    if not commit:
        return False, "缺少 hash"

    if not force:
        out, err, code = run_git(["status", "--porcelain=v1"])
        if code != 0:
            return False, (err or "无法检测工作区状态")
        if (out or "").strip():
            return False, "工作区有未提交更改，请先提交/撤回/暂存（stash）后再执行还原"

    _, err, code = run_git(["reset", "--hard", commit], timeout=120)
    if code != 0:
        return False, (err or "还原工作区失败")
    run_git(["clean", "-fd"], timeout=120)
    return True, ""


def revert_commit(commit: str):
    if not REPO_PATH:
        return False, "未打开仓库", ""

    commit = (commit or "").strip()
    if not commit:
        return False, "缺少 hash", ""

    # 要求工作区干净，避免 Revert 过程中出现混乱
    out, err, code = run_git(["status", "--porcelain=v1"])
    if code != 0:
        return False, (err or "无法检测工作区状态"), ""
    if (out or "").strip():
        return False, "工作区有未提交更改，请先提交/撤回/暂存（stash）后再执行 Revert", ""

    # 解析为完整提交哈希
    full_out, full_err, full_code = run_git(["rev-parse", commit])
    if full_code != 0:
        return False, (full_err or "无法解析提交哈希"), ""
    full = (full_out or "").strip()

    # 执行 Revert，自动使用提交消息
    _, rerr, rcode = run_git(["revert", "--no-edit", full], timeout=600)
    if rcode != 0:
        # 尝试中止以恢复现场
        run_git(["revert", "--abort"], timeout=120)
        return False, (rerr or "Revert 失败"), ""

    # 返回新的 HEAD
    head_out, head_err, head_code = run_git(["rev-parse", "HEAD"]) 
    new_head = (head_out or "").strip() if head_code == 0 else ""
    return True, "", new_head


def _git_apply_reverse_patch(patch_text: str):
    if not patch_text:
        return False, "空 patch"

    attempts = [
        ["apply", "-R", "--whitespace=nowarn", "--recount"],
        ["apply", "-R", "--whitespace=nowarn", "--recount", "-C0"],
        ["apply", "-R", "--whitespace=nowarn", "--recount", "--unidiff-zero"],
        ["apply", "-R", "--whitespace=nowarn", "--recount", "--unidiff-zero", "-C0"],
        ["apply", "-R", "--whitespace=nowarn", "--recount", "--3way"],
        ["apply", "-R", "--whitespace=nowarn", "--recount", "--3way", "-C0"],
        ["apply", "-R", "--whitespace=nowarn", "--recount", "--3way", "--unidiff-zero"],
        ["apply", "-R", "--whitespace=nowarn", "--recount", "--ignore-space-change"],
        ["apply", "-R", "--whitespace=nowarn", "--recount", "--ignore-whitespace"],
    ]

    last_msg = ""
    for args in attempts:
        out, err, code = run_git(args, input_data=patch_text, timeout=120)
        if code == 0:
            return True, ""
        last_msg = (err or out or "git apply 失败")

    return False, last_msg


def revert_file(filepath: str, status: str):
    if not REPO_PATH:
        return False, "未打开仓库"
    filepath = (filepath or "").replace("\\", "/").lstrip("/")
    if not filepath:
        return False, "缺少 path"

    full = _safe_repo_abspath(filepath)
    if not full:
        return False, "非法路径"

    # Unified undo semantics (A+B+C):
    # - If file exists in HEAD, restore it from HEAD (works for modified/deleted/staged-deleted).
    # - Otherwise, treat as untracked: remove it if present.
    _, _, code_head = run_git(["cat-file", "-e", f"HEAD:{filepath}"], timeout=30)
    if code_head == 0:
        _, err, code = run_git(["checkout", "HEAD", "--", filepath], timeout=120)
        if code != 0:
            return False, (err or "撤回文件失败")
        return True, ""

    try:
        if os.path.exists(full) and (not os.path.isdir(full)):
            os.remove(full)
        return True, ""
    except Exception as e:
        return False, str(e)


def rename_file(old_path: str, new_path: str):
    if not REPO_PATH:
        return False, "未打开仓库"

    old_rel = (old_path or "").replace("\\", "/").lstrip("/")
    new_rel = (new_path or "").replace("\\", "/").lstrip("/")
    if not old_rel or not new_rel:
        return False, "缺少 path"
    if old_rel == new_rel:
        return False, "新旧路径相同"

    old_abs = _safe_repo_abspath(old_rel)
    new_abs = _safe_repo_abspath(new_rel)
    if not old_abs or not new_abs:
        return False, "非法路径"
    if not os.path.exists(old_abs) or os.path.isdir(old_abs):
        return False, "原文件不存在"
    if os.path.exists(new_abs):
        return False, "目标路径已存在"

    parent = os.path.dirname(new_abs)
    if parent and (not os.path.exists(parent)):
        os.makedirs(parent, exist_ok=True)

    # Prefer git mv for tracked files
    out_t, err_t, code_t = run_git(["ls-files", "--error-unmatch", "--", old_rel], timeout=30)
    if code_t == 0:
        out, err, code = run_git(["mv", "--", old_rel, new_rel], timeout=120)
        if code != 0:
            return False, (err or out or "git mv 失败")
        return True, ""

    # Untracked file: filesystem rename
    try:
        os.replace(old_abs, new_abs)
        return True, ""
    except Exception as e:
        return False, str(e)


def mkdir_repo(rel_path: str):
    if not REPO_PATH:
        return False, "未打开仓库"
    rp = str(rel_path or "").strip().replace("\\", "/").lstrip("/")
    rp = rp.rstrip("/")
    if not rp:
        return False, "缺少 path"
    full = _safe_repo_abspath(rp)
    if not full:
        return False, "非法路径"
    try:
        Path(full).mkdir(parents=True, exist_ok=False)
        return True, ""
    except FileExistsError:
        return False, "目录已存在"
    except Exception as e:
        return False, str(e)


def revert_hunk(filepath: str, hunk_idx: int, status: str, ctx_lines: int = 5):
    raw_patch, err = _get_raw_file_diff_patch(filepath, status, ctx_lines=ctx_lines)
    if err:
        return False, err
    if raw_patch is None or (not raw_patch.strip()):
        return False, "无可撤回变更"

    file_header, hunks = _parse_unified_patch_with_mapping(raw_patch)
    if hunk_idx < 0 or hunk_idx >= len(hunks):
        return False, "hunk_index 越界"

    h = hunks[hunk_idx]
    include = set(range(len(h.get("raw_lines") or [])))
    patch = _build_partial_hunk_patch(file_header, h, include)
    if not patch:
        return False, "构建 patch 失败"
    ok, msg = _git_apply_reverse_patch(patch)
    if ok:
        return True, ""
    # Retry with zero-context diff to reduce context mismatches
    raw_patch2, err2 = _get_raw_file_diff_patch(filepath, status, ctx_lines=0)
    if err2 or (not raw_patch2) or (not raw_patch2.strip()):
        return False, msg
    file_header2, hunks2 = _parse_unified_patch_with_mapping(raw_patch2)
    if hunk_idx < 0 or hunk_idx >= len(hunks2):
        return False, msg
    h2 = hunks2[hunk_idx]
    include2 = set(range(len(h2.get("raw_lines") or [])))
    patch2 = _build_partial_hunk_patch(file_header2, h2, include2)
    if not patch2:
        return False, msg
    return _git_apply_reverse_patch(patch2)


def revert_line(filepath: str, hunk_idx: int, line_idx: int, status: str, ctx_lines: int = 5, signature: dict | None = None):
    if not REPO_PATH:
        return False, "未打开仓库"

    filepath = (filepath or "").replace("\\", "/").lstrip("/")
    if not filepath:
        return False, "缺少 path"

    full = _safe_repo_abspath(filepath)
    if not full:
        return False, "非法路径"

    st = (status or "").strip().upper() or "M"
    if st == "U":
        return False, "未跟踪文件不支持按行撤回"

    try:
        hunk_idx = int(hunk_idx)
        line_idx = int(line_idx)
    except Exception:
        return False, "参数非法"

    hunks, err = get_file_diff(filepath, st, ctx_lines)
    if err:
        return False, err
    if not hunks:
        return False, "无可撤回变更"
    if hunk_idx < 0 or hunk_idx >= len(hunks):
        return False, "hunk_index 越界"

    h = hunks[hunk_idx] or {}
    lines = h.get("lines") or []

    def _to_int(x):
        try:
            return int(x)
        except Exception:
            return None

    def _sig_match(dl0: dict, sig0: dict) -> bool:
        if not dl0 or not sig0:
            return False
        t0 = (dl0.get("type") or "").lower()
        ts = (sig0.get("type") or "").lower()
        if t0 != ts:
            return False
        if (dl0.get("text") or "") != (sig0.get("text") or ""):
            return False
        if (dl0.get("old_text") or "") != (sig0.get("old_text") or ""):
            return False
        n0 = _to_int(dl0.get("new"))
        ns = _to_int(sig0.get("new"))
        o0 = _to_int(dl0.get("old"))
        os = _to_int(sig0.get("old"))
        if ns is not None and n0 != ns:
            return False
        if os is not None and o0 != os:
            return False
        return True

    if signature:
        found_idx = None
        for i, dli in enumerate(lines):
            try:
                if _sig_match(dli or {}, signature):
                    found_idx = i
                    break
            except Exception as e:
                logger.debug(f"Exception ignored: {e}")
        if found_idx is not None:
            line_idx = found_idx

    if line_idx < 0 or line_idx >= len(lines):
        return False, "line_index 越界"
    dl = lines[line_idx] or {}

    full = _safe_repo_abspath(filepath)
    if not full:
        return False, "非法路径"

    # Read file with encoding detection and preserve original encoding.
    # If the file was decoded with replacement (lossy), forbids revert to avoid corrupting bytes.
    try:
        with _file_encoding_lock:
            lossy = _file_decode_lossy.get(filepath)
            enc_hint = _file_last_encoding.get(filepath) or "unknown"
        if lossy:
            return False, f"该文件读取时发生编码替换（{enc_hint} + replace），无法安全执行撤回写入。请先用正确编码修复文件后再撤回。"
    except Exception as e:
        logger.debug(f"Exception ignored: {e}")

    try:
        content, enc = get_file_content(filepath, return_encoding=True)
        if content is None:
            return False, "无法读取文件内容"
        enc_used = enc or "utf-8"
    except Exception as e:
        return False, str(e)

    # Keep EOL style based on original content
    normalized = str(content).replace("\r\n", "\n").replace("\r", "\n")
    eol = "\r\n" if "\r\n" in str(content) else "\n"
    cur_lines = normalized.split("\n")

    def _ensure_eol(s: str) -> str:
        if s is None:
            s = ""
        s = str(s)
        return s

    t = (dl.get("type") or "").lower()
    if t == "context":
        return False, "无法撤回上下文行"

    if t in ("added", "modified"):
        new_no = dl.get("new")
        if new_no is None:
            return False, "缺少 new 行号"
        try:
            new_no = int(new_no)
        except Exception:
            return False, "new 行号非法"
        idx0 = new_no - 1
        if idx0 < 0 or idx0 >= len(cur_lines):
            return False, "目标行号越界"

        if t == "added":
            cur_lines.pop(idx0)
        else:
            old_text = dl.get("old_text")
            if old_text is None:
                return False, "缺少 old_text"
            cur_lines[idx0] = str(old_text)

    elif t == "removed":
        ins_text = dl.get("text")
        if ins_text is None:
            return False, "缺少 text"
        ins_line = str(ins_text)

        insert_at = None
        for j in range(line_idx + 1, len(lines)):
            nxt = lines[j] or {}
            n_new = nxt.get("new")
            if n_new is not None:
                try:
                    insert_at = max(0, int(n_new) - 1)
                    break
                except Exception as e:
                    logger.debug(f"Exception ignored: {e}")
        if insert_at is None:
            for j in range(line_idx - 1, -1, -1):
                prv = lines[j] or {}
                p_new = prv.get("new")
                if p_new is not None:
                    try:
                        insert_at = max(0, int(p_new))
                        break
                    except Exception as e:
                        logger.debug(f"Exception ignored: {e}")
        if insert_at is None:
            insert_at = len(cur_lines)
        if insert_at > len(cur_lines):
            insert_at = len(cur_lines)
        cur_lines.insert(insert_at, ins_line)

    else:
        return False, "不支持的行类型"

    try:
        out_txt = "\n".join(cur_lines)
        out_txt = out_txt.replace("\n", eol)
        ok, msg = save_file_content(filepath, out_txt, force_encoding=enc_used)
        if not ok:
            return False, msg
    except Exception as e:
        return False, str(e)

    return True, ""


def revert_multi_lines(filepath: str, hunk_idx: int, start_line_idx: int, end_line_idx: int, status: str, ctx_lines: int = 5):
    if not REPO_PATH:
        return False, "未打开仓库"

    filepath = (filepath or "").replace("\\", "/").lstrip("/")
    if not filepath:
        return False, "缺少 path"

    st = (status or "").strip().upper() or "M"
    if st == "U":
        return False, "未跟踪文件不支持按行撤回"

    try:
        hunk_idx = int(hunk_idx)
        start_line_idx = int(start_line_idx)
        end_line_idx = int(end_line_idx)
    except Exception:
        return False, "参数非法"

    if start_line_idx > end_line_idx:
        start_line_idx, end_line_idx = end_line_idx, start_line_idx

    hunks, err = get_file_diff(filepath, st, ctx_lines)
    if err:
        return False, err
    if not hunks:
        return False, "无可撤回变更"
    if hunk_idx < 0 or hunk_idx >= len(hunks):
        return False, "hunk_index 越界"

    h = hunks[hunk_idx] or {}
    lines = h.get("lines") or []
    if not lines:
        return False, "无可撤回变更"

    start_line_idx = max(0, start_line_idx)
    end_line_idx = min(len(lines) - 1, end_line_idx)
    if start_line_idx > end_line_idx:
        return False, "line_index 越界"

    targets = []
    for li in range(start_line_idx, end_line_idx + 1):
        dl = lines[li] or {}
        t = (dl.get("type") or "").lower()
        if not t or t == "context":
            continue
        sig = {
            "type": dl.get("type"),
            "new": dl.get("new"),
            "old": dl.get("old"),
            "text": dl.get("text"),
            "old_text": dl.get("old_text"),
        }
        pos = dl.get("new") if t in ("added", "modified") else dl.get("old")
        try:
            pos = int(pos) if pos is not None else None
        except Exception:
            pos = None
        targets.append((li, pos, sig))

    if not targets:
        return False, "无可撤回变更"

    first_type = (targets[0][2].get("type") or "").lower()
    for _, __, sig in targets:
        if (sig.get("type") or "").lower() != first_type:
            return False, "多行撤回要求所有行类型相同"

    for li, pos, sig in sorted(targets, key=lambda x: (x[1] is None, x[1] if x[1] is not None else x[0]), reverse=True):
        ok, msg = revert_line(filepath, hunk_idx, li, st, ctx_lines, sig)
        if not ok:
            return False, msg

    return True, ""


def get_log(limit: int = 50):
    """获取提交历史"""
    logger.debug("获取提交历史")
    try:
        limit_i = int(limit)
    except Exception:
        limit_i = 50
    if limit_i < 1:
        limit_i = 1
    if limit_i > 200:
        limit_i = 200
    fmt = "--pretty=format:%H%x00%an%x00%ae%x00%ad%x00%s"
    out, _, code = run_git(["log", fmt, "--date=iso", f"-{limit_i}"])
    if code != 0:
        logger.warning("获取提交历史失败")
        return []
    
    commits = []
    for line in out.splitlines():
        parts = line.split("\x00")
        if len(parts) < 5:
            continue
        full_hash = parts[0]
        commits.append({
            "hash":      full_hash[:7],  # 短hash用于显示
            "full_hash": full_hash,      # 完整hash用于操作
            "author":    parts[1],
            "email":     parts[2],
            "date":      parts[3],
            "message":   parts[4]
        })
    
    logger.info(f"成功获取提交历史，共 {len(commits)} 条")
    return commits


def _ai_config_path():
    # Backward-compatible alias: we now use a single unified config file.
    return _hivo_cfg_path()


def _hivo_ai_pick_latest_history_path(base_dir: Path) -> Path | None:
    try:
        if not base_dir or not base_dir.exists() or not base_dir.is_dir():
            return None
        items = list(base_dir.glob("hivo_ai_chat_history_*.json"))
        if not items:
            return None
        items.sort(key=lambda p: p.stat().st_mtime if p.exists() else 0, reverse=True)
        return items[0]
    except Exception:
        return None


def _hivo_ai_history_path():
    repo = (REPO_PATH or "")
    repo_key = "global"
    try:
        if repo:
            repo_key = hashlib.sha1(repo.encode("utf-8", errors="ignore")).hexdigest()[:10]
    except Exception:
        repo_key = "global"
    base = _hivo_ai_data_dir()
    p0 = base / f"hivo_ai_chat_history_{repo_key}.json"
    if not repo:
        try:
            if not (p0.exists() and p0.is_file() and p0.stat().st_size > 0):
                p1 = _hivo_ai_pick_latest_history_path(base)
                if p1:
                    return p1
        except Exception:
            p1 = _hivo_ai_pick_latest_history_path(base)
            if p1:
                return p1
    return p0


_AI_GLOBAL_PROFILE_ID = "__global__"


def _hivo_ai_get_tool_call_id(tc: dict) -> str:
    """从 tool_call 对象中提取 id，兼容多种字段名（id/call_id/tool_call_id）。
    某些 OpenAI 兼容 API（如百度千帆）可能使用非标准的 id 字段名。
    """
    if not isinstance(tc, dict):
        return ""
    for key in ("id", "call_id", "tool_call_id"):
        val = tc.get(key)
        if val is not None:
            s = str(val).strip()
            if s:
                return s
    return ""


def _hivo_ai_clean_tool_calls(tc_list) -> list:
    """清洗原生 tool_calls 列表，标准化格式并补全缺失的 id。
    返回标准化后的 tool_calls 列表，每个元素包含 id/type/function{name,arguments}。
    """
    if not isinstance(tc_list, list) or not tc_list:
        return []
    cleaned = []
    for idx, tc in enumerate(tc_list):
        if not isinstance(tc, dict):
            continue
        tc_id = _hivo_ai_get_tool_call_id(tc)
        if not tc_id:
            tc_id = f"call_{idx}"
        tc_type = str(tc.get("type") or "function")
        tc_func = tc.get("function") if isinstance(tc.get("function"), dict) else {}
        tc_name = str(tc_func.get("name") or "")
        tc_args = str(tc_func.get("arguments") or "")
        cleaned.append({
            "id": tc_id,
            "type": tc_type,
            "function": {"name": tc_name, "arguments": tc_args},
        })
    return cleaned


def _hivo_ai_normalize_message(m: dict) -> dict | None:
    """将单条消息标准化为核心格式，保留 role/content/tool_calls/tool_call_id/name 等关键字段。
    返回标准化后的消息字典，若消息无效则返回 None。
    用于消除多处历史消息加载逻辑中的重复代码。
    """
    if not isinstance(m, dict):
        return None
    role = str(m.get("role") or "").strip()
    if role not in ("user", "assistant", "tool"):
        return None
    c0 = m.get("content")
    if isinstance(c0, list):
        content = c0
    else:
        content = str(c0 or "")
    # 无有效内容且非 assistant/tool 消息则返回 None
    if not _hivo_ai_msg_has_valid_content({"content": content}) and role not in ("assistant", "tool"):
        return None
    item = {"role": role, "content": content}
    if role == "assistant":
        try:
            tc_list = m.get("tool_calls")
            cleaned_tc = _hivo_ai_clean_tool_calls(tc_list)
            if cleaned_tc:
                item["tool_calls"] = cleaned_tc
        except Exception:
            pass
    elif role == "tool":
        try:
            tcid = str(m.get("tool_call_id") or "").strip()
            if tcid:
                item["tool_call_id"] = tcid
        except Exception:
            pass
    return item


def _hivo_ai_fix_tool_call_ids(messages: list):
    """修复 tool 消息的 tool_call_id，确保与前一条 assistant 消息的 tool_calls id 一一对应。
    以 assistant 消息中的 tool_calls id 为准，强制对齐，确保工具调用链路 ID 完全匹配。
    可以修复旧版历史数据中 tool 消息缺少 tool_call_id 或 id 不一致的问题。
    """
    if not isinstance(messages, list) or not messages:
        return messages
    try:
        fixed = []
        pending_tc = []
        pending_idx = 0
        for m in messages:
            if not isinstance(m, dict):
                continue
            role = str(m.get("role") or "")
            if role == "assistant":
                tc_list = m.get("tool_calls")
                if isinstance(tc_list, list) and tc_list:
                    pending_tc = [_hivo_ai_get_tool_call_id(tc) for tc in tc_list if isinstance(tc, dict)]
                    pending_idx = 0
                else:
                    pending_tc = []
                    pending_idx = 0
                fixed.append(m)
            elif role == "tool":
                if pending_tc and pending_idx < len(pending_tc):
                    expected_id = pending_tc[pending_idx]
                    if expected_id:
                        m["tool_call_id"] = expected_id
                    pending_idx += 1
                fixed.append(m)
            else:
                pending_tc = []
                pending_idx = 0
                fixed.append(m)
        return fixed
    except Exception:
        return messages


def _hivo_ai_clean_att_meta(meta_in, include_path: bool = False) -> list:
    """清洗 _attMeta 列表，标准化为 name/mime/size/url(可选 path) 格式。
    用于消除多处 _attMeta 处理的重复代码。
    """
    if not isinstance(meta_in, list) or not meta_in:
        return []
    meta_out = []
    for a in meta_in:
        if not isinstance(a, dict):
            continue
        o = {
            "name": str(a.get("name") or "file"),
            "mime": str(a.get("mime") or ""),
            "size": int(a.get("size") or 0),
        }
        u = str(a.get("url") or "").strip()
        if u:
            o["url"] = u
        if include_path:
            pth = str(a.get("path") or "").strip()
            if pth:
                o["path"] = pth.replace("\\", "/")
        meta_out.append(o)
    return meta_out


def _hivo_ai_clean_messages(messages: list, limit: int):
    if not isinstance(messages, list):
        return []
    cleaned = []
    for m in messages:
        if not isinstance(m, dict):
            continue
        role = str(m.get("role") or "").strip()
        if role not in ("user", "assistant", "tool"):
            continue
        # 处理 content：支持字符串和多模态数组两种格式
        raw_content = m.get("content")
        if isinstance(raw_content, list):
            if role == "tool":
                # tool 消息的多模态 content：剥离图片数据，只保留文本（避免历史记录臃肿）
                text_parts = []
                for p in raw_content:
                    if isinstance(p, dict):
                        pt = str(p.get("type") or "")
                        if pt == "text":
                            text_parts.append({"type": "text", "text": str(p.get("text") or "")})
                if len(text_parts) == 1:
                    content = text_parts[0]["text"]
                elif text_parts:
                    content = text_parts
                else:
                    content = ""
            else:
                content = raw_content
        else:
            content = str(raw_content or "")
        # 无有效内容且非 assistant/tool 消息则跳过
        if not _hivo_ai_msg_has_valid_content({"content": content}) and role not in ("assistant", "tool"):
            continue
        item = {"role": role, "content": content}
        if role == "user":
            try:
                mid = str(m.get("id") or "").strip()
                if mid:
                    item["id"] = mid
            except Exception:
                pass
            ug = str(m.get("undo_gid") or "").strip()
            if ug:
                item["undo_gid"] = ug
            # persist sanitized _attMeta for reliable client render after reload (preserve url if present)
            try:
                meta_out = _hivo_ai_clean_att_meta(m.get("_attMeta"))
                if meta_out:
                    item["_attMeta"] = meta_out
            except Exception:
                pass
        elif role == "assistant":
            try:
                mid = str(m.get("id") or "").strip()
                if mid:
                    item["id"] = mid
            except Exception:
                pass
            # 保留原生 tool_calls（含 id/type/function），用于会话恢复后继续多轮工具调用
            # 注意：必须完整保留所有 tool_call，不过滤，确保与后续 tool 消息一一对应
            try:
                tc_list = m.get("tool_calls")
                cleaned_tc = _hivo_ai_clean_tool_calls(tc_list)
                if cleaned_tc:
                    item["tool_calls"] = cleaned_tc
            except Exception:
                pass
            # 保留 tool_receipts（含 request/data），过滤 attachments，做长度截断
            try:
                tr = m.get("tool_receipts")
                if isinstance(tr, list) and tr:
                    # 读取配置阈值
                    try:
                        cfg0 = _hivo_load_cfg()
                        lim0 = cfg0.get("limits") if isinstance(cfg0.get("limits"), dict) else {}
                        req_max = int(lim0.get("tool_receipt_request_max") or 2000)
                        dat_max = int(lim0.get("tool_receipt_data_max") or 4000)
                    except Exception:
                        req_max, dat_max = 2000, 4000
                    out = []
                    for r in tr[:50]:
                        if not isinstance(r, dict):
                            continue
                        t = str(r.get("type") or "")
                        if t == "attachments":
                            continue
                        ok1 = bool(r.get("ok"))
                        msg1 = str(r.get("msg") or "")[:800]
                        req = r.get("request") if isinstance(r.get("request"), dict) else {}
                        dat = r.get("data") if isinstance(r.get("data"), dict) else {}
                        def _trim_json(o, mx):
                            try:
                                s = json.dumps(o, ensure_ascii=False)
                                if len(s) > mx:
                                    s = s[:mx]
                                return json.loads(s)
                            except Exception:
                                return {}
                        out.append({
                            "type": t,
                            "ok": ok1,
                            "msg": msg1,
                            "request": _trim_json(req, req_max),
                            "data": _trim_json(dat, dat_max),
                        })
                    if out:
                        item["tool_receipts"] = out
            except Exception:
                pass
            # stage/can_continue 信息
            try:
                st = str(m.get("stage") or "").strip()
                if st:
                    item["stage"] = st
                item["can_continue"] = bool(m.get("can_continue", False))
            except Exception:
                pass
            # 思考过程
            try:
                tc = str(m.get("thinking_content") or "").strip()
                if tc:
                    item["thinking_content"] = tc
            except Exception:
                pass
            # 持久化最小化 _attMeta（渲染附件缩略图）与 __continue_snapshot（用于 Continue 恢复）
            try:
                meta_out = _hivo_ai_clean_att_meta(m.get("_attMeta"), include_path=True)
                if meta_out:
                    item["_attMeta"] = meta_out
            except Exception:
                pass
            try:
                snap = m.get("__continue_snapshot")
                if isinstance(snap, dict) and isinstance(snap.get("body"), dict):
                    b = snap["body"]
                    body = {
                        "profile_id": str(b.get("profile_id") or ""),
                        "session_id": str(b.get("session_id") or ""),
                        "user_text": str(b.get("user_text") or ""),
                        "messages": b.get("messages") if isinstance(b.get("messages"), list) else [],
                        "dyn_context": str(b.get("dyn_context") or ""),
                        "context_refs": b.get("context_refs") if isinstance(b.get("context_refs"), list) else [],
                        "undo_gid": str(b.get("undo_gid") or ""),
                        "attachments": [
                            {
                                "name": str(a.get("name") or "file"),
                                "mime": str(a.get("mime") or ""),
                                "size": int(a.get("size") or 0),
                                **({"url": str(a.get("url") or "")} if str(a.get("url") or "").strip() else {}),
                                **({"path": str(a.get("path") or "").replace("\\", "/")} if str(a.get("path") or "").strip() else {}),
                            }
                            for a in (b.get("attachments") if isinstance(b.get("attachments"), list) else [])
                            if isinstance(a, dict)
                        ],
                        "env_observation": str(b.get("env_observation") or ""),
                    }
                    meta_snap = _hivo_ai_clean_att_meta(snap.get("_attMeta"), include_path=True)
                    item["__continue_snapshot"] = {"body": body, "_attMeta": meta_snap}
            except Exception:
                pass
        elif role == "tool":
            try:
                tcid = str(m.get("tool_call_id") or "").strip()
                if tcid:
                    item["tool_call_id"] = tcid
            except Exception:
                pass
        cleaned.append(item)

    cleaned = _hivo_ai_fix_tool_call_ids(cleaned)

    lim = int(limit) if str(limit).strip() else 80
    if lim < 10:
        lim = 10
    if lim > 500:
        lim = 500
    return cleaned[-lim:]


def _hivo_ai_guess_session_title(messages: list):
    try:
        if not isinstance(messages, list):
            return ""
        for m in messages:
            if not isinstance(m, dict):
                continue
            if str(m.get("role") or "").strip() != "user":
                continue
            txt = str(m.get("content") or "").strip()
            if not txt:
                continue
            txt = " ".join(txt.split())
            if len(txt) > 26:
                txt = txt[:26].rstrip() + "…"
            return txt
        return ""
    except Exception:
        return ""


def _hivo_ai_load_history_data():
    p = _hivo_ai_history_path()
    try:
        if not p.exists():
            return {}
        data = json.loads(p.read_text(encoding="utf-8") or "{}")
        return data if isinstance(data, dict) else {}
    except Exception:
        return {}


def _hivo_ai_write_history_data(data: dict):
    p = _hivo_ai_history_path()
    _atomic_write_json(p, data, indent=2)


def _hivo_ai_ensure_profile_node(data: dict, profile_id: str):
    byp = data.get("by_profile")
    if not isinstance(byp, dict):
        byp = {}
        data["by_profile"] = byp
    node = byp.get(profile_id)

    # Migration: old format stored list directly
    if isinstance(node, list):
        sid = uuid.uuid4().hex
        node = {
            "active_session_id": sid,
            "default_session_id": sid,
            "session_order": [sid],
            "sessions": [
                {
                    "id": sid,
                    "title": "会话",
                    "updated_at": time.time(),
                    "messages": node,
                }
            ],
        }
        byp[profile_id] = node

    if not isinstance(node, dict):
        sid = uuid.uuid4().hex
        node = {
            "active_session_id": sid,
            "default_session_id": sid,
            "session_order": [sid],
            "sessions": [
                {
                    "id": sid,
                    "title": "会话",
                    "updated_at": time.time(),
                    "messages": [],
                }
            ],
        }
        byp[profile_id] = node

    sess = node.get("sessions")
    if not isinstance(sess, list):
        sess = []
        node["sessions"] = sess

    # Ensure default session id exists
    default_sid = str(node.get("default_session_id") or "").strip()
    if not default_sid:
        default_sid = str(sess[0].get("id") or "").strip() if sess else ""
        if default_sid:
            node["default_session_id"] = default_sid

    # Ensure session order list exists
    order = node.get("session_order")
    if not isinstance(order, list):
        order = []
        node["session_order"] = order
    # Normalize order: only keep existing ids
    existing_ids = [str(x.get("id") or "").strip() for x in sess if isinstance(x, dict) and str(x.get("id") or "").strip()]
    exist_set = set(existing_ids)
    norm = [str(x).strip() for x in order if str(x).strip() in exist_set]
    for sid0 in existing_ids:
        if sid0 not in norm:
            norm.append(sid0)
    node["session_order"] = norm

    active = str(node.get("active_session_id") or "").strip()
    if active and any(isinstance(x, dict) and str(x.get("id") or "").strip() == active for x in sess):
        return node

    if sess:
        node["active_session_id"] = str(sess[0].get("id") or "").strip()
    else:
        sid = uuid.uuid4().hex
        node["active_session_id"] = sid
        node["sessions"] = [{"id": sid, "title": "会话", "updated_at": time.time(), "messages": []}]
    return node


def list_ai_sessions(profile_id: str):
    pid = _AI_GLOBAL_PROFILE_ID
    data = _hivo_ai_load_history_data()
    node = _hivo_ai_ensure_profile_node(data, pid)
    sessions = []
    by_id = {}
    for s in (node.get("sessions") or []):
        if not isinstance(s, dict):
            continue
        sid = str(s.get("id") or "").strip()
        if not sid:
            continue
        item = {
            "id": sid,
            "title": str(s.get("title") or "会话"),
            "updated_at": float(s.get("updated_at") or 0.0),
        }
        by_id[sid] = item
        sessions.append(item)

    order = node.get("session_order") if isinstance(node.get("session_order"), list) else []

    # Keep by stored order; fallback to updated_at desc.
    if order:
        ordered = []
        for sid in [str(x).strip() for x in order if str(x).strip()]:
            if sid in by_id:
                ordered.append(by_id[sid])
        # Append any missing sessions to the end (do NOT reorder them to the front).
        ordered_ids = set([y.get("id") for y in ordered])
        missing = [x for x in sessions if x.get("id") not in ordered_ids]
        sessions = ordered + missing
    else:
        sessions.sort(key=lambda x: x.get("updated_at") or 0.0, reverse=True)
    return {"active_session_id": str(node.get("active_session_id") or ""), "sessions": sessions}


def reorder_ai_sessions(profile_id: str, session_ids: list[str]):
    pid = _AI_GLOBAL_PROFILE_ID
    data = _hivo_ai_load_history_data()
    node = _hivo_ai_ensure_profile_node(data, pid)
    sess = node.get("sessions")
    if not isinstance(sess, list):
        return False, "会话数据异常"
    existing = [str(x.get("id") or "").strip() for x in sess if isinstance(x, dict) and str(x.get("id") or "").strip()]
    exist_set = set(existing)
    ordered = []
    for sid in (session_ids or []):
        s = str(sid or "").strip()
        if not s or s not in exist_set:
            continue
        if s not in ordered:
            ordered.append(s)
    final_order = list(ordered)
    for sid in existing:
        if sid not in final_order:
            final_order.append(sid)
    node["session_order"] = final_order
    _hivo_ai_write_history_data(data)
    return True, ""


def set_ai_active_session(profile_id: str, session_id: str):
    sid = str(session_id or "").strip()
    pid = _AI_GLOBAL_PROFILE_ID
    if not sid:
        return False, "缺少参数"

    data = _hivo_ai_load_history_data()
    node = _hivo_ai_ensure_profile_node(data, pid)
    sess = node.get("sessions") or []
    if not any(isinstance(x, dict) and str(x.get("id") or "").strip() == sid for x in sess):
        return False, "会话不存在"
    node["active_session_id"] = sid
    _hivo_ai_write_history_data(data)
    return True, ""


def create_ai_session(profile_id: str, title: str | None = None):
    pid = _AI_GLOBAL_PROFILE_ID
    data = _hivo_ai_load_history_data()
    node = _hivo_ai_ensure_profile_node(data, pid)
    sid = uuid.uuid4().hex
    t = str(title or "会话").strip() or "会话"
    sess = node.get("sessions")
    if not isinstance(sess, list):
        sess = []
        node["sessions"] = sess
    sess.append({"id": sid, "title": t, "updated_at": time.time(), "messages": []})
    node["active_session_id"] = sid
    # Append new session to order so it appears at the end by default.
    order = node.get("session_order")
    if not isinstance(order, list):
        order = []
    sid_s = str(sid or "").strip()
    if sid_s and sid_s not in [str(x).strip() for x in order if str(x).strip()]:
        order.append(sid_s)
    node["session_order"] = order
    _hivo_ai_write_history_data(data)
    return True, "", sid


def delete_ai_session(profile_id: str, session_id: str):
    sid = str(session_id or "").strip()
    pid = _AI_GLOBAL_PROFILE_ID
    if not sid:
        return False, "缺少参数"

    data = _hivo_ai_load_history_data()
    node = _hivo_ai_ensure_profile_node(data, pid)
    sess = node.get("sessions")
    if not isinstance(sess, list):
        return False, "会话不存在"
    # Collect attachments to delete for this session before removing it (best-effort)
    try:
        target_session = next((x for x in sess if isinstance(x, dict) and str(x.get("id") or "").strip() == sid), None)
        att_paths = []
        if isinstance(target_session, dict):
            msgs = target_session.get("messages")
            if isinstance(msgs, list):
                for m in msgs:
                    if not isinstance(m, dict):
                        continue
                    if str(m.get("role") or "").strip() != "user":
                        continue
                    meta = m.get("_attMeta")
                    if not isinstance(meta, list):
                        continue
                    for a in meta:
                        if not isinstance(a, dict):
                            continue
                        u = str(a.get("url") or "").strip()
                        if not u:
                            continue
                        # Extract attachments relative path
                        rel = ""
                        if "/api/ai_attachment" in u and "path=" in u:
                            try:
                                from urllib.parse import urlparse, parse_qs
                                pr = urlparse(u)
                                qs = parse_qs(pr.query or "")
                                p0 = qs.get("path")
                                if isinstance(p0, list) and p0:
                                    rel = str(p0[0] or "").strip()
                            except Exception:
                                rel = ""
                        elif u.startswith("attachments/"):
                            rel = u
                        if not rel:
                            continue
                        # Normalize and ensure it's inside the attachments directory
                        base_dir = os.path.join(str(_hivo_ai_data_dir()), "attachments")
                        full = os.path.abspath(os.path.join(str(_hivo_ai_data_dir()), rel))
                        base_abs = os.path.abspath(base_dir)
                        if full.startswith(base_abs):
                            att_paths.append(full)
        # Delete files (dedup)
        seen = set()
        removed = 0
        for pth in att_paths:
            if not pth or pth in seen:
                continue
            seen.add(pth)
            try:
                if os.path.isfile(pth):
                    os.remove(pth)
                    removed += 1
            except Exception:
                pass
        if removed:
            logger.info(f"删除会话 {sid} 同步清理附件文件 {removed} 个")
    except Exception:
        # ignore cleanup errors
        pass
    new_sess = [x for x in sess if isinstance(x, dict) and str(x.get("id") or "").strip() != sid]
    if len(new_sess) == len(sess):
        return False, "会话不存在"
    node["sessions"] = new_sess

    # 清理内存中的会话状态，防止内存泄漏
    try:
        _hivo_session_mem.pop(sid, None)
    except Exception:
        pass

    # Update session_order and pick a reasonable next active session.
    try:
        order0 = node.get("session_order")
        if not isinstance(order0, list):
            order0 = []
        order_ids0 = [str(x).strip() for x in order0 if str(x).strip()]
        existing_ids = [str(x.get("id") or "").strip() for x in new_sess if isinstance(x, dict) and str(x.get("id") or "").strip()]
        existing_set = set(existing_ids)

        # Normalize order: keep existing only, and ensure deleted sid is removed.
        order_ids = [x for x in order_ids0 if x and x != sid and x in existing_set]
        # Append any sessions not in order to the end.
        for x in existing_ids:
            if x and x not in order_ids:
                order_ids.append(x)
        node["session_order"] = order_ids

        if str(node.get("active_session_id") or "").strip() == sid:
            # Prefer the previous session in the stored order; fallback to next.
            next_active = ""
            try:
                idx0 = -1
                if sid in order_ids0:
                    idx0 = order_ids0.index(sid)
                if idx0 >= 0:
                    for j in range(idx0 - 1, -1, -1):
                        cand = str(order_ids0[j] or "").strip()
                        if cand and cand in existing_set:
                            next_active = cand
                            break
                    if not next_active:
                        for j in range(idx0 + 1, len(order_ids0)):
                            cand = str(order_ids0[j] or "").strip()
                            if cand and cand in existing_set:
                                next_active = cand
                                break
            except Exception:
                next_active = ""
            if not next_active:
                next_active = existing_ids[0] if existing_ids else ""
            node["active_session_id"] = next_active
    except Exception:
        if str(node.get("active_session_id") or "").strip() == sid:
            node["active_session_id"] = str(new_sess[0].get("id") or "").strip() if new_sess else ""
    _hivo_ai_write_history_data(data)
    return True, ""


def load_ai_chat_history(profile_id: str, limit: int = 40, session_id: str | None = None):
    pid = _AI_GLOBAL_PROFILE_ID
    data = _hivo_ai_load_history_data()
    node = _hivo_ai_ensure_profile_node(data, pid)
    sid = str(session_id or "").strip() or str(node.get("active_session_id") or "").strip()
    if not sid:
        return [], None

    for s in (node.get("sessions") or []):
        if not isinstance(s, dict):
            continue
        if str(s.get("id") or "").strip() != sid:
            continue
        arr = s.get("messages")
        if not isinstance(arr, list):
            return [], None
        out = []
        last_assistant = None
        for m in arr[-max(1, int(limit)):]:
            if not isinstance(m, dict):
                continue
            role = str(m.get("role") or "").strip()
            # 处理 content：支持字符串和多模态数组两种格式
            raw_content = m.get("content")
            if isinstance(raw_content, list):
                content = raw_content
            else:
                content = str(raw_content or "")
            if role not in ("user", "assistant", "system", "tool"):
                continue
            # Skip empty content except we still allow tracking for last assistant (for session-level continue)
            if not _hivo_ai_msg_has_valid_content({"content": content}) and role not in ("assistant", "tool"):
                continue
            item = {"role": role, "content": content}
            # Preserve common identifiers when present
            mid = str(m.get("id") or "").strip()
            if mid:
                item["id"] = mid
            if role == "user":
                ug = str(m.get("undo_gid") or "").strip()
                if ug:
                    item["undo_gid"] = ug
                # include stored _attMeta so client can render without local cache
                meta_in = m.get("_attMeta")
                if isinstance(meta_in, list) and meta_in:
                    try:
                        meta_out = []
                        for a in meta_in:
                            if not isinstance(a, dict):
                                continue
                            o = {
                                "name": str(a.get("name") or "file"),
                                "mime": str(a.get("mime") or ""),
                                "size": int(a.get("size") or 0),
                            }
                            u = str(a.get("url") or "").strip()
                            if u:
                                o["url"] = u
                            meta_out.append(o)
                        if meta_out:
                            item["_attMeta"] = meta_out
                    except Exception:
                        pass
            elif role == "assistant":
                # Preserve light assistant fields for compatibility only
                stage = m.get("stage")
                if isinstance(stage, str) and stage:
                    item["stage"] = stage
                ug = str(m.get("undo_gid") or "").strip()
                if ug:
                    item["undo_gid"] = ug
                # 保留思考过程
                try:
                    tc = str(m.get("thinking_content") or "").strip()
                    if tc:
                        item["thinking_content"] = tc
                except Exception:
                    pass
                # Preserve tool_receipts for cross-session display
                try:
                    tr = m.get("tool_receipts")
                    if isinstance(tr, list) and tr:
                        item["tool_receipts"] = tr
                except Exception:
                    pass
                # 保留原生 tool_calls（含 id），确保会话恢复后工具调用链 ID 匹配
                # 注意：必须完整保留所有 tool_call，不过滤，确保与后续 tool 消息一一对应
                try:
                    tc_list = m.get("tool_calls")
                    cleaned_tc = _hivo_ai_clean_tool_calls(tc_list)
                    if cleaned_tc:
                        item["tool_calls"] = cleaned_tc
                except Exception:
                    pass
                # Track for session-level continue computation
                last_assistant = m
            elif role == "tool":
                try:
                    tcid = str(m.get("tool_call_id") or "").strip()
                    if tcid:
                        item["tool_call_id"] = tcid
                except Exception:
                    pass
            out.append(item)
        # Derive session-level continue strictly from last assistant
        cont_obj = None
        try:
            if isinstance(last_assistant, dict):
                stage = str(last_assistant.get("stage") or "").strip()
                can_c = bool(last_assistant.get("can_continue"))
                snap = last_assistant.get("__continue_snapshot") if isinstance(last_assistant.get("__continue_snapshot"), dict) else None
                if snap and (can_c or stage in ("error", "aborted")):
                    cont_obj = {
                        "message_id": str(last_assistant.get("id") or ""),
                        "status": ("active" if can_c else ("error" if stage == "error" else "aborted")),
                        "snapshot": snap,
                        "updated_at": time.time(),
                    }
        except Exception:
            cont_obj = None
        out = _hivo_ai_fix_tool_call_ids(out)
        return out, cont_obj
    return [], None


def save_ai_chat_history(profile_id: str, messages: list, limit: int = 80, session_id: str | None = None):
    pid = _AI_GLOBAL_PROFILE_ID
    # Strip any continue-related fields before persisting (runtime-only)
    def _strip_continue_fields(arr: list):
        out = []
        for m in arr or []:
            if not isinstance(m, dict):
                continue
            mm = dict(m)
            mm.pop("can_continue", None)
            mm.pop("can_continue_reason", None)
            mm.pop("__continue_snapshot", None)
            out.append(mm)
        return out

    cleaned = _hivo_ai_clean_messages(_strip_continue_fields(messages), limit)

    data = _hivo_ai_load_history_data()
    node = _hivo_ai_ensure_profile_node(data, pid)
    sid = str(session_id or "").strip() or str(node.get("active_session_id") or "").strip()
    if not sid:
        ok, msg, sid = create_ai_session(pid, title="会话")
        if not ok:
            return False, msg
        data = _hivo_ai_load_history_data()
        node = _hivo_ai_ensure_profile_node(data, pid)

    sess = node.get("sessions")
    if not isinstance(sess, list):
        return False, "会话数据异常"

    # If client-provided session_id is stale (e.g. deleted), fall back to current active session.
    if session_id:
        if not any(isinstance(x, dict) and str(x.get("id") or "").strip() == sid for x in sess):
            sid2 = str(node.get("active_session_id") or "").strip()
            if sid2 and any(isinstance(x, dict) and str(x.get("id") or "").strip() == sid2 for x in sess):
                sid = sid2
            else:
                ok, msg, new_sid = create_ai_session(pid, title="会话")
                if not ok:
                    return False, msg
                sid = str(new_sid or "").strip()
                data = _hivo_ai_load_history_data()
                node = _hivo_ai_ensure_profile_node(data, pid)
                sess = node.get("sessions")
                if not isinstance(sess, list):
                    return False, "会话数据异常"
    for s in sess:
        if not isinstance(s, dict):
            continue
        if str(s.get("id") or "").strip() == sid:
            s["messages"] = cleaned
            s["updated_at"] = time.time()
            title = str(s.get("title") or "").strip()
            auto_titled = bool(s.get("auto_titled"))
            if (not auto_titled) and title == "会话":
                t = _hivo_ai_guess_session_title(cleaned)
                if t:
                    s["title"] = t
                    s["auto_titled"] = True
            node["active_session_id"] = sid
            _hivo_ai_write_history_data(data)
            return True, ""
    return False, "会话不存在"


def load_hivo_ai_config():
    try:
        data = _hivo_load_cfg()
        if not isinstance(data, dict):
            data = {}
        # Strict new-structure requirements
        if not isinstance(data.get("profiles"), list):
            data["profiles"] = []
        if not isinstance(data.get("features"), dict):
            data["features"] = {"web_search_enabled": False, "ai_stream": True}
        if not isinstance(data.get("web_search"), dict):
            data["web_search"] = {"active_profile_id": "", "profiles": []}

        # Normalize active_profile_id to a valid profile id when possible.
        active = str(data.get("active_profile_id") or "").strip()
        profiles = data.get("profiles") if isinstance(data.get("profiles"), list) else []
        if active and any(isinstance(x, dict) and str(x.get("id") or "").strip() == active for x in profiles):
            return data
        if profiles and isinstance(profiles[0], dict) and str(profiles[0].get("id") or "").strip():
            data["active_profile_id"] = str(profiles[0].get("id") or "").strip()
            return data
        data["active_profile_id"] = ""
        return data
    except Exception:
        return {"active_profile_id": "", "profiles": [], "features": {"web_search_enabled": False, "ai_stream": True}, "web_search": {"active_profile_id": "", "profiles": []}}


def get_workspace_context(max_entries: int = 80):
    if not REPO_PATH:
        return ""
    try:
        repo_root = os.path.abspath(REPO_PATH)
    except Exception:
        repo_root = str(REPO_PATH)

    branch = ""
    try:
        out, _, code = run_git(["rev-parse", "--abbrev-ref", "HEAD"])
        if code == 0:
            branch = (out or "").strip()
    except Exception:
        branch = ""

    entries = []
    try:
        names = sorted(os.listdir(repo_root))
        for n in names:
            if n == ".git":
                continue
            if n.startswith(".") and n not in (".github", ".gitignore"):
                continue
            p = os.path.join(repo_root, n)
            entries.append(n + ("/" if os.path.isdir(p) else ""))
            if len(entries) >= int(max_entries):
                break
    except Exception:
        entries = []

    tree = "\n".join(["- " + e for e in entries])
    if entries and len(entries) >= int(max_entries):
        tree += "\n- …"
    return (
        f"Repo: {repo_root}\n"
        + (f"Branch: {branch}\n" if branch else "")
        + "Hints: If a file path is uncertain, use find_files(name) or list_dir_tree(path,depth) before calling open_file/save_file.\n"
        + "Top-level:\n"
        + (tree or "- (empty)")
    )


def find_files_by_name(name: str, max_results: int = 20):
    q = str(name or "").strip()
    if not REPO_PATH:
        return []
    ql = q.lower()
    repo_root = os.path.abspath(REPO_PATH)
    out = []
    try:
        for root, dirs, files in os.walk(repo_root):
            if ".git" in dirs:
                dirs.remove(".git")
            dirs[:] = [d for d in dirs if not d.startswith(".")]
            try:
                dirs.sort()
            except Exception as e:
                logger.debug(f"Exception ignored: {e}")
            try:
                files.sort()
            except Exception as e:
                logger.debug(f"Exception ignored: {e}")
            for fn in files:
                if fn.startswith("."):
                    continue
                fcl = fn.lower()
                if (not q) or (fcl == ql) or (ql in fcl):
                    abs_path = os.path.join(root, fn)
                    rel = os.path.relpath(abs_path, repo_root).replace("\\", "/")
                    out.append(rel)
                    if len(out) >= int(max_results):
                        return out
    except Exception:
        return out
    return out


def save_hivo_ai_config(cfg: dict):
    if not isinstance(cfg, dict):
        return False, "配置格式非法"

    features_in = cfg.get("features") if isinstance(cfg.get("features"), dict) else {}
    web_search_in = cfg.get("web_search") if isinstance(cfg.get("web_search"), dict) else {}

    base_cfg = _hivo_load_cfg()
    if not isinstance(base_cfg, dict):
        base_cfg = {}

    # Accept both new and legacy formats.
    profiles_in = cfg.get("profiles")
    if isinstance(profiles_in, list):
        profiles = []
        for p0 in profiles_in:
            if not isinstance(p0, dict):
                continue
            pid = str(p0.get("id") or "").strip()
            name = str(p0.get("name") or "").strip()
            base_url = str(p0.get("base_url") or p0.get("endpoint") or "").strip()
            api_key = str(p0.get("api_key") or "")
            model = str(p0.get("model") or "").strip()
            supports_vision = bool(p0.get("supports_vision", False))
            if not pid:
                pid = uuid.uuid4().hex
            if not name:
                name = model or "Profile"
            profiles.append({
                "id": pid,
                "name": name,
                "base_url": base_url,
                "api_key": api_key,
                "model": model,
                "supports_vision": supports_vision,
            })

        active = str(cfg.get("active_profile_id") or "").strip()
        if active and any(x.get("id") == active for x in profiles):
            active_id = active
        else:
            active_id = profiles[0]["id"] if profiles else ""
        base_cfg["active_profile_id"] = active_id
        base_cfg["profiles"] = profiles
    else:
        base_url = str(cfg.get("base_url") or cfg.get("endpoint") or "").strip()
        api_key = str(cfg.get("api_key") or "")
        model = str(cfg.get("model") or "").strip()
        pid = str(cfg.get("id") or "").strip() or uuid.uuid4().hex
        base_cfg["active_profile_id"] = pid
        base_cfg["profiles"] = [{
            "id": pid,
            "name": str(cfg.get("name") or model or "Default"),
            "base_url": base_url,
            "api_key": api_key,
            "model": model,
            "supports_vision": bool(cfg.get("supports_vision", False)),
        }]

    # Persist web search config together with ai config (same file).
    # IMPORTANT: merge-only behavior to avoid resetting settings when client submits partial config.
    try:
        if "features" not in base_cfg or not isinstance(base_cfg.get("features"), dict):
            base_cfg["features"] = {}
        if isinstance(features_in, dict) and "web_search_enabled" in features_in:
            base_cfg["features"]["web_search_enabled"] = bool(features_in.get("web_search_enabled"))
        if isinstance(features_in, dict) and "ai_stream" in features_in:
            base_cfg["features"]["ai_stream"] = bool(features_in.get("ai_stream"))
    except Exception as e:
        logger.debug(f"Exception ignored: {e}")

    try:
        # Only overwrite web_search when client provided it.
        if isinstance(cfg.get("web_search"), dict):
            ws_profiles = web_search_in.get("profiles") if isinstance(web_search_in.get("profiles"), list) else []
            ws_profiles2 = []
            for p0 in ws_profiles:
                if not isinstance(p0, dict):
                    continue
                ws_profiles2.append({
                    "id": str(p0.get("id") or "").strip(),
                    "provider": str(p0.get("provider") or "").strip(),
                    "api_key": str(p0.get("api_key") or ""),
                    "base_url": str(p0.get("base_url") or "").strip(),
                    "timeout": int(p0.get("timeout") or 20),
                    "top_k": int(p0.get("top_k") or 5),
                })
            ws_active = str(web_search_in.get("active_profile_id") or "").strip()
            if ws_active and any(x.get("id") == ws_active for x in ws_profiles2):
                active2 = ws_active
            else:
                active2 = str(ws_profiles2[0].get("id") or "").strip() if ws_profiles2 else ""
            base_cfg["web_search"] = {"active_profile_id": active2, "profiles": ws_profiles2}
    except Exception as e:
        logger.debug(f"Exception ignored: {e}")

    ok, msg = _hivo_save_cfg(base_cfg)
    if not ok:
        return False, msg
    return True, ""


def _ai_load_web_search_cfg():
    """Return (enabled, web_search_dict) from unified hivo_ai_config.json."""
    try:
        aic = load_hivo_ai_config()
        if isinstance(aic, dict):
            feat = aic.get("features") if isinstance(aic.get("features"), dict) else {}
            ws = aic.get("web_search") if isinstance(aic.get("web_search"), dict) else None
            if ws is not None:
                return bool(feat.get("web_search_enabled", False)), ws
    except Exception:
        return False, {}
    return False, {}


def _ai_build_chat_url(base_url: str):
    u = (base_url or "").strip()
    if not u:
        return ""
    u = u.rstrip("/")
    # If user provides full path, use it directly.
    if u.endswith("/chat/completions"):
        return u
    # Extract path after scheme+host to determine URL depth.
    try:
        parsed = urllib.parse.urlparse(u)
        path = parsed.path or ""
    except Exception:
        path = u
    # Count non-empty path segments.
    segments = [s for s in path.split("/") if s]
    # If URL has any path (version-only like /v1, or non-standard like /v2/tokenplan/personal),
    # just append /chat/completions without forcing /v1 prefix.
    if len(segments) >= 1:
        return u + "/chat/completions"
    # Bare domain with no path: use standard /v1/chat/completions.
    return u + "/v1/chat/completions"


def _ai_build_models_url(base_url: str):
    u = (base_url or "").strip()
    if not u:
        return ""
    u = u.rstrip("/")
    # If the user provided the full chat completions URL, strip it before building the models URL.
    # This handles non-standard paths like /zen/v1/chat/completions as well as standard /v1/chat/completions.
    _chat_suffix = "/chat/completions"
    if u.endswith(_chat_suffix):
        u = u[:-len(_chat_suffix)]
    if u.endswith("/models"):
        return u
    return u + "/models"


def _parse_upstream_error(status_code: int, body: str) -> str:
    """Parse upstream API error response into a human-readable message.

    Handles OpenAI-compatible JSON error format ({"error":{"message":"...","code":"..."}})
    and Cloudflare/HTML error pages (error code:1010 etc.).

    Returns empty string when no known error pattern is detected (caller should fall back
    to a generic message for non-2xx status codes).
    """
    body = (body or "").strip()
    if not body:
        return ""

    # Try to parse as OpenAI-compatible JSON error first.
    try:
        data = json.loads(body)
        err = data.get("error")
        if isinstance(err, dict):
            parts = []
            em = str(err.get("message") or err.get("msg") or "").strip()
            ec = str(err.get("code") or "")
            et = str(err.get("type") or "").strip()
            if em:
                parts.append(em)
            if ec:
                parts.append(f"[code={ec}]")
            if et and et not in (ec, em):
                parts.append(f"[type={et}]")
            if parts:
                return f"上游返回错误 (HTTP {status_code}): " + " ".join(parts)
        elif isinstance(err, str) and err.strip():
            return f"上游返回错误 (HTTP {status_code}): {err.strip()}"
    except Exception:
        pass

    # Detect Cloudflare error pages (HTML body).
    body_lower = body.lower()
    if "error code:1010" in body_lower or "cf-error" in body_lower or "cloudflare" in body_lower:
        return (
            f"上游返回错误 (HTTP {status_code}): 请求被 Cloudflare 安全策略拦截（Error 1010）。\n"
            "可能原因：User-Agent 被屏蔽或 IP 被限制，已自动添加浏览器 UA 头重试。\n"
            "如仍失败，请联系上游服务商开放 API 访问权限。"
        )

    # Generic: include status code and first 400 chars of body (only for error status codes).
    if status_code >= 400:
        snip = body.strip()
        if len(snip) > 400:
            snip = snip[:400] + "…"
        return f"上游返回错误 (HTTP {status_code}): {snip}"

    return ""




def _ai_build_opener():
    """Build a urllib opener with custom SSL context (TLS 1.2+)."""
    https_handler = urllib.request.HTTPSHandler(context=_AI_SSL_CONTEXT) if _AI_SSL_CONTEXT else urllib.request.HTTPSHandler()
    return urllib.request.build_opener(https_handler)


def _ai_extract_content(msg: dict) -> str:
    """Extract text content from an API response message, supporting multiple formats.

    Handles:
    - Standard OpenAI: content="text"
    - Anthropic-compatible: content=[{"type":"text","text":"..."}]
    - Reasoning models: reasoning/reasoning_details fields
    - Empty/null content + reasoning fallback
    """
    content = msg.get("content")
    if isinstance(content, str):
        text = content.strip()
        if text:
            return text
    elif isinstance(content, list):
        # Anthropic-style content array: [{"type":"text","text":"Hello"}]
        parts = []
        for it in content:
            if isinstance(it, dict):
                if str(it.get("type") or "").strip() == "text":
                    parts.append(str(it.get("text") or ""))
            elif isinstance(it, str):
                parts.append(it)
        text = "".join(parts).strip()
        if text:
            return text

    # Fallback: reasoning / thinking models that put output in non-standard fields
    reasoning = str(msg.get("reasoning") or "").strip()
    if not reasoning:
        rd = msg.get("reasoning_details")
        if isinstance(rd, list):
            reasoning = "".join(str(it.get("text") or "") for it in rd if isinstance(it, dict)).strip()
    if reasoning:
        return reasoning

    # Last resort: text field (some APIs use this)
    text_val = str(msg.get("text") or "").strip()
    if text_val:
        return text_val

    return ""


def _ai_parse_model_list(data):
    """Flexible model list parser supporting OpenAI, Ollama, vLLM, LiteLLM, etc.

    Returns a sorted, deduplicated list of model ID strings.
    """
    if not isinstance(data, dict):
        return []

    arr = data.get("data")
    if not isinstance(arr, list):
        arr = []

    out = []
    for it in arr:
        if not isinstance(it, (dict, str)):
            continue
        mid = None
        if isinstance(it, dict):
            mid = str(it.get("id") or it.get("name") or "").strip()
        else:
            mid = str(it).strip()
        if not mid:
            continue
        out.append(mid)
        if len(out) >= 300:
            break

    # Fallback: some APIs (Ollama, Gemini via proxy) use "models" key with model name strings
    if not out:
        alt = data.get("models")
        if isinstance(alt, list):
            for it in alt:
                if isinstance(it, str):
                    mid = it.strip()
                    if mid:
                        out.append(mid)
                elif isinstance(it, dict):
                    mid = str(it.get("id") or it.get("name") or "").strip()
                    if mid:
                        out.append(mid)
                if len(out) >= 300:
                    break

    return sorted(list(dict.fromkeys(out)))


_ai_models_cache = {
    "items": {},
    "ttl_s": 60.0,
}


def ai_list_models(base_url: str, api_key: str | None = None):
    base_url = str(base_url or "").strip()
    api_key = str(api_key or "")
    url = _ai_build_models_url(base_url)
    if not url:
        return False, "未配置 API Base URL", []

    cache_key = (url, "1" if api_key else "0")
    try:
        ent = _ai_models_cache.get("items", {}).get(cache_key)
        if isinstance(ent, dict) and (time.time() - float(ent.get("ts") or 0.0)) <= float(_ai_models_cache.get("ttl_s") or 60.0):
            models = ent.get("models")
            if isinstance(models, list):
                return True, "", models
    except Exception as e:
        logger.debug(f"Exception ignored: {e}")

    headers = {
        "Accept": "application/json",
        "User-Agent": "GitManager/1.0",
    }
    if api_key:
        headers["Authorization"] = f"Bearer {api_key}"
        # Some APIs (e.g. Anthropic, Together) use x-api-key instead of Bearer.
        headers["x-api-key"] = api_key

    req = urllib.request.Request(url, headers=headers, method="GET")
    try:
        with _ai_build_opener().open(req, timeout=30) as resp:
            raw = resp.read().decode("utf-8", errors="replace")
        data = json.loads(raw or "{}")
    except urllib.error.HTTPError as e:
        try:
            body = e.read().decode("utf-8", errors="replace")
        except Exception:
            body = ""
        msg = _parse_upstream_error(e.code, body)
        return False, msg, []
    except Exception as e:
        return False, str(e), []

    # ── Flexible model-list parsing: supports OpenAI, Ollama, vLLM, LiteLLM, etc. ──
    out = _ai_parse_model_list(data)
    if not out:
        return False, "响应格式不支持或未返回任何模型", []

    try:
        items = _ai_models_cache.get("items")
        if not isinstance(items, dict):
            items = {}
            _ai_models_cache["items"] = items
        items[cache_key] = {"ts": time.time(), "models": out}
    except Exception as e:
        logger.debug(f"Exception ignored: {e}")
    return True, "", out


def get_capabilities_spec():
    cfg0 = _hivo_load_cfg()
    cap = cfg0.get("capabilities") if isinstance(cfg0, dict) else None
    cap = cap if isinstance(cap, dict) else {}
    dis_ep = cap.get("disable_endpoints")
    dis_ep = dis_ep if isinstance(dis_ep, list) else []
    dis_tool = cap.get("disable_tools")
    dis_tool = dis_tool if isinstance(dis_tool, list) else []
    ep_extra = cap.get("endpoints_extra")
    ep_extra = ep_extra if isinstance(ep_extra, list) else []
    tool_extra = cap.get("agent_tools_extra")
    tool_extra = tool_extra if isinstance(tool_extra, list) else []

    endpoints = [
        # Repo/Workspace
        {"method": "GET", "path": "/api/status"},
        {"method": "GET", "path": "/api/sync_status", "query": {"fetch": "0|1"}},
        {"method": "POST", "path": "/api/open_repo", "body": {"path": "<abs path>"}},
        {"method": "POST", "path": "/api/set_origin", "body": {"url": "<git remote url>"}},
        {"method": "GET", "path": "/api/staged_files"},
        {"method": "POST", "path": "/api/unstage_file", "body": {"path": "<rel path>"}},
        {"method": "POST", "path": "/api/discard_staged_file", "body": {"path": "<rel path>"}},
        {"method": "POST", "path": "/api/unstage_all_staged"},
        {"method": "POST", "path": "/api/discard_all_staged"},
        {"method": "GET", "path": "/api/files", "query": {"max_age": "0-10"}},
        {"method": "GET", "path": "/api/diff", "query": {"path": "<rel>", "status": "M/A/D/R/U", "ctx": "0-200"}},
        {"method": "GET", "path": "/api/file_content", "query": {"path": "<rel>"}},
        {"method": "GET", "path": "/api/raw_file", "query": {"path": "<rel>"}},
        {"method": "GET", "path": "/api/read_file_range", "query": {"path": "<rel>", "start": "1+", "end": "1+"}},
        {"method": "GET", "path": "/api/list_dir_tree", "query": {"path": "<rel>", "depth": "1-10", "max_entries": "1-2000"}},
        {"method": "GET", "path": "/api/search_code", "query": {"query": "<text>", "case_sensitive": "0/1", "max_results": "1-200"}},
        {"method": "GET", "path": "/api/find_files", "query": {"name": "<filename>"}},
        {"method": "GET", "path": "/api/workspace_context"},

        # Write / File ops
        {"method": "POST", "path": "/api/save_file", "body": {"path": "<rel>", "content": "<text>"}},
        {"method": "POST", "path": "/api/delete_file", "body": {"path": "<rel>"}},
        {"method": "POST", "path": "/api/file_content_head", "body": {"path": "<rel>"}},
        {"method": "POST", "path": "/api/rename_file", "body": {"old_path": "<rel>", "new_path": "<rel>"}},
        {"method": "POST", "path": "/api/revert_file", "body": {"path": "<rel>", "status": "M/A/D/R/U"}},
        {"method": "POST", "path": "/api/revert_hunk", "body": {"path": "<rel>", "hunk_index": 0, "status": "M/A/D/R/U"}},
        {"method": "POST", "path": "/api/revert_line", "body": {"path": "<rel>", "hunk_index": 0, "line_index": 0, "status": "M/A/D/R/U"}},
        {"method": "POST", "path": "/api/revert_multi_lines", "body": {"path": "<rel>", "hunk_index": 0, "start_line_index": 0, "end_line_index": 0, "status": "M/A/D/R/U"}},
        {"method": "POST", "path": "/api/revert_all", "body": {"paths": ["<rel>"]}},

        # Git
        {"method": "GET", "path": "/api/commits"},
        {"method": "GET", "path": "/api/log"},
        {"method": "GET", "path": "/api/commit_detail", "query": {"hash": "<sha>"}},
        {"method": "GET", "path": "/api/commit_file_diff", "query": {"hash": "<sha>", "path": "<rel>"}},
        {"method": "GET", "path": "/api/branches"},
        {"method": "GET", "path": "/api/commit_push_status", "query": {"hash": "<sha>"}},
        {"method": "GET", "path": "/api/file_log", "query": {"path": "<rel>", "limit": "1-200"}},
        {"method": "POST", "path": "/api/stage_file", "body": {"path": "<rel>"}},
        {"method": "POST", "path": "/api/commit", "body": {"message": "<text>", "files": ["<rel>"]}},
        {"method": "POST", "path": "/api/commit_hunks", "body": {"message": "<text>", "path": "<rel>", "status": "M/A/D/R/U", "hunks": [], "ctx": 0}},
        {"method": "POST", "path": "/api/commit_lines", "body": {"message": "<text>", "path": "<rel>", "status": "M/A/D/R/U", "lines": [], "ctx": 0}},
        {"method": "POST", "path": "/api/pull", "body": {}},
        {"method": "POST", "path": "/api/push", "body": {}},
        {"method": "POST", "path": "/api/stash_and_pull", "body": {}},
        {"method": "POST", "path": "/api/commit_and_pull", "body": {"message": "<text>", "files": ["<rel>"]}},
        {"method": "POST", "path": "/api/switch_branch", "body": {"branch": "<name>"}},
        {"method": "POST", "path": "/api/stash_and_switch", "body": {"branch": "<name>"}},
        {"method": "POST", "path": "/api/commit_and_switch", "body": {"branch": "<name>", "message": "<text>", "files": ["<rel>"]}},
        {"method": "GET", "path": "/api/staged_files"},
        {"method": "POST", "path": "/api/unstage_file", "body": {"path": "<rel>"}},
        {"method": "POST", "path": "/api/unstage_all_staged", "body": {}},
        {"method": "POST", "path": "/api/discard_staged_file", "body": {"path": "<rel>"}},
        {"method": "POST", "path": "/api/discard_all_staged", "body": {}},
        {"method": "GET", "path": "/api/stash_list"},
        {"method": "POST", "path": "/api/stash_pop", "body": {"ref": "stash@{0}"}},
        {"method": "POST", "path": "/api/stash_drop", "body": {"ref": "stash@{0}"}},
        {"method": "POST", "path": "/api/fetch_remotes", "body": {}},
        {"method": "POST", "path": "/api/pull_file", "body": {"path": "<rel>"}},
        {"method": "POST", "path": "/api/restore_file", "body": {"hash": "<sha>", "path": "<rel>"}},
        {"method": "POST", "path": "/api/restore_workspace", "body": {"hash": "<sha>"}},
        {"method": "POST", "path": "/api/revert_commit", "body": {"hash": "<sha>"}},
        {"method": "POST", "path": "/api/soft_reset_commit", "body": {"hash": "<sha>"}},

        # AI
        {"method": "GET", "path": "/api/ai_config"},
        {"method": "POST", "path": "/api/ai_config", "body": {"config": {"active_profile_id": "", "profiles": []}}},
        {"method": "GET", "path": "/api/hivo_ai_config"},
        {"method": "POST", "path": "/api/hivo_ai_config", "body": {"config": {}}},
        {"method": "GET", "path": "/api/ai_models", "query": {"base_url": "<url>", "api_key": "<key>"}},
        {"method": "POST", "path": "/api/ai_chat", "body": {"messages": [], "profile_id": "<id>", "session_id": "<id>"}},
        {"method": "GET", "path": "/api/ai_chat_history", "query": {"profile_id": "<id>", "session_id": "<id>", "limit": "1-200"}},
        {"method": "POST", "path": "/api/ai_chat_history", "body": {"profile_id": "<id>", "session_id": "<id>", "messages": []}},
        {"method": "POST", "path": "/api/hivo_agent", "body": {"profile_id": "<id>", "session_id": "<id>", "user_text": "<text>", "messages": [], "dyn_context": "<text>", "env_observation": "<text?>"}},
        {"method": "GET", "path": "/api/ai_sessions", "query": {"profile_id": "<id>"}},
        {"method": "POST", "path": "/api/ai_sessions", "body": {"profile_id": "<id>", "action": "create/delete/set_active/reorder", "session_id": "<id>"}},
        {"method": "POST", "path": "/api/ai_cache_clear", "body": {}},

        # System
        {"method": "POST", "path": "/api/open_file", "body": {"name": "<name or rel path>", "view": "editor/split/change/unified"}},
        {"method": "POST", "path": "/api/verify_python", "body": {"paths": ["<rel>"]}},
        {"method": "POST", "path": "/api/run_cmd", "body": {"cmd": "<single line>", "timeout": 30, "cwd": "<rel?>"}},
        {"method": "POST", "path": "/api/browser_tool", "body": {"action": "open/click/type/eval/text/screenshot/wait/close", "session_id": "<id>", "params": {}}},
        {"method": "POST", "path": "/api/undo", "body": {}},
        {"method": "GET", "path": "/api/undo_stats"},
        {"method": "GET", "path": "/api/capabilities"},
    ]

    try:
        for e in ep_extra:
            if isinstance(e, dict) and e.get("method") and e.get("path"):
                endpoints.append(e)
    except Exception as e:
        logger.debug(f"Exception ignored: {e}")

    # Agent Tools 定义
    agent_tools = [
        {
            "type": "save_file",
            "desc": "创建新文件或覆盖已有文件的内容（同时支持新建和修改）。",
            "required": ["path", "content"],
            "properties": {
                "path": {"type": "string", "desc": "相对路径，必须在仓库内，例如 src/main.py"},
                "content": {"type": "string", "desc": "完整文件内容"},
            },
            "example": {"type": "save_file", "path": "README.md", "content": "# Title"},
        },
        {
            "type": "edit_file",
            "desc": "局部编辑文件内容（将文件中指定的 old_string 替换为 new_string）。适用于修改大文件的一小部分，比 save_file 更高效。old_string 必须在文件中唯一匹配，否则返回错误。",
            "required": ["path", "old_string", "new_string"],
            "properties": {
                "path": {"type": "string", "desc": "相对路径，必须在仓库内"},
                "old_string": {"type": "string", "desc": "要替换的旧文本（必须在文件中唯一存在）"},
                "new_string": {"type": "string", "desc": "新文本内容"},
            },
            "example": {"type": "edit_file", "path": "README.md", "old_string": "# Old Title", "new_string": "# New Title"},
        },
        {
            "type": "find_files",
            "desc": "按文件名模糊搜索工作区内的文件（子串匹配，不区分大小写），返回匹配文件的相对路径列表。用户只给文件名或部分名称时，先用此工具定位真实路径。",
            "required": ["name"],
            "properties": {
                "name": {"type": "string", "desc": "文件名或部分名称（支持模糊匹配，如 server / .py / 123 均可）"},
            },
            "example": {"type": "find_files", "name": "server"},
        },
        {
            "type": "delete_file",
            "desc": "删除文件。",
            "required": ["path"],
            "properties": {
                "path": {"type": "string", "desc": "相对路径"},
            },
            "example": {"type": "delete_file", "path": "tmp.txt"},
        },
        {
            "type": "rename_file",
            "desc": "重命名文件（移动路径）。",
            "required": ["old_path", "new_path"],
            "properties": {
                "old_path": {"type": "string", "desc": "原相对路径"},
                "new_path": {"type": "string", "desc": "新相对路径"},
            },
            "example": {"type": "rename_file", "old_path": "a.txt", "new_path": "b.txt"},
        },
        {
            "type": "copy_file",
            "desc": "复制文件或目录（支持复制到另一个目录，或复制为新文件名）。",
            "required": ["source", "destination"],
            "properties": {
                "source": {"type": "string", "desc": "源相对路径（文件或目录）"},
                "destination": {"type": "string", "desc": "目标相对路径（文件或目录）"},
            },
            "example": {"type": "copy_file", "source": "a.txt", "destination": "b.txt"},
        },
        {
            "type": "mkdir",
            "desc": "创建目录（递归创建父目录）。",
            "required": ["path"],
            "properties": {
                "path": {"type": "string", "desc": "相对路径，例如 src/utils"},
            },
            "example": {"type": "mkdir", "path": "docs"},
        },
        {
            "type": "browser_open",
            "desc": "（可选）Playwright：打开网页。未安装 Playwright 时返回 ok=false。",
            "required": ["url"],
            "properties": {
                "session_id": {"type": "string", "desc": "可选，会话 id（默认 global）"},
                "url": {"type": "string", "desc": "http/https URL"},
                "wait_ms": {"type": "number", "desc": "可选，打开后等待毫秒"},
            },
            "example": {"type": "browser_open", "session_id": "s1", "url": "https://example.com", "wait_ms": 500},
        },
        {
            "type": "browser_click",
            "desc": "（可选）Playwright：点击元素（CSS selector）。未安装 Playwright 时返回 ok=false。",
            "required": ["selector"],
            "properties": {
                "session_id": {"type": "string", "desc": "可选，会话 id（默认 global）"},
                "selector": {"type": "string", "desc": "CSS selector"},
                "wait_ms": {"type": "number", "desc": "可选，点击后等待毫秒"},
            },
            "example": {"type": "browser_click", "session_id": "s1", "selector": "button[type=submit]", "wait_ms": 500},
        },
        {
            "type": "browser_type",
            "desc": "（可选）Playwright：在输入框输入文本。未安装 Playwright 时返回 ok=false。",
            "required": ["selector", "text"],
            "properties": {
                "session_id": {"type": "string", "desc": "可选，会话 id（默认 global）"},
                "selector": {"type": "string", "desc": "CSS selector"},
                "text": {"type": "string", "desc": "输入文本"},
                "clear": {"type": "boolean", "desc": "可选，是否先清空（默认 true）"},
                "wait_ms": {"type": "number", "desc": "可选，输入后等待毫秒"},
            },
            "example": {"type": "browser_type", "session_id": "s1", "selector": "#q", "text": "hello", "clear": True},
        },
        {
            "type": "browser_eval",
            "desc": "（可选）Playwright：执行页面 JS 并返回结果。未安装 Playwright 时返回 ok=false。",
            "required": ["script"],
            "properties": {
                "session_id": {"type": "string", "desc": "可选，会话 id（默认 global）"},
                "script": {"type": "string", "desc": "JS 表达式/函数体（evaluate）"},
            },
            "example": {"type": "browser_eval", "session_id": "s1", "script": "() => document.title"},
        },
        {
            "type": "browser_text",
            "desc": "（可选）Playwright：获取页面内容（selector 则取 textContent，否则返回 html 截断）。未安装 Playwright 时返回 ok=false。",
            "required": [],
            "properties": {
                "session_id": {"type": "string", "desc": "可选，会话 id（默认 global）"},
                "selector": {"type": "string", "desc": "可选，CSS selector"},
            },
            "example": {"type": "browser_text", "session_id": "s1", "selector": "body"},
        },
        {
            "type": "browser_screenshot",
            "desc": "（可选）Playwright：截图（base64 png）。未安装 Playwright 时返回 ok=false。",
            "required": [],
            "properties": {
                "session_id": {"type": "string", "desc": "可选，会话 id（默认 global）"},
                "full_page": {"type": "boolean", "desc": "可选，是否整页（默认 true）"},
            },
            "example": {"type": "browser_screenshot", "session_id": "s1", "full_page": True},
        },
        {
            "type": "browser_wait",
            "desc": "（可选）Playwright：等待毫秒。未安装 Playwright 时返回 ok=false。",
            "required": ["ms"],
            "properties": {
                "session_id": {"type": "string", "desc": "可选，会话 id（默认 global）"},
                "ms": {"type": "number", "desc": "等待毫秒"},
            },
            "example": {"type": "browser_wait", "session_id": "s1", "ms": 1000},
        },
        {
            "type": "browser_close",
            "desc": "（可选）Playwright：关闭会话页面。未安装 Playwright 时返回 ok=false。",
            "required": [],
            "properties": {
                "session_id": {"type": "string", "desc": "可选，会话 id（默认 global）"},
            },
            "example": {"type": "browser_close", "session_id": "s1"},
        },
        {
            "type": "revert_hunk",
            "desc": "撤回某个文件的指定 hunk。",
            "required": ["path", "hunk_index"],
            "properties": {
                "path": {"type": "string", "desc": "相对路径"},
                "hunk_index": {"type": "number", "desc": "hunk 索引（从 0 开始）"},
                "status": {"type": "string", "desc": "可选：M/A/D/R/U（默认 M）"},
            },
            "example": {"type": "revert_hunk", "path": "server.py", "hunk_index": 0, "status": "M"},
        },
        {
            "type": "revert_line",
            "desc": "撤回某个文件的指定行（位于某个 hunk 内）。",
            "required": ["path", "hunk_index", "line_index"],
            "properties": {
                "path": {"type": "string", "desc": "相对路径"},
                "hunk_index": {"type": "number", "desc": "hunk 索引（从 0 开始）"},
                "line_index": {"type": "number", "desc": "行索引（从 0 开始）"},
                "status": {"type": "string", "desc": "可选：M/A/D/R/U（默认 M）"},
            },
            "example": {"type": "revert_line", "path": "server.py", "hunk_index": 0, "line_index": 0, "status": "M"},
        },
        {
            "type": "revert_multi_lines",
            "desc": "撤回某个文件的多行（位于某个 hunk 内）。",
            "required": ["path", "hunk_index", "start_line_index", "end_line_index"],
            "properties": {
                "path": {"type": "string", "desc": "相对路径"},
                "hunk_index": {"type": "number", "desc": "hunk 索引（从 0 开始）"},
                "start_line_index": {"type": "number", "desc": "起始行索引（从 0 开始）"},
                "end_line_index": {"type": "number", "desc": "结束行索引（从 0 开始，包含）"},
                "status": {"type": "string", "desc": "可选：M/A/D/R/U（默认 M）"},
            },
            "example": {"type": "revert_multi_lines", "path": "server.py", "hunk_index": 0, "start_line_index": 0, "end_line_index": 3, "status": "M"},
        },
        {
            "type": "revert_file",
            "desc": "撤回整个文件的修改（恢复为 HEAD 版本或删除未跟踪文件）。",
            "required": ["path"],
            "properties": {
                "path": {"type": "string", "desc": "相对路径"},
                "status": {"type": "string", "desc": "可选：M/A/D/R/U（默认 M）"},
            },
            "example": {"type": "revert_file", "path": "server.py", "status": "M"},
        },
        {
            "type": "revert_all",
            "desc": "撤回多个文件的所有修改（整文件级别）。",
            "required": ["paths"],
            "properties": {
                "paths": {"type": "array", "desc": "相对路径数组"},
            },
            "example": {"type": "revert_all", "paths": ["server.py", "README.md"]},
        },
        {
            "type": "open_file",
            "desc": "在 IDE 中打开文件（必须走 /api/open_file 调度）。支持多种视图：editor=文本视图，split=并排对比（双栏），change=变更视图(单栏)，unified=统一 diff 视图(单栏)。",
            "required": ["name"],
            "properties": {
                "name": {"type": "string", "desc": "文件名或相对路径"},
                "view": {"type": "string", "desc": "可选视图：editor/split/change/unified（默认 editor）"},
            },
            "example": {"type": "open_file", "name": "README.md", "view": "change"},
        },
        {
            "type": "status",
            "desc": "获取当前仓库/工作区状态（包含仓库是否打开、分支、变更概览等）。",
            "required": [],
            "properties": {},
            "example": {"type": "status"},
        },
        {
            "type": "open_repo",
            "desc": "打开/切换仓库（有副作用，会改变当前工作区仓库）。",
            "required": ["path"],
            "properties": {
                "path": {"type": "string", "desc": "仓库路径（绝对路径或可展开的用户路径）"},
            },
            "example": {"type": "open_repo", "path": "C:/path/to/repo"},
        },
        {
            "type": "set_origin",
            "desc": "设置当前仓库的远程 origin 地址（等同 git remote set-url origin <url>）。",
            "required": ["url"],
            "properties": {
                "url": {"type": "string", "desc": "新的远程 git 地址，例如 https://gitee.com/xxx/yyy.git"},
            },
            "example": {"type": "set_origin", "url": "https://gitee.com/gzgdata/ntos-test.git"},
        },
        {
            "type": "workspace_context",
            "desc": "获取工作区上下文摘要（用于让 AI 了解当前工作区、文件树摘要等）。",
            "required": [],
            "properties": {},
            "example": {"type": "workspace_context"},
        },
        {
            "type": "pull_safe",
            "desc": "拉取远程更新（git pull --no-edit），并返回：是否有冲突、是否本地修改会被覆盖、冲突文件列表。",
            "required": [],
            "properties": {},
            "example": {"type": "pull_safe"},
        },
        {
            "type": "stash_and_pull",
            "desc": "暂存本地修改（git stash push -u）后执行更新（git pull --no-edit）。完成后可用 stash_pop 恢复修改。",
            "required": [],
            "properties": {},
            "example": {"type": "stash_and_pull"},
        },
        {
            "type": "commit_and_pull",
            "desc": "提交本地修改（git add -A + git commit）后执行更新（git pull --no-edit）。",
            "required": ["message"],
            "properties": {
                "message": {"type": "string", "desc": "提交信息"},
            },
            "example": {"type": "commit_and_pull", "message": "WIP: save changes before pull"},
        },
        {
            "type": "switch_branch_safe",
            "desc": "切换分支（包含覆盖检测）。若工作区有未提交修改且会被覆盖，则返回 needs_handling/affected_files。",
            "required": ["branch"],
            "properties": {
                "branch": {"type": "string", "desc": "目标分支名（支持 remotes/origin/foo 或 origin/foo 或本地分支名）"},
                "detached": {"type": "boolean", "desc": "可选：远端分支时是否仅切到远端（detached HEAD）"},
            },
            "example": {"type": "switch_branch_safe", "branch": "main"},
        },
        {
            "type": "stash_and_switch",
            "desc": "暂存本地修改（git stash push -u）后切换分支。完成后可用 stash_pop 恢复修改。",
            "required": ["branch"],
            "properties": {
                "branch": {"type": "string", "desc": "目标分支名"},
            },
            "example": {"type": "stash_and_switch", "branch": "main"},
        },
        {
            "type": "commit_and_switch",
            "desc": "提交本地修改（git add -A + git commit）后切换分支。",
            "required": ["branch", "message"],
            "properties": {
                "branch": {"type": "string", "desc": "目标分支名"},
                "message": {"type": "string", "desc": "提交信息"},
            },
            "example": {"type": "commit_and_switch", "branch": "main", "message": "save changes"},
        },
        {
            "type": "list_files",
            "desc": "列出仓库文件变更列表（用于了解哪些文件有修改/新增/删除）。",
            "required": [],
            "properties": {
                "max_age": {"type": "number", "desc": "可选：缓存最大年龄（秒），0 表示强制刷新"},
            },
            "example": {"type": "list_files", "max_age": 0},
        },
        {
            "type": "diff_file",
            "desc": "获取单个文件的 diff（用于审阅变更）。",
            "required": ["path"],
            "properties": {
                "path": {"type": "string", "desc": "相对路径"},
                "status": {"type": "string", "desc": "可选：M/A/D/R/U（默认 M）"},
                "ctx_lines": {"type": "number", "desc": "可选：diff 上下文行数（0-200）"},
            },
            "example": {"type": "diff_file", "path": "server.py", "status": "M", "ctx_lines": 80},
        },
        {
            "type": "get_file",
            "desc": "统一文件接口：文本直接返回 content；媒体/二进制直接返回 data_url（base64，受限大小），不返回本地URL。路径不确定时，先用 find_files 按文件名搜索定位。",
            "required": ["path"],
            "properties": {
                "path": {"type": "string", "desc": "相对路径（路径不确定请先用 find_files 搜索定位）"}
            },
            "example": {"type": "get_file", "path": "README.md"},
        },
        {
            "type": "file_stat",
            "desc": "获取文件或目录的元信息（是否存在、大小、修改时间、类型等）。不返回文件内容，仅用于确认文件状态。",
            "required": ["path"],
            "properties": {
                "path": {"type": "string", "desc": "相对路径"}
            },
            "example": {"type": "file_stat", "path": "README.md"},
        },
        {
            "type": "read_pdf",
            "desc": "读取 PDF 文件内容。将 PDF 文本提取为纯文本返回。支持读取部分页面。",
            "required": ["path"],
            "properties": {
                "path": {"type": "string", "desc": "PDF 文件相对路径"},
                "page_num": {"type": "integer", "desc": "指定页码（1 开始），不指定则读取全部"},
                "max_pages": {"type": "integer", "desc": "最多读取页数（默认 50 页）"},
            },
            "example": {"type": "read_pdf", "path": "docs/report.pdf", "max_pages": 10},
        },
        {
            "type": "read_docx",
            "desc": "读取 Word（.docx）文件内容。将文档文本提取为纯文本返回。",
            "required": ["path"],
            "properties": {
                "path": {"type": "string", "desc": "docx 文件相对路径"},
            },
            "example": {"type": "read_docx", "path": "docs/report.docx"},
        },
        {
            "type": "read_xlsx",
            "desc": "读取 Excel（.xlsx）文件内容。将表格数据提取为文本或 JSON 返回。",
            "required": ["path"],
            "properties": {
                "path": {"type": "string", "desc": "xlsx 文件相对路径"},
                "sheet_name": {"type": "string", "desc": "指定工作表名称，不指定则读取第一个"},
                "max_rows": {"type": "integer", "desc": "最多读取行数（默认 200 行）"},
            },
            "example": {"type": "read_xlsx", "path": "data/sheet.xlsx", "max_rows": 50},
        },
        {
            "type": "undo_stats",
            "desc": "获取 Undo 统计信息（用于判断当前是否可撤销/撤销栈状态）。",
            "required": [],
            "properties": {},
            "example": {"type": "undo_stats"},
        },
        {
            "type": "ai_cache_clear",
            "desc": "清理 AI 缓存（例如模型列表缓存等）。",
            "required": [],
            "properties": {},
            "example": {"type": "ai_cache_clear"},
        },
        {
            "type": "branches",
            "desc": "列出本地/远端分支信息（只读）。",
            "required": [],
            "properties": {},
            "example": {"type": "branches"},
        },
        {
            "type": "commits",
            "desc": "列出最近提交记录（只读）。",
            "required": [],
            "properties": {
                "limit": {"type": "number", "desc": "可选：最大条数"},
            },
            "example": {"type": "commits", "limit": 30},
        },
        {
            "type": "commit_detail",
            "desc": "获取单个提交详情（只读）。",
            "required": ["hash"],
            "properties": {
                "hash": {"type": "string", "desc": "提交 hash（短或长）"},
            },
            "example": {"type": "commit_detail", "hash": "abc1234"},
        },
        {
            "type": "commit_file_diff",
            "desc": "获取某次提交中某个文件的 diff（只读）。",
            "required": ["hash", "path"],
            "properties": {
                "hash": {"type": "string", "desc": "提交 hash"},
                "path": {"type": "string", "desc": "文件相对路径"},
            },
            "example": {"type": "commit_file_diff", "hash": "abc1234", "path": "server.py"},
        },
        {
            "type": "commit_push_status",
            "desc": "查询某次提交的 push 状态（是否已推送/远端是否存在等，只读）。",
            "required": ["hash"],
            "properties": {
                "hash": {"type": "string", "desc": "提交 hash"},
            },
            "example": {"type": "commit_push_status", "hash": "abc1234"},
        },
        {
            "type": "api_request",
            "desc": "调用 /api/* HTTP 接口。**仅在无更具体工具时使用**，优先使用本列表中的专用工具（如 file_content, diff_file 等）。",
            "required": ["method", "path"],
            "properties": {
                "method": {"type": "string", "desc": "GET 或 POST"},
                "path": {"type": "string", "desc": "/api/ 开头的路径，例如 /api/workspace_context"},
                "query": {"type": "object", "desc": "可选，URL 查询参数"},
                "body": {"type": "object", "desc": "可选，POST 请求的 JSON body"},
            },
            "example": {"type": "api_request", "method": "GET", "path": "/api/branches"},
        },
        {
            "type": "undo_last_turn",
            "desc": "撤销上一轮操作（文件修改、重命名、AI 配置变更等）。",
            "required": [],
            "properties": {},
            "example": {"type": "undo_last_turn"},
        },
        {
            "type": "run_cmd",
            "desc": "执行底层终端命令（Windows: CMD/PowerShell；Linux/macOS: Shell）。这是**系统命令执行入口**，不是普通业务功能。**仅当没有任何其他工具能完成任务时使用**。禁止换行、管道、重定向、复合命令（如 `cmd1 && cmd2`）。跨平台要求：生成命令时必须根据当前运行平台选择对应指令；若命令不兼容，必须明确说明当前平台并给出正确平台命令。返回 result.output（stdout+stderr）。",
            "required": ["cmd"],
            "properties": {
                "cmd": {"type": "string", "desc": "单行命令，例如 'git log -1'"},
                "timeout": {"type": "number", "desc": "超时秒数，默认 30，最大 600"},
                "cwd": {"type": "string", "desc": "可选，工作目录（相对仓库根目录）"},
            },
            "example": {"type": "run_cmd", "cmd": "git status", "timeout": 30},
        },
        {
            "type": "list_processes",
            "desc": "列出系统进程和端口占用情况。支持按名称过滤进程或按端口查找占用进程。",
            "required": [],
            "properties": {
                "filter": {"type": "string", "desc": "可选，按进程名或端口号过滤"},
                "limit": {"type": "integer", "desc": "最多返回条数（默认 50）"},
            },
            "example": {"type": "list_processes", "filter": "python"},
        },
        {
            "type": "env_info",
            "desc": "获取当前运行环境信息（Python版本、操作系统、已安装的关键依赖包版本等）。",
            "required": [],
            "properties": {
                "type": {"type": "string", "desc": "可选，all/python/system/packages，默认 all"},
            },
            "example": {"type": "env_info", "type": "all"},
        },
        {
            "type": "web_search",
            "desc": "联网搜索（RAG 检索入口）。当问题需要最新信息/外部资料时先调用此工具；基于返回的 items 生成答案并给出 URL 引用。若不可用/失败需降级：明确说明无法联网并改用本地上下文回答。",
            "required": ["query"],
            "properties": {
                "query": {"type": "string", "desc": "搜索关键词/问题"},
                "top_k": {"type": "number", "desc": "可选，返回条数 1-10（默认 5）"},
                "timeout": {"type": "number", "desc": "可选，超时秒数（默认 20）"},
            },
            "example": {"type": "web_search", "query": "Python 3.13 release date", "top_k": 5},
        },
        {
            "type": "web_fetch",
            "desc": "抓取网页正文（与 web_search 配合的 RAG 阅读器）。只用于从已知 URL 获取内容；返回纯文本（会做基础去标签，且有长度上限）。不可用时必须降级说明。",
            "required": ["url"],
            "properties": {
                "url": {"type": "string", "desc": "http/https URL"},
                "timeout": {"type": "number", "desc": "可选，超时秒数（默认 20）"},
            },
            "example": {"type": "web_fetch", "url": "https://example.com"},
        },
        {
            "type": "search_code",
            "desc": "在仓库代码中搜索（只读）。",
            "required": ["query"],
            "properties": {
                "query": {"type": "string", "desc": "搜索关键词/正则"},
                "case_sensitive": {"type": "boolean", "desc": "可选"},
                "max_results": {"type": "number", "desc": "可选"},
            },
            "example": {"type": "search_code", "query": "ai_chat", "case_sensitive": False, "max_results": 50},
        },
        {
            "type": "read_file_range",
            "desc": "读取文件片段。",
            "required": ["path", "start", "end"],
            "properties": {
                "path": {"type": "string", "desc": "相对路径"},
                "start": {"type": "number", "desc": "起始行(1-indexed)"},
                "end": {"type": "number", "desc": "结束行(1-indexed)"},
            },
            "example": {"type": "read_file_range", "path": "server.py", "start": 1, "end": 120},
        },
        {
            "type": "list_dir_tree",
            "desc": "列目录树。",
            "required": ["path"],
            "properties": {
                "path": {"type": "string", "desc": "相对路径或空"},
                "depth": {"type": "number", "desc": "可选"},
                "max_entries": {"type": "number", "desc": "可选"},
            },
            "example": {"type": "list_dir_tree", "path": "", "depth": 3, "max_entries": 500},
        },
        {
            "type": "verify_python",
            "desc": "编译校验 Python 文件。",
            "required": ["paths"],
            "properties": {
                "paths": {"type": "array", "desc": "相对路径列表"},
            },
            "example": {"type": "verify_python", "paths": ["server.py"]},
        },
        {
            "type": "stage_file",
            "desc": "Git 暂存单个文件。",
            "required": ["path"],
            "properties": {
                "path": {"type": "string", "desc": "相对路径"},
            },
            "example": {"type": "stage_file", "path": "README.md"},
        },
        {
            "type": "staged_files",
            "desc": "列出暂存区（index）中的文件列表（git diff --cached --name-status）。",
            "required": [],
            "properties": {},
            "example": {"type": "staged_files"},
        },
        {
            "type": "unstage_file",
            "desc": "将单个文件从暂存区恢复到工作区（git restore --staged）。",
            "required": ["path"],
            "properties": {
                "path": {"type": "string", "desc": "相对路径"},
            },
            "example": {"type": "unstage_file", "path": "README.md"},
        },
        {
            "type": "unstage_all_staged",
            "desc": "取消全部暂存（保留工作区改动）。",
            "required": [],
            "properties": {},
            "example": {"type": "unstage_all_staged"},
        },
        {
            "type": "discard_staged_file",
            "desc": "丢弃单个文件在暂存区和工作区的改动，恢复到 HEAD（危险）。",
            "required": ["path"],
            "properties": {
                "path": {"type": "string", "desc": "相对路径"},
            },
            "example": {"type": "discard_staged_file", "path": "README.md"},
        },
        {
            "type": "discard_all_staged",
            "desc": "丢弃全部暂存区内容，并同步丢弃相关工作区改动，恢复到 HEAD（危险）。",
            "required": [],
            "properties": {},
            "example": {"type": "discard_all_staged"},
        },
        {
            "type": "commit",
            "desc": "Git 提交。",
            "required": ["message"],
            "properties": {
                "message": {"type": "string", "desc": "提交信息"},
                "files": {"type": "array", "desc": "可选：要提交的文件列表（相对路径）"},
            },
            "example": {"type": "commit", "message": "fix: ...", "files": ["README.md"]},
        },
        {
            "type": "pull",
            "desc": "Git pull 拉取远端。",
            "required": [],
            "properties": {},
            "example": {"type": "pull"},
        },
        {
            "type": "push",
            "desc": "Git push 推送远端。",
            "required": [],
            "properties": {},
            "example": {"type": "push"},
        },
        {
            "type": "switch_branch",
            "desc": "切换分支。",
            "required": ["branch"],
            "properties": {
                "branch": {"type": "string", "desc": "分支名"},
            },
            "example": {"type": "switch_branch", "branch": "main"},
        },
        {
            "type": "stash_pop",
            "desc": "应用 stash（stash pop）。",
            "required": [],
            "properties": {
                "ref": {"type": "string", "desc": "可选：指定 stash 引用，例如 stash@{0}（默认最新）"},
            },
            "example": {"type": "stash_pop", "ref": "stash@{0}"},
        },
        {
            "type": "stash_list",
            "desc": "列出 stash 列表（git stash list）。",
            "required": [],
            "properties": {},
            "example": {"type": "stash_list"},
        },
        {
            "type": "stash_drop",
            "desc": "丢弃指定 stash（git stash drop <ref>）。",
            "required": ["ref"],
            "properties": {
                "ref": {"type": "string", "desc": "stash 引用，例如 stash@{0}"},
            },
            "example": {"type": "stash_drop", "ref": "stash@{0}"},
        },
        {
            "type": "fetch_remotes",
            "desc": "拉取最新远端引用（git fetch --all --prune），用于刷新远端分支列表。",
            "required": [],
            "properties": {},
            "example": {"type": "fetch_remotes"},
        },
    ]

    try:
        for t in tool_extra:
            if isinstance(t, dict) and str(t.get("type") or "").strip():
                agent_tools.append(t)
    except Exception as e:
        logger.debug(f"Exception ignored: {e}")

    try:
        dis_ep_set = set(str(x) for x in dis_ep if x is not None and str(x).strip())
        if dis_ep_set:
            endpoints = [e for e in endpoints if (str(e.get("path") or "") not in dis_ep_set)]
    except Exception as e:
        logger.debug(f"Exception ignored: {e}")

    try:
        dis_tool_set = set(str(x) for x in dis_tool if x is not None and str(x).strip())
        if dis_tool_set:
            agent_tools = [t for t in agent_tools if (str(t.get("type") or "") not in dis_tool_set)]
    except Exception as e:
        logger.debug(f"Exception ignored: {e}")

    lines = []
    lines.append("系统接口索引（后端单一真源，自动生成）：")
    lines.append("")
    lines.append("Action 联动卡片（Browser Flow）协议：")
    lines.append("- 你可以在回复中输出一个 language=json 的代码块，前端会自动识别并渲染为可执行卡片（Browser Flow）。")
    lines.append("- 代码块内容必须是严格 JSON（不要写注释/尾逗号/多余文字）。")
    lines.append("- 触发条件：必须使用 Markdown 围栏代码块 ```json ... ```（小写 json），并且代码块内只放 JSON。")
    lines.append("- 两种格式均可：")
    lines.append("  1) [{\"action\":\"open\",\"params\":{...}}, ...]")
    lines.append("  2) {\"session_id\":\"s1\",\"steps\":[{\"action\":\"open\",\"params\":{...}}, ...]}")
    lines.append("- action 允许值：open/click/type/eval/text/screenshot/wait/close（会映射到 /api/browser_tool）。")
    lines.append("- 最小可用示例（可直接照抄）：")
    lines.append("  ```json")
    lines.append("  [")
    lines.append("    {\"action\":\"open\",\"params\":{\"url\":\"https://example.com\"}},")
    lines.append("    {\"action\":\"screenshot\",\"params\":{}}")
    lines.append("  ]")
    lines.append("  ```")
    lines.append("")
    lines.append("重要约束：")
    lines.append("- 优先使用 /api/* 完成业务 CRUD；只有接口确实无法覆盖时才使用 /api/run_cmd")
    lines.append("- run_cmd 是底层终端命令执行入口（不是普通业务功能），必须是单行命令，不允许换行/管道/重定向/复合命令")
    lines.append("- 跨平台：生成命令必须匹配当前运行平台（Windows vs Linux/macOS）；不兼容时必须说明平台并提供正确命令")
    lines.append("- 工具选择：只允许从下方 Agent Tools 列表中选择；输出 JSON 时严格遵守 required 字段与类型")
    lines.append("- 参数不确定：先澄清/补全（例如 path 不明确、缺少 content），不要瞎猜")
    lines.append("")
    lines.append("Agent Tools（前端可执行的工具 JSON schema 摘要）：")
    for t in agent_tools:
        tt = str(t.get("type") or "").strip()
        if not tt:
            continue
        req = t.get("required") or []
        req_s = ",".join([str(x) for x in req]) if isinstance(req, list) else ""
        lines.append(f"- {tt} required=[{req_s}]")
    lines.append("")
    for e in endpoints:
        method = str(e.get("method") or "").upper()
        path = str(e.get("path") or "")
        if not method or not path:
            continue
        lines.append(f"- {method} {path}")
    text = "\n".join(lines)

    return {
        "version": "1",
        "generated_at": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
        "endpoints": endpoints,
        "agent_tools": agent_tools,
        "text": text,
    }


_AI_SSL_CONTEXT = None
try:
    _AI_SSL_CONTEXT = ssl.create_default_context()
    try:
        _AI_SSL_CONTEXT.minimum_version = ssl.TLSVersion.TLSv1_2
    except Exception as e:
        logger.debug(f"Exception ignored: {e}")
except Exception:
    _AI_SSL_CONTEXT = None


def _ai_pick_profile(cfg: dict, profile_id: str | None):
    if not isinstance(cfg, dict):
        return None
    profiles = cfg.get("profiles")
    if not isinstance(profiles, list):
        return None

    pid = (profile_id or "").strip()
    if pid:
        for p0 in profiles:
            if isinstance(p0, dict) and str(p0.get("id") or "").strip() == pid:
                return p0

    active = str(cfg.get("active_profile_id") or "").strip()
    if active:
        for p0 in profiles:
            if isinstance(p0, dict) and str(p0.get("id") or "").strip() == active:
                return p0

    for p0 in profiles:
        if isinstance(p0, dict):
            return p0
    return None


def _ai_estimate_text_tokens(text: str):
    if isinstance(text, list):
        try:
            s2 = ""
            for p in text:
                if not isinstance(p, dict):
                    continue
                if str(p.get("type") or "") == "text":
                    s2 += str(p.get("text") or "")
            text = s2
        except Exception:
            text = ""
    s = str(text or "")
    if not s:
        return 0
    try:
        cjk = 0      # 中日韩文字
        digits = 0    # 数字
        punct = 0     # 标点符号
        whitespace = 0
        other = 0
        for ch in s:
            o = ord(ch)
            if (
                (0x4E00 <= o <= 0x9FFF)
                or (0x3400 <= o <= 0x4DBF)
                or (0x20000 <= o <= 0x2A6DF)
                or (0x2A700 <= o <= 0x2B73F)
                or (0x2B740 <= o <= 0x2B81F)
                or (0x2B820 <= o <= 0x2CEAF)
                or (0xF900 <= o <= 0xFAFF)
                or (0x2F800 <= o <= 0x2FA1F)
                or (0x3040 <= o <= 0x30FF)   # 日文假名
                or (0xAC00 <= o <= 0xD7AF)   # 韩文
            ):
                cjk += 1
            elif ch.isdigit():
                digits += 1
            elif ch.isalpha() or ch == '_':
                other += 1
            elif ch.isspace():
                whitespace += 1
            else:
                punct += 1
        # CJK: ~1 token/字; 英文字母: ~4字符/token; 数字: ~3字符/token; 标点: ~2字符/token
        return int(cjk + digits / 3.0 + other / 4.0 + punct / 2.0 + whitespace / 4.0)
    except Exception:
        return max(1, int(len(s) / 4))


def _ai_estimate_messages_tokens(messages: list):
    try:
        total = 0
        arr = messages if isinstance(messages, list) else []
        for m in arr:
            if not isinstance(m, dict):
                continue
            total += 4
            total += _ai_estimate_text_tokens(m.get("role"))
            c = m.get("content")
            if isinstance(c, str):
                total += _ai_estimate_text_tokens(c)
            elif isinstance(c, list):
                # OpenAI 多模态格式：数组 content
                for p in c:
                    if not isinstance(p, dict):
                        continue
                    pt = str(p.get("type") or "")
                    if pt == "text":
                        total += _ai_estimate_text_tokens(p.get("text"))
                    elif pt in ("image_url", "image", "input_image"):
                        # 图片 token 估算（base64 不计入文本 token，固定开销）
                        total += 85  # GPT-4o 每张图片约 85-170 tokens
            # 估算 tool_calls 的 token 开销
            role = str(m.get("role") or "")
            if role == "assistant":
                tc_list = m.get("tool_calls")
                if isinstance(tc_list, list):
                    for tc in tc_list:
                        if not isinstance(tc, dict):
                            continue
                        total += _ai_estimate_text_tokens(tc.get("id"))
                        total += _ai_estimate_text_tokens(tc.get("type"))
                        func = tc.get("function") if isinstance(tc.get("function"), dict) else {}
                        total += _ai_estimate_text_tokens(func.get("name"))
                        total += _ai_estimate_text_tokens(func.get("arguments"))
            elif role == "tool":
                total += _ai_estimate_text_tokens(m.get("tool_call_id"))
        total += 8
        return int(total)
    except Exception:
        return 0


def _ai_truncate_text_to_tokens(text: str, keep_tokens: int):
    try:
        s = str(text or "")
        t = int(keep_tokens or 0)
        if t <= 0:
            return ""
        if _ai_estimate_text_tokens(s) <= t:
            return s
        try:
            max_chars = max(0, int(t * 4))
            if len(s) <= max_chars:
                return s
            return s[-max_chars:]
        except Exception:
            return s
    except Exception:
        return str(text or "")


def _hivo_normalize_tool_data(obj):
    """递归地把工具返回数据中的反斜杠路径替换为正斜杠，避免 markdown-it 转义破坏。"""
    try:
        if isinstance(obj, str):
            # 不处理 URL 协议（http://、https://、data: 等）
            if re.match(r'^[a-zA-Z]+://', obj):
                return obj
            return obj.replace('\\', '/')
        if isinstance(obj, list):
            return [_hivo_normalize_tool_data(item) for item in obj]
        if isinstance(obj, dict):
            return {k: _hivo_normalize_tool_data(v) for k, v in obj.items()}
        return obj
    except Exception:
        return obj


def _hivo_ai_msg_has_valid_content(m):
    """判断消息是否有有效内容（支持字符串和多模态数组两种 content 格式）。"""
    try:
        if not isinstance(m, dict):
            return False
        c = m.get("content")
        if isinstance(c, str):
            return bool(c.strip())
        if isinstance(c, list):
            for p in c:
                if isinstance(p, dict):
                    pt = str(p.get("type") or "")
                    if pt == "text" and str(p.get("text") or "").strip():
                        return True
                    if pt in ("image_url", "image", "input_image"):
                        return True
            return False
        return False
    except Exception:
        return False


def _ai_trim_messages_to_budget(messages: list, max_total_tokens: int, reserve_output_tokens: int = 0):
    try:
        if not isinstance(messages, list):
            return []
        max_total = int(max_total_tokens or 0)
        reserve = max(0, int(reserve_output_tokens or 0))
        if max_total <= 0:
            return [m for m in messages if isinstance(m, dict)]

        arr = [m for m in messages if isinstance(m, dict)]
        if not arr:
            return []

        def is_dup(a, b):
            if not a or not b:
                return False
            if str(a.get("role") or "") != str(b.get("role") or ""):
                return False
            ca = str(a.get("content") or "").strip()
            cb = str(b.get("content") or "").strip()
            if not ca or not cb:
                return False
            if ca == cb:
                return True
            if len(ca) > 20 and len(cb) > 20 and ca[:80] == cb[:80]:
                return True
            return False

        cleaned = []
        last = None
        for m in arr:
            r = str(m.get("role") or "")
            if r not in ("system", "user", "assistant", "tool"):
                continue
            # tool 消息和带 tool_calls 的 assistant 消息即使 content 为空也要保留（工具调用链路完整性）
            if r == "tool":
                pass
            elif r == "assistant":
                tc = m.get("tool_calls")
                has_tc = isinstance(tc, list) and len(tc) > 0
                if not has_tc and not _hivo_ai_msg_has_valid_content(m):
                    continue
            else:
                if not _hivo_ai_msg_has_valid_content(m):
                    continue
            if is_dup(last, m):
                continue
            if r == "assistant" and _hivo_ai_msg_has_valid_content(m) and str(m.get("content") or "").strip() in ("对话内容过长，已超过限制",):
                continue
            cleaned.append(m)
            last = m

        sys_msgs = [m for m in cleaned if str(m.get("role") or "") == "system"]
        non_sys = [m for m in cleaned if str(m.get("role") or "") != "system"]

        sys_text = "\n\n".join([
            str(m.get("content") or "").strip()
            for m in sys_msgs
            if isinstance(m.get("content"), str)
        ])
        base = []
        if sys_text:
            base.append({"role": "system", "content": sys_text})

        def total_tokens(msgs):
            return _ai_estimate_messages_tokens(msgs)

        budget = max_total - reserve
        if total_tokens(base + non_sys) <= budget:
            return base + non_sys

        turns = []
        i = 0
        while i < len(non_sys):
            r = str(non_sys[i].get("role") or "")
            if r == "user":
                turn_msgs = [non_sys[i]]
                i += 1
                if i < len(non_sys) and str(non_sys[i].get("role") or "") == "assistant":
                    turn_msgs.append(non_sys[i])
                    i += 1
                    while i < len(non_sys) and str(non_sys[i].get("role") or "") == "tool":
                        turn_msgs.append(non_sys[i])
                        i += 1
                turns.append(tuple(turn_msgs))
            elif r == "assistant":
                turn_msgs = [non_sys[i]]
                i += 1
                while i < len(non_sys) and str(non_sys[i].get("role") or "") == "tool":
                    turn_msgs.append(non_sys[i])
                    i += 1
                turns.append(tuple(turn_msgs))
            else:
                turns.append((non_sys[i],))
                i += 1

        recent = []
        used = total_tokens(base)
        for t in reversed(turns):
            cand = list(t) + recent
            if used + total_tokens(cand) <= budget:
                recent = cand
                used = used + total_tokens(t)
            else:
                break

        # 找到 recent 对应的起始位置，确保不会从中间截断一个 turn
        recent_start_idx = 0
        if recent:
            first_recent_msg = recent[0]
            for idx, m in enumerate(non_sys):
                if m is first_recent_msg:
                    recent_start_idx = idx
                    break
        older_msgs = non_sys[:recent_start_idx] if recent_start_idx > 0 else []
        has_older = len(older_msgs) > 0

        summary_txt = ""
        if has_older:
            parts = []
            keep_turns = 6
            turn_count = 0
            last_tc_map = {}
            for m in older_msgs:
                role = str(m.get("role") or "")
                if role == "user":
                    turn_count += 1
                    if turn_count > keep_turns:
                        break
                if role == "assistant":
                    tc = m.get("tool_calls")
                    if isinstance(tc, list):
                        for t in tc:
                            if isinstance(t, dict):
                                tid = str(t.get("id") or "")
                                tname = str((t.get("function") or {}).get("name") or "")
                                if tid and tname:
                                    last_tc_map[tid] = tname
                c = m.get("content")
                if isinstance(c, str):
                    txt = c
                elif isinstance(c, list):
                    txt_parts = []
                    for p in c:
                        if isinstance(p, dict) and str(p.get("type") or "") == "text":
                            txt_parts.append(str(p.get("text") or ""))
                    txt = " ".join(txt_parts)
                else:
                    txt = ""
                if role == "user":
                    parts.append("U:" + txt[:160])
                elif role == "assistant":
                    tc = m.get("tool_calls")
                    has_tc = isinstance(tc, list) and len(tc) > 0
                    if has_tc:
                        tnames = [str((t.get("function") or {}).get("name") or "?") for t in tc if isinstance(t, dict)]
                        parts.append("A:[tool:" + ",".join(tnames[:3]) + "] " + txt[:120])
                    else:
                        parts.append("A:" + txt[:220])
                elif role == "tool":
                    tcid = str(m.get("tool_call_id") or "")
                    tname = last_tc_map.get(tcid) or "tool"
                    parts.append("T(" + tname + "):" + txt[:80])
                else:
                    parts.append(txt[:120])
            summary_blob = "\n".join(parts)
            sum_allow = max(64, int(budget * 0.25))
            summary_txt = _ai_truncate_text_to_tokens(summary_blob, sum_allow)
            summary_msg = {"role": "system", "content": "历史摘要：\n" + summary_txt}
            composed = base + [summary_msg] + recent
        else:
            composed = base + recent

        if total_tokens(composed) <= budget:
            return composed

        if base:
            sys_allow = max(32, int(budget * 0.25))
            base = [{"role": "system", "content": _ai_truncate_text_to_tokens(base[0].get("content"), sys_allow)}]
        else:
            base = []

        if has_older:
            sum_allow2 = max(32, int(budget * 0.2))
            composed = base + [{"role": "system", "content": _ai_truncate_text_to_tokens("历史摘要：\n" + summary_txt, sum_allow2)}] + recent
        else:
            composed = base + recent

        # 从前往后按 turn 粒度裁剪，确保不会从中间截断 assistant + tool 消息对
        while total_tokens(composed) > budget and len(recent) > 1:
            # 找到下一个 turn 的起始位置（user 或 assistant 消息）
            next_start = 1
            while next_start < len(recent):
                r = str(recent[next_start].get("role") or "")
                if r in ("user", "assistant"):
                    break
                next_start += 1
            if next_start >= len(recent):
                break
            recent = recent[next_start:]
            if has_older:
                composed = base + [{"role": "system", "content": _ai_truncate_text_to_tokens("历史摘要：\n" + summary_txt, sum_allow2)}] + recent
            else:
                composed = base + recent

        if total_tokens(composed) <= budget:
            return composed

        first_sys = base[0] if base else None
        last_user = None
        for m in reversed(non_sys):
            if str(m.get("role") or "") == "user":
                last_user = m
                break
        final_out = []
        if first_sys and isinstance(first_sys.get("content"), str):
            sys_allow3 = max(24, int(budget * 0.3))
            final_out.append({"role": "system", "content": _ai_truncate_text_to_tokens(first_sys.get("content"), sys_allow3)})
        if last_user:
            lu_content = last_user.get("content")
            remain3 = max(24, budget - _ai_estimate_messages_tokens(final_out) - 8)
            if isinstance(lu_content, str):
                final_out.append({"role": "user", "content": _ai_truncate_text_to_tokens(lu_content, remain3)})
            elif isinstance(lu_content, list):
                # 多模态消息：保留 text 部分和尽量多的 image_url
                parts = []
                for p in lu_content:
                    if not isinstance(p, dict):
                        continue
                    pt = str(p.get("type") or "")
                    if pt == "text":
                        parts.append({"type": "text", "text": _ai_truncate_text_to_tokens(str(p.get("text") or ""), max(24, remain3 - 100))})
                    elif pt in ("image_url", "image", "input_image") and remain3 > 100:
                        parts.append(p)
                        remain3 -= 85
                if parts:
                    final_out.append({"role": "user", "content": parts})

        # 简化超大的 tool 消息，避免 base64 图片等大内容撑爆上下文
        if final_out and max_total_tokens > 0:
            tool_content_limit = max(512, int(max_total_tokens * 0.05))
            for m in final_out:
                if not isinstance(m, dict):
                    continue
                if str(m.get("role") or "") != "tool":
                    continue
                c = m.get("content")
                if not isinstance(c, str):
                    continue
                # 粗略估算：长度 / 4 ≈ token 数
                est_tokens = len(c) // 4
                if est_tokens <= tool_content_limit:
                    continue
                # 检测是否是 base64 图片
                is_b64_img = c.startswith("data:image") or (len(c) > 100 and all(ch in "ABCDEFGHIJKLMNOPQRSTUVWXYZabcdefghijklmnopqrstuvwxyz0123456789+/=" for ch in c[:80]))
                if is_b64_img:
                    m["content"] = "[image content omitted - too large to fit in context window]"
                else:
                    # 长文本：截断前 200 字符 + 摘要
                    head = c[:200]
                    m["content"] = f"[truncated - content too large ({est_tokens} tokens)]\n{head}..."

        if final_out:
            return final_out
        return None
    except Exception:
        return messages if isinstance(messages, list) else []


# ════════════════════════════════════════════════════════
#  多模型 Provider 适配层
# ════════════════════════════════════════════════════════

class BaseLLMProvider:
    """LLM Provider 抽象基类。

    子类负责：
    - 将内部统一格式（OpenAI 风格 messages/tools）转换为对应提供商的请求格式
    - 将提供商的响应转换为内部统一格式

    统一内部响应格式:
    {
        "content": str,               # 文本内容
        "tool_calls": [               # 可选：工具调用列表（原生 function calling）
            {"id": str, "type": "function", "function": {"name": str, "arguments": str}}
        ],
        "finish_reason": str,         # stop / tool_calls / length / ...
        "raw": dict,                  # 原始响应（用于调试）
    }
    """

    def __init__(self, profile: dict):
        self.profile = profile or {}
        self.base_url = str(self.profile.get("base_url") or self.profile.get("endpoint") or "").strip()
        self.api_key = str(self.profile.get("api_key") or "")
        self.model = str(self.profile.get("model") or "").strip()

    def build_request(self, messages, tools=None, temperature=None, max_tokens=None, stream=False):
        """构建请求。返回 (url, headers, body)。"""
        raise NotImplementedError

    def parse_response(self, raw_body: str) -> dict:
        """解析非流式响应。返回统一格式。"""
        raise NotImplementedError

    def parse_stream_chunk(self, chunk_text: str) -> dict:
        """解析流式单个 chunk。返回统一格式的增量部分（content/delta），或 None 表示跳过。"""
        return None

    def stream_is_done(self, chunk_text: str) -> bool:
        """判断流式是否结束。"""
        return chunk_text.strip() == "[DONE]"

    def create_stream_state(self) -> dict:
        """创建流式累积状态对象。"""
        return {
            "content_parts": [],
            "reasoning_parts": [],
            "finish_reason": "",
            "tool_calls": {},
        }

    def accumulate_chunk(self, state: dict, chunk_text: str) -> dict:
        """累积一个流式 chunk，返回当前增量（用于前端显示）。

        返回格式: {"content_delta": str, "reasoning_delta": str}
        """
        parsed = self.parse_stream_chunk(chunk_text)
        if not parsed:
            return {}

        content_delta = parsed.get("content_delta")
        reasoning_delta = parsed.get("reasoning_delta")

        if content_delta:
            state["content_parts"].append(str(content_delta))
        if reasoning_delta:
            state["reasoning_parts"].append(str(reasoning_delta))

        fr = parsed.get("finish_reason")
        if fr:
            state["finish_reason"] = fr

        tc_delta = parsed.get("tool_calls_delta")
        if isinstance(tc_delta, list):
            for d in tc_delta:
                if not isinstance(d, dict):
                    continue
                idx = d.get("index")
                if idx is None:
                    continue
                idx = int(idx)
                if idx not in state["tool_calls"]:
                    state["tool_calls"][idx] = {
                        "id": "", "type": "function",
                        "function": {"name": "", "arguments": ""}
                    }
                acc = state["tool_calls"][idx]
                # 兼容多种 id 字段名（id/call_id/tool_call_id）
                delta_id = ""
                for k in ("id", "call_id", "tool_call_id"):
                    if d.get(k):
                        delta_id = str(d[k])
                        break
                if delta_id:
                    acc["id"] = delta_id
                if d.get("type"):
                    acc["type"] = str(d["type"])
                func_d = d.get("function") or {}
                if isinstance(func_d, dict):
                    if func_d.get("name"):
                        acc["function"]["name"] = str(func_d["name"])
                    if func_d.get("arguments"):
                        acc["function"]["arguments"] += str(func_d["arguments"])

        return {
            "content_delta": content_delta,
            "reasoning_delta": reasoning_delta,
        }

    def get_stream_result(self, state: dict) -> dict:
        """获取流式最终结果，返回统一格式。"""
        content = "".join(state["content_parts"])
        reasoning_content = "".join(state["reasoning_parts"])
        if not content:
            content = reasoning_content

        tool_calls_list = []
        if state["tool_calls"]:
            for idx in sorted(state["tool_calls"].keys()):
                tool_calls_list.append(state["tool_calls"][idx])

        result = {
            "content": content,
            "reasoning_content": reasoning_content,
            "raw": {"stream": True},
            "finish_reason": state["finish_reason"],
        }
        if tool_calls_list:
            result["tool_calls"] = tool_calls_list
        return result


class OpenAIProvider(BaseLLMProvider):
    """OpenAI 兼容格式（默认）。适用于 OpenAI、Azure OpenAI、vLLM、Ollama 等。"""

    PROVIDER_ID = "openai"

    def build_request(self, messages, tools=None, temperature=None, max_tokens=None, stream=False):
        url = _ai_build_chat_url(self.base_url)
        headers = {
            "Content-Type": "application/json",
            "User-Agent": "GitManager/1.0",
        }
        if self.api_key:
            headers["Authorization"] = f"Bearer {self.api_key}"
            headers["x-api-key"] = self.api_key

        body = {
            "model": self.model,
            "messages": messages,
        }
        if temperature is not None:
            body["temperature"] = float(temperature)
        if max_tokens and max_tokens > 0:
            body["max_tokens"] = max_tokens
            body["max_completion_tokens"] = max_tokens
        if stream:
            body["stream"] = True
        if tools:
            body["tools"] = tools
            body["tool_choice"] = "auto"

        return url, headers, body

    def parse_response(self, raw_body: str) -> dict:
        data = json.loads(raw_body or "{}")
        msg = (data.get("choices") or [{}])[0].get("message", {})
        content = _ai_extract_content(msg)
        finish_reason = str((data.get("choices") or [{}])[0].get("finish_reason") or "")

        reasoning_content = ""
        try:
            rc = msg.get("reasoning_content") or msg.get("reasoning")
            if isinstance(rc, str):
                reasoning_content = rc
            elif isinstance(rc, list):
                reasoning_content = "".join(
                    str(it.get("text") or "") if isinstance(it, dict) else str(it or "")
                    for it in rc
                )
        except Exception:
            reasoning_content = ""

        tool_calls = msg.get("tool_calls")
        result = {
            "content": content,
            "reasoning_content": reasoning_content,
            "finish_reason": finish_reason,
            "raw": data,
        }
        if isinstance(tool_calls, list) and tool_calls:
            result["tool_calls"] = tool_calls
        return result

    def parse_stream_chunk(self, chunk_text: str) -> dict:
        try:
            j = json.loads(chunk_text)
        except Exception:
            return None
        choice = (j.get("choices") or [{}])[0]
        delta = (choice.get("delta") or {})

        delta_content = delta.get("content")
        if not delta_content:
            delta_content = delta.get("text")
        if isinstance(delta_content, list):
            delta_content = "".join(
                str(it.get("text") or "") if isinstance(it, dict) else str(it or "")
                for it in delta_content
            )

        delta_reasoning = delta.get("reasoning") or delta.get("reasoning_content")
        if not delta_reasoning:
            rd = delta.get("reasoning_details")
            if isinstance(rd, list):
                delta_reasoning = "".join(str(it.get("text") or "") for it in rd if isinstance(it, dict))

        result = {}
        if delta_content:
            result["content_delta"] = str(delta_content)
        if delta_reasoning:
            result["reasoning_delta"] = str(delta_reasoning)

        fr = str(choice.get("finish_reason") or "")
        if fr:
            result["finish_reason"] = fr

        tc_delta = delta.get("tool_calls")
        if tc_delta is not None:
            result["tool_calls_delta"] = tc_delta

        return result


class AnthropicProvider(BaseLLMProvider):
    """Anthropic Claude 原生 API 格式。"""

    PROVIDER_ID = "anthropic"

    def build_request(self, messages, tools=None, temperature=None, max_tokens=None, stream=False):
        base = self.base_url.rstrip("/")
        if not base.endswith("/v1"):
            base = base + "/v1"
        url = base + "/messages"

        headers = {
            "Content-Type": "application/json",
            "x-api-key": self.api_key,
            "anthropic-version": "2023-06-01",
        }

        system_prompt = ""
        user_messages = []
        for m in messages:
            if not isinstance(m, dict):
                continue
            role = str(m.get("role") or "")
            content = m.get("content")
            if role == "system":
                if isinstance(content, str):
                    if system_prompt:
                        system_prompt += "\n\n" + content
                    else:
                        system_prompt = content
            elif role in ("user", "assistant"):
                user_messages.append({"role": role, "content": content})

        body = {
            "model": self.model,
            "messages": user_messages,
        }
        if system_prompt:
            body["system"] = system_prompt
        if temperature is not None:
            body["temperature"] = float(temperature)
        if max_tokens and max_tokens > 0:
            body["max_tokens"] = max_tokens
        if stream:
            body["stream"] = True
        if tools:
            anthropic_tools = []
            for t in tools:
                if not isinstance(t, dict):
                    continue
                func = t.get("function", {}) if t.get("type") == "function" else t
                name = str(func.get("name") or "")
                desc = str(func.get("description") or "")
                params = func.get("parameters", {})
                if not name:
                    continue
                anthropic_tools.append({"name": name, "description": desc, "input_schema": params})
            if anthropic_tools:
                body["tools"] = anthropic_tools

        return url, headers, body

    def parse_response(self, raw_body: str) -> dict:
        data = json.loads(raw_body or "{}")
        content = data.get("content")
        text_parts = []
        reasoning_parts = []
        tool_calls = []
        if isinstance(content, list):
            for idx, block in enumerate(content):
                if not isinstance(block, dict):
                    continue
                btype = str(block.get("type") or "")
                if btype == "text":
                    text_parts.append(str(block.get("text") or ""))
                elif btype == "thinking":
                    reasoning_parts.append(str(block.get("thinking") or block.get("text") or ""))
                elif btype == "tool_use":
                    tool_calls.append({
                        "id": str(block.get("id") or f"toolu_{idx}"),
                        "type": "function",
                        "function": {
                            "name": str(block.get("name") or ""),
                            "arguments": json.dumps(block.get("input", {}), ensure_ascii=False),
                        },
                    })

        finish_reason = str(data.get("stop_reason") or "")
        if finish_reason == "tool_use":
            finish_reason = "tool_calls"

        result = {
            "content": "".join(text_parts).strip(),
            "reasoning_content": "".join(reasoning_parts),
            "finish_reason": finish_reason,
            "raw": data,
        }
        if tool_calls:
            result["tool_calls"] = tool_calls
        return result

    def parse_stream_chunk(self, chunk_text: str) -> dict:
        try:
            j = json.loads(chunk_text)
        except Exception:
            return None
        etype = str(j.get("type") or "")
        result = {}

        if etype == "content_block_delta":
            delta = j.get("delta", {})
            dtype = str(delta.get("type") or "")
            if dtype == "text_delta":
                text = str(delta.get("text") or "")
                if text:
                    result["content_delta"] = text
            elif dtype == "input_json_delta":
                partial = str(delta.get("partial_json") or "")
                if partial:
                    result["tool_arguments_delta"] = partial
        elif etype == "message_delta":
            delta = j.get("delta", {})
            sr = str(delta.get("stop_reason") or "")
            if sr:
                result["finish_reason"] = "tool_calls" if sr == "tool_use" else sr

        return result

    def stream_is_done(self, chunk_text: str) -> bool:
        try:
            j = json.loads(chunk_text)
            return str(j.get("type") or "") == "message_stop"
        except Exception:
            return False

    def create_stream_state(self) -> dict:
        state = super().create_stream_state()
        state["anthropic_current_tool_idx"] = -1
        return state

    def accumulate_chunk(self, state: dict, chunk_text: str) -> dict:
        try:
            j = json.loads(chunk_text or "{}")
        except Exception:
            return {}

        etype = str(j.get("type") or "")
        content_delta = None
        reasoning_delta = None

        if etype == "content_block_start":
            cb = j.get("content_block") or {}
            cb_idx = int(j.get("index") or 0)
            cb_type = str(cb.get("type") or "")
            if cb_type == "tool_use":
                state["anthropic_current_tool_idx"] = cb_idx
                tool_id = str(cb.get("id") or f"toolu_{cb_idx}")
                tool_name = str(cb.get("name") or "")
                state["tool_calls"][cb_idx] = {
                    "id": tool_id,
                    "type": "function",
                    "function": {"name": tool_name, "arguments": ""}
                }

        elif etype == "content_block_delta":
            delta = j.get("delta") or {}
            dtype = str(delta.get("type") or "")
            if dtype == "text_delta":
                text = str(delta.get("text") or "")
                if text:
                    content_delta = text
                    state["content_parts"].append(text)
            elif dtype == "input_json_delta":
                partial = str(delta.get("partial_json") or "")
                if partial:
                    idx = state.get("anthropic_current_tool_idx", -1)
                    if idx >= 0 and idx in state["tool_calls"]:
                        state["tool_calls"][idx]["function"]["arguments"] += partial

        elif etype == "message_delta":
            delta = j.get("delta") or {}
            sr = str(delta.get("stop_reason") or "")
            if sr:
                state["finish_reason"] = "tool_calls" if sr == "tool_use" else sr

        return {
            "content_delta": content_delta,
            "reasoning_delta": reasoning_delta,
        }


class GeminiProvider(BaseLLMProvider):
    """Google Gemini 原生 API 格式（v1beta / v1 generateContent）。"""

    PROVIDER_ID = "gemini"

    def build_request(self, messages, tools=None, temperature=None, max_tokens=None, stream=False):
        base = self.base_url.rstrip("/")
        if not base.endswith("/v1") and not base.endswith("/v1beta"):
            base = base + "/v1"
        url = f"{base}/models/{self.model}:generateContent"
        if stream:
            url = f"{base}/models/{self.model}:streamGenerateContent"

        headers = {
            "Content-Type": "application/json",
            "x-goog-api-key": self.api_key,
        }

        system_instruction = None
        contents = []
        for m in messages:
            if not isinstance(m, dict):
                continue
            role = str(m.get("role") or "")
            content = m.get("content")
            if role == "system":
                if isinstance(content, str) and content.strip():
                    system_instruction = {"parts": [{"text": content}]}
            elif role in ("user", "model"):
                gemini_role = "model" if role == "assistant" else "user"
                parts = []
                if isinstance(content, str):
                    parts.append({"text": content})
                elif isinstance(content, list):
                    for p in content:
                        if isinstance(p, dict):
                            ptype = str(p.get("type") or "")
                            if ptype == "text":
                                parts.append({"text": str(p.get("text") or "")})
                            elif ptype in ("image_url", "image"):
                                img = p.get("image_url", {}) if isinstance(p.get("image_url"), dict) else p
                                url = str(img.get("url") or "")
                                if url.startswith("data:"):
                                    idx = url.find(",")
                                    if idx > 0:
                                        mime = url[5:idx].split(";")[0]
                                        b64 = url[idx + 1:]
                                        parts.append({"inline_data": {"mime_type": mime, "data": b64}})
                if parts:
                    contents.append({"role": gemini_role, "parts": parts})

        body = {"contents": contents}
        if system_instruction:
            body["system_instruction"] = system_instruction

        gen_cfg = {}
        if temperature is not None:
            gen_cfg["temperature"] = float(temperature)
        if max_tokens and max_tokens > 0:
            gen_cfg["maxOutputTokens"] = max_tokens
        if gen_cfg:
            body["generationConfig"] = gen_cfg

        if tools:
            gemini_funcs = []
            for t in tools:
                if not isinstance(t, dict):
                    continue
                func = t.get("function", {}) if t.get("type") == "function" else t
                name = str(func.get("name") or "")
                desc = str(func.get("description") or "")
                params = func.get("parameters", {})
                if not name:
                    continue
                gemini_funcs.append({"name": name, "description": desc, "parameters": params})
            if gemini_funcs:
                body["tools"] = [{"function_declarations": gemini_funcs}]

        return url, headers, body

    def parse_response(self, raw_body: str) -> dict:
        data = json.loads(raw_body or "{}")
        candidates = data.get("candidates")
        if not isinstance(candidates, list) or not candidates:
            return {"content": "", "finish_reason": "", "raw": data}

        cand = candidates[0]
        content = cand.get("content", {})
        parts = content.get("parts", []) if isinstance(content, dict) else []
        if not isinstance(parts, list):
            parts = []

        text_parts = []
        tool_calls = []
        for idx, part in enumerate(parts):
            if not isinstance(part, dict):
                continue
            if "text" in part:
                text_parts.append(str(part.get("text") or ""))
            elif "functionCall" in part:
                fc = part.get("functionCall", {})
                tool_calls.append({
                    "id": f"gemini_{idx}",
                    "type": "function",
                    "function": {
                        "name": str(fc.get("name") or ""),
                        "arguments": json.dumps(fc.get("args", {}), ensure_ascii=False),
                    },
                })

        finish_reason = str(cand.get("finishReason") or "")
        fr_lower = finish_reason.lower()
        if fr_lower == "stop":
            finish_reason = "stop"
        elif fr_lower in ("function_call", "tool_call"):
            finish_reason = "tool_calls"
        elif fr_lower == "max_tokens":
            finish_reason = "length"

        result = {
            "content": "".join(text_parts).strip(),
            "reasoning_content": "",
            "finish_reason": finish_reason,
            "raw": data,
        }
        if tool_calls:
            result["tool_calls"] = tool_calls
        return result

    def parse_stream_chunk(self, chunk_text: str) -> dict:
        try:
            j = json.loads(chunk_text)
        except Exception:
            return None
        result = {}

        candidates = j.get("candidates")
        if isinstance(candidates, list) and candidates:
            cand = candidates[0]
            content = cand.get("content", {})
            parts = content.get("parts", []) if isinstance(content, dict) else []
            if isinstance(parts, list):
                for part in parts:
                    if isinstance(part, dict) and "text" in part:
                        t = str(part.get("text") or "")
                        if t:
                            result["content_delta"] = t
                            break

            fr = str(cand.get("finishReason") or "")
            if fr:
                fr_lower = fr.lower()
                if fr_lower == "stop":
                    result["finish_reason"] = "stop"
                elif fr_lower in ("function_call", "tool_call"):
                    result["finish_reason"] = "tool_calls"
                elif fr_lower == "max_tokens":
                    result["finish_reason"] = "length"

        return result

    def stream_is_done(self, chunk_text: str) -> bool:
        return False  # Gemini 通过 HTTP 流结束自动终止

    def accumulate_chunk(self, state: dict, chunk_text: str) -> dict:
        try:
            j = json.loads(chunk_text or "{}")
        except Exception:
            return {}

        content_delta = None
        reasoning_delta = None

        candidates = j.get("candidates")
        if isinstance(candidates, list) and candidates:
            cand = candidates[0]
            content = cand.get("content", {})
            parts = content.get("parts", []) if isinstance(content, dict) else []
            if isinstance(parts, list):
                for part in parts:
                    if isinstance(part, dict):
                        if "text" in part:
                            t = str(part.get("text") or "")
                            if t:
                                content_delta = t
                                state["content_parts"].append(t)
                        elif "functionCall" in part:
                            fc = part.get("functionCall", {})
                            name = str(fc.get("name") or "")
                            args = fc.get("args", {})
                            if name and isinstance(args, dict):
                                args_str = json.dumps(args, ensure_ascii=False)
                                found = False
                                for idx, tc in state["tool_calls"].items():
                                    if tc["function"]["name"] == name and tc["function"]["arguments"] == args_str:
                                        found = True
                                        break
                                if not found:
                                    idx = len(state["tool_calls"])
                                    state["tool_calls"][idx] = {
                                        "id": f"gemini_{idx}",
                                        "type": "function",
                                        "function": {"name": name, "arguments": args_str}
                                    }

            fr = str(cand.get("finishReason") or "")
            if fr:
                fr_lower = fr.lower()
                if fr_lower == "stop":
                    state["finish_reason"] = "stop"
                elif fr_lower in ("function_call", "tool_call"):
                    state["finish_reason"] = "tool_calls"
                elif fr_lower == "max_tokens":
                    state["finish_reason"] = "length"

        return {
            "content_delta": content_delta,
            "reasoning_delta": reasoning_delta,
        }


_PROVIDER_REGISTRY = {
    OpenAIProvider.PROVIDER_ID: OpenAIProvider,
    AnthropicProvider.PROVIDER_ID: AnthropicProvider,
    GeminiProvider.PROVIDER_ID: GeminiProvider,
}


def _get_provider(profile: dict) -> BaseLLMProvider:
    """根据 profile 配置获取对应的 Provider 实例。"""
    if not isinstance(profile, dict):
        return OpenAIProvider({})
    provider_id = str(profile.get("provider") or "").strip().lower()
    if not provider_id:
        provider_id = OpenAIProvider.PROVIDER_ID
    cls = _PROVIDER_REGISTRY.get(provider_id, OpenAIProvider)
    return cls(profile)


def ai_chat(messages: list, temperature: float | None = None, profile_id: str | None = None, timeout_s: int = 60, stream: bool = False, on_delta=None, tools: list | None = None):
    with ai_config_lock:
        cfg = load_hivo_ai_config()
    prof = _ai_pick_profile(cfg, profile_id)
    if not prof:
        return False, "未配置可用模型", None

    limits = cfg.get("limits") if isinstance(cfg, dict) else None
    limits = limits if isinstance(limits, dict) else {}
    try:
        max_input_tokens = int(limits.get("ai_max_input_tokens") or 0)
    except Exception:
        max_input_tokens = 0
    try:
        max_output_tokens = int(limits.get("ai_max_output_tokens") or 0)
    except Exception:
        max_output_tokens = 0

    provider = _get_provider(prof)

    if not provider.base_url:
        return False, "未配置 API Base URL", None
    if not provider.model:
        return False, "未配置 Model", None
    if not isinstance(messages, list) or not messages:
        return False, "messages 为空", None

    msgs2 = _ai_trim_messages_to_budget(messages, max_total_tokens=max_input_tokens, reserve_output_tokens=max_output_tokens)
    if msgs2 is None:
        try:
            budget = max(128, int(max_input_tokens or 0) - int(max_output_tokens or 0))
        except Exception:
            budget = 512
        budget = max(128, budget)
        sys_txt = ""
        last_user = None
        for m in messages:
            if str(m.get("role") or "") == "system" and isinstance(m.get("content"), str):
                if sys_txt:
                    sys_txt += "\n\n" + m.get("content")
                else:
                    sys_txt = m.get("content")
        for m in reversed(messages):
            if str(m.get("role") or "") == "user" and isinstance(m.get("content"), str) and m.get("content").strip():
                last_user = m
                break
        out = []
        if sys_txt:
            sys_allow = max(32, int(budget * 0.3))
            out.append({"role": "system", "content": _ai_truncate_text_to_tokens(sys_txt, sys_allow)})
        if last_user:
            remain = max(32, budget - _ai_estimate_messages_tokens(out) - 8)
            out.append({"role": "user", "content": _ai_truncate_text_to_tokens(last_user.get("content"), remain)})
        msgs2 = out if out else []
    if not msgs2:
        return False, "messages 为空", None

    try:
        tmo = int(timeout_s or 0)
    except Exception:
        tmo = 60
    tmo = max(5, min(120, int(tmo)))

    def _build_req():
        url, headers, body = provider.build_request(msgs2, tools=tools, temperature=temperature, max_tokens=max_output_tokens, stream=stream)
        req = urllib.request.Request(
            url,
            data=json.dumps(body, ensure_ascii=False).encode("utf-8"),
            headers=headers,
            method="POST",
        )
        return url, headers, body, req

    url, headers, payload, req = _build_req()
    if not url:
        return False, "API Base URL 非法", None

    def _do_request(req_):
        with _ai_build_opener().open(req_, timeout=tmo) as resp_:
            if stream:
                stream_state = provider.create_stream_state()
                while True:
                    try:
                        line_b = resp_.readline()
                    except Exception:
                        break
                    if not line_b:
                        break
                    line = line_b.decode("utf-8", errors="replace").strip()
                    if not line:
                        continue
                    if not line.startswith("data:"):
                        continue
                    chunk = line[len("data:"):].strip()
                    if provider.stream_is_done(chunk):
                        break
                    delta = provider.accumulate_chunk(stream_state, chunk)
                    if not delta:
                        continue
                    delta_content = delta.get("content_delta")
                    delta_reasoning = delta.get("reasoning_delta")
                    if delta_content:
                        try:
                            if callable(on_delta):
                                on_delta({"delta": str(delta_content), "delta_type": "content"})
                        except Exception as e:
                            logger.debug(f"Exception ignored: {e}")
                    if delta_reasoning:
                        try:
                            if callable(on_delta):
                                on_delta({"delta": str(delta_reasoning), "delta_type": "reasoning"})
                        except Exception as e:
                            logger.debug(f"Exception ignored: {e}")
                result = provider.get_stream_result(stream_state)
                return True, "", result
            else:
                raw = resp_.read().decode("utf-8", errors="replace")
                try:
                    parsed = provider.parse_response(raw)
                except Exception as e:
                    snip = (raw or "").strip()
                    if snip:
                        cf_msg = _parse_upstream_error(200, snip)
                        if cf_msg:
                            return False, cf_msg, None
                    if len(snip) > 400:
                        snip = snip[:400] + "…"
                    return False, f"上游响应不是有效 JSON: {e}; snippet={snip}", None
                return True, "", parsed

    try:
        return _do_request(req)
    except urllib.error.HTTPError as e:
        try:
            body = e.read().decode("utf-8", errors="replace")
        except Exception:
            body = ""
        msg = _parse_upstream_error(e.code, body)
        if e.code == 429 or (500 <= e.code < 600):
            for retry_i in range(3):
                try:
                    wait_s = min(2 ** retry_i, 8)
                    time.sleep(wait_s)
                    logger.info(f"[ai_chat] 自动重试 {retry_i + 1}/3 (HTTP {e.code})")
                    _, _, _, req2 = _build_req()
                    ok, err_msg, result = _do_request(req2)
                    if ok:
                        logger.info(f"[ai_chat] 重试成功 (第 {retry_i + 1} 次)")
                        return ok, err_msg, result
                except Exception:
                    continue
        return False, msg, None
    except Exception as e:
        err_str = str(e).lower()
        if "timeout" in err_str or "timed out" in err_str or "connection" in err_str:
            try:
                logger.info(f"[ai_chat] 网络错误自动重试: {err_str[:100]}")
                time.sleep(2)
                _, _, _, req3 = _build_req()
                ok, err_msg, result = _do_request(req3)
                if ok:
                    logger.info("[ai_chat] 网络错误重试成功")
                    return ok, err_msg, result
            except Exception:
                pass
        return False, str(e), None


# ════════════════════════════════════════════════════════
#  Agent Action Framework（Action Registry）
# ════════════════════════════════════════════════════════
# Unified registry for agent actions: registration, discovery,
# validation, system-prompt generation, and frontend serialization.
# Actions are host-application-level operations independent of
# Tool Calls / MCP Tools. Scope support allows exposing only
# relevant actions per page/module/workspace context.

class ActionRegistry:
    """可扩展的 Agent Action 注册中心。

    统一管理 Action 的定义、参数、权限、确认策略和作用域（Scope）。
    支持注册、注销、发现、校验、System Prompt 动态生成及前端序列化。
    新增/删除 Action 只需调用 register()/unregister()，无需修改核心流程。
    """

    def __init__(self):
        self._actions: dict = {}  # action_type → definition dict
        self._lock = threading.Lock()

    # ── Registration ──────────────────────────────────────────

    def register(self, action_type: str, desc: str,
                 params: dict | None = None,
                 scope: list | None = None):
        """注册（或更新）一个 Agent Action 定义。

        Args:
            action_type: 动作唯一标识（如 clear_session）
            desc:        动作的人类可读描述
            params:      {参数名: {"type":"string","desc":"..."}} 或 {}
            scope:       作用域列表，如 ["global","ai_chat","editor"]
                         Action 仅在匹配的 scope 下暴露给 AI
        """
        with self._lock:
            self._actions[action_type] = {
                "action_type": action_type,
                "desc": str(desc or "").strip(),
                "params": params if isinstance(params, dict) else {},
                "requiresConfirmation": False,
                "scope": list(scope) if scope else ["global"],
            }

    def unregister(self, action_type: str):
        """注销一个 Action，后续将不再暴露给 AI 和前端。"""
        with self._lock:
            self._actions.pop(action_type, None)

    # ── Query ─────────────────────────────────────────────────

    def get(self, action_type: str) -> dict | None:
        """获取单个 Action 定义。"""
        with self._lock:
            return self._actions.get(action_type)

    def get_all(self, scope: str | None = None) -> list:
        """获取所有 Action 定义，可按 scope 过滤。"""
        with self._lock:
            actions = list(self._actions.values())
        if scope:
            actions = [a for a in actions if scope in (a.get("scope") or [])]
        return actions

    def validate(self, action: dict) -> bool:
        """校验一个 action 对象是否为已注册的 Action。"""
        at = str(action.get("action_type") or "").strip()
        if not at:
            return False
        with self._lock:
            return at in self._actions

    @property
    def count(self) -> int:
        with self._lock:
            return len(self._actions)

    @property
    def action_types(self) -> list:
        with self._lock:
            return list(self._actions.keys())

    # ── System Prompt ─────────────────────────────────────────

    def build_system_prompt_text(self, scope: str | None = None) -> str:
        """动态生成 System Prompt 中的可用 Action 列表。

        Args:
            scope: 仅暴露匹配该 scope 的 Action（None=全部暴露）
        """
        actions = self.get_all(scope=scope)
        if not actions:
            return "（当前上下文无可用的 Agent 动作）"
        lines = []
        for a in actions:
            at = a.get("action_type", "")
            desc = str(a.get("desc", "") or "").strip()
            params = a.get("params") if isinstance(a.get("params"), dict) else {}
            param_str = ""
            if params:
                parts = []
                for pk, pv in params.items():
                    ptype = str(pv.get("type") or "").strip() if isinstance(pv, dict) else ""
                    parts.append(f"{pk}:{ptype}")
                if parts:
                    param_str = "（参数：" + ", ".join(parts) + "）"
            lines.append(f"- {at}{param_str}：{desc}")
        return "\n".join(lines)

    # ── Serialization ─────────────────────────────────────────

    def to_serializable(self, scope: str | None = None) -> list:
        """将 Action 定义序列化为前端可消费的 JSON 列表。

        Returns: [{action_type, desc, params, requiresConfirmation, scope}, ...]
        """
        actions = self.get_all(scope=scope)
        return [{
            "action_type": a["action_type"],
            "desc": a["desc"],
            "params": a["params"],
            "requiresConfirmation": a["requiresConfirmation"],
            "scope": a["scope"],
        } for a in actions]


# ── Global Action Registry singleton ─────────────────────────
# 所有 Agent Action 通过此实例统一管理。
# 新增 Action：AGENT_ACTION_REGISTRY.register("xxx", ...)
# 删除 Action：AGENT_ACTION_REGISTRY.unregister("xxx")
AGENT_ACTION_REGISTRY = ActionRegistry()


def _ai_build_system_context_text(cfg=None, use_native_tools=True):
    """构建系统提示词（强约束）。可传入 cfg 避免重复 _hivo_load_cfg()。
    use_native_tools: 是否使用原生 function calling（原生模式下精简工具 JSON 描述，避免与 tools 参数冲突）。
    """
    try:
        spec = get_capabilities_spec()
    except Exception:
        spec = {}
    try:
        tools0 = (spec.get("agent_tools") or []) if isinstance(spec, dict) else []
        eps0 = (spec.get("endpoints") or []) if isinstance(spec, dict) else []

        slim_tools = []
        for t in (tools0[:200] if isinstance(tools0, list) else []):
            if not isinstance(t, dict):
                continue
            tt = str(t.get("type") or "").strip()
            if not tt:
                continue
            desc = str(t.get("desc") or "").strip()
            if len(desc) > 160:
                desc = desc[:160] + "…"
            req = t.get("required") if isinstance(t.get("required"), list) else []
            props_in = t.get("properties") if isinstance(t.get("properties"), dict) else {}
            props = {}
            for k, v in list(props_in.items())[:30]:
                if not isinstance(v, dict):
                    continue
                ptype = str(v.get("type") or "").strip()
                pdesc = str(v.get("desc") or "").strip()
                if len(pdesc) > 120:
                    pdesc = pdesc[:120] + "…"
                props[str(k)] = {"type": ptype, "desc": pdesc}
            slim_tools.append({"type": tt, "desc": desc, "required": req, "properties": props})

        slim_eps = []
        for e in (eps0[:200] if isinstance(eps0, list) else []):
            if not isinstance(e, dict):
                continue
            m = str(e.get("method") or "").upper().strip()
            p = str(e.get("path") or "").strip()
            if not m or not p:
                continue
            slim_eps.append({"method": m, "path": p})

        slim = {
            "version": spec.get("version") if isinstance(spec, dict) else "",
            "generated_at": spec.get("generated_at") if isinstance(spec, dict) else "",
            "agent_tools": slim_tools,
            "endpoints": slim_eps,
        }
    except Exception:
        slim = {"version": "", "generated_at": "", "agent_tools": [], "endpoints": []}

    if use_native_tools:
        slim_registry = {
            "version": slim.get("version", ""),
            "generated_at": slim.get("generated_at", ""),
            "endpoints": slim.get("endpoints", []),
            "note": "工具定义通过 function calling 传递，此处仅列出 API 端点供参考",
        }
        tool_registry = "TOOL_REGISTRY_JSON:\n" + json.dumps(slim_registry, ensure_ascii=False, indent=2)
    else:
        tool_registry = "TOOL_REGISTRY_JSON:\n" + json.dumps(slim, ensure_ascii=False, indent=2)

    try:
        cfg0 = cfg if isinstance(cfg, dict) and cfg else _hivo_load_cfg()
        extra = str((cfg0.get("system_context_extra") if isinstance(cfg0, dict) else "") or "").strip()
    except Exception:
        extra = ""

    tool_call_section = (
        "# 工具调用\n"
        "- 仅使用系统提供的工具，不要臆造工具名或参数\n"
        "- 每轮最多同时调用 3 个工具\n"
        "- 工具调用后等待结果返回再继续下一步\n"
        "- 文件路径不明确时先用 find_files 模糊搜索，支持部分文件名匹配\n"
        "- 修改大文件时优先用 edit_file（局部替换），比 save_file 更高效可靠\n"
        "- 只需确认文件是否存在或获取元信息时用 file_stat，不要用 get_file\n"
        "- 所有路径使用正斜杠 /，不要用反斜杠 \\\n"
        "- 需要读取文档文件时用 read_pdf/read_docx/read_xlsx\n"
        "- 排查端口冲突或查看运行进程时用 list_processes\n"
        "- 排查环境问题时用 env_info 查看 Python 版本和依赖\n\n"
        if use_native_tools else
        "# 工具调用\n"
        "- 工具定义见 TOOL_REGISTRY_JSON\n"
        "- 只使用已定义的工具，确保 required 参数完整\n"
        "- 每轮最多 3 个工具调用，末尾独占一行\n"
        "- 工具调用后等待回执再进行下一步\n"
        "- 文件路径不明确时先用 find_files 模糊搜索，支持部分文件名匹配\n"
        "- 修改大文件时优先用 edit_file（局部替换），比 save_file 更高效可靠\n"
        "- 只需确认文件是否存在或获取元信息时用 file_stat，不要用 get_file\n"
        "- 所有路径使用正斜杠 /，不要用反斜杠 \\\n"
        "- 需要读取文档文件时用 read_pdf/read_docx/read_xlsx\n"
        "- 排查端口冲突或查看运行进程时用 list_processes\n"
        "- 排查环境问题时用 env_info 查看 Python 版本和依赖\n\n"
    )

    strong = (
        "你是 Hivo，Git Manager 系统的智能助手，负责帮助用户管理 Git 仓库和文件操作。\n\n"
        "# 核心原则\n"
        "1. 工具驱动：通过调用工具完成实际操作，而非仅描述步骤\n"
        "2. 结果导向：基于工具返回的真实结果回复用户，不臆造结果\n"
        "3. 清晰沟通：用自然语言解释操作结果，隐藏技术细节\n\n"
        "# 术语\n"
        "- 暂存（stash）：git stash 暂存栈，恢复用 stash_pop\n"
        "- 暂存区（staged/index）：git add 后的区域\n\n"
        "# 交互约束\n"
        "- Git 操作前先用 status 确认仓库已打开\n"
        "- pull/switch 遇到未提交修改时，先让用户选择处理方式再执行\n"
        "- 用户未明确选择前，禁止执行改变仓库状态的操作\n"
        "- 复合操作后说明状态并给出后续选项\n\n"
        "# 输出规范\n"
        "- 使用 Markdown 格式\n"
        "- 可执行命令代码块统一用 `run-` 前缀标签，前端自动渲染运行按钮\n"
        "- 收到【用户已执行命令】消息后基于输出内容分析，不重复执行\n\n"
        + tool_call_section +
        "# Agent 动作\n"
        "- 使用 action_type 而非 type 区分\n"
        "- 可用动作：\n" + AGENT_ACTION_REGISTRY.build_system_prompt_text(scope="ai_chat") + "\n"
        "- 当用户要求 UI 级操作（清空会话/切换主题等）时直接输出 action JSON\n\n"
    )

    if extra:
        return strong + "\n\n" + extra + "\n\n" + tool_registry
    return strong + "\n\n" + tool_registry


def _hivo_scan_json_objects(text: str, max_objects: int = 10):
    out = []
    s = str(text or "")
    if not s:
        return out
    depth = 0
    start = -1
    in_str = False
    esc = False
    for i, ch in enumerate(s):
        if in_str:
            if esc:
                esc = False
            elif ch == "\\":
                esc = True
            elif ch == '"':
                in_str = False
            continue
        if ch == '"':
            in_str = True
            continue
        if ch == '{':
            if depth == 0:
                start = i
            depth += 1
            continue
        if ch == '}':
            if depth > 0:
                depth -= 1
                if depth == 0 and start >= 0:
                    seg = s[start:i + 1]
                    out.append(seg)
                    if len(out) >= max_objects:
                        break
                    start = -1
            continue
    return out


def _hivo_tools_to_openai_format(agent_tools: list) -> list:
    """将内部 agent_tools 格式转换为 OpenAI function calling 格式。

    内部格式: {"type": "tool_name", "desc": "...", "required": [...], "properties": {"name": {"type": "...", "desc": "..."}}}
    OpenAI 格式: {"type": "function", "function": {"name": "...", "description": "...", "parameters": {"type": "object", "properties": {...}, "required": [...]}}}
    """
    if not isinstance(agent_tools, list):
        return []
    out = []
    for t in agent_tools:
        if not isinstance(t, dict):
            continue
        name = str(t.get("type") or "").strip()
        if not name:
            continue
        desc = str(t.get("desc") or t.get("description") or "")
        props_in = t.get("properties") if isinstance(t.get("properties"), dict) else {}
        required_in = t.get("required") if isinstance(t.get("required"), list) else []

        properties = {}
        for k, v in props_in.items():
            if not isinstance(v, dict):
                continue
            ptype = str(v.get("type") or "string").strip()
            pdesc = str(v.get("desc") or v.get("description") or "")
            prop = {"type": ptype}
            if pdesc:
                prop["description"] = pdesc
            enum = v.get("enum")
            if isinstance(enum, list) and enum:
                prop["enum"] = enum
            properties[str(k)] = prop

        parameters = {"type": "object", "properties": properties}
        required = [str(x) for x in required_in if str(x).strip() and str(x) in properties]
        if required:
            parameters["required"] = required

        out.append({
            "type": "function",
            "function": {
                "name": name,
                "description": desc,
                "parameters": parameters,
            },
        })
    return out


def _hivo_parse_native_tool_calls(tool_calls: list) -> list:
    """将原生 function calling 返回的 tool_calls 转换为内部调用格式。

    原生格式: [{"id": "...", "type": "function", "function": {"name": "...", "arguments": "{...}"}}]
    内部格式: [{"type": "tool_name", ...arg_fields...}]
    """
    if not isinstance(tool_calls, list):
        return []
    out = []
    for tc in tool_calls:
        if not isinstance(tc, dict):
            continue
        func = tc.get("function", {}) if tc.get("type") == "function" else tc
        name = str(func.get("name") or "").strip()
        if not name:
            continue
        args_str = str(func.get("arguments") or "")
        try:
            args = json.loads(args_str) if args_str.strip() else {}
        except Exception:
            args = {}
        if not isinstance(args, dict):
            args = {}
        call = {"type": name}
        call.update(args)
        out.append(call)
    return out


def _hivo_supports_native_tools(profile: dict) -> bool:
    """判断 profile 配置的模型是否支持原生 function calling。

    目前策略：OpenAI 兼容格式默认支持；Anthropic/Gemini 也支持。
    只有显式标记 native_tools=false 时才不支持。
    """
    if not isinstance(profile, dict):
        return True
    nt = profile.get("native_tools")
    if nt is not None:
        return bool(nt)
    return True


def _hivo_extract_tool_calls(text: str, max_calls: int = 3):
    calls = []
    s = str(text or "").strip()
    if not s:
        return calls

    # 1) Fast path: whole output is a JSON object or JSON array.
    #    Note: system prompt允许"单个 JSON 对象或 JSON 数组"，这里必须支持数组，否则多指令不会被执行。
    try:
        whole = json.loads(s)
        if isinstance(whole, dict):
            if str(whole.get("type") or "").strip():
                return [whole]
        if isinstance(whole, list):
            for it in whole:
                if isinstance(it, dict) and str(it.get("type") or "").strip():
                    calls.append(it)
                    if len(calls) >= max_calls:
                        break
            if calls:
                return calls
    except Exception as e:
        logger.debug(f"Exception ignored: {e}")

    # 2) Fallback: scan embedded JSON objects in mixed text.
    for seg in _hivo_scan_json_objects(s, max_objects=max_calls * 2):
        try:
            obj = json.loads(seg)
        except Exception:
            continue
        if isinstance(obj, dict) and str(obj.get("type") or "").strip():
            calls.append(obj)
            if len(calls) >= max_calls:
                break
    return calls


# ── Agent Action Registry ──────────────────────────────────────────────
# All agent actions registered through the unified ActionRegistry.
# To add a new action, call AGENT_ACTION_REGISTRY.register(...).
# To remove, call AGENT_ACTION_REGISTRY.unregister("action_type").

AGENT_ACTION_REGISTRY.register(
    "clear_session", "清空当前会话的所有消息（保留会话本身）",
    scope=["global", "ai_chat"],
)
AGENT_ACTION_REGISTRY.register(
    "close_session", "关闭当前会话，切换到新会话",
    scope=["global", "ai_chat"],
)
AGENT_ACTION_REGISTRY.register(
    "open_settings", "打开 AI 设置面板",
    scope=["global", "ai_chat"],
)
AGENT_ACTION_REGISTRY.register(
    "switch_theme", "切换界面主题（dark 或 light）",
    params={"theme": {"type": "string", "desc": "目标主题：dark 或 light"}},
    scope=["global", "ai_chat"],
)
AGENT_ACTION_REGISTRY.register(
    "focus_input", "聚焦 AI 输入框，方便用户快速输入",
    scope=["global", "ai_chat"],
)
AGENT_ACTION_REGISTRY.register(
    "toggle_sidebar", "展开或收起侧边栏",
    params={"side": {"type": "string", "desc": "要操作的侧边栏：left 或 right"}},
    scope=["global", "ai_chat"],
)
AGENT_ACTION_REGISTRY.register(
    "undo_last_turn", "撤销上一轮 AI 的操作",
    scope=["global", "ai_chat"],
)


def _hivo_extract_agent_actions(text: str, max_actions: int = 3):
    """Extract agent action calls from LLM response text.
    An agent action uses 'action_type' (not 'type') as the discriminator key."""
    actions = []
    s = str(text or "").strip()
    if not s:
        return actions
    # Scan for JSON objects containing "action_type"
    for seg in _hivo_scan_json_objects(s, max_objects=max_actions * 2):
        try:
            obj = json.loads(seg)
        except Exception:
            continue
        if isinstance(obj, dict) and str(obj.get("action_type") or "").strip():
            actions.append(obj)
            if len(actions) >= max_actions:
                break
    return actions


def _hivo_ws_emit_action(run_id: str, session_id: str, action: dict):
    """Forward an agent action request to the frontend for execution."""
    try:
        payload = {
            "type": "ai_agent_action",
            "run_id": str(run_id or ""),
            "session_id": str(session_id or ""),
            "action": action if isinstance(action, dict) else {},
            "ts": time.time(),
        }
        broadcast_to_clients(payload)
    except Exception as e:
        logger.debug(f"Exception ignored: {e}")


def _hivo_ws_emit_action_registry(run_id: str = "", session_id: str = ""):
    """Send the full Action Registry definitions to frontend for dynamic consumption.

    Frontend uses these definitions to drive confirmation dialogs,
    parameter validation, and capability discovery — without hardcoding
    action metadata.
    """
    try:
        defs = AGENT_ACTION_REGISTRY.to_serializable()
        payload = {
            "type": "ai_action_registry",
            "run_id": str(run_id or ""),
            "session_id": str(session_id or ""),
            "actions": defs,
            "ts": time.time(),
        }
        broadcast_to_clients(payload)
    except Exception as e:
        logger.debug(f"Exception ignored: {e}")


def _hivo_tool_receipt_line(name: str, ok: bool, detail: str = ""):
    s = f"- {name}: {'ok' if ok else 'fail'}"
    if detail:
        s += f" ({detail})"
    return s


def _hivo_exec_tool(tool: dict, undo_gid: str = "", run_id: str = "", agent_deadline: float = 0.0):
    t = str((tool or {}).get("type") or "").strip()
    if not t:
        return False, "missing type", {}

    # System/Repo ops
    if t == "status":
        try:
            ok, msg, st0 = get_repo_status()
            if not ok:
                return False, msg or "status failed", {}
            return True, "", {"status": st0}
        except Exception as e:
            return False, str(e), {}

    if t == "open_repo":
        path0 = str(tool.get("path") or "").strip()
        if not path0:
            return False, "缺少 path", {}
        ok, msg = open_repo(path0)
        if not ok:
            return False, msg or "open_repo failed", {"path": path0}
        return True, "", {"path": path0}

    if t == "set_origin":
        url0 = str(tool.get("url") or "").strip()
        if not url0:
            return False, "缺少 url", {}
        ok, msg, origin_url = set_origin_url(url0)
        if not ok:
            return False, msg or "set_origin failed", {"url": url0}
        return True, "", {"origin_url": origin_url}

    if t == "workspace_context":
        try:
            txt = workspace_context_text()
            return True, "", {"content": txt}
        except Exception as e:
            return False, str(e), {}

    if t == "open_file":
        name = str(tool.get("name") or tool.get("path") or "").strip()
        view = str(tool.get("view") or "editor").strip().lower() or "editor"
        if not name:
            return False, "缺少 name", {}
        if view not in ("editor", "change", "split", "unified"):
            view = "editor"

        # If looks like a path, try as-is first.
        pick = ""
        cand = name.replace("\\", "/").lstrip("/")
        full = _safe_repo_abspath(cand)
        if full and os.path.isfile(full):
            pick = cand
        else:
            items = find_files_by_name(name, max_results=20)
            if not items:
                return False, "未找到文件", {"name": name}
            pick = str(items[0] or "")
            nm_lower = name.lower()
            for rel in items:
                bn = str(rel or "").replace("\\", "/").split("/")[-1]
                if bn.lower() == nm_lower:
                    pick = str(rel or "")
                    break

        content, encoding = get_file_content(pick, return_encoding=True)
        if content is None:
            return False, "无法读取文件内容", {"file": {"path": pick, "view": view}}
        try:
            broadcast_to_clients({
                "type": "open_file",
                "file": {
                    "path": pick,
                    "file_name": os.path.basename(pick),
                    "view": view,
                    "encoding": encoding,
                    "content": content,
                },
            })
        except Exception as e:
            logger.debug(f"Exception ignored: {e}")
        return True, "", {"file": {"path": pick, "view": view, "encoding": encoding, "content": content}}

    # File/Workspace
    if t in ("get_file",):
        rp_in = str(tool.get("path") or tool.get("name") or "").strip()
        rp, full = _hivo_resolve_full_path(rp_in)
        if not (full and os.path.exists(full)) or os.path.isdir(full):
            return False, "文件不存在", {"file": {"path": rp}}
        # 判定媒体 or 文本
        try:
            mime0, _enc0 = mimetypes.guess_type(full)
        except Exception:
            mime0 = None
        # 预判常见结构化类型（优先于 is_media 判定）
        mime_s = str(mime0 or "").lower()
        ext_s = os.path.splitext(full)[1].lower()
        is_pdf = (mime_s == "application/pdf" or ext_s == ".pdf")
        is_docx = (mime_s == "application/vnd.openxmlformats-officedocument.wordprocessingml.document" or ext_s == ".docx")
        is_xlsx = (mime_s == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet" or ext_s == ".xlsx")
        is_pptx = (mime_s == "application/vnd.openxmlformats-officedocument.presentationml.presentation" or ext_s == ".pptx")
        # Legacy MS Office formats
        is_doc = (mime_s == "application/msword" or ext_s == ".doc")
        is_xls = (mime_s == "application/vnd.ms-excel" or ext_s == ".xls")
        is_ppt = (mime_s == "application/vnd.ms-powerpoint" or ext_s == ".ppt")
        is_zip = (mime_s in ("application/zip", "application/x-zip-compressed") or ext_s == ".zip")
        is_media = False
        if not (is_pdf or is_docx or is_xlsx or is_pptx or is_doc or is_xls or is_ppt or is_zip):
            if mime0 and (mime_s.startswith("image/") or mime_s.startswith("audio/") or mime_s.startswith("video/")):
                is_media = True
            else:
                try:
                    with open(full, "rb") as fchk:
                        sample = fchk.read(4096)
                    if b"\x00" in sample:
                        is_media = True
                except Exception:
                    pass
        bn = os.path.basename(rp.replace("\\", "/"))
        if is_media:
            # 返回 data_url（base64），大小上限从配置读取；超过阈值优先尝试按比例缩放到阈值内，失败才回退为截断。
            try:
                size0 = int(os.path.getsize(full))
            except Exception:
                size0 = 0
            ftype = mime0 or "application/octet-stream"
            try:
                cfg0 = _hivo_load_cfg()
                lim0 = cfg0.get("limits") if isinstance(cfg0.get("limits"), dict) else {}
                max_bytes = int(lim0.get("media_base64_max_bytes") or 256*1024)
            except Exception:
                max_bytes = 256*1024

            b: bytes = b""
            truncated = False
            downscaled = False
            width, height = (0, 0)
            if size0 > max_bytes and ftype.startswith("image/"):
                # 兜底：尝试用 Pillow 压缩/缩放到阈值内
                try:
                    from PIL import Image  # type: ignore
                    import io
                    with Image.open(full) as im:
                        im = im.convert("RGB")
                        try:
                            width, height = im.size
                        except Exception:
                            width, height = (0, 0)
                        # 估算按面积缩放比，避免过度缩小
                        import math
                        ratio = math.sqrt(max(0.1, min(1.0, float(max_bytes) / float(max(1, size0)))))
                        w, h = im.size
                        nw = max(1, int(w * ratio))
                        nh = max(1, int(h * ratio))
                        if nw < w or nh < h:
                            im = im.resize((nw, nh))
                            width, height = (nw, nh)
                        # 逐步降低质量以贴近阈值
                        buf = io.BytesIO()
                        q = 85
                        for q in (85, 75, 65, 55, 45, 35):
                            buf.seek(0); buf.truncate(0)
                            im.save(buf, format="JPEG", quality=q, optimize=True)
                            if buf.tell() <= max_bytes:
                                break
                        b = buf.getvalue()
                        ftype = "image/jpeg"
                        downscaled = True
                except Exception:
                    downscaled = False
                    b = b""

            if not b:
                # 常规路径：读取原始字节；若超限且未能缩放，则截断
                try:
                    with open(full, "rb") as fb:
                        if size0 > max_bytes and not downscaled:
                            b = fb.read(max_bytes)
                            truncated = True
                        else:
                            b = fb.read()
                except Exception:
                    b = b""; truncated = False

            import base64
            b64 = ""
            data_url = ""
            try:
                if b:
                    b64 = base64.b64encode(b).decode("ascii", errors="ignore")
                    if b64:
                        data_url = f"data:{ftype};base64,{b64}"
            except Exception:
                b64 = ""
                data_url = ""
            if ftype.startswith("image/"):
                return True, "", {
                    "file": {
                        "path": rp,
                        "file_type": ftype,
                        "size": size0,
                        "data_url": data_url,
                        "width": int(width),
                        "height": int(height),
                        "truncated": bool(truncated) and (not downscaled),
                    }
                }
            else:
                # 非图片媒体：对音频/视频暂不实现解析，返回受支持标记
                if ftype.startswith("audio/") or ftype.startswith("video/"):
                    return True, "", {
                        "file": {
                            "path": rp,
                            "file_type": ftype,
                            "size": size0,
                            "supported": False,
                        }
                    }
                else:
                    # 未知二进制：返回基础元数据（避免巨大 data_url）
                    try:
                        import hashlib
                        h = hashlib.sha256()
                        with open(full, "rb") as fb2:
                            for chunk in iter(lambda: fb2.read(1024*1024), b""):
                                h.update(chunk)
                        sha256 = h.hexdigest()
                    except Exception:
                        sha256 = ""
                    return True, "", {
                        "file": {
                            "path": rp,
                            "file_type": ftype or "application/octet-stream",
                            "size": size0,
                            "sha256": sha256,
                            "description": "未知二进制文件，暂不支持解析",
                        }
                    }
        # 文本：读取内容（区分附件与仓库文件）
        try:
            base_data = str(_hivo_ai_data_dir())
        except Exception:
            base_data = ""
        is_under_attachments = False
        try:
            if base_data:
                ap = os.path.abspath(full)
                is_under_attachments = (os.path.commonpath([base_data, ap]) == base_data)
        except Exception:
            is_under_attachments = False

        if is_under_attachments:
            # 直接读取绝对路径并进行编码探测/解码
            try:
                with open(full, "rb") as ftxt:
                    data = ftxt.read()
            except Exception:
                data = b""
            if not data:
                return True, "", {"file": {"path": rp, "file_type": "text/plain", "encoding": "utf-8", "content": ""}}
            encoding = None
            try:
                try:
                    content = data.decode("utf-8")
                    encoding = "utf-8"
                except UnicodeDecodeError:
                    try:
                        content = data.decode("utf-8-sig")
                        encoding = "utf-8-sig"
                    except UnicodeDecodeError:
                        ok = False
                        for enc0 in ("gb18030", "gbk", "cp936", "gb2312"):
                            try:
                                content = data.decode(enc0)
                                encoding = enc0
                                ok = True
                                break
                            except UnicodeDecodeError:
                                continue
                        if not ok:
                            try:
                                enc = _detect_text_encoding_from_bytes(data)
                                try:
                                    content = data.decode(enc)
                                    encoding = enc
                                except UnicodeDecodeError:
                                    content = data.decode(enc, errors="replace")
                                    encoding = enc
                            except Exception:
                                content = data.decode("utf-8", errors="replace")
                                encoding = "utf-8"
            except Exception:
                content = ""
                encoding = "utf-8"
            ftype_txt = None
            try:
                ftype_txt, _ = mimetypes.guess_type(bn)
            except Exception:
                ftype_txt = None
            return True, "", {"file": {"path": rp, "file_type": (ftype_txt or "text/plain"), "encoding": encoding, "content": content}}
        else:
            # 仓库内文件：按相对路径读取
            rp_norm = str(rp or "").replace("\\", "/")
            content, encoding = get_file_content(rp_norm, return_encoding=True)
            if content is None:
                return False, "无法读取文件内容", {"file": {"path": rp}}
            ftype_txt = None
            try:
                ftype_txt, _ = mimetypes.guess_type(bn)
            except Exception:
                ftype_txt = None
            return True, "", {"file": {"path": rp, "file_type": (ftype_txt or "text/plain"), "encoding": encoding, "content": content}}

    if t == "file_stat":
        rp_in = str(tool.get("path") or "").strip()
        rp, full = _hivo_resolve_full_path(rp_in)
        if not full:
            return False, "非法路径", {"path": rp, "exists": False}
        if not os.path.exists(full):
            return True, "", {"path": rp, "exists": False}
        try:
            st = os.stat(full)
            is_dir = os.path.isdir(full)
            ftype = None
            if not is_dir:
                try:
                    ftype, _ = mimetypes.guess_type(full)
                except Exception:
                    ftype = None
            return True, "", {
                "path": rp,
                "exists": True,
                "is_dir": is_dir,
                "is_file": not is_dir,
                "size": int(st.st_size),
                "modified_time": int(st.st_mtime),
                "file_type": ftype or ("inode/directory" if is_dir else "application/octet-stream"),
            }
        except Exception as e:
            return False, str(e), {"path": rp, "exists": True}

    # 文档文件读取工具（PDF/Word/Excel）
    if t == "read_pdf":
        rp = str(tool.get("path") or "").strip()
        if not rp:
            return False, "path 不能为空", {}
        full = _safe_repo_abspath(rp)
        if not full or not os.path.exists(full) or os.path.isdir(full):
            return False, "文件不存在", {}
        try:
            import PyPDF2
            page_num = tool.get("page_num")
            max_pages = int(tool.get("max_pages") or 50)
            texts = []
            with open(full, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                total = len(reader.pages)
                if page_num is not None:
                    idx = int(page_num) - 1
                    if 0 <= idx < total:
                        texts.append(f"--- Page {page_num} ---\n" + (reader.pages[idx].extract_text() or ""))
                    else:
                        return False, f"页码超出范围（共 {total} 页）", {}
                else:
                    for i in range(min(total, max_pages)):
                        texts.append(f"--- Page {i + 1} ---\n" + (reader.pages[i].extract_text() or ""))
            return True, "", {"path": rp, "total_pages": total, "content": "\n\n".join(texts)}
        except ImportError:
            return False, "未安装 PyPDF2，请先执行 pip install PyPDF2", {}
        except Exception as e:
            return False, str(e), {}

    if t == "read_docx":
        rp = str(tool.get("path") or "").strip()
        if not rp:
            return False, "path 不能为空", {}
        full = _safe_repo_abspath(rp)
        if not full or not os.path.exists(full) or os.path.isdir(full):
            return False, "文件不存在", {}
        try:
            import docx
            doc = docx.Document(full)
            texts = []
            for para in doc.paragraphs:
                if para.text:
                    texts.append(para.text)
            return True, "", {"path": rp, "content": "\n".join(texts)}
        except ImportError:
            return False, "未安装 python-docx，请先执行 pip install python-docx", {}
        except Exception as e:
            return False, str(e), {}

    if t == "read_xlsx":
        rp = str(tool.get("path") or "").strip()
        if not rp:
            return False, "path 不能为空", {}
        full = _safe_repo_abspath(rp)
        if not full or not os.path.exists(full) or os.path.isdir(full):
            return False, "文件不存在", {}
        try:
            import openpyxl
            wb = openpyxl.load_workbook(full, data_only=True)
            sheet_name = str(tool.get("sheet_name") or "").strip()
            if sheet_name and sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
            else:
                ws = wb.active
            max_rows = int(tool.get("max_rows") or 200)
            rows = []
            for idx, row in enumerate(ws.iter_rows(values_only=True), start=1):
                if idx > max_rows:
                    break
                rows.append([str(cell) if cell is not None else "" for cell in row])
            return True, "", {"path": rp, "sheet": ws.title, "rows": rows}
        except ImportError:
            return False, "未安装 openpyxl，请先执行 pip install openpyxl", {}
        except Exception as e:
            return False, str(e), {}

    if t in ("read_file_range",):
        rp = _hivo_resolve_path_if_needed(str(tool.get("path") or "").strip())
        start = int(tool.get("start") or 1)
        end = int(tool.get("end") or start)
        content, encoding = get_file_content(rp, return_encoding=True)
        if content is None:
            return False, "无法读取文件内容", {"path": rp}
        lines = str(content or "").splitlines()
        start_i = max(1, start)
        end_i = max(start_i, end)
        if start_i > len(lines):
            return True, "", {"file": {"path": rp, "encoding": encoding, "content": ""}}
        seg = lines[start_i - 1:min(len(lines), end_i)]
        return True, "", {"file": {"path": rp, "encoding": encoding, "content": "\n".join(seg), "range": {"start": start_i, "end": min(len(lines), end_i)}}}

    if t in ("list_dir_tree",):
        p0 = str(tool.get("path") or "").strip() or "."
        depth = int(tool.get("depth") or 4)
        max_entries = int(tool.get("max_entries") or 800)
        out0, err0 = list_dir_tree(p0, max_depth=depth, max_entries=max_entries)
        if err0:
            return False, err0, {"path": p0}
        return True, "", {"path": p0, "depth": depth, "max_entries": max_entries, "entries": out0}

    if t in ("search_code",):
        q = str(tool.get("query") or "").strip()
        case_sensitive = bool(tool.get("case_sensitive"))
        max_results = int(tool.get("max_results") or 50)
        hits, err = search_code(q, case_sensitive=case_sensitive, max_results=max_results)
        if err:
            return False, str(err or "search_code failed"), {}
        return True, "", {"query": q, "case_sensitive": case_sensitive, "max_results": max_results, "hits": hits if isinstance(hits, list) else []}

    if t in ("find_files",):
        name = str(tool.get("name") or "").strip()
        hits = find_files_by_name(name, max_results=50)
        return True, "", {"name": name, "files": hits}

    if t in ("list_files",):
        try:
            max_age = float(tool.get("max_age") or 0)
        except Exception:
            max_age = 0.0
        if max_age < 0:
            max_age = 0.0
        if max_age > 10:
            max_age = 10.0
        files0 = get_changed_files_cached(max_age_sec=max_age)
        return True, "", {"max_age": max_age, "files": files0}

    if t in ("diff_file",):
        rp = _hivo_resolve_path_if_needed(str(tool.get("path") or "").strip())
        if not rp:
            return False, "缺少 path", {}
        st = str(tool.get("status") or "M").strip() or "M"
        try:
            ctx_lines = int(tool.get("ctx_lines") if (tool.get("ctx_lines") is not None) else (tool.get("ctx") if (tool.get("ctx") is not None) else 80))
        except Exception:
            ctx_lines = 80
        diff0, err0 = get_file_diff(rp, st, ctx_lines=ctx_lines)
        if err0:
            return False, err0, {"path": rp}
        return True, "", {"file": {"path": rp, "diff": diff0, "status": st, "ctx_lines": ctx_lines}}

    if t in ("staged_files",):
        files0, err0 = get_staged_files()
        if err0:
            return False, err0, {}
        return True, "", {"files": files0}

    if t == "stash_list":
        out, err, code = run_git(["stash", "list"], timeout=30)
        if code != 0:
            return False, (err or out or "stash list failed"), {"output": (out or "").strip()}
        raw = (out or "").strip()
        stashes = [ln.strip() for ln in raw.splitlines() if ln.strip()]
        top = stashes[0] if stashes else ""
        return True, "", {"has_stash": bool(stashes), "stash_count": len(stashes), "stash_top": top, "stashes": stashes}

    if t == "stash_drop":
        ref = str(tool.get("ref") or "").strip()
        if not ref:
            return False, "缺少 ref", {}
        out, err, code = run_git(["stash", "drop", ref], timeout=60)
        if code != 0:
            return False, (err or out or "stash drop failed"), {"output": (out or "").strip()}
        try:
            invalidate_changed_files_cache()
            notify_files_updated()
        except Exception as e:
            logger.debug(f"Exception ignored: {e}")
        return True, "", {"ref": ref, "output": (out or "").strip()}

    if t == "unstage_file":
        rp = str(tool.get("path") or "").strip()
        if not rp:
            return False, "缺少 path", {}
        ok, msg = unstage_file(rp)
        if not ok:
            return False, msg or "unstage_file failed", {"path": rp}
        invalidate_changed_files_cache()
        notify_files_updated()
        return True, "", {"path": rp}

    if t == "discard_staged_file":
        rp = str(tool.get("path") or "").strip()
        if not rp:
            return False, "缺少 path", {}
        ok, msg = discard_staged_file(rp)
        if not ok:
            return False, msg or "discard_staged_file failed", {"path": rp}
        invalidate_changed_files_cache()
        notify_files_updated()
        return True, "", {"path": rp}

    if t == "unstage_all_staged":
        ok, msg = unstage_all_staged()
        if not ok:
            return False, msg or "unstage_all_staged failed", {}
        invalidate_changed_files_cache()
        notify_files_updated()
        return True, "", {}

    if t == "discard_all_staged":
        ok, msg = discard_all_staged()
        if not ok:
            return False, msg or "discard_all_staged failed", {}
        invalidate_changed_files_cache()
        notify_files_updated()
        return True, "", {}


    # Write
    if t in ("save_file",):
        rp = str(tool.get("path") or "").strip()
        content = str(tool.get("content") or "")
        try:
            cur = get_file_content(rp)
            if isinstance(cur, str) and cur == content:
                return True, "", {"path": rp, "no_change": True}
        except Exception as e:
            logger.debug(f"Exception ignored: {e}")
        if undo_gid:
            snap = _undo_capture_file_snapshot(rp)
            if snap is not None:
                op0 = "add" if (isinstance(snap, dict) and (snap.get("exists") is False)) else "modify"
                _undo_record(undo_gid, {"type": "file_snapshot", "op": op0, "snapshot": snap})
        ok, msg = save_file_content(rp, content)
        if ok:
            invalidate_changed_files_cache()
            notify_files_updated()
        return bool(ok), msg or "", {"path": rp}

    if t == "edit_file":
        rp = str(tool.get("path") or "").strip()
        old_str = str(tool.get("old_string") or "")
        new_str = str(tool.get("new_string") or "")
        if not rp:
            return False, "path 不能为空", {}
        if not old_str:
            return False, "old_string 不能为空", {}
        try:
            cur = get_file_content(rp)
            if not isinstance(cur, str):
                return False, "文件不是文本文件", {}
        except Exception as e:
            return False, f"读取文件失败: {e}", {}
        # 检查 old_string 是否唯一匹配
        count = cur.count(old_str)
        if count == 0:
            return False, "未找到匹配的 old_string", {"found": 0}
        if count > 1:
            return False, f"old_string 不唯一，匹配到 {count} 处，请提供更多上下文", {"found": count}
        # 无变化则直接返回
        if old_str == new_str:
            return True, "", {"path": rp, "no_change": True}
        # 保存快照用于 undo
        if undo_gid:
            snap = _undo_capture_file_snapshot(rp)
            if snap is not None:
                op0 = "modify"
                _undo_record(undo_gid, {"type": "file_snapshot", "op": op0, "snapshot": snap})
        new_content = cur.replace(old_str, new_str, 1)
        ok, msg = save_file_content(rp, new_content)
        if ok:
            invalidate_changed_files_cache()
            notify_files_updated()
        return bool(ok), msg or "", {"path": rp, "replaced": 1}

    if t == "delete_file":
        rp = str(tool.get("path") or "").strip()
        if undo_gid:
            snap = _undo_capture_file_snapshot(rp)
            if snap is not None:
                _undo_record(undo_gid, {"type": "file_snapshot", "op": "delete", "snapshot": snap})
        full = _safe_repo_abspath(rp)
        if not full:
            return False, "非法路径", {}
        try:
            if os.path.isdir(full):
                return False, "目标是目录", {}
            if os.path.exists(full):
                os.remove(full)
            invalidate_changed_files_cache()
            notify_files_updated()
            return True, "", {"path": rp}
        except Exception as e:
            return False, str(e), {}

    if t == "rename_file":
        oldp = str(tool.get("old_path") or "").strip()
        newp = str(tool.get("new_path") or "").strip()
        if undo_gid:
            snap = _undo_capture_file_snapshot(oldp)
            if snap is not None:
                _undo_record(undo_gid, {"type": "file_snapshot", "snapshot": snap})
        ok1, msg1 = rename_file(oldp, newp)
        if ok1:
            invalidate_changed_files_cache()
            notify_files_updated()
        return bool(ok1), msg1 or "", {"old_path": oldp, "new_path": newp}

    if t == "copy_file":
        src = str(tool.get("source") or "").strip()
        dst = str(tool.get("destination") or "").strip()
        if not src or not dst:
            return False, "source 和 destination 不能为空", {}
        full_src = _safe_repo_abspath(src)
        full_dst = _safe_repo_abspath(dst)
        if not full_src or not full_dst:
            return False, "非法路径", {}
        if not os.path.exists(full_src):
            return False, "源文件不存在", {}
        try:
            # 确保目标目录存在
            dst_dir = os.path.dirname(full_dst)
            if dst_dir and not os.path.exists(dst_dir):
                os.makedirs(dst_dir, exist_ok=True)
            # 如果目标是已存在的目录，则将源复制到目录内
            if os.path.isdir(full_dst):
                dst_name = os.path.basename(full_src)
                full_dst = os.path.join(full_dst, dst_name)
            if os.path.isdir(full_src):
                import shutil
                shutil.copytree(full_src, full_dst)
                is_dir = True
            else:
                import shutil
                shutil.copy2(full_src, full_dst)
                is_dir = False
            invalidate_changed_files_cache()
            notify_files_updated()
            return True, "", {"source": src, "destination": dst, "is_dir": is_dir}
        except Exception as e:
            return False, str(e), {}

    if t == "mkdir":
        rp = str(tool.get("path") or "").strip()
        ok1, msg1 = mkdir_repo(rp)
        if ok1:
            invalidate_changed_files_cache()
            notify_files_updated()
        return bool(ok1), msg1 or "", {"path": rp}

    if t == "browser_open":
        sid = str(tool.get("session_id") or run_id or "global").strip() or "global"
        u = str(tool.get("url") or "").strip()
        wait_ms = int(tool.get("wait_ms") or 0)
        return _pw_open(sid, u, wait_ms=wait_ms)

    if t == "browser_click":
        sid = str(tool.get("session_id") or run_id or "global").strip() or "global"
        sel = str(tool.get("selector") or "").strip()
        wait_ms = int(tool.get("wait_ms") or 0)
        return _pw_click(sid, sel, wait_ms=wait_ms)

    if t == "browser_type":
        sid = str(tool.get("session_id") or run_id or "global").strip() or "global"
        sel = str(tool.get("selector") or "").strip()
        txt = str(tool.get("text") or "")
        clear = bool(tool.get("clear", True))
        wait_ms = int(tool.get("wait_ms") or 0)
        return _pw_type(sid, sel, txt, clear=clear, wait_ms=wait_ms)

    if t == "browser_eval":
        sid = str(tool.get("session_id") or run_id or "global").strip() or "global"
        js = str(tool.get("script") or "")
        return _pw_eval(sid, js)

    if t == "browser_text":
        sid = str(tool.get("session_id") or run_id or "global").strip() or "global"
        sel = str(tool.get("selector") or "").strip()
        return _pw_text(sid, selector=sel)

    if t == "browser_screenshot":
        sid = str(tool.get("session_id") or run_id or "global").strip() or "global"
        full_page = bool(tool.get("full_page", True))
        return _pw_screenshot(sid, full_page=full_page)

    if t == "browser_wait":
        sid = str(tool.get("session_id") or run_id or "global").strip() or "global"
        ms = int(tool.get("ms") or 0)
        return _pw_wait(sid, ms=ms)

    if t == "browser_close":
        sid = str(tool.get("session_id") or run_id or "global").strip() or "global"
        if not _pw_is_available():
            return False, "Playwright 未安装（可选能力）。", {}
        _pw_close_session(sid)
        return True, "", {"session_id": sid}

    if t == "revert_file":
        rp = str(tool.get("path") or "").strip()
        st = str(tool.get("status") or "M").strip() or "M"
        ok, msg = revert_file(rp, st)
        if ok:
            invalidate_changed_files_cache()
            notify_files_updated()
        return bool(ok), msg or "", {"path": rp}

    if t == "revert_hunk":
        rp = str(tool.get("path") or "").strip()
        idx = int(tool.get("hunk_index") or 0)
        st = str(tool.get("status") or "M").strip() or "M"
        ok, msg = revert_hunk(rp, idx, st)
        if ok:
            invalidate_changed_files_cache()
            notify_files_updated()
        return bool(ok), msg or "", {"path": rp, "hunk_index": idx}

    if t == "revert_line":
        rp = str(tool.get("path") or "").strip()
        idx = int(tool.get("hunk_index") or 0)
        li = int(tool.get("line_index") or 0)
        st = str(tool.get("status") or "M").strip() or "M"
        ok, msg = revert_line(rp, idx, li, st)
        if ok:
            invalidate_changed_files_cache()
            notify_files_updated()
        return bool(ok), msg or "", {"path": rp, "hunk_index": idx, "line_index": li}

    if t == "revert_multi_lines":
        rp = str(tool.get("path") or "").strip()
        idx = int(tool.get("hunk_index") or 0)
        s0 = int(tool.get("start_line_index") or 0)
        e0 = int(tool.get("end_line_index") or 0)
        st = str(tool.get("status") or "M").strip() or "M"
        ok, msg = revert_multi_lines(rp, idx, s0, e0, st)
        if ok:
            invalidate_changed_files_cache()
            notify_files_updated()
        return bool(ok), msg or "", {"path": rp, "hunk_index": idx, "start": s0, "end": e0}

    if t == "revert_all":
        arr = tool.get("paths")
        if not isinstance(arr, list):
            arr = []
        ok, msg = revert_all(arr)
        if ok:
            invalidate_changed_files_cache()
            notify_files_updated()
        return bool(ok), msg or "", {"paths": arr}

    # Git
    if t == "branches":
        try:
            brs, cur = get_branches()
            return True, "", {"branches": brs, "current": cur}
        except Exception as e:
            return False, str(e), {}

    if t == "commits":
        try:
            limit = int(tool.get("limit") or 50)
        except Exception:
            limit = 50
        out0 = get_log(limit=limit)
        return True, "", {"limit": limit, "commits": out0}

    if t == "commit_detail":
        h0 = str(tool.get("hash") or "").strip()
        if not h0:
            return False, "缺少 hash", {}
        try:
            detail = get_commit_detail(h0)
            if not isinstance(detail, dict) or ("error" in detail):
                return False, str((detail or {}).get("error") or "获取提交详情失败"), {}
            return True, "", {"commit": detail}
        except Exception as e:
            return False, str(e), {}

    if t == "commit_file_diff":
        h0 = str(tool.get("hash") or "").strip()
        p0 = str(tool.get("path") or "").strip().replace("\\", "/").lstrip("/")
        if not h0:
            return False, "缺少 hash", {}
        if not p0:
            return False, "缺少 path", {}
        try:
            hunks = get_commit_file_diff(h0, p0)
            return True, "", {"hash": h0, "path": p0, "hunks": hunks}
        except Exception as e:
            return False, str(e), {}

    if t == "commit_push_status":
        h0 = str(tool.get("hash") or "").strip()
        if not h0:
            return False, "缺少 hash", {}
        try:
            pushed, branches, err = is_commit_pushed(h0)
            return True, "", {"hash": h0, "pushed": pushed, "branches": branches, "error": err}
        except Exception as e:
            return False, str(e), {}

    if t == "stage_file":
        rp = str(tool.get("path") or "").strip()
        out, err, code = run_git(["add", "--", rp])
        if code != 0:
            return False, (err or out or "stage failed"), {}
        invalidate_changed_files_cache()
        notify_files_updated()
        return True, "", {"path": rp}

    if t == "commit":
        msg = str(tool.get("message") or "").strip()
        files0 = tool.get("files")
        if isinstance(files0, list) and files0:
            for rp in files0:
                run_git(["add", "--", str(rp)])
        out, err, code = run_git(["commit", "-m", msg])
        if code != 0:
            return False, (err or out or "commit failed"), {}
        invalidate_changed_files_cache()
        notify_files_updated()
        return True, "", {"output": out}

    if t in ("fetch_remotes",):
        ok, msg, out = fetch_remotes()
        if not ok:
            return False, msg or "fetch_remotes failed", {"output": (out or "")}
        try:
            invalidate_changed_files_cache()
            notify_files_updated()
        except Exception as e:
            logger.debug(f"Exception ignored: {e}")
        return True, "", {"output": (out or "")}

    if t in ("pull_safe",):
        out, err, code = run_git(["pull", "--no-edit"], timeout=300)
        output = (out or "")
        error = (err or "")
        local_changes_conflict = False
        affected_files = []
        try:
            if code != 0 and (
                ("would be overwritten" in error.lower()) or
                ("will be overwritten" in error.lower()) or
                ("本地修改" in error) or
                ("覆盖" in error)
            ):
                local_changes_conflict = True
                lines = error.split("\n")
                in_file_list = False
                for line in lines:
                    line = (line or "").strip()
                    if ("would be overwritten" in line.lower()) or ("will be overwritten" in line.lower()):
                        in_file_list = True
                        continue
                    if in_file_list:
                        if line.startswith("Please") or line.startswith("Aborting") or (not line):
                            break
                        if line and (not line.startswith("error:")) and (not line.startswith("hint:")):
                            affected_files.append(line.strip())
        except Exception as e:
            logger.debug(f"Exception ignored: {e}")

        conflict_files, _ = get_unmerged_files()
        has_conflicts = bool(conflict_files)
        ok = (code == 0) and (not has_conflicts)
        if ok:
            try:
                invalidate_changed_files_cache()
                notify_files_updated()
            except Exception as e:
                logger.debug(f"Exception ignored: {e}")
        return True, "", {
            "ok": ok,
            "output": output.strip(),
            "error": error.strip(),
            "has_conflicts": has_conflicts,
            "conflict_files": conflict_files,
            "local_changes_conflict": local_changes_conflict,
            "affected_files": affected_files,
        }

    if t == "stash_and_pull":
        stash_out, stash_err, stash_code = run_git(["stash", "push", "-u", "-m", "Auto stash before pull"], timeout=60)
        if stash_code != 0:
            return False, (stash_err or stash_out or "git stash 失败"), {"output": (stash_out or "").strip(), "error": (stash_err or "").strip()}
        stashed = "No local changes to save" not in (stash_out or "")

        pull_out, pull_err, pull_code = run_git(["pull", "--no-edit"], timeout=300)
        conflict_files, _ = get_unmerged_files()
        has_conflicts = bool(conflict_files)
        ok = (pull_code == 0) and (not has_conflicts)
        try:
            invalidate_changed_files_cache()
            notify_files_updated()
        except Exception as e:
            logger.debug(f"Exception ignored: {e}")
        return True, "", {
            "ok": ok,
            "stashed": stashed,
            "pending_pop": bool(stashed),
            "pull_output": (pull_out or "").strip(),
            "pull_error": (pull_err or "").strip(),
            "has_conflicts": has_conflicts,
            "conflict_files": conflict_files,
            "message": "成功暂存修改并更新" if ok else ("更新未完成：你的本地修改已暂存（可稍后手动恢复）" if stashed else "更新未完成"),
        }

    if t == "commit_and_pull":
        commit_msg = str(tool.get("message") or "").strip()
        if not commit_msg:
            return False, "提交信息不能为空", {}

        add_out, add_err, add_code = run_git(["add", "-A"], timeout=60)
        if add_code != 0:
            return False, (add_err or add_out or "暂存文件失败"), {}

        commit_out, commit_err, commit_code = run_git(["commit", "-m", commit_msg], timeout=60)
        if commit_code != 0:
            return False, (commit_err or commit_out or "提交失败"), {}

        pull_out, pull_err, pull_code = run_git(["pull", "--no-edit"], timeout=300)
        conflict_files, _ = get_unmerged_files()
        has_conflicts = bool(conflict_files)
        ok = (pull_code == 0) and (not has_conflicts)
        try:
            invalidate_changed_files_cache()
            notify_files_updated()
        except Exception as e:
            logger.debug(f"Exception ignored: {e}")
        return True, "", {
            "ok": ok,
            "commit_output": (commit_out or "").strip(),
            "pull_output": (pull_out or "").strip(),
            "pull_error": (pull_err or "").strip(),
            "has_conflicts": has_conflicts,
            "conflict_files": conflict_files,
            "message": "成功提交并更新" if ok else "提交成功但更新时发生冲突",
        }

    if t == "push":
        out, err, code = run_git(["push"])
        if code != 0:
            return False, (err or out or "push failed"), {}
        return True, "", {"output": out}

    if t == "stash_pop":
        ref = str(tool.get("ref") or "").strip()
        cmd = ["stash", "pop"]
        if ref:
            cmd.append(ref)
        out, err, code = run_git(cmd)
        if code != 0:
            list_out, list_err, list_code = run_git(["stash", "list"], timeout=30)
            stash_kept = False
            top = ""
            try:
                raw = (list_out or "").strip()
                if raw:
                    stash_kept = True
                    top = raw.splitlines()[0].strip()
            except Exception:
                stash_kept = False
                top = ""
            conflict_files, _ = get_unmerged_files()
            return False, (err or out or "stash pop failed"), {
                "has_conflict": bool(conflict_files) or ("CONFLICT" in (out or "")) or ("CONFLICT" in (err or "")),
                "conflict_files": conflict_files,
                "stash_kept": stash_kept,
                "stash_top": top,
                "output": (out or "").strip(),
                "error_detail": (err or "").strip(),
            }
        try:
            invalidate_changed_files_cache()
            notify_files_updated()
        except Exception as e:
            logger.debug(f"Exception ignored: {e}")
        return True, "", {"output": out}

    if t in ("switch_branch_safe",):
        target_branch = str(tool.get("branch") or "").strip()
        if not target_branch:
            return False, "未指定目标分支", {}

        # 获取当前分支
        current_out, current_err, current_code = run_git(["branch", "--show-current"], timeout=30)
        current_branch = (current_out or "").strip()
        if current_code != 0:
            return False, (current_err or "获取当前分支失败"), {}
        if target_branch == current_branch:
            return True, "", {"ok": True, "current": current_branch, "message": "已在目标分支上"}

        status_out, _, _ = run_git(["status", "--porcelain"], timeout=30)
        has_changes = bool((status_out or "").strip())

        is_remote, remote_ref, _raw = _normalize_remote_ref(target_branch)
        want_detached = bool(tool.get("detached"))

        if not has_changes:
            if is_remote and want_detached:
                remote_ref = remote_ref or target_branch.replace("remotes/", "", 1)
                out, err, code = run_git(["switch", "--detach", remote_ref], timeout=120)
                if code != 0:
                    out, err, code = run_git(["checkout", remote_ref], timeout=120)
                if code != 0:
                    return True, "", {"ok": False, "error": err or "切换到远端分支失败", "output": out or ""}
                _, cur = get_branches()
                try:
                    invalidate_changed_files_cache()
                    notify_files_updated()
                except Exception as e:
                    logger.debug(f"Exception ignored: {e}")
                return True, "", {"ok": True, "current": cur, "message": f"成功切换到分支 {cur}"}

            ok2, cur2, err_msg, out_msg, _safe_err = switch_branch(target_branch)
            if ok2:
                try:
                    invalidate_changed_files_cache()
                    notify_files_updated()
                except Exception as e:
                    logger.debug(f"Exception ignored: {e}")
                return True, "", {"ok": True, "current": cur2, "message": f"成功切换到分支 {cur2}"}
            return True, "", {"ok": False, "error": err_msg or "切换分支失败", "output": out_msg or ""}

        # 有本地修改：尝试切换，失败时解析"会覆盖"的受影响文件列表
        if is_remote and (not want_detached):
            return True, "", {
                "ok": False,
                "needs_handling": True,
                "has_uncommitted_changes": True,
                "affected_files": [],
                "error": "工作区有未提交的修改，请先提交/暂存后再从远端分支创建/切换本地分支",
                "message": "工作区有未提交的修改，请先处理后再切换远端分支",
            }

        test_out, test_err, test_code = run_git(["checkout", target_branch], timeout=60)
        if test_code == 0:
            try:
                invalidate_changed_files_cache()
                notify_files_updated()
            except Exception as e:
                logger.debug(f"Exception ignored: {e}")
            return True, "", {
                "ok": True,
                "current": target_branch,
                "has_uncommitted_changes": True,
                "message": f"成功切换到分支 {target_branch}，未提交的修改已保留",
            }

        error_msg = (test_err or "").lower()
        if ("would be overwritten" in error_msg) or ("overwritten by checkout" in error_msg):
            affected_files = []
            try:
                lines = (test_err or "").split("\n")
                in_file_list = False
                for line in lines:
                    line = (line or "").strip()
                    if ("would be overwritten" in line.lower()) or ("overwritten by checkout" in line.lower()):
                        in_file_list = True
                        continue
                    if in_file_list:
                        if line.startswith("Please") or line.startswith("Aborting") or (not line):
                            break
                        if line and (not line.startswith("error:")) and (not line.startswith("hint:")):
                            affected_files.append(line.strip())
            except Exception:
                affected_files = []
            return True, "", {
                "ok": False,
                "needs_handling": True,
                "has_uncommitted_changes": True,
                "affected_files": affected_files,
                "error": "工作区有未提交的修改会被覆盖",
                "message": "切换分支会覆盖当前未提交的修改，请先处理这些修改",
            }

        return True, "", {"ok": False, "error": (test_err or "切换分支失败"), "output": (test_out or "")}

    if t == "stash_and_switch":
        target_branch = str(tool.get("branch") or "").strip()
        if not target_branch:
            return False, "未指定目标分支", {}

        stash_out, stash_err, stash_code = run_git(
            ["stash", "push", "-u", "-m", f"Auto stash before switching to {target_branch}"],
            timeout=60,
        )
        if stash_code != 0:
            return False, (stash_err or stash_out or "git stash 失败"), {"output": (stash_out or "").strip(), "error": (stash_err or "").strip()}
        stashed = "No local changes to save" not in (stash_out or "")

        dirty2, detail2 = _has_worktree_changes()
        if dirty2:
            if stashed:
                try:
                    run_git(["stash", "pop"], timeout=60)
                except Exception as e:
                    logger.debug(f"Exception ignored: {e}")
            return False, "暂存后工作区仍存在未提交更改，无法安全切换分支（可能包含未被 stash 的变更）", {"output": detail2 or ""}

        ok2, cur2, err_msg, out_msg, _safe_err = switch_branch(target_branch)
        try:
            invalidate_changed_files_cache()
            notify_files_updated()
        except Exception as e:
            logger.debug(f"Exception ignored: {e}")
        if ok2:
            return True, "", {"ok": True, "current": cur2, "stashed": stashed, "has_stash": bool(stashed), "message": f"成功切换到分支 {cur2}"}
        return True, "", {"ok": False, "error": err_msg or "切换分支失败", "output": out_msg or "", "stashed": stashed, "has_stash": bool(stashed)}

    if t == "commit_and_switch":
        target_branch = str(tool.get("branch") or "").strip()
        commit_msg = str(tool.get("message") or "").strip()
        if not target_branch:
            return False, "未指定目标分支", {}
        if not commit_msg:
            return False, "提交信息不能为空", {}

        add_out, add_err, add_code = run_git(["add", "-A"], timeout=60)
        if add_code != 0:
            return False, (add_err or add_out or "暂存文件失败"), {}

        commit_out, commit_err, commit_code = run_git(["commit", "-m", commit_msg], timeout=60)
        if commit_code != 0:
            return False, (commit_err or commit_out or "提交失败"), {}

        ok2, cur2, err_msg, out_msg, _safe_err = switch_branch(target_branch)
        try:
            invalidate_changed_files_cache()
            notify_files_updated()
        except Exception as e:
            logger.debug(f"Exception ignored: {e}")
        if ok2:
            return True, "", {"ok": True, "current": cur2, "commit_output": (commit_out or "").strip(), "message": f"成功提交并切换到分支 {cur2}"}
        return True, "", {"ok": False, "error": err_msg or "切换分支失败", "output": out_msg or "", "commit_output": (commit_out or "").strip()}

    # Undo
    if t in ("undo_last_turn", "undo"):
        gid, actions = _undo_pop_latest_group()
        if not actions:
            return False, "无可撤销操作", {}
        ok, msg = _undo_apply_actions(actions)
        if not ok:
            return False, msg or "撤销失败", {"group_id": gid}
        invalidate_changed_files_cache()
        notify_files_updated()
        return True, "", {"group_id": gid}

    if t == "verify_python":
        arr = tool.get("paths") or tool.get("files") or []
        if not isinstance(arr, list):
            arr = []
        paths = [str(x or "").strip().lstrip("@") for x in arr]
        paths = [p for p in paths if p]
        if not paths:
            paths = ["server.py"]
        checked = []
        for rp in paths[:50]:
            safe = _safe_repo_abspath(rp)
            if not safe or (not os.path.isfile(safe)):
                return False, f"非法路径或文件不存在: {rp}", {"files": checked}
            checked.append(rp.replace("\\", "/"))
            try:
                py_compile.compile(safe, doraise=True)
            except py_compile.PyCompileError as e:
                return False, str(e), {"files": checked}
        return True, "", {"files": checked}

    if t == "run_cmd":
        cmd = str(tool.get("cmd") or "").strip()
        timeout = int(tool.get("timeout") or 30)
        cwd = str(tool.get("cwd") or "").strip()
        ok1, msg1, out1 = _run_cmd_simple(cmd, timeout=timeout, cwd=cwd, run_id=run_id, agent_deadline=agent_deadline, undo_gid=undo_gid)
        if not ok1:
            return False, msg1 or "run_cmd failed", {"result": {"cmd": cmd, "timeout": timeout, "cwd": cwd, "output": out1}}
        return True, "", {"result": {"cmd": cmd, "timeout": timeout, "cwd": cwd, "output": out1}}

    if t == "list_processes":
        try:
            import psutil
            flt = str(tool.get("filter") or "").strip().lower()
            limit = int(tool.get("limit") or 50)
            procs = []
            for p in psutil.process_iter(["pid", "name", "exe", "cmdline", "connections", "status"]):
                try:
                    info = p.info
                    name = str(info.get("name") or "").lower()
                    pid = int(info.get("pid") or 0)
                    # 按名称过滤
                    if flt and flt not in name and flt != str(pid):
                        continue
                    # 获取端口
                    ports = []
                    try:
                        for conn in p.connections(kind="inet"):
                            if conn.laddr:
                                ports.append(str(conn.laddr.port))
                    except Exception:
                        pass
                    cmdline = ""
                    try:
                        cl = info.get("cmdline")
                        if cl:
                            cmdline = " ".join(str(c) for c in cl)
                    except Exception:
                        pass
                    exe = str(info.get("exe") or "")
                    procs.append({
                        "pid": pid,
                        "name": str(info.get("name") or ""),
                        "exe": exe,
                        "cmdline": cmdline,
                        "ports": ports,
                        "status": str(info.get("status") or ""),
                    })
                except Exception:
                    pass
                if len(procs) >= limit:
                    break
            return True, "", {"processes": procs, "total": len(procs)}
        except ImportError:
            return False, "未安装 psutil，请先执行 pip install psutil", {}
        except Exception as e:
            return False, str(e), {}

    if t == "env_info":
        try:
            etype = str(tool.get("type") or "all").strip().lower()
            import sys
            import platform
            result = {}
            if etype in ("all", "python"):
                result["python_version"] = sys.version
                result["python_executable"] = sys.executable
                result["python_path"] = sys.path
            if etype in ("all", "system"):
                result["os"] = platform.platform()
                result["machine"] = platform.machine()
                result["processor"] = platform.processor()
            if etype in ("all", "packages"):
                try:
                    import pkg_resources
                    installed = []
                    for dist in pkg_resources.working_set:
                        installed.append({"name": dist.project_name, "version": dist.version})
                    result["packages"] = sorted(installed, key=lambda x: x["name"].lower())[:100]
                except Exception:
                    result["packages"] = "无法获取（缺少 setuptools）"
            return True, "", result
        except Exception as e:
            return False, str(e), {}

    if t == "web_search":
        q = str(tool.get("query") or "").strip()
        try:
            top_k = int(tool.get("top_k") or 5)
        except Exception:
            top_k = 5
        try:
            timeout = int(tool.get("timeout") or 20)
        except Exception:
            timeout = 20
        ok1, msg1, data1 = _hivo_web_search(q, top_k=top_k, timeout=timeout)
        if not ok1:
            return False, msg1 or "web_search failed", data1 if isinstance(data1, dict) else {"result": {"query": q, "items": []}}
        return True, "", data1 if isinstance(data1, dict) else {"result": {"query": q, "items": []}}

    if t == "web_fetch":
        u = str(tool.get("url") or "").strip()
        try:
            timeout = int(tool.get("timeout") or 20)
        except Exception:
            timeout = 20
        ok1, msg1, data1 = _hivo_web_fetch(u, timeout=timeout)
        if not ok1:
            return False, msg1 or "web_fetch failed", data1 if isinstance(data1, dict) else {"result": {"url": u, "text": ""}}
        return True, "", data1 if isinstance(data1, dict) else {"result": {"url": u, "text": ""}}

    return False, f"unsupported tool: {t}", {}


def hivo_agent_run(run_id: str, profile_id: str, session_id: str, user_text: str, history_messages: list | None = None, dyn_context: str = "", undo_gid: str = "", attachments: list | None = None, resume: dict | None = None):
    cfg = _hivo_load_cfg()
    agent_conf = _hivo_agent_conf(cfg)
    prof = _ai_pick_profile(cfg, profile_id)
    mem_conf = _hivo_mem_conf(cfg)
    feat = cfg.get("features") if isinstance(cfg, dict) else None
    feat = feat if isinstance(feat, dict) else {}
    memory_enabled = bool(feat.get("memory_enabled", True))
    tool_memory_enabled = bool(feat.get("tool_memory_enabled", True))
    tool_cache_enabled = bool(feat.get("tool_cache_enabled", True))
    topic_iso_enabled = bool(feat.get("topic_isolation_enabled", True))
    max_rounds = int(cfg.get("max_rounds") or 12)
    max_calls = int(cfg.get("max_tool_calls_per_round") or 3)
    max_hist = int(cfg.get("max_visible_history") or 80)

    try:
        cap_spec = get_capabilities_spec()
        agent_tools = (cap_spec.get("agent_tools") or []) if isinstance(cap_spec, dict) else []
    except Exception:
        agent_tools = []

    try:
        agent_timeout_s = int(cfg.get("agent_timeout_s") or 0)
    except Exception:
        agent_timeout_s = 0
    if agent_timeout_s <= 0:
        agent_timeout_s = 300

    # Agent 超时上限由配置决定，不再硬编码
    agent_deadline = time.time() + max(30, agent_timeout_s)

    try:
        llm_timeout_s = int(cfg.get("llm_timeout_s") or 0)
    except Exception:
        llm_timeout_s = 0
    if llm_timeout_s <= 0:
        llm_timeout_s = 45
    llm_timeout_s = max(10, min(90, int(llm_timeout_s)))
    repeat_conf = cfg.get("repeat_block") if isinstance(cfg.get("repeat_block"), dict) else {}
    rep_window = int(repeat_conf.get("window") or 3)
    rep_max_same = int(repeat_conf.get("max_same") or 2)
    rep_sig_mode = str(repeat_conf.get("signature") or "tool_types")
    rep_escalation_limit = int(repeat_conf.get("escalation_limit") or 3)

    undo_gid_eff = str(undo_gid or "").strip()
    if not undo_gid_eff:
        undo_gid_eff = f"ai-{session_id}-{run_id}" if session_id else f"ai-{run_id}"

    st = _hivo_get_session_state(session_id) if memory_enabled else {"summary": "", "chat": [], "tool_log": [], "tool_cache": OrderedDict(), "topic_id": 1, "topic_archive": []}

    # Topic isolation: if user explicitly switches topic, start a new scoped memory.
    try:
        if memory_enabled and topic_iso_enabled and _hivo_detect_topic_switch(user_text):
            _hivo_start_new_topic(st, reason=str(user_text or "")[:200])
            iso_note = "【话题已切换 / 上下文隔离】\n从这一轮开始，你必须将其视为一个全新的问题域：\n- 禁止引用上一话题的工具回执、推理结论或未相关上下文\n- 若需要旧信息，必须要求用户重新提供，或通过工具重新获取\n"
            dyn_context = (iso_note + ("\n" + str(dyn_context or "") if dyn_context else "")).strip()
    except Exception as e:
        logger.debug(f"Exception ignored: {e}")
    try:
        long_summary = str(st.get("summary") or "").strip()
    except Exception:
        long_summary = ""
    try:
        tool_log0 = st.get("tool_log") if isinstance(st.get("tool_log"), list) else []
    except Exception:
        tool_log0 = []
    try:
        tool_mem_block = _hivo_format_tool_memory_block(tool_log0, limit=int(mem_conf.get("tool_log_items") or 80)) if tool_memory_enabled else ""
    except Exception:
        tool_mem_block = ""
    try:
        tc0 = st.get("tool_cache")
        if not isinstance(tc0, OrderedDict):
            tc0 = OrderedDict(tc0 or {})
            st["tool_cache"] = tc0
        tool_cache = tc0
    except Exception:
        tool_cache = OrderedDict()

    hist = []
    try:
        base_mem = st.get("chat") if (memory_enabled and isinstance(st, dict)) else None
        if isinstance(base_mem, list) and base_mem:
            for m in base_mem:
                nm = _hivo_ai_normalize_message(m)
                if nm:
                    hist.append(nm)
        else:
            base = history_messages if isinstance(history_messages, list) else []
            for m in base:
                if isinstance(m, dict) and bool(m.get("pending")) and str(m.get("role") or "").strip() == "assistant":
                    continue
                nm = _hivo_ai_normalize_message(m)
                if nm:
                    hist.append(nm)
    except Exception:
        hist = []

    # 确保 tool 消息的 tool_call_id 与 assistant 消息的 tool_calls id 一致
    hist = _hivo_ai_fix_tool_call_ids(hist)

    # Short-term memory window: use max_visible_history (capped to avoid token overflow)
    keep_n = max_hist if max_hist > 0 else 80
    # Ensure a sane upper bound (max 160 messages = 80 turns) to prevent token overflow
    keep_n = max(6, min(160, int(keep_n)))
    if keep_n > 0 and len(hist) > keep_n:
        hist = hist[-keep_n:]

    def _format_text_attachments(att0: list):
        try:
            blocks = []
            for a in (att0 or []):
                if not isinstance(a, dict):
                    continue
                t = a.get("text")
                if t is None:
                    continue
                nm = str(a.get("name") or "").strip() or "file"
                mime = str(a.get("mime") or "").strip() or "text/plain"
                s = str(t or "")
                if not s.strip():
                    continue
                # avoid huge prompt injection
                if len(s) > 8000:
                    s = s[:8000] + "\n...（已截断）..."
                blocks.append(f"【文件：{nm} / {mime}】\n{s}")
            return "\n\n".join(blocks).strip()
        except Exception:
            return ""

    # Normalize attachments once and detect if any meaningful content is present
    atts = attachments if isinstance(attachments, list) else []
    def _has_meaningful_attachments(arr:list):
        try:
            for a in (arr or []):
                if not isinstance(a, dict):
                    continue
                u = str(a.get("url") or a.get("data_url") or "").strip()
                pth = str(a.get("path") or "").strip()
                t = a.get("text")
                if pth or u or (isinstance(t, str) and t.strip()):
                    return True
        except Exception:
            pass
        return False
    attachments_present = _has_meaningful_attachments(atts)

    # 自动提取文档附件（PDF/Word/Excel）的文本内容，回填到 text 字段
    # 主流 AI Agent 做法：上传文档后 AI 直接看到内容，无需额外工具调用
    def _extract_attachment_doc(a: dict) -> str:
        """尝试从附件中提取文档文本内容，失败返回空字符串。"""
        try:
            mime = str(a.get("mime") or "")
            if not any(k in mime for k in ["pdf", "wordprocessingml", "spreadsheetml"]):
                return ""
            # 解析文件路径
            pth = str(a.get("path") or "").strip()
            if not pth:
                # 从 url 解析 path
                u = str(a.get("url") or "").strip()
                if u and "/api/ai_attachment" in u and "path=" in u:
                    from urllib.parse import urlparse, parse_qs
                    pr = urlparse(u)
                    qs = parse_qs(pr.query or "")
                    p0 = (qs.get("path") or [""])[0]
                    pth = str(p0 or "").strip()
            if not pth:
                return ""
            _, full = _hivo_resolve_full_path(pth)
            if not full or not os.path.exists(full) or os.path.isdir(full):
                return ""
            if "pdf" in mime:
                import PyPDF2
                reader = PyPDF2.PdfReader(full)
                pages = []
                for i in range(min(len(reader.pages), 50)):
                    t = reader.pages[i].extract_text() or ""
                    pages.append(t)
                content = "\n".join(pages)
            elif "wordprocessingml" in mime:
                import docx
                doc = docx.Document(full)
                paras = [p.text for p in doc.paragraphs if p.text]
                content = "\n".join(paras)
            elif "spreadsheetml" in mime:
                import openpyxl
                wb = openpyxl.load_workbook(full, data_only=True)
                ws = wb.active
                rows = []
                for idx, row in enumerate(ws.iter_rows(values_only=True), start=1):
                    if idx > 200:
                        break
                    rows.append("\t".join(str(c) if c is not None else "" for c in row))
                content = "\n".join(rows)
            else:
                return ""
            if content.strip():
                return content[:8000] + ("\n...（已截断）..." if len(content) > 8000 else "")
        except ImportError:
            pass  # 依赖未安装，跳过
        except Exception:
            pass
        return ""
    try:
        for a in (atts or []):
            if not isinstance(a, dict):
                continue
            if a.get("text") is not None:
                continue  # 已有文本内容
            extracted = _extract_attachment_doc(a)
            if extracted:
                a["text"] = extracted
    except Exception:
        pass

    # 预置回执：即便模型未调用任何工具，但用户本轮携带了附件，也生成一条摘要回执便于前端展示"工具摘要"按钮
    pre_receipts: list[dict] = []
    try:
        if attachments_present:
            cnt = 0
            try:
                cnt = sum(1 for a in (atts or []) if isinstance(a, dict))
            except Exception:
                cnt = len(atts or [])
            pre_receipts.append({"type": "attachments", "ok": True, "msg": f"{cnt} 个附件", "data": {"count": cnt}})
    except Exception:
        pre_receipts = []

    if user_text:
        if atts:
            # 方案 B：OpenAI 多模态格式
            # 图片用 image_url 注入 content 数组；文本文件内容拼入 text 部分；文档文件自动提取文本
            # 构造 OpenAI 多模态 content 数组
            content_parts = []
            # 文本部分：用户文本 + 文本附件内容 + 文档附件内容（已回填到 text） + 其他附件提示
            txt_files = _format_text_attachments(atts)
            joined_text = []
            if str(user_text):
                joined_text.append(str(user_text))
            if txt_files:
                joined_text.append(txt_files)
            # 其他附件（非图片、无 text 内容）：告知 AI 文件存在
            other_files = []
            for a in (atts or []):
                if not isinstance(a, dict):
                    continue
                mime = str(a.get("mime") or "")
                name = str(a.get("name") or "").strip()
                size = a.get("size")
                if not name:
                    continue
                if mime.startswith("image/"):
                    continue  # 图片已单独注入 image_url
                if a.get("text") is not None and str(a.get("text") or "").strip():
                    continue  # 已有文本内容（原始文本或自动提取的文档内容）
                sz_str = f" ({size} bytes)" if isinstance(size, (int, float)) and size > 0 else ""
                other_files.append(f"{name} / {mime}{sz_str}")
            if other_files:
                joined_text.append("【上传的文件】\n" + "\n".join(other_files))
            text_content = "\n\n".join([x for x in joined_text if x]).strip()
            if text_content:
                content_parts.append({"type": "text", "text": text_content})

            # 图片部分：用 image_url 格式注入 base64 data_url
            for a in (atts or []):
                if not isinstance(a, dict):
                    continue
                mime = str(a.get("mime") or "")
                if not mime.startswith("image/"):
                    continue
                # 优先使用 data_url（base64），其次从 path 读取文件转 base64
                data_url = str(a.get("data_url") or "").strip()
                # 从 url 解析 path 作为 fallback
                if not data_url:
                    pth = str(a.get("path") or "").replace("\\", "/").strip()
                    # 若前端未提供 path，从 url 解析
                    if not pth:
                        try:
                            u = str(a.get("url") or "").strip()
                            if u and ("/api/ai_attachment" in u) and ("path=" in u):
                                from urllib.parse import urlparse, parse_qs
                                pr = urlparse(u)
                                if pr and pr.path and pr.path.endswith("/api/ai_attachment"):
                                    qs = parse_qs(pr.query or "")
                                    p0 = (qs.get("path") or [""])[0]
                                    pth = str(p0 or "").replace("\\", "/").lstrip("/")
                        except Exception:
                            pth = pth
                    if pth:
                        try:
                            if not (pth.startswith("attachments/") or pth.startswith("hivo_ai_data/attachments/")):
                                full = _safe_repo_abspath(pth)
                            else:
                                base_data = str(_hivo_ai_data_dir())
                                if pth.startswith("hivo_ai_data/attachments/"):
                                    suffix = pth.split("hivo_ai_data/", 1)[1]
                                else:
                                    suffix = pth
                                full = os.path.abspath(os.path.join(base_data, suffix))
                            if full and os.path.exists(full) and not os.path.isdir(full):
                                import base64 as _b64
                                with open(full, "rb") as _f:
                                    b64 = _b64.b64encode(_f.read()).decode("ascii")
                                data_url = f"data:{mime};base64,{b64}"
                        except Exception:
                            data_url = ""
                if data_url:
                    content_parts.append({
                        "type": "image_url",
                        "image_url": {"url": data_url}
                    })

            if content_parts:
                hist.append({"role": "user", "content": content_parts})
            else:
                hist.append({"role": "user", "content": str(user_text)})
        else:
            hist.append({"role": "user", "content": str(user_text)})

    sys_text = ""
    try:
        sys_text = str(agent_conf.get("system_prompt") or "")
    except Exception:
        sys_text = ""

    use_native_tools = _hivo_supports_native_tools(prof) if prof else False

    # Always append tool registry so custom prompts never hide available tools (e.g. run_cmd)
    try:
        tool_ctx = _ai_build_system_context_text(cfg, use_native_tools)
    except Exception:
        tool_ctx = ""
    if sys_text.strip() and tool_ctx:
        sys_text = sys_text.strip() + "\n\n" + tool_ctx
    elif not sys_text.strip():
        sys_text = tool_ctx

    # 追加工具调用相关提示
    if use_native_tools:
        sys_text = sys_text.rstrip() + (
            "\n\n路径不明确时先 find_files/search_code；读取大文件建议用 read_file_range 分段。"
            "用户上传的图片已直接内嵌在消息中，无需再调用工具读取。"
            "用户上传的文档（PDF/Word/Excel）已自动提取内容并显示在消息中，无需再调用工具读取。"
        )
    else:
        sys_text = sys_text.rstrip() + "\n\n【工具调用格式】务必输出严格 JSON（不要多余文本/注释）。" \
            "路径不明确时先 find_files/search_code；读取大文件建议用 read_file_range 分段。" \
            "用户上传的图片已直接内嵌在消息中，无需再调用工具读取。" \
            "用户上传的文档（PDF/Word/Excel）已自动提取内容并显示在消息中，无需再调用工具读取。" \
            "工具名称和参数格式参见上方工具描述。"
    sys0 = {"role": "system", "content": sys_text}
    msgs = [sys0]
    # One-shot resume guidance (runtime-only): when client hints a resume from last assistant error/aborted
    try:
        r = resume if isinstance(resume, dict) else None
        r_reason = str((r or {}).get("reason") or "").strip().lower()
        r_mid = str((r or {}).get("from_message_id") or "").strip()
        if r and r_mid and r_reason in ("error", "aborted"):
            guide = (
                "【续跑提示】检测到上轮中断（" + ("错误" if r_reason == "error" else "已中止") + 
                "）。请从中断点继续，避免重复已完成步骤；如需引用上轮结论或回执，请先简述已完成内容再继续。"
            )
            msgs.append({"role": "system", "content": guide})
    except Exception:
        pass
    if long_summary:
        msgs.append({"role": "system", "content": long_summary})
    if tool_mem_block:
        msgs.append({"role": "system", "content": tool_mem_block})
    if dyn_context and str(dyn_context).strip():
        # 限制 dyn_context 长度，避免占满输入 token 预算
        _dyn = str(dyn_context).strip()
        if len(_dyn) > 20000:
            _dyn = _dyn[:20000] + "\n...（动态上下文已截断）..."
        msgs.append({"role": "system", "content": _dyn})
    msgs.extend(hist)

    _hivo_ws_emit(run_id, session_id, "sending", _hivo_status_message(cfg, "sending"))

    # ── Broadcast Action Registry to frontend ──
    # Ensures frontend has the latest action definitions for dynamic
    # confirmation, parsing, and permission checks.
    try:
        _hivo_ws_emit_action_registry(run_id, session_id)
    except Exception as e:
        logger.debug(f"Exception ignored: {e}")

    try:
        feat = cfg.get("features") if isinstance(cfg, dict) else None
        feat = feat if isinstance(feat, dict) else {}
        ai_stream_enabled = bool(feat.get("ai_stream", True))
    except Exception:
        ai_stream_enabled = True

    if str(agent_conf.get("mode") or "").strip() == "fallback_chat":
        _hivo_ws_emit(run_id, session_id, "thinking", _hivo_status_message(cfg, "thinking"))
        rem = max(5, int(agent_deadline - time.time()))
        ok, msg, result = ai_chat(
            msgs,
            temperature=None,
            profile_id=profile_id,
            timeout_s=min(llm_timeout_s, rem),
            stream=ai_stream_enabled,
            on_delta=(lambda d: _hivo_ws_emit_delta(run_id, session_id, d.get("delta", ""), d.get("delta_type", "content"))) if ai_stream_enabled else None,
        )
        if not ok:
            err = msg or "对话失败"
            _hivo_ws_emit(run_id, session_id, "error", err)
            if _hivo_is_timeout_error(err):
                _hivo_ws_emit_final(run_id, session_id, err, ok=False, extra={"rounds": 1, "can_continue": True, "continue_reason": "timeout"})
            else:
                _hivo_ws_emit_final(run_id, session_id, err, ok=False, extra={"rounds": 1})
            return False, err, run_id
        content = str((result or {}).get("content") or "")
        if not content.strip():
            err = "模型未返回任何内容，请重试或更换模型。"
            _hivo_ws_emit(run_id, session_id, "error", err)
            _hivo_ws_emit_final(run_id, session_id, err, ok=False, extra={"rounds": 1})
            return False, err, run_id
        try:
            if user_text:
                st_chat = st.get("chat") if isinstance(st.get("chat"), list) else []
                st_chat.append({"role": "user", "content": str(user_text)})
                st_chat.append({"role": "assistant", "content": content})
                if memory_enabled:
                    st["chat"] = st_chat[-max_hist:]
                # roll older into summary when needed
                if memory_enabled and len(st_chat) > keep_n:
                    older = st_chat[:-keep_n]
                    s2 = _hivo_summarize_for_long_term(older, max_chars=int(mem_conf.get("long_term_summary_chars") or 3500), profile_id=profile_id)
                    if s2:
                        st["summary"] = s2
        except Exception as e:
            logger.debug(f"Exception ignored: {e}")
        _hivo_ws_emit(run_id, session_id, "done", _hivo_status_message(cfg, "done"))
        _hivo_ws_emit_final(run_id, session_id, content, ok=True, extra={"rounds": 1, "tool_receipts": pre_receipts, "finish_reason": (result or {}).get("finish_reason", ""), "reasoning_content": (result or {}).get("reasoning_content", "")})
        return True, content, run_id

    tool_sig_hist = []
    last_tool_receipts = []
    rep_escalations = 0
    reasoning_content_hist = ""
    openai_tools = _hivo_tools_to_openai_format(agent_tools) if (use_native_tools and agent_tools) else None
    for round_i in range(max(1, max_rounds)):
        if _hivo_agent_is_cancelled(run_id):
            final = "已取消执行。"
            _hivo_ws_emit(run_id, session_id, "error", "用户取消")
            _hivo_ws_emit_final(run_id, session_id, final, ok=False, extra={"rounds": round_i, "can_continue": True, "continue_reason": "cancelled"})
            return False, final, run_id
        if time.time() > agent_deadline:
            final = "执行超时，已中止。"
            _hivo_ws_emit(run_id, session_id, "error", "执行超时")
            _hivo_ws_emit_final(run_id, session_id, final, ok=False, extra={"rounds": round_i, "can_continue": True, "continue_reason": "timeout"})
            return False, final, run_id

        _hivo_ws_emit(run_id, session_id, "thinking", _hivo_status_message(cfg, "thinking", round_i=round_i + 1, max_rounds=max_rounds))
        if _hivo_agent_is_cancelled(run_id):
            final = "已取消执行。"
            _hivo_ws_emit(run_id, session_id, "error", "用户取消")
            _hivo_ws_emit_final(run_id, session_id, final, ok=False, extra={"rounds": round_i + 1, "can_continue": True, "continue_reason": "cancelled"})
            return False, final, run_id
        rem = max(5, int(agent_deadline - time.time()))

        ok, msg, result = ai_chat(
            msgs,
            temperature=None,
            profile_id=profile_id,
            timeout_s=min(llm_timeout_s, rem),
            stream=ai_stream_enabled,
            on_delta=(lambda d: _hivo_ws_emit_delta(run_id, session_id, d.get("delta", ""), d.get("delta_type", "content"))) if ai_stream_enabled else None,
            tools=openai_tools,
        )
        if not ok:
            err = msg or "对话失败"
            _hivo_ws_emit(run_id, session_id, "error", err)
            if _hivo_is_timeout_error(err):
                _hivo_ws_emit_final(run_id, session_id, err, ok=False, extra={"rounds": round_i + 1, "can_continue": True, "continue_reason": "timeout"})
            else:
                _hivo_ws_emit_final(run_id, session_id, err, ok=False, extra={"rounds": round_i + 1})
            return False, err, run_id

        content = str((result or {}).get("content") or "")
        finish_reason = str((result or {}).get("finish_reason") or "")
        native_tool_calls = (result or {}).get("tool_calls")
        # 累积所有轮次的思考过程，贴合主流 Agent 的单段思考体验
        _cur_reasoning = str((result or {}).get("reasoning_content") or "")
        if _cur_reasoning:
            if reasoning_content_hist:
                reasoning_content_hist += "\n\n" + _cur_reasoning
            else:
                reasoning_content_hist = _cur_reasoning

        calls = []
        if use_native_tools and isinstance(native_tool_calls, list) and native_tool_calls:
            calls = _hivo_parse_native_tool_calls(native_tool_calls)
            if not calls:
                msgs.append({"role": "assistant", "content": content})
        else:
            msgs.append({"role": "assistant", "content": content})

        if not content.strip() and not calls:
            err = "模型未返回任何内容，请重试或更换模型。"
            _hivo_ws_emit(run_id, session_id, "error", err)
            _hivo_ws_emit_final(run_id, session_id, err, ok=False, extra={"rounds": round_i + 1})
            return False, err, run_id
        # ── finish_reason 感知：根据 LLM 返回的结束原因增强终端判断 ──
        if finish_reason == "length" and round_i == 0:
            # 首轮即因 token 限制被截断 → 几乎肯定丢失了工具调用或关键内容
            trunc_err = "模型回复因 token 限制被截断（finish_reason=length）。请缩短上下文或提高 max_tokens 后重试。"
            _hivo_ws_emit(run_id, session_id, "error", trunc_err)
            _hivo_ws_emit_final(run_id, session_id, trunc_err, ok=False,
                extra={"rounds": round_i + 1, "can_continue": True, "continue_reason": "truncated", "finish_reason": finish_reason})
            return False, trunc_err, run_id
        if finish_reason == "length" and round_i > 0:
            # 后续轮次被截断：可能丢失了工具调用链，但注入提示引导模型收束
            _trunc_warn = "\n\n[系统提示] 上一轮回复因 token 限制被截断，若干工具调用可能不完整。请在当前轮优先给出阶段性结论，必要时使用更少的工具（≤2个）完成剩余任务。"
            content += _trunc_warn
        if not calls:
            calls = _hivo_extract_tool_calls(content, max_calls=max_calls)
        # ── Agent Actions: extract and forward to frontend ──
        agent_actions = _hivo_extract_agent_actions(content, max_actions=max_calls)
        if agent_actions:
            # Validate against registry; only emit registered actions
            valid_actions = [a for a in agent_actions if AGENT_ACTION_REGISTRY.validate(a)]
            if valid_actions:
                try:
                    _hivo_ws_emit(run_id, session_id, "executing", _hivo_status_message(cfg, "executing", tool_count=len(valid_actions)))
                except Exception as e:
                    logger.debug(f"Exception ignored: {e}")
                for a in valid_actions:
                    _hivo_ws_emit_action(run_id, session_id, a)
                # Save chat memory and end round (actions are terminal for this round)
                try:
                    if user_text:
                        st_chat = st.get("chat") if isinstance(st.get("chat"), list) else []
                        st_chat.append({"role": "user", "content": str(user_text)})
                        st_chat.append({"role": "assistant", "content": content})
                        if memory_enabled:
                            st["chat"] = st_chat[-max_hist:]
                        if memory_enabled and len(st_chat) > keep_n:
                            older = st_chat[:-keep_n]
                            s2 = _hivo_summarize_for_long_term(older, max_chars=int(mem_conf.get("long_term_summary_chars") or 3500), profile_id=profile_id)
                            if s2:
                                st["summary"] = s2
                except Exception as e:
                    logger.debug(f"Exception ignored: {e}")
                _hivo_ws_emit(run_id, session_id, "done", _hivo_status_message(cfg, "done"))
                _hivo_ws_emit_final(run_id, session_id, content, ok=True,
                    extra={"rounds": round_i + 1, "tool_receipts": (pre_receipts + last_tool_receipts) if pre_receipts else last_tool_receipts,
                           "agent_actions": valid_actions, "finish_reason": finish_reason, "reasoning_content": reasoning_content_hist})
                return True, content, run_id
            # Invalid actions: treat as no-ops and continue with tool calls
            _trunc_warn = "\n\n[系统提示] 检测到未注册的 Agent Action，已忽略。请使用已注册的动作类型。"
            content += _trunc_warn
        # ── finish_reason 二次校验：LLM 表示有工具调用但未提取到有效调用 ──
        if not calls and finish_reason == "tool_calls":
            parse_err = "模型请求了工具调用 (finish_reason=tool_calls)，但未能解析出有效的工具调用 JSON。这可能是因为模型输出格式不符合预期，请重试。"
            _hivo_ws_emit(run_id, session_id, "error", parse_err)
            _hivo_ws_emit_final(run_id, session_id, parse_err, ok=False,
                extra={"rounds": round_i + 1, "can_continue": True, "continue_reason": "tool_calls_parse_error", "finish_reason": finish_reason})
            return False, parse_err, run_id
        # 允许在首轮即按需调用"查看文件"等工具读取附件内容，避免因抑制导致无法查看附件。
        allow = agent_conf.get("tools")
        if isinstance(allow, list) and allow:
            allow_set = set(str(x).strip() for x in allow if str(x).strip())
            filtered_calls = [c for c in calls if str(c.get("type") or "").strip() in allow_set]
            if len(filtered_calls) < len(calls):
                dropped = [str(c.get("type") or "?") for c in calls if str(c.get("type") or "").strip() not in allow_set]
                logger.info(f"[Agent] 工具白名单过滤: 丢弃 {len(calls) - len(filtered_calls)} 个未授权工具调用: {dropped}")
            calls = filtered_calls
        if not calls:
            try:
                if user_text:
                    st_chat = st.get("chat") if isinstance(st.get("chat"), list) else []
                    st_chat.append({"role": "user", "content": str(user_text)})
                    st_chat.append({"role": "assistant", "content": content})
                    if memory_enabled:
                        st["chat"] = st_chat[-max_hist:]
                    if memory_enabled and len(st_chat) > keep_n:
                        older = st_chat[:-keep_n]
                        s2 = _hivo_summarize_for_long_term(older, max_chars=int(mem_conf.get("long_term_summary_chars") or 3500), profile_id=profile_id)
                        if s2:
                            st["summary"] = s2
            except Exception as e:
                logger.debug(f"Exception ignored: {e}")
            _hivo_ws_emit(run_id, session_id, "done", _hivo_status_message(cfg, "done"))
            _hivo_ws_emit_final(run_id, session_id, content, ok=True, extra={"rounds": round_i + 1, "tool_receipts": (pre_receipts + last_tool_receipts) if pre_receipts else last_tool_receipts, "finish_reason": finish_reason, "reasoning_content": reasoning_content_hist})
            return True, content, run_id

        sig = _hivo_repeat_signature(calls, mode=rep_sig_mode)
        tool_sig_hist.append(sig)
        if rep_window > 0 and len(tool_sig_hist) >= rep_window:
            recent = tool_sig_hist[-rep_window:]
            same = sum(1 for x in recent if x == recent[-1])
            if same > rep_max_same:
                if rep_escalations < rep_escalation_limit:
                    rep_escalations += 1
                    guide = _hivo_repeat_escalation_prompt(rep_escalations, last_sig=sig)
                    _hivo_ws_emit(run_id, session_id, "thinking", f"检测到重复趋势，注入修正策略 (level {rep_escalations}/{rep_escalation_limit})")
                    msgs.append({"role": "assistant", "content": content})
                    msgs.append({"role": "user", "content": guide})
                    continue

                final = "检测到重复工具调用，已中止以避免死循环。请你调整需求或补充关键信息后重试。"
                _hivo_ws_emit(run_id, session_id, "error", "检测到重复行为")
                _hivo_ws_emit_final(run_id, session_id, final, ok=False, extra={"rounds": round_i + 1, "can_continue": True, "continue_reason": "repeat"})
                return False, final, run_id

        try:
            _hivo_ws_emit(run_id, session_id, "executing", _hivo_status_message(cfg, "executing", tool_count=len(calls)))
        except Exception as e:
            logger.debug(f"Exception ignored: {e}")

        ws_enabled = None
        try:
            ws_enabled, _sc0 = _ai_load_web_search_cfg()
        except Exception:
            ws_enabled = False
        receipts = []
        for i_tool, c in enumerate(calls or []):
            if _hivo_agent_is_cancelled(run_id):
                final = "已取消执行。"
                _hivo_ws_emit(run_id, session_id, "error", "用户取消")
                _hivo_ws_emit_final(run_id, session_id, final, ok=False, extra={"rounds": round_i + 1, "can_continue": True, "continue_reason": "cancelled"})
                return False, final, run_id
            if time.time() > agent_deadline:
                final = "执行超时，已中止。"
                _hivo_ws_emit(run_id, session_id, "error", "执行超时")
                _hivo_ws_emit_final(run_id, session_id, final, ok=False, extra={"rounds": round_i + 1, "can_continue": True, "continue_reason": "timeout"})
                return False, final, run_id

            name = str(c.get("type") or "")
            try:
                msg0 = _hivo_status_message(cfg, "executing", tool_count=len(calls))
                if ws_enabled and name in ("web_search", "web_fetch"):
                    msg0 = "搜索中..."
                _hivo_ws_emit(
                    run_id,
                    session_id,
                    "executing",
                    msg0,
                    extra={"tool": name, "tool_i": i_tool + 1, "tool_n": len(calls), "can_cancel": True},
                )
            except Exception as e:
                logger.debug(f"Exception ignored: {e}")

            c2 = c
            try:
                remain_s = int(max(1, agent_deadline - time.time()))
            except Exception:
                remain_s = 1
            try:
                if isinstance(c, dict):
                    c2 = dict(c)
                    tt = str(c2.get("type") or "")
                    if tt in ("run_cmd", "web_search", "web_fetch"):
                        if "timeout" in c2:
                            try:
                                t0 = int(c2.get("timeout") or 0)
                            except Exception:
                                t0 = 0
                            if t0 <= 0:
                                t0 = 30
                            c2["timeout"] = int(max(1, min(t0, remain_s)))
            except Exception:
                c2 = c
            cache_key = _hivo_tool_call_sig(c)
            cached = None
            try:
                if tool_cache_enabled and cache_key and isinstance(tool_cache, OrderedDict):
                    cached = tool_cache.get(cache_key)
            except Exception:
                cached = None

            if isinstance(cached, dict) and "ok" in cached:
                ok1 = bool(cached.get("ok"))
                msg1 = str(cached.get("msg") or "")
                data1 = cached.get("data") if isinstance(cached.get("data"), dict) else {}
                try:
                    if isinstance(data1, dict) and "cache_hit" not in data1:
                        data1 = dict(data1)
                        data1["cache_hit"] = True
                except Exception as e:
                    logger.debug(f"Exception ignored: {e}")
            else:
                ok1, msg1, data1 = _hivo_exec_tool(c2, undo_gid=undo_gid_eff, run_id=run_id, agent_deadline=agent_deadline)
                # 写操作后清理工具缓存，避免返回旧内容
                _tool_type = str(c2.get("type") or "").strip()
                if _tool_type in ("save_file", "delete_file", "rename_file", "unstage_file", "discard_staged_file", "unstage_all_staged", "discard_all_staged") and isinstance(tool_cache, OrderedDict):
                    try:
                        tool_cache.clear()
                    except Exception:
                        pass
                try:
                    if tool_cache_enabled and cache_key and isinstance(tool_cache, OrderedDict):
                        tool_cache[cache_key] = {"ok": bool(ok1), "msg": str(msg1 or ""), "data": data1 if isinstance(data1, dict) else {}}
                        while len(tool_cache) > int(mem_conf.get("tool_cache_items") or 120):
                            try:
                                tool_cache.popitem(last=False)
                            except Exception:
                                break
                except Exception as e:
                    logger.debug(f"Exception ignored: {e}")

            try:
                if tool_memory_enabled and memory_enabled:
                    rec = _hivo_tool_log_record(name, c, bool(ok1), str(msg1 or ""), data=data1 if isinstance(data1, dict) else {})
                    tool_log0.append(rec)
                    while len(tool_log0) > int(mem_conf.get("tool_log_items") or 80):
                        tool_log0.pop(0)
                    st["tool_log"] = tool_log0
            except Exception as e:
                logger.debug(f"Exception ignored: {e}")

            receipts.append({"type": name, "ok": bool(ok1), "msg": msg1 or "", "request": c2 if isinstance(c2, dict) else {}, "data": data1})

        # Auto browser screenshot: only when the round includes a state-changing browser action
        # (open/click/type) but no screenshot was taken, append one extra receipt.
        try:
            try:
                featS = cfg.get("features") if isinstance(cfg, dict) and isinstance(cfg.get("features"), dict) else {}
                auto_enabled = bool(featS.get("browser_auto_screenshot", False))
            except Exception:
                auto_enabled = False
            if not auto_enabled:
                raise Exception("auto_screenshot_disabled")

            had_browser_action = False
            had_shot = False
            sid_auto = ""
            for r in receipts:
                nm0 = str(r.get("type") or "")
                if nm0 in ("browser_screenshot", "browser_screenshot_auto"):
                    had_shot = True
                if nm0 in ("browser_open", "browser_click", "browser_type"):
                    had_browser_action = True
                    if not sid_auto:
                        d0 = r.get("data") if isinstance(r.get("data"), dict) else {}
                        sid_auto = str(d0.get("session_id") or "").strip()
            if had_browser_action and (not had_shot):
                if not sid_auto:
                    sid_auto = str(session_id or "").strip() or str(run_id or "").strip() or "global"
                okS, msgS, dataS = _pw_screenshot(sid_auto, full_page=True)
                receipts.append({"type": "browser_screenshot_auto", "ok": bool(okS), "msg": msgS or "", "data": dataS if isinstance(dataS, dict) else {}})
        except Exception as e:
            if str(e) != "auto_screenshot_disabled":
                logger.debug(f"Exception ignored: {e}")

        receipt_lines = ["【工具回执(真实执行结果)】"]
        total_payload_chars = 0
        # Per-tool payload caps (chars) to prevent any single tool flooding the receipt
        _TOOL_PAYLOAD_CAPS = {
            "file_content": 20000,
            "read_file_range": 20000,
            "get_file": 12000,
            "run_cmd": 15000,
            "web_search": 16000,
            "web_fetch": 12000,
            "browser_text": 12000,
            "browser_screenshot": 8000,
            "browser_eval": 8000,
        }
        try:
            limits = (cfg or {}).get("limits") if isinstance(cfg, dict) else {}
            if limits and isinstance(limits, dict):
                server_input = int(limits.get("ai_max_input_tokens") or 0)
                if server_input > 0:
                    max_total_payload_chars = max(12000, min(80000, int(server_input * 0.25)))
                else:
                    max_total_payload_chars = 24000
            else:
                max_total_payload_chars = 24000
        except Exception:
            max_total_payload_chars = 24000
        for r in receipts:
            nm = str(r.get("type") or "")
            ok1 = bool(r.get("ok"))
            detail = str(r.get("msg") or "")
            receipt_lines.append(_hivo_tool_receipt_line(nm, ok1, detail))

            # For read/inspect tools, include returned content so the model can actually see it.
            try:
                data1 = r.get("data") if isinstance(r.get("data"), dict) else {}
                if not ok1:
                    continue
                payload = ""
                if nm in ("file_content", "read_file_range"):
                    f0 = data1.get("file") if isinstance(data1.get("file"), dict) else {}
                    payload = str(f0.get("content") or "")
                    path1 = str(f0.get("path") or "")
                    if payload:
                        receipt_lines.append(f"  path: {path1}")
                elif nm in ("raw_file",):
                    f0 = data1.get("file") if isinstance(data1.get("file"), dict) else {}
                    if f0:
                        payload = json.dumps({"file": f0}, ensure_ascii=False, indent=2)
                elif nm in ("get_file",):
                    f0 = data1.get("file") if isinstance(data1.get("file"), dict) else {}
                    if f0:
                        # 文本：直接附带内容并给出 path；媒体：仅提示类型与是否截断，避免刷屏
                        if "content" in f0:
                            payload = str(f0.get("content") or "")
                            path1 = str(f0.get("path") or "")
                            if payload:
                                receipt_lines.append(f"  path: {path1}")
                        elif "data_url" in f0:
                            mt = str(f0.get("file_type") or "")
                            trunc = bool(f0.get("truncated"))
                            payload = f"image_base64({mt}) - {'truncated' if trunc else 'full'}"
                elif nm in ("run_cmd",):
                    res1 = data1.get("result") if isinstance(data1.get("result"), dict) else {}
                    payload = str(res1.get("output") or data1.get("output") or "")
                elif nm in ("web_search",):
                    res1 = data1.get("result") if isinstance(data1.get("result"), dict) else {}
                    if res1:
                        payload = json.dumps(res1, ensure_ascii=False, indent=2)
                elif nm in ("web_fetch",):
                    res1 = data1.get("result") if isinstance(data1.get("result"), dict) else {}
                    payload = str(res1.get("text") or "")

                elif nm.startswith("browser_"):
                    if nm in ("browser_text",):
                        payload = str(data1.get("text") or data1.get("html") or "")
                        try:
                            if len(payload) > 12000:
                                payload = payload[:12000]
                        except Exception:
                            pass
                    elif nm in ("browser_screenshot",):
                        b64 = str(data1.get("image_b64") or "")
                        if b64:
                            try:
                                if len(b64) > 8000:
                                    b64 = b64[:8000]
                            except Exception:
                                pass
                            payload = "image_b64(png, truncated):\n" + b64
                    elif nm in ("browser_eval",):
                        try:
                            payload = json.dumps({"result": data1.get("result"), "obs": data1.get("obs")}, ensure_ascii=False, indent=2)
                        except Exception:
                            payload = str(data1.get("result") or "")
                    else:
                        try:
                            obs0 = data1.get("obs") if isinstance(data1.get("obs"), dict) else {}
                            if obs0:
                                payload = json.dumps({"obs": obs0}, ensure_ascii=False, indent=2)
                        except Exception:
                            payload = ""

                if payload:
                    # Apply per-tool-type payload cap if one exists
                    tool_cap = _TOOL_PAYLOAD_CAPS.get(nm, 0)
                    if tool_cap > 0 and len(payload) > tool_cap:
                        payload = payload[:tool_cap] + f"\n...（已截断，原文{len(payload)}字）..."
                    if total_payload_chars >= max_total_payload_chars:
                        continue
                    room = max_total_payload_chars - total_payload_chars
                    chunk = payload[:room]
                    total_payload_chars += len(chunk)
                    receipt_lines.append("  ---")
                    receipt_lines.append(chunk)
            except Exception as e:
                logger.debug(f"Exception ignored: {e}")
        receipt_text = "\n".join(receipt_lines)
        # 统一把工具摘要中的反斜杠路径替换为正斜杠（避免 markdown-it 转义破坏）
        receipt_text = _hivo_normalize_tool_data(receipt_text)
        # 累积所有轮次的工具回执（多轮工具调用时，每轮都追加，而不是覆盖）
        try:
            if isinstance(last_tool_receipts, list) and isinstance(receipts, list):
                last_tool_receipts.extend(receipts)
            else:
                last_tool_receipts = receipts[:] if isinstance(receipts, list) else []
        except Exception:
            last_tool_receipts = receipts

        if use_native_tools and isinstance(native_tool_calls, list) and native_tool_calls:
            # 统一清洗 tool_calls：补全缺失的 id，确保 assistant 消息和 tool 消息的 id 完全一致
            cleaned_tool_calls = _hivo_ai_clean_tool_calls(native_tool_calls)
            if cleaned_tool_calls:
                # 原生工具调用：先添加 assistant 消息（带 tool_calls），再添加对应 tool 消息
                msgs.append({
                    "role": "assistant",
                    "content": content or "",
                    "tool_calls": cleaned_tool_calls,
                })
                for idx, tc in enumerate(cleaned_tool_calls):
                    tc_id = tc["id"]
                    receipt_content = receipt_text
                    image_data_url = ""
                    if idx < len(receipts):
                        r = receipts[idx]
                        ok1 = bool(r.get("ok"))
                        detail = str(r.get("msg") or "")
                        data1 = r.get("data") if isinstance(r.get("data"), dict) else {}
                        # 检查是否是图片返回（get_file 返回 data_url）
                        try:
                            f0 = data1.get("file") if isinstance(data1.get("file"), dict) else {}
                            if f0:
                                du = str(f0.get("data_url") or "").strip()
                                if du and du.startswith("data:image/"):
                                    image_data_url = du
                        except Exception:
                            pass
                        receipt_content = json.dumps({
                            "ok": ok1,
                            "msg": detail,
                            "data": _hivo_normalize_tool_data(data1),
                        }, ensure_ascii=False)
                    # 构造 tool 消息内容：有图片则用多模态数组，否则用纯文本
                    if image_data_url:
                        tool_content = [
                            {"type": "text", "text": receipt_content},
                            {"type": "image_url", "image_url": {"url": image_data_url}},
                        ]
                    else:
                        tool_content = receipt_content
                    msgs.append({
                        "role": "tool",
                        "tool_call_id": tc_id,
                        "content": tool_content,
                    })
        else:
            msgs.append({"role": "assistant", "content": content})
            msgs.append({
                "role": "system",
                "content": "【工具回执】\n" + receipt_text
                + "\n\n请基于上述真实结果生成面向用户的自然语言回复：\n"
                + "- 不要输出回执原文或内部字段（如 'ok', 'msg', 'data'）\n"
                + "- 若任务完成，直接给出结论；若还需继续，输出下一步的工具 JSON\n"
                + "- 若回执中有错误或失败，向用户说明原因并建议替代方案"
            })

    final = "已达到最大轮次，任务仍未完成。请你缩小问题范围或补充信息后重试。"
    _hivo_ws_emit(run_id, session_id, "error", "达到最大轮次")
    _hivo_ws_emit_final(run_id, session_id, final, ok=False, extra={"rounds": max_rounds, "can_continue": True, "continue_reason": "max_rounds"})
    return False, final, run_id


def _run_cmd_simple(cmd: str, timeout: int = 30, cwd: str = "", run_id: str = "", agent_deadline: float = 0.0, undo_gid: str = ""):
    if not REPO_PATH:
        return False, "未打开仓库", ""
    c = str(cmd or "").strip()
    if not c:
        return False, "缺少 cmd", ""
    if "\n" in c or "\r" in c:
        return False, "cmd 必须是单行命令", ""
    rid = str(run_id or "").strip()
    try:
        workdir = REPO_PATH
        if cwd:
            safe = _safe_repo_abspath(cwd)
            if safe and os.path.isdir(safe):
                workdir = safe

        to_s = int(timeout or 30)
        to_s = max(1, min(3600, to_s))
        if agent_deadline and agent_deadline > 0:
            try:
                remain = int(max(1, agent_deadline - time.time()))
                to_s = max(1, min(to_s, remain))
            except Exception as e:
                logger.debug(f"Exception ignored: {e}")

        pre_map, pre_snaps = _undo_prepare_cmd_snapshots(undo_gid)

        start_t = time.time()
        # Use shlex to safely split command into args list (avoid shell=True injection)
        try:
            args = shlex.split(c)
        except ValueError:
            return False, "cmd 格式错误：无法解析命令", ""
        p = subprocess.Popen(
            args,
            cwd=workdir,
            shell=False,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True,
            encoding="utf-8",
            errors="replace",
        )
        if rid:
            try:
                with _hivo_agent_proc_lock:
                    _hivo_agent_run_proc[rid] = p
            except Exception as e:
                logger.debug(f"Exception ignored: {e}")

        out_chunks = []
        err_chunks = []
        while True:
            if rid and _hivo_agent_is_cancelled(rid):
                try:
                    p.terminate()
                except Exception as e:
                    logger.debug(f"Exception ignored: {e}")
                return False, "已取消", ""
            if agent_deadline and agent_deadline > 0 and time.time() > agent_deadline:
                try:
                    p.terminate()
                except Exception as e:
                    logger.debug(f"Exception ignored: {e}")
                return False, "执行超时", ""
            if (time.time() - start_t) > to_s:
                try:
                    p.terminate()
                except Exception as e:
                    logger.debug(f"Exception ignored: {e}")
                return False, "命令超时", ""

            rc = p.poll()
            if rc is not None:
                try:
                    o, e = p.communicate(timeout=1)
                except Exception:
                    o, e = "", ""
                if o:
                    out_chunks.append(o)
                if e:
                    err_chunks.append(e)
                out = ("".join(out_chunks) + ("\n" + "".join(err_chunks) if err_chunks else "")).strip()
                if rc != 0:
                    return False, f"exit={rc}", out
                try:
                    _undo_finalize_cmd_snapshots(undo_gid, pre_map, pre_snaps)
                except Exception as e:
                    logger.debug(f"Exception ignored: {e}")
                return True, "", out
            time.sleep(0.2)
    except Exception as e:
        return False, str(e), ""
    finally:
        try:
            if rid:
                _hivo_agent_clear_proc(rid)
        except Exception as e:
            logger.debug(f"Exception ignored: {e}")


_rl_lock = threading.Lock()
_rl_windows: dict[str, deque] = {}


def _rl_make_key(handler, profile_id: str | None = None, session_id: str | None = None):
    ip = ""
    try:
        ip = str((handler.client_address or [""])[0] or "")
    except Exception:
        ip = ""
    pid = str(profile_id or "").strip()
    sid = str(session_id or "").strip()
    base = sid or pid or ip or "anon"
    repo_key = "global"
    try:
        if REPO_PATH:
            repo_key = hashlib.sha1(str(REPO_PATH).encode("utf-8", errors="ignore")).hexdigest()[:10]
    except Exception:
        repo_key = "global"
    return f"{repo_key}::{base}::{ip}"


def _rl_conf(level: str):
    lvl = str(level or "").lower()
    if lvl == "modify":
        return 60.0, 5, 8
    if lvl == "search":
        return 60.0, 10, 15
    return 60.0, 20, 30


def _rl_check_and_get_delay(key: str, level: str):
    window_s, soft_n, hard_n = _rl_conf(level)
    now = time.time()
    wkey = f"{level}::{key}"
    with _rl_lock:
        q = _rl_windows.get(wkey)
        if q is None:
            q = deque()
            _rl_windows[wkey] = q
        cutoff = now - window_s
        while q and q[0] < cutoff:
            q.popleft()
        cur = len(q)
        if cur >= hard_n:
            retry_after = int(max(1.0, (q[0] + window_s) - now)) if q else int(window_s)
            return False, 0.0, retry_after
        delay = 0.0
        if cur >= soft_n:
            delay = min(3.0, 0.25 * (cur - soft_n + 1))
        q.append(now)
        return True, delay, 0


_ai_cache_lock = threading.Lock()
_ai_cache: "OrderedDict[str, dict]" = OrderedDict()


def _ai_norm_query(text: str):
    s = str(text or "").strip().lower()
    s = re.sub(r"\s+", " ", s)
    s = re.sub(r"[\u0000-\u001f]", " ", s)
    s = re.sub(r"[\t\r\n]", " ", s)
    s = re.sub(r"[\"'`~!@#$%^&*()\-_=+\[\]{};:,./<>?\\|]", " ", s)
    s = re.sub(r"\s+", " ", s).strip()
    return s


def _ai_keywords(text: str):
    s = _ai_norm_query(text)
    if not s:
        return set()
    toks = re.findall(r"[a-z0-9]{2,}|[\u4e00-\u9fff]{2,}", s)
    return set(toks)


def _ai_jaccard(a: set, b: set):
    if not a or not b:
        return 0.0
    inter = len(a & b)
    if inter <= 0:
        return 0.0
    uni = len(a | b)
    return float(inter) / float(uni or 1)


def _ai_cache_get(query: str, profile_id: str | None = None):
    now = time.time()
    nq = _ai_norm_query(query)
    if not nq:
        return None
    kset = _ai_keywords(nq)
    repo_key = "global"
    try:
        if REPO_PATH:
            repo_key = hashlib.sha1(str(REPO_PATH).encode("utf-8", errors="ignore")).hexdigest()[:10]
    except Exception:
        repo_key = "global"
    prefix = repo_key + "::" + str(profile_id or "").strip() + "::"

    with _ai_cache_lock:
        dead = []
        for ck, ent in list(_ai_cache.items()):
            ttl = float(ent.get("ttl") or 0)
            ts = float(ent.get("ts") or 0)
            if ttl > 0 and now - ts > ttl:
                dead.append(ck)
        for ck in dead:
            try:
                _ai_cache.pop(ck, None)
            except Exception as e:
                logger.debug(f"Exception ignored: {e}")

        exact_key = prefix + nq
        ent = _ai_cache.get(exact_key)
        if ent is not None:
            ent["hits"] = int(ent.get("hits") or 0) + 1
            ent["ts"] = now
            try:
                _ai_cache.move_to_end(exact_key)
            except Exception as e:
                logger.debug(f"Exception ignored: {e}")
            return ent

        best_key = ""
        best_score = 0.0
        scan = 0
        for ck, e in reversed(list(_ai_cache.items())):
            if prefix and not ck.startswith(prefix):
                continue
            scan += 1
            if scan > 60:
                break
            ks = e.get("kw") or set()
            try:
                score = _ai_jaccard(kset, set(ks))
            except Exception:
                score = 0.0
            if score > best_score:
                best_score = score
                best_key = ck
        if best_key and best_score >= 0.92:
            e = _ai_cache.get(best_key)
            if e is not None:
                e["hits"] = int(e.get("hits") or 0) + 1
                e["ts"] = now
                try:
                    _ai_cache.move_to_end(best_key)
                except Exception as e:
                    logger.debug(f"Exception ignored: {e}")
                return e
    return None


def _ai_cache_put(query: str, response: str, profile_id: str | None = None):
    now = time.time()
    nq = _ai_norm_query(query)
    if not nq:
        return
    repo_key = "global"
    try:
        if REPO_PATH:
            repo_key = hashlib.sha1(str(REPO_PATH).encode("utf-8", errors="ignore")).hexdigest()[:10]
    except Exception:
        repo_key = "global"
    prefix = repo_key + "::" + str(profile_id or "").strip() + "::"
    key = prefix + nq
    kw = _ai_keywords(nq)
    ent = {
        "q": nq,
        "kw": kw,
        "resp": str(response or ""),
        "ts": now,
        "hits": 0,
        "ttl": 20.0 * 60.0,
    }
    with _ai_cache_lock:
        _ai_cache[key] = ent
        try:
            _ai_cache.move_to_end(key)
        except Exception as e:
            logger.debug(f"Exception ignored: {e}")
        while len(_ai_cache) > 220:
            try:
                _ai_cache.popitem(last=False)
            except Exception:
                break


def _ai_cache_clear():
    with _ai_cache_lock:
        _ai_cache.clear()


def get_file_log(filepath: str, limit: int = 50):
    """Get git log for a single file."""
    if not REPO_PATH:
        return []
    fp = (filepath or "").replace("\\", "/").lstrip("/")
    if not fp:
        return []

    try:
        limit = int(limit)
    except Exception:
        limit = 50
    if limit <= 0:
        limit = 50
    if limit > 200:
        limit = 200

    fmt = "--pretty=format:%H%x00%an%x00%ae%x00%ad%x00%s"
    out, _, code = run_git(["log", fmt, "--date=iso", f"-{limit}", "--", fp])
    if code != 0:
        return []

    commits = []
    for line in (out or "").splitlines():
        parts = line.split("\x00")
        if len(parts) < 5:
            continue
        full_hash = parts[0]
        commits.append({
            "hash": full_hash[:7],
            "full_hash": full_hash,
            "author": parts[1],
            "email": parts[2],
            "date": parts[3],
            "message": parts[4],
        })
    return commits


def get_commit_detail(commit_hash):
    """获取提交详情"""
    logger.debug(f"获取提交详情: {commit_hash}")
    fmt = "--pretty=format:%H%x00%an%x00%ae%x00%ad%x00%s%x00%b"

    out, _, code = run_git(["show", "--no-patch", fmt, "--date=iso", commit_hash])
    parts = out.strip().split("\x00")
    if len(parts) < 5:
        logger.error(f"解析提交详情失败: {commit_hash}")
        return {"error":"解析失败"}

    info = {
        "hash":      parts[0][:7],  # 短hash
        "full_hash": parts[0],      # 完整hash
        "author":    parts[1],
        "email":     parts[2],
        "date":      parts[3],
        "subject":   parts[4],
        "message":   parts[4],
        "body":      parts[5] if len(parts) > 5 else ""
    }

    # 获取文件变更详情（使用numstat获取准确的增删行数）
    stat_out, _, _ = run_git(["show", "--numstat", "--format=", commit_hash])
    files = []
    
    for line in stat_out.splitlines():
        line = line.strip()
        if not line:
            continue
        
        # numstat格式: <added>\t<removed>\t<filepath>
        parts = line.split('\t')
        if len(parts) >= 3:
            added_str = parts[0]
            removed_str = parts[1]
            filepath = parts[2]
            
            # 处理二进制文件(显示为'-')
            try:
                added = int(added_str) if added_str != '-' else 0
                removed = int(removed_str) if removed_str != '-' else 0
            except ValueError:
                added = 0
                removed = 0
            
            # 使用name-status获取准确的文件状态
            status_out, _, _ = run_git(["show", "--name-status", "--format=", commit_hash, "--", filepath])
            status = 'M'  # 默认修改
            
            for status_line in status_out.splitlines():
                if filepath in status_line:
                    if status_line.startswith('A'):
                        status = 'A'
                    elif status_line.startswith('D'):
                        status = 'D'
                    elif status_line.startswith('M'):
                        status = 'M'
                    elif status_line.startswith('R'):
                        status = 'R'
                    break
            
            files.append({
                "path": filepath,
                "status": status,
                "added": added,
                "removed": removed
            })

    info["files"] = files
    logger.info(f"成功获取提交详情: {commit_hash}, 文件数: {len(files)}")
    return info


def get_commit_file_diff(commit_hash, filepath):
    """获取提交中某个文件的 diff"""
    logger.debug(f"获取提交文件 diff: {commit_hash} - {filepath}")
    
    # 使用git show获取该commit中该文件的diff
    # 添加--format=选项来只输出diff内容，不包含提交信息
    out, err, code = run_git(["show", "--format=", "--unified=5", f"{commit_hash}", "--", filepath])
    if code != 0:
        logger.error(f"获取提交文件 diff 失败: {commit_hash} - {filepath}, 错误: {err}")
        return []

    # 检查是否有diff内容
    if not out or not out.strip():
        logger.warning(f"提交文件 diff 为空: {commit_hash} - {filepath}")
        return []

    logger.info(f"成功获取提交文件 diff: {commit_hash} - {filepath}, diff长度: {len(out)}")
    return parse_diff(out)

# ════════════════════════════════════════════════════════
#  分支操作
# ════════════════════════════════════════════════════════

def get_branches():
    """获取分支列表"""
    logger.debug("获取分支列表")
    out, _, code = run_git(["branch", "-a"])
    if code != 0:
        logger.warning("获取分支列表失败")
        return [], None
    
    branches = []
    current  = None
    for line in out.splitlines():
        line = line.strip()
        if line.startswith("*"):
            current = line[2:].strip()
            branches.append(current)
        elif line:
            branches.append(line)
    
    logger.info(f"成功获取分支列表，共 {len(branches)} 个分支，当前分支: {current}")
    return branches, current


def fetch_remotes():
    if not REPO_PATH:
        return False, "未打开仓库", ""
    out, err, code = run_git(["fetch", "--all", "--prune"], timeout=300)
    if code != 0:
        return False, (err or out or "fetch 失败"), (out or "")
    return True, "", (out or "")


def get_sync_status(do_fetch: bool = False):
    """Return (ok: bool, data: dict, error: str|None).

    Compute ahead/behind between local HEAD and upstream remote-tracking ref.
    Prefer @{u} (configured upstream). Fallback to origin/<branch> when possible.
    """
    if not REPO_PATH:
        return False, {}, "未打开仓库"

    if do_fetch:
        ok, msg, _out = fetch_remotes()
        if not ok:
            return False, {}, (msg or "fetch 失败")

    bout, berr, bcode = run_git(["branch", "--show-current"], timeout=30)
    branch = (bout or "").strip() if bcode == 0 else ""
    if not branch:
        # Detached HEAD or unborn branch: treat as a valid state.
        # Sync-status is not meaningful without an upstream branch, but it shouldn't be a 400.
        return True, {
            "branch": "",
            "upstream": None,
            "ahead": 0,
            "behind": 0,
            "has_upstream": False,
            "detached": True,
        }, None

    upstream = ""
    uout, uerr, ucode = run_git(["rev-parse", "--abbrev-ref", "--symbolic-full-name", "@{u}"], timeout=30)
    if ucode == 0:
        upstream = (uout or "").strip()
    else:
        cand = f"origin/{branch}"
        if _remote_ref_exists(cand):
            upstream = cand

    if not upstream:
        return True, {
            "branch": branch,
            "upstream": None,
            "ahead": 0,
            "behind": 0,
            "has_upstream": False,
        }, None

    out, err, code = run_git(["rev-list", "--left-right", "--count", f"HEAD...{upstream}"], timeout=60)
    if code != 0:
        return False, {}, (err or out or "无法计算同步状态")

    ahead = 0
    behind = 0
    try:
        parts = (out or "").strip().split()
        if len(parts) >= 2:
            ahead = int(parts[0])
            behind = int(parts[1])
    except Exception:
        ahead = 0
        behind = 0

    return True, {
        "branch": branch,
        "upstream": upstream,
        "ahead": max(0, ahead),
        "behind": max(0, behind),
        "has_upstream": True,
    }, None

def is_commit_pushed(commit_hash):
    """判断某个提交是否已经存在于任意远端分支（基于本地 remote refs）。"""
    commit_hash = (commit_hash or "").strip()
    if not commit_hash:
        return False, [], "缺少 hash"

    out, err, code = run_git(["branch", "-r", "--contains", commit_hash])
    if code != 0:
        return False, [], (err or "查询远端分支失败")

    branches = []
    for ln in (out or "").splitlines():
        b = (ln or "").strip()
        if not b:
            continue
        # ignore symbolic refs like: origin/HEAD -> origin/main
        if "->" in b:
            continue
        branches.append(b)

    return bool(branches), branches, None

def _has_worktree_changes():
    """Return (dirty: bool, detail: str)."""
    out, err, code = run_git(["status", "--porcelain=v1"])
    if code != 0:
        return True, (err or "无法检测工作区状态")
    dirty = bool((out or "").strip())
    return dirty, (out or "")

def _branch_exists_local(branch_name: str):
    branch_name = (branch_name or "").strip()
    if not branch_name:
        return False
    _, _, code = run_git(["show-ref", "--verify", "--quiet", f"refs/heads/{branch_name}"])
    return code == 0

def _remote_ref_exists(remote_ref: str):
    remote_ref = (remote_ref or "").strip()
    if not remote_ref:
        return False
    # remote_ref example: origin/foo
    _, _, code = run_git(["show-ref", "--verify", "--quiet", f"refs/remotes/{remote_ref}"])
    return code == 0

def _normalize_remote_ref(branch_name: str):
    """Return (is_remote: bool, remote_ref: str|None, raw: str).

    Accepts:
    - remotes/origin/foo  (from `git branch -a`)
    - origin/foo          (sometimes used as remote ref)
    """
    raw = (branch_name or "").strip()
    if not raw:
        return False, None, raw
    if raw.startswith("remotes/"):
        rr = raw.replace("remotes/", "", 1)
        return True, rr, raw
    # If user passes origin/foo directly and it's a real remote-tracking ref, treat it as remote.
    if "/" in raw and _remote_ref_exists(raw):
        return True, raw, raw
    return False, None, raw

def switch_branch(branch_name: str):
    """Switch to a branch.

    - If working tree is dirty, refuse (consistent with safe UI behavior).
    - If selecting a remote branch (remotes/origin/foo), create a local tracking branch.
    """
    if not REPO_PATH:
        return False, None, "未打开仓库", "", True

    branch_name = (branch_name or "").strip()
    if not branch_name:
        return False, None, "缺少分支名", "", True

    dirty, detail = _has_worktree_changes()
    if dirty:
        return False, None, "工作区有未提交更改，请先提交/撤回/暂存（stash）后再切换分支", detail, True

    # Normalize
    is_remote, remote_ref, raw = _normalize_remote_ref(branch_name)

    # Prefer git switch, fallback to checkout
    if is_remote and remote_ref:
        # remote_ref like origin/foo
        local_name = remote_ref
        if "/" in remote_ref:
            local_name = remote_ref.split("/", 1)[1]

        if _branch_exists_local(local_name):
            out, err, code = run_git(["switch", local_name], timeout=120)
            if code != 0:
                out, err, code = run_git(["checkout", local_name], timeout=120)

            # Ensure upstream tracking is set (avoid detached-ish state and make pulls/pushes predictable)
            if code == 0:
                _, _, ucode = run_git(["branch", "--set-upstream-to", remote_ref, local_name], timeout=60)
                # ignore upstream set failures (e.g. remote ref missing), switch is already done
        else:
            out, err, code = run_git(["switch", "-c", local_name, "--track", remote_ref], timeout=120)
            if code != 0:
                out, err, code = run_git(["checkout", "-b", local_name, "--track", remote_ref], timeout=120)
    else:
        out, err, code = run_git(["switch", raw], timeout=120)
        if code != 0:
            out, err, code = run_git(["checkout", raw], timeout=120)

    if code != 0:
        return False, None, (err or "切换分支失败"), (out or ""), False

    # Refresh current branch
    _, cur = get_branches()
    return True, cur, "", (out or ""), False

# ════════════════════════════════════════════════════════
#  HTTP 处理器
# ════════════════════════════════════════════════════════

class Handler(BaseHTTPRequestHandler):
    
    # 禁用默认的日志输出
    def log_message(self, format, *args):
        # 使用我们的 logger 记录 HTTP 请求
        rid = getattr(self, "_req_id", None)
        if rid:
            logger.info(f"[{rid}] {self.address_string()} - {format % args}")
        else:
            logger.info(f"{self.address_string()} - {format % args}")

    def _ensure_req_id(self):
        try:
            rid = getattr(self, "_req_id", None)
            if rid:
                return rid
            rid = (self.headers.get("X-Req-Id") or "").strip()
            if not rid:
                rid = f"{int(time.time() * 1000)}-{os.getpid()}-{threading.get_ident()}"
            self._req_id = rid
            return rid
        except Exception:
            return None

    def _require_repo(self):
        if not REPO_PATH:
            self.send_json({"error": "未打开仓库"}, 400)
            return False
        return True

    def send_json(self, data, status=200):
        """发送 JSON 响应"""
        try:
            rid = self._ensure_req_id()

            payload = data
            if isinstance(data, dict):
                try:
                    if "req_id" not in data and rid:
                        data["req_id"] = rid
                except Exception as e:
                    logger.debug(f"Exception ignored: {e}")

                need_wrap = not ("ok" in data and "data" in data and "error" in data)
                if need_wrap:
                    err = data.get("error")
                    if err is None and data.get("ok") is False:
                        err = data.get("msg") or "请求失败"
                    ok_val = data.get("ok")
                    if ok_val is None:
                        ok_val = (err is None)
                    payload = {
                        "ok": bool(ok_val),
                        "data": data,
                        "error": err,
                        "req_id": rid,
                    }
                    payload.update(data)
                else:
                    if rid and ("req_id" not in data or not data.get("req_id")):
                        data["req_id"] = rid
                    payload = data

            body = json.dumps(payload, ensure_ascii=False).encode("utf-8")
            self.send_response(status)
            self.send_header("Content-Type", "application/json; charset=utf-8")
            self.send_header("Content-Length", len(body))
            self.send_header("Access-Control-Allow-Origin", "*")
            self.end_headers()
            self.wfile.write(body)
            
            # 记录响应
            if status >= 400:
                if rid:
                    logger.warning(f"[{rid}] 响应错误 {status}: {self.path} - {payload}")
                else:
                    logger.warning(f"响应错误 {status}: {self.path} - {payload}")
            else:
                if rid:
                    logger.debug(f"[{rid}] 响应成功 {status}: {self.path}")
                else:
                    logger.debug(f"响应成功 {status}: {self.path}")
        except (ConnectionAbortedError, BrokenPipeError, ConnectionResetError) as e:
            # 客户端主动中止/断开连接（例如页面刷新/取消请求）
            try:
                rid = self._ensure_req_id()
            except Exception:
                rid = ""
            if rid:
                logger.debug(f"[{rid}] 客户端已断开，响应发送中止: {e}")
            else:
                logger.debug(f"客户端已断开，响应发送中止: {e}")
        except Exception as e:
            logger.error(f"发送 JSON 响应失败: {e}", exc_info=True)

    def send_html(self):
        """发送 HTML 页面"""
        html_path = Path(__file__).parent / "index.html"
        try:
            data = html_path.read_bytes()
            self.send_response(200)
            self.send_header("Content-Type", "text/html; charset=utf-8")
            self.send_header("Content-Length", len(data))
            self.end_headers()
            self.wfile.write(data)
            logger.debug("发送 HTML 页面: index.html")
        except Exception as e:
            logger.error(f"发送 HTML 页面失败: {e}")
            body = f"找不到 index.html: {e}".encode()
            self.send_response(500)
            self.send_header("Content-Length", len(body))
            self.end_headers()
            self.wfile.write(body)

    def do_OPTIONS(self):
        """处理 OPTIONS 请求"""
        logger.debug(f"OPTIONS 请求: {self.path}")
        self.send_response(200)
        self.send_header("Access-Control-Allow-Origin", "*")
        self.send_header("Access-Control-Allow-Methods", "GET,POST,OPTIONS")
        self.send_header("Access-Control-Allow-Headers", "Content-Type, X-Req-Id")
        self.end_headers()

    def do_GET(self):
        """处理 GET 请求"""
        self._ensure_req_id()
        rid = getattr(self, "_req_id", None)
        if rid:
            logger.info(f"[{rid}] 收到 GET 请求: {self.path}")
        else:
            logger.info(f"收到 GET 请求: {self.path}")
        parsed = urlparse(self.path)
        p  = parsed.path
        qs = parse_qs(parsed.query)

        # 关键：unquote 解码 URL 编码（处理中文/空格路径）
        def qget(key):
            value = unquote(qs.get(key, [""])[0])
            logger.debug(f"查询参数 {key} = {value}")
            return value

        if p in ("/", "/index.html"):
            self.send_html()
            return

        elif p == "/api/undo_stats":
            try:
                sid = qget("session_id")
            except Exception:
                sid = ""
            payload = {"ok": True, "undo_steps": _undo_get_steps()}
            if sid:
                payload["session_id"] = sid
                payload["session_undo_steps"] = _undo_get_steps_for_session(sid)
            self.send_json(payload)
            return

        elif p == "/api/undo_history":
            try:
                sid = qget("session_id")
            except Exception:
                sid = ""
            try:
                limit = int(qget("limit") or "20")
            except Exception:
                limit = 20
            limit = max(1, min(200, int(limit)))
            prefix = ("ai-" + str(sid).strip() + "-") if str(sid or "").strip() else ""
            items = []
            try:
                with undo_lock:
                    order = list(_undo_group_order)
                    for gid in reversed(order):
                        g = str(gid or "").strip()
                        if not g:
                            continue
                        if prefix and (not g.startswith(prefix)):
                            continue
                        acts = _undo_groups.get(g)
                        if not isinstance(acts, list):
                            continue
                        types = []
                        created_ts = None
                        add_n = 0
                        del_n = 0
                        mod_n = 0
                        unknown_n = 0
                        file_paths = []
                        for a in acts:
                            if not isinstance(a, dict):
                                continue
                            tp = str(a.get("type") or "").strip()
                            if tp:
                                types.append(tp)
                            try:
                                ts0 = a.get("ts")
                                if ts0 is not None:
                                    tsf = float(ts0)
                                    if (created_ts is None) or (tsf < created_ts):
                                        created_ts = tsf
                            except Exception as e:
                                logger.debug(f"Exception ignored: {e}")
                            if tp == "file_snapshot":
                                op = str(a.get("op") or "").strip()
                                try:
                                    snap = a.get("snapshot")
                                except Exception:
                                    snap = None
                                if not op:
                                    try:
                                        if isinstance(snap, dict) and (snap.get("exists") is False):
                                            op = "add"
                                        else:
                                            op = ""
                                    except Exception:
                                        op = ""
                                if op == "add":
                                    add_n += 1
                                elif op == "delete":
                                    del_n += 1
                                elif op == "modify":
                                    mod_n += 1
                                else:
                                    unknown_n += 1
                                try:
                                    if isinstance(snap, dict):
                                        rp = str(snap.get("path") or "").strip().replace("\\", "/").lstrip("/")
                                        if rp:
                                            file_paths.append(rp)
                                except Exception as e:
                                    logger.debug(f"Exception ignored: {e}")
                            elif tp == "rename":
                                try:
                                    oldp = str(a.get("old_path") or "").strip().replace("\\", "/").lstrip("/")
                                    newp = str(a.get("new_path") or "").strip().replace("\\", "/").lstrip("/")
                                    if oldp:
                                        file_paths.append(oldp)
                                    if newp:
                                        file_paths.append(newp)
                                except Exception as e:
                                    logger.debug(f"Exception ignored: {e}")
                                try:
                                    snaps = a.get("snapshots")
                                    if isinstance(snaps, list):
                                        for s0 in snaps:
                                            if not isinstance(s0, dict):
                                                continue
                                            rp = str(s0.get("path") or "").strip().replace("\\", "/").lstrip("/")
                                            if rp:
                                                file_paths.append(rp)
                                except Exception as e:
                                    logger.debug(f"Exception ignored: {e}")
                        uniq_files = []
                        try:
                            seen = set()
                            for rp in file_paths:
                                if rp in seen:
                                    continue
                                seen.add(rp)
                                uniq_files.append(rp)
                        except Exception:
                            uniq_files = []
                        title_path = ""
                        try:
                            if uniq_files:
                                title_path = uniq_files[0]
                        except Exception:
                            title_path = ""
                        items.append({
                            "group_id": g,
                            "action_count": len(acts),
                            "types": types[:6],
                            "first_type": (types[0] if types else ""),
                            "last_type": (types[-1] if types else ""),
                            "created_ts": created_ts,
                            "add": add_n,
                            "delete": del_n,
                            "modify": mod_n,
                            "unknown": unknown_n,
                            "files": uniq_files[:6],
                            "file_count": len(uniq_files),
                            "title_path": title_path,
                        })
                        if len(items) >= limit:
                            break
            except Exception:
                items = []
            payload = {"ok": True, "items": items, "limit": limit}
            if sid:
                payload["session_id"] = sid
                payload["session_undo_steps"] = _undo_get_steps_for_session(sid)
            self.send_json(payload)
            return

        elif p == "/api/status":
            logger.debug("处理 /api/status 请求")
            origin_url = ""
            has_staged = False
            staged_count = 0
            has_stash = False
            stash_count = 0
            stash_top = ""
            if REPO_PATH:
                try:
                    out, _, code = run_git(["config", "--get", "remote.origin.url"])
                    if code == 0:
                        origin_url = (out or "").strip()
                except Exception:
                    origin_url = ""
                try:
                    sf, err_sf = get_staged_files()
                    if not err_sf and isinstance(sf, list):
                        staged_count = len(sf)
                        has_staged = staged_count > 0
                except Exception:
                    has_staged = False
                    staged_count = 0
                try:
                    list_out, _, list_code = run_git(["stash", "list"], timeout=30)
                    if list_code == 0:
                        raw = (list_out or "").strip()
                        if raw:
                            lines = [x for x in raw.splitlines() if x.strip()]
                            stash_count = len(lines)
                            has_stash = stash_count > 0
                            stash_top = lines[0].strip() if lines else ""
                except Exception:
                    has_stash = False
                    stash_count = 0
                    stash_top = ""
            self.send_json({
                "repo":  REPO_PATH,
                "valid": bool(REPO_PATH and
                              os.path.isdir(os.path.join(REPO_PATH, ".git"))),
                "ws_port": WS_PORT if WEBSOCKET_AVAILABLE else None,
                "origin_url": origin_url,
                "has_staged_changes": has_staged,
                "staged_count": staged_count,
                "has_stash": has_stash,
                "stash_count": stash_count,
                "stash_top": stash_top,
            })

        elif p == "/api/stash_list":
            logger.debug("处理 /api/stash_list 请求")
            if not self._require_repo():
                return
            list_out, list_err, list_code = run_git(["stash", "list"], timeout=30)
            if list_code != 0:
                self.send_json({"ok": False, "error": (list_err or "获取 stash 列表失败"), "output": (list_out or "").strip()}, 400)
                return
            raw = (list_out or "").strip()
            lines = [x for x in raw.splitlines() if x.strip()] if raw else []
            self.send_json({
                "ok": True,
                "has_stash": bool(lines),
                "stash_count": len(lines),
                "stash_top": lines[0].strip() if lines else "",
                "stashes": lines,
            })

        elif p == "/api/sync_status":
            logger.debug("处理 /api/sync_status 请求")
            if not self._require_repo():
                return
            try:
                fetch_flag = str(qget("fetch") or "0").strip()
                do_fetch = fetch_flag in ("1", "true", "True")
            except Exception:
                do_fetch = False
            ok, data, err = get_sync_status(do_fetch=do_fetch)
            if not ok:
                self.send_json({"ok": False, "error": err or "获取同步状态失败"}, 400)
                return
            self.send_json({"ok": True, **(data or {})})

        elif p == "/api/capabilities":
            logger.debug("处理 /api/capabilities 请求")
            try:
                spec = get_capabilities_spec()
                self.send_json({
                    "ok": True,
                    "version": spec.get("version"),
                    "generated_at": spec.get("generated_at"),
                    "text": spec.get("text"),
                    "endpoints": spec.get("endpoints"),
                    "agent_tools": spec.get("agent_tools"),
                })
            except Exception as e:
                logger.error(f"/api/capabilities failed: {e}")
                self.send_json({"ok": False, "error": str(e)}, 500)

        elif p == "/api/files":
            logger.debug("处理 /api/files 请求")
            if not self._require_repo():
                return
            try:
                max_age = float(qget("max_age") or "")
            except Exception:
                max_age = 1.0
            if max_age < 0:
                max_age = 0.0
            if max_age > 10:
                max_age = 10.0
            self.send_json({"files": get_changed_files_cached(max_age_sec=max_age)})

        elif p == "/api/staged_files":
            logger.debug("处理 /api/staged_files 请求")
            if not self._require_repo():
                return
            files2, err2 = get_staged_files()
            if err2:
                self.send_json({"ok": False, "error": err2}, 400)
                return
            self.send_json({"ok": True, "files": files2})

        elif p == "/api/diff":
            logger.debug("处理 /api/diff 请求")
            if not self._require_repo():
                return
            fp     = qget("path")
            status = qget("status") or "M"
            ctx = qget("ctx") or "5"
            hunks, err = get_file_diff(fp, status, ctx)
            self.send_json({"hunks": hunks, "error": err})

        elif p == "/api/file_content":
            logger.debug("处理 /api/file_content GET 请求")
            if not self._require_repo():
                return
            fp = qget("path")
            content, encoding = get_file_content(fp, return_encoding=True)
            if content is not None:
                self.send_json({"ok": True, "content": content, "encoding": encoding})
            else:
                self.send_json({"ok": False, "error": "无法读取文件内容"}, 404)

        elif p == "/api/raw_file":
            logger.debug("处理 /api/raw_file GET 请求")
            if not self._require_repo():
                return
            fp = qget("path")
            if not fp:
                self.send_json({"ok": False, "error": "缺少 path"}, 400)
                return
            full = _safe_repo_abspath(fp)
            if not full or (not os.path.exists(full)) or os.path.isdir(full):
                self.send_json({"ok": False, "error": "文件不存在"}, 404)
                return
            try:
                file_size = os.path.getsize(full)
            except Exception:
                file_size = 0

            # Hard cap for raw media streaming (avoid excessive bandwidth/abuse).
            MAX_RAW_FILE_BYTES = 256 * 1024 * 1024  # 256MB
            if file_size > MAX_RAW_FILE_BYTES:
                self.send_json({"ok": False, "error": f"文件过大（>{MAX_RAW_FILE_BYTES} bytes）"}, 413)
                return

            ctype, _enc = mimetypes.guess_type(full)
            if not ctype:
                ctype = "application/octet-stream"

            range_header = (self.headers.get("Range") or "").strip()
            start = 0
            end = max(0, file_size - 1)
            is_range = False
            if range_header.lower().startswith("bytes="):
                try:
                    spec = range_header.split("=", 1)[1].strip()
                    # Only support a single range: start-end
                    if "," in spec:
                        spec = spec.split(",", 1)[0].strip()
                    if "-" in spec:
                        a, b = spec.split("-", 1)
                        a = a.strip()
                        b = b.strip()
                        if a == "":
                            # suffix range: -N (last N bytes)
                            n = int(b or "0")
                            if n <= 0:
                                raise ValueError("invalid suffix")
                            start = max(0, file_size - n)
                            end = max(0, file_size - 1)
                        else:
                            start = int(a)
                            end = int(b) if b != "" else max(0, file_size - 1)
                        if start < 0:
                            start = 0
                        if end < start:
                            end = start
                        if file_size > 0:
                            end = min(end, file_size - 1)
                        is_range = True
                except Exception:
                    # Invalid Range
                    self.send_response(416)
                    self.send_header("Content-Range", f"bytes */{file_size}")
                    self.send_header("Access-Control-Allow-Origin", "*")
                    self.end_headers()
                    return

            if file_size <= 0:
                self.send_response(200)
                self.send_header("Content-Type", ctype)
                self.send_header("Content-Length", "0")
                self.send_header("Cache-Control", "no-store")
                self.send_header("Access-Control-Allow-Origin", "*")
                self.send_header("Accept-Ranges", "bytes")
                self.end_headers()
                return

            length = max(0, end - start + 1)
            try:
                self.send_response(206 if is_range else 200)
                self.send_header("Content-Type", ctype)
                self.send_header("Cache-Control", "no-store")
                self.send_header("Access-Control-Allow-Origin", "*")
                self.send_header("Accept-Ranges", "bytes")
                if is_range:
                    self.send_header("Content-Range", f"bytes {start}-{end}/{file_size}")
                self.send_header("Content-Length", str(length if is_range else file_size))
                self.end_headers()

                CHUNK = 64 * 1024
                with open(full, "rb") as f:
                    if is_range:
                        f.seek(start)
                        remain = length
                        while remain > 0:
                            buf = f.read(min(CHUNK, remain))
                            if not buf:
                                break
                            self.wfile.write(buf)
                            remain -= len(buf)
                    else:
                        while True:
                            buf = f.read(CHUNK)
                            if not buf:
                                break
                            self.wfile.write(buf)
            except Exception as e:
                logger.error(f"发送 raw_file 失败: {fp} - {e}", exc_info=True)
            return

        elif p == "/api/log":
            logger.debug("处理 /api/log 请求")
            if not self._require_repo():
                return
            self.send_json({"log": get_log(limit=50)})

        elif p == "/api/commits":
            logger.debug("处理 /api/commits 请求")
            if not self._require_repo():
                return
            try:
                limit = int(qget("limit") or "50")
            except Exception:
                limit = 50
            self.send_json({"ok": True, "commits": get_log(limit=limit), "limit": limit})

        elif p == "/api/file_log":
            logger.debug("处理 /api/file_log 请求")
            if not self._require_repo():
                return
            fp = qget("path")
            if not fp:
                self.send_json({"error": "缺少 path"}, 400)
                return
            try:
                limit = int(qget("limit") or "50")
            except Exception:
                limit = 50
            self.send_json({"log": get_file_log(fp, limit=limit)})

        elif p == "/api/commit_detail":
            logger.debug("处理 /api/commit_detail 请求")
            if not self._require_repo():
                return
            commit = qget("hash")
            if not commit:
                self.send_json({"error":"缺少 hash"}, 400)
                return
            detail = get_commit_detail(commit)
            if "error" not in detail:
                # 确保返回完整hash和短hash
                if "full_hash" not in detail and "hash" in detail:
                    full_hash = detail.get("hash", "")
                    if len(full_hash) > 7:
                        detail["full_hash"] = full_hash
                        detail["hash"] = full_hash[:7]
                head_out, _, head_code = run_git(["rev-parse", "HEAD"])
                head_full = (head_out or "").strip() if head_code == 0 else ""
                detail["is_head"] = bool(head_full and (detail.get("full_hash") == head_full or detail.get("hash") == head_full))
            self.send_json(detail)

        elif p == "/api/commit_file_diff":
            logger.debug("处理 /api/commit_file_diff 请求")
            if not self._require_repo():
                return
            commit = qget("hash")
            fp     = qget("path")
            hunks  = get_commit_file_diff(commit, fp)
            self.send_json({"hunks": hunks})

        elif p == "/api/branches":
            logger.debug("处理 /api/branches 请求")
            if not self._require_repo():
                return
            branches, current = get_branches()
            self.send_json({"branches": branches, "current": current})

        elif p == "/api/commit_push_status":
            logger.debug("处理 /api/commit_push_status 请求")
            if not self._require_repo():
                return
            commit = qget("hash")
            pushed, branches, err = is_commit_pushed(commit)
            self.send_json({"pushed": pushed, "branches": branches, "error": err})

        elif p == "/api/ai_config":
            logger.debug("处理 /api/ai_config 请求")
            with ai_config_lock:
                self.send_json({"config": load_hivo_ai_config()})

        elif p == "/api/hivo_ai_config":
            logger.debug("处理 /api/hivo_ai_config 请求")
            cfg0 = _hivo_load_cfg()
            self.send_json({"ok": True, "config": cfg0})

        elif p == "/api/workspace_context":
            logger.debug("处理 /api/workspace_context 请求")
            if not self._require_repo():
                return
            ctx = get_workspace_context(max_entries=80)
            self.send_json({"ok": True, "context": ctx})

        elif p == "/api/find_files":
            logger.debug("处理 /api/find_files 请求")
            if not self._require_repo():
                return
            name = qget("name") or qget("q") or ""
            items = find_files_by_name(name, max_results=200)
            self.send_json({"ok": True, "files": items})

        elif p == "/api/list_dir_tree":
            logger.debug("处理 /api/list_dir_tree 请求")
            if not self._require_repo():
                return
            rel = qget("path") or qget("dir") or ""
            try:
                depth = int(qget("depth") or 3)
            except Exception:
                depth = 3
            try:
                cap = int(qget("max_entries") or 500)
            except Exception:
                cap = 500
            tree, err = list_dir_tree(rel, max_depth=depth, max_entries=cap)
            if tree is None:
                self.send_json({"ok": False, "msg": err or "生成失败"}, 400)
                return
            self.send_json({"ok": True, "tree": tree})

        elif p == "/api/read_file_range":
            logger.debug("处理 /api/read_file_range 请求")
            if not self._require_repo():
                return
            fp = qget("path") or ""
            try:
                start = int(qget("start") or 1)
            except Exception:
                start = 1
            try:
                end = int(qget("end") or (start + 200))
            except Exception:
                end = start + 200
            data, err = read_file_range(fp, start=start, end=end)
            if data is None:
                self.send_json({"ok": False, "msg": err or "读取失败"}, 400)
                return
            self.send_json({"ok": True, "data": data})

        elif p == "/api/search_code":
            logger.debug("处理 /api/search_code 请求")
            if not self._require_repo():
                return
            try:
                pid = qget("profile_id") or ""
                sid = qget("session_id") or ""
                key = _rl_make_key(self, profile_id=pid, session_id=sid)
                ok_rl, delay_s, retry_after = _rl_check_and_get_delay(key, "search")
                if not ok_rl:
                    self.send_json({"ok": False, "msg": f"请求过于频繁，请在 {retry_after} 秒后重试", "retry_after": retry_after}, 429)
                    return
                if delay_s > 0:
                    time.sleep(delay_s)
            except Exception as e:
                logger.debug(f"Exception ignored: {e}")
            q = qget("query") or qget("q") or ""
            cs = str(qget("case_sensitive") or "").strip().lower() in ("1", "true", "yes", "y")
            try:
                cap = int(qget("max_results") or 50)
            except Exception:
                cap = 50
            hits, err = search_code(q, case_sensitive=cs, max_results=cap)
            if hits is None:
                self.send_json({"ok": False, "msg": err or "搜索失败"}, 400)
                return
            self.send_json({"ok": True, "hits": hits})

        elif p == "/api/ai_chat_history":
            pid = qget("profile_id")
            sid = qget("session_id")
            try:
                limit = int(qget("limit") or "40")
            except Exception:
                limit = 40
            logger.info(f"处理 /api/ai_chat_history GET 请求 (profile_id={pid or ''}, session_id={sid or ''}, limit={limit})")
            with ai_history_lock:
                hist_pack = load_ai_chat_history(pid, limit=limit, session_id=sid)
                hist = hist_pack[0]
                info = list_ai_sessions(pid) if pid else {"active_session_id": "", "sessions": []}
            self.send_json({"ok": True, "messages": hist, "active_session_id": info.get("active_session_id"), "sessions": info.get("sessions")})

        elif p == "/api/ai_sessions":
            pid = qget("profile_id")
            logger.info(f"处理 /api/ai_sessions GET 请求 (profile_id={pid or ''})")
            with ai_history_lock:
                info = list_ai_sessions(pid)
            self.send_json({"ok": True, "active_session_id": info.get("active_session_id"), "sessions": info.get("sessions")})

        elif p == "/api/ai_attachment":
            # Serve files saved by /api/ai_store_attachments
            try:
                rel = qget("path") or ""
                rel = str(rel or "").replace("\\", "/").lstrip("/")
                if not rel.startswith("attachments/"):
                    self.send_json({"error": "非法路径"}, 400)
                    return
                base = _hivo_ai_data_dir()
                full = os.path.join(str(base), rel)
                full = os.path.abspath(full)
                # Ensure inside data dir
                if not full.startswith(os.path.abspath(str(base))):
                    self.send_json({"error": "非法路径"}, 400)
                    return
                if (not os.path.exists(full)) or (not os.path.isfile(full)):
                    self.send_json({"error": "文件不存在"}, 404)
                    return
                ctype, _ = mimetypes.guess_type(full)
                ctype = ctype or "application/octet-stream"
                try:
                    with open(full, "rb") as f:
                        data = f.read()
                    self.send_response(200)
                    self.send_header("Content-Type", ctype)
                    self.send_header("Content-Length", str(len(data)))
                    self.end_headers()
                    self.wfile.write(data)
                except Exception as e:
                    self.send_json({"error": str(e)}, 500)
            except Exception:
                self.send_json({"error": "读取失败"}, 500)

        elif p == "/api/tags":
            # 列出所有 tag
            if not REPO_PATH:
                self.send_json({"ok": False, "error": "未打开仓库"}, 400)
                return
            try:
                # 格式：refname|objectname|taggername|taggerdate|subject
                fmt = "%(refname:short)|%(objectname)|%(taggername)|%(taggerdate:short)|%(subject)"
                out, err, code = run_git(["for-each-ref", "--format=" + fmt, "refs/tags"], timeout=60)
                tags = []
                if code == 0 and out:
                    for line in out.splitlines():
                        parts = line.split("|", 4)
                        if len(parts) >= 5:
                            tags.append({
                                "name": parts[0],
                                "commit": parts[1],
                                "tagger": parts[2],
                                "date": parts[3],
                                "message": parts[4]
                            })
                # 同时获取指向的提交短 hash
                for t in tags:
                    try:
                        o2, e2, c2 = run_git(["rev-parse", "--short", t["commit"]], timeout=30)
                        if c2 == 0 and o2:
                            t["short"] = o2.strip()
                    except Exception:
                        t["short"] = (t["commit"] or "")[:7]
                self.send_json({"ok": True, "tags": tags})
            except Exception as e:
                logger.error(f"处理 /api/tags 异常: {e}", exc_info=True)
                self.send_json({"ok": False, "error": str(e)}, 500)

        elif p == "/api/conflicts":
            # 获取冲突文件列表 + 每个文件的冲突块
            if not REPO_PATH:
                self.send_json({"ok": False, "error": "未打开仓库"}, 400)
                return
            try:
                conflict_files, _ = get_unmerged_files()
                files_info = []
                for f in conflict_files:
                    info = {"path": f, "hunks": []}
                    try:
                        # 读取文件内容，提取 <<<<<<< / ======= / >>>>>>> 冲突块
                        full = os.path.join(REPO_PATH, f)
                        if os.path.exists(full):
                            with open(full, "r", encoding="utf-8", errors="replace") as fp:
                                content = fp.read()
                            info["hunks"] = _extract_conflict_hunks(content)
                    except Exception as e3:
                        info["error"] = str(e3)
                    files_info.append(info)
                self.send_json({"ok": True, "files": files_info})
            except Exception as e:
                logger.error(f"处理 /api/conflicts 异常: {e}", exc_info=True)
                self.send_json({"ok": False, "error": str(e)}, 500)

        else:
            logger.warning(f"未知的 GET 请求路径: {p}")
            self.send_json({"error":"Not found"}, 404)

    def do_POST(self):
        """处理 POST 请求"""
        global REPO_PATH
        self._ensure_req_id()
        rid = getattr(self, "_req_id", None)
        if rid:
            logger.info(f"[{rid}] 收到 POST 请求: {self.path}")
        else:
            logger.info(f"收到 POST 请求: {self.path}")
        p  = urlparse(self.path).path
        try:
            length = int(self.headers.get('Content-Length', '0'))
            try:
                cfg0 = _hivo_load_cfg()
                limits0 = cfg0.get("limits") if isinstance(cfg0, dict) else None
                limits0 = limits0 if isinstance(limits0, dict) else {}
                # Special-case larger body for attachment uploads (base64 inflates ~4/3)
                if p == "/api/ai_store_attachments":
                    max_body = int(limits0.get("max_upload_request_bytes") or (32 * 1024 * 1024))
                else:
                    max_body = int(limits0.get("max_post_body_bytes") or 0)
            except Exception:
                max_body = 0
            if max_body > 0 and length > max_body:
                logger.warning(f"[{rid}] 请求体过大: {length} bytes > {max_body} bytes")
                self.send_json({"ok": False, "error": f"附件或请求内容过大（>{max_body/1024/1024:.1f} MB），请减少附件数量或压缩图片"}, 413)
                return
            body = self.rfile.read(length).decode('utf-8')
            data = json.loads(body) if body else {}
            logger.debug(f"POST 请求数据: {json.dumps(data, ensure_ascii=False)[:200]}...")
        except Exception as e:
            logger.error(f"解析 JSON 失败: {e}")
            self.send_json({"ok": False, "error": "请求体 JSON 解析失败"}, 400)
            return

        try:
            if p == "/api/open_repo":
                logger.info("处理 /api/open_repo 请求")
                raw = data.get("path", "").strip()
                if not raw:
                    logger.warning("打开仓库失败: 路径为空")
                    self.send_json({"error":"路径为空"}, 400)
                    return
                logger.info(f"尝试打开仓库: {raw}")
                ok, msg = open_repo(raw)
                if not ok:
                    logger.error(f"打开仓库失败: {msg} - {raw}")
                    self.send_json({"error": msg or "打开仓库失败"}, 400)
                    return
                # 获取当前分支
                branch = get_current_branch()
                # 获取仓库远程地址
                try:
                    _, _, st0 = get_repo_status()
                    origin_url = st0.get("origin_url") if isinstance(st0, dict) else ""
                except Exception:
                    origin_url = ""
                self.send_json({"repo": REPO_PATH, "branch": branch, "origin_url": origin_url})

            elif p == "/api/set_origin":
                logger.info("处理 /api/set_origin 请求")
                url0 = ""
                try:
                    url0 = str((data or {}).get("url") or "").strip()
                except Exception:
                    url0 = ""
                ok, msg, origin_url = set_origin_url(url0)
                if not ok:
                    self.send_json({"error": msg or "设置远程地址失败"}, 400)
                    return
                self.send_json({"ok": True, "origin_url": origin_url})

            elif p == "/api/unstage_file":
                logger.info("处理 /api/unstage_file 请求")
                if not self._require_repo():
                    return
                fp = (data.get("path") or "").strip()
                if not fp:
                    self.send_json({"error": "缺少 path"}, 400)
                    return
                _, err, code = run_git(["restore", "--staged", "--", fp], timeout=120)
                if code != 0:
                    self.send_json({"ok": False, "error": err or "取消暂存失败"}, 400)
                    return
                try:
                    invalidate_changed_files_cache()
                    notify_files_updated()
                except Exception as e:
                    logger.debug(f"Exception ignored: {e}")
                self.send_json({"ok": True})

            elif p == "/api/discard_staged_file":
                logger.info("处理 /api/discard_staged_file 请求")
                if not self._require_repo():
                    return
                fp = (data.get("path") or "").strip()
                if not fp:
                    self.send_json({"error": "缺少 path"}, 400)
                    return
                _, err, code = run_git(["restore", "--staged", "--worktree", "--source=HEAD", "--", fp], timeout=120)
                if code != 0:
                    self.send_json({"ok": False, "error": err or "丢弃暂存区失败"}, 400)
                    return
                try:
                    invalidate_changed_files_cache()
                    notify_files_updated()
                except Exception as e:
                    logger.debug(f"Exception ignored: {e}")
                self.send_json({"ok": True})

            elif p == "/api/unstage_all_staged":
                logger.info("处理 /api/unstage_all_staged 请求")
                if not self._require_repo():
                    return
                _, err, code = run_git(["restore", "--staged", "."], timeout=180)
                if code != 0:
                    self.send_json({"ok": False, "error": err or "取消全部暂存失败"}, 400)
                    return
                try:
                    invalidate_changed_files_cache()
                    notify_files_updated()
                except Exception as e:
                    logger.debug(f"Exception ignored: {e}")
                self.send_json({"ok": True})

            elif p == "/api/discard_all_staged":
                logger.info("处理 /api/discard_all_staged 请求")
                if not self._require_repo():
                    return
                _, err, code = run_git(["restore", "--staged", "--worktree", "--source=HEAD", "."], timeout=180)
                if code != 0:
                    self.send_json({"ok": False, "error": err or "丢弃全部暂存失败"}, 400)
                    return
                try:
                    invalidate_changed_files_cache()
                    notify_files_updated()
                except Exception as e:
                    logger.debug(f"Exception ignored: {e}")
                self.send_json({"ok": True})

            elif p == "/api/revert_hunk":
                logger.info("处理 /api/revert_hunk 请求")
                if not self._require_repo():
                    return
                fp  = data.get("path", "")
                idx = data.get("hunk_index", -1)
                st  = data.get("status", "M")
                with _apply_revert_lock:
                    ok, msg = revert_hunk(fp, int(idx), st)
                    if ok:
                        invalidate_changed_files_cache()
                        notify_files_updated()
                self.send_json({"ok": ok, "msg": msg})

            elif p == "/api/revert_line":
                logger.info("处理 /api/revert_line 请求")
                if not self._require_repo():
                    return
                fp  = data.get("path", "")
                h_idx = data.get("hunk_index", -1)
                l_idx = data.get("line_index", -1)
                st  = data.get("status", "M")
                ctx = data.get("ctx", 5)
                sig = data.get("signature")
                with _apply_revert_lock:
                    ok, msg = revert_line(fp, int(h_idx), int(l_idx), st, ctx, sig)
                    if ok:
                        invalidate_changed_files_cache()
                        notify_files_updated()
                self.send_json({"ok": ok, "msg": msg})
            
            elif p == "/api/revert_multi_lines":
                logger.info("处理 /api/revert_multi_lines 请求")
                if not self._require_repo():
                    return
                fp  = data.get("path", "")
                h_idx = data.get("hunk_index", -1)
                start_l_idx = data.get("start_line_index", -1)
                end_l_idx = data.get("end_line_index", -1)
                st  = data.get("status", "M")
                ctx = data.get("ctx", 5)
                with _apply_revert_lock:
                    ok, msg = revert_multi_lines(fp, int(h_idx), int(start_l_idx), int(end_l_idx), st, ctx)
                    if ok:
                        invalidate_changed_files_cache()
                        notify_files_updated()
                self.send_json({"ok": ok, "msg": msg})
            
            elif p == "/api/file_content":
                logger.info("处理 /api/file_content 请求")
                if not self._require_repo():
                    return
                fp = data.get("path", "")
                content, encoding = get_file_content(fp, return_encoding=True)
                if content is not None:
                    self.send_json({"ok": True, "content": content, "encoding": encoding})
                else:
                    self.send_json({"ok": False, "error": "无法读取文件内容"})

            elif p == "/api/file_content_head":
                logger.info("处理 /api/file_content_head 请求")
                if not self._require_repo():
                    return
                fp = data.get("path", "")
                content = get_head_file_content(fp)
                if content is not None:
                    self.send_json({"ok": True, "content": content})
                else:
                    self.send_json({"ok": False, "content": None})

            elif p == "/api/open_file":
                logger.info("处理 /api/open_file 请求")
                if not self._require_repo():
                    return
                name = (data.get("name") or data.get("path") or "").strip()
                view = (data.get("view") or "editor").strip().lower()
                if not name:
                    self.send_json({"ok": False, "msg": "缺少 name"}, 400)
                    return

                # Normalize view
                if view not in ("editor", "change", "split", "unified"):
                    view = "editor"

                pick = ""
                # If looks like a path, try as-is first.
                cand = name.replace("\\", "/").lstrip("/")
                full = _safe_repo_abspath(cand)
                if full and os.path.isfile(full):
                    pick = cand
                else:
                    # Find by name (basename or partial)
                    items = find_files_by_name(name, max_results=20)
                    if not items:
                        self.send_json({"ok": False, "msg": "未找到文件"}, 404)
                        return
                    pick = str(items[0] or "")
                    nm_lower = name.lower()
                    for rel in items:
                        bn = str(rel or "").replace("\\", "/").split("/")[-1]
                        if bn.lower() == nm_lower:
                            pick = str(rel or "")
                            break

                content, encoding = get_file_content(pick, return_encoding=True)
                if content is None:
                    self.send_json({"ok": False, "msg": "无法读取文件内容", "path": pick, "view": view}, 400)
                    return
                self.send_json({
                    "ok": True,
                    "path": pick,
                    "view": view,
                    "content": content,
                    "encoding": encoding,
                })

            elif p == "/api/verify_python":
                logger.info("处理 /api/verify_python 请求")
                if not self._require_repo():
                    return
                arr = data.get("paths") or data.get("files") or []
                if not isinstance(arr, list):
                    arr = []
                paths = [str(x or "").strip().lstrip("@") for x in arr]
                paths = [p for p in paths if p]
                if not paths:
                    paths = ["server.py"]
                if len(paths) > 50:
                    paths = paths[:50]

                checked = []
                for rp in paths:
                    safe = _safe_repo_abspath(rp)
                    if not safe or (not os.path.isfile(safe)):
                        self.send_json({"ok": False, "files": checked, "error": f"非法路径或文件不存在: {rp}"}, 400)
                        return
                    checked.append(rp.replace("\\", "/"))
                    try:
                        py_compile.compile(safe, doraise=True)
                    except py_compile.PyCompileError as e:
                        self.send_json({"ok": False, "files": checked, "error": str(e)}, 400)
                        return
                    except Exception as e:
                        self.send_json({"ok": False, "files": checked, "error": str(e)}, 500)
                        return
                self.send_json({"ok": True, "files": checked})

            elif p == "/api/save_file":
                logger.info("处理 /api/save_file 请求")
                if not self._require_repo():
                    return
                try:
                    key = _rl_make_key(self)
                    ok_rl, delay_s, retry_after = _rl_check_and_get_delay(key, "modify")
                    if not ok_rl:
                        self.send_json({"ok": False, "msg": f"请求过于频繁，请在 {retry_after} 秒后重试", "retry_after": retry_after}, 429)
                        return
                    if delay_s > 0:
                        time.sleep(delay_s)
                except Exception as e:
                    logger.debug(f"Exception ignored: {e}")
                idem_key = (self.headers.get("X-Idempotency-Key") or "").strip()
                if idem_key:
                    ent = _idempotency_get(idem_key)
                    if ent and isinstance(ent.get("payload"), dict):
                        self.send_json(ent.get("payload"), int(ent.get("code") or 200))
                        return
                fp = data.get("path", "")
                content = data.get("content", "")
                undo_gid = (self.headers.get("X-Undo-Group") or "").strip()
                try:
                    cur = get_file_content(fp)
                    if isinstance(cur, str) and isinstance(content, str) and cur == content:
                        payload = {"ok": True, "msg": "", "path": str(fp or "").replace("\\", "/").lstrip("/"), "no_change": True}
                        if idem_key:
                            _idempotency_set(idem_key, payload, 200)
                        self.send_json(payload)
                        return
                except Exception as e:
                    logger.debug(f"Exception ignored: {e}")
                if undo_gid:
                    snap = _undo_capture_file_snapshot(fp)
                    if snap is not None:
                        op0 = "add" if (isinstance(snap, dict) and (snap.get("exists") is False)) else "modify"
                        _undo_record(undo_gid, {"type": "file_snapshot", "op": op0, "snapshot": snap})
                ok, msg = save_file_content(fp, content)
                if ok:
                    invalidate_changed_files_cache()
                    notify_files_updated()
                try:
                    norm_path = str(fp or "").replace("\\", "/").lstrip("/")
                    if norm_path.startswith("./"):
                        norm_path = norm_path[2:]
                except Exception:
                    norm_path = str(fp or "")
                payload = {"ok": ok, "msg": msg, "path": norm_path}
                if idem_key:
                    _idempotency_set(idem_key, payload, 200)
                self.send_json(payload)

            elif p == "/api/mkdir":
                logger.info("处理 /api/mkdir 请求")
                if not self._require_repo():
                    return
                try:
                    key = _rl_make_key(self)
                    ok_rl, delay_s, retry_after = _rl_check_and_get_delay(key, "modify")
                    if not ok_rl:
                        self.send_json({"ok": False, "msg": f"请求过于频繁，请在 {retry_after} 秒后重试", "retry_after": retry_after}, 429)
                        return
                    if delay_s > 0:
                        time.sleep(delay_s)
                except Exception as e:
                    logger.debug(f"Exception ignored: {e}")
                idem_key = (self.headers.get("X-Idempotency-Key") or "").strip()
                if idem_key:
                    ent = _idempotency_get(idem_key)
                    if ent and isinstance(ent.get("payload"), dict):
                        self.send_json(ent.get("payload"), int(ent.get("code") or 200))
                        return
                fp = data.get("path", "")
                ok, msg = mkdir_repo(fp)
                if ok:
                    invalidate_changed_files_cache()
                    notify_files_updated()
                try:
                    norm_path = str(fp or "").replace("\\", "/").lstrip("/")
                    if norm_path.startswith("./"):
                        norm_path = norm_path[2:]
                except Exception:
                    norm_path = str(fp or "")
                payload = {"ok": ok, "msg": msg, "path": norm_path}
                if idem_key:
                    _idempotency_set(idem_key, payload, 200)
                self.send_json(payload)

            elif p == "/api/browser_tool":
                logger.info("处理 /api/browser_tool 请求")
                action = str((data.get("action") if isinstance(data, dict) else "") or "").strip().lower()
                sid = str((data.get("session_id") if isinstance(data, dict) else "") or "").strip() or "global"
                params = (data.get("params") if isinstance(data, dict) else None)
                params = params if isinstance(params, dict) else {}

                if not action:
                    self.send_json({"ok": False, "msg": "缺少 action"}, 400)
                    return

                if action == "open":
                    ok, msg, out = _pw_open(sid, str(params.get("url") or ""), wait_ms=int(params.get("wait_ms") or 0))
                    self.send_json({"ok": bool(ok), "msg": msg, **(out if isinstance(out, dict) else {})})
                    return
                if action == "click":
                    ok, msg, out = _pw_click(sid, str(params.get("selector") or ""), wait_ms=int(params.get("wait_ms") or 0))
                    self.send_json({"ok": bool(ok), "msg": msg, **(out if isinstance(out, dict) else {})})
                    return
                if action == "type":
                    ok, msg, out = _pw_type(
                        sid,
                        str(params.get("selector") or ""),
                        str(params.get("text") or ""),
                        clear=bool(params.get("clear", True)),
                        wait_ms=int(params.get("wait_ms") or 0),
                    )
                    self.send_json({"ok": bool(ok), "msg": msg, **(out if isinstance(out, dict) else {})})
                    return
                if action == "eval":
                    ok, msg, out = _pw_eval(sid, str(params.get("script") or ""))
                    self.send_json({"ok": bool(ok), "msg": msg, **(out if isinstance(out, dict) else {})})
                    return
                if action == "text":
                    ok, msg, out = _pw_text(sid, selector=str(params.get("selector") or "").strip())
                    self.send_json({"ok": bool(ok), "msg": msg, **(out if isinstance(out, dict) else {})})
                    return
                if action == "screenshot":
                    ok, msg, out = _pw_screenshot(sid, full_page=bool(params.get("full_page", True)))
                    self.send_json({"ok": bool(ok), "msg": msg, **(out if isinstance(out, dict) else {})})
                    return
                if action == "wait":
                    ok, msg, out = _pw_wait(sid, ms=int(params.get("ms") or 0))
                    self.send_json({"ok": bool(ok), "msg": msg, **(out if isinstance(out, dict) else {})})
                    return
                if action == "close":
                    if not _pw_is_available():
                        self.send_json({"ok": False, "msg": "Playwright 未安装（可选能力）。"}, 400)
                        return
                    _pw_close_session(sid)
                    self.send_json({"ok": True, "msg": "", "session_id": sid})
                    return

                self.send_json({"ok": False, "msg": f"不支持的 action: {action}"}, 400)
                return

            elif p == "/api/delete_file":
                logger.info("处理 /api/delete_file 请求")
                if not self._require_repo():
                    return
                idem_key = (self.headers.get("X-Idempotency-Key") or "").strip()
                if idem_key:
                    ent = _idempotency_get(idem_key)
                    if ent and isinstance(ent.get("payload"), dict):
                        self.send_json(ent.get("payload"), int(ent.get("code") or 200))
                        return
                fp = (data.get("path") or "").strip()
                undo_gid = (self.headers.get("X-Undo-Group") or "").strip()
                if undo_gid:
                    snap = _undo_capture_file_snapshot(fp)
                    if snap is not None:
                        _undo_record(undo_gid, {"type": "file_snapshot", "op": "delete", "snapshot": snap})
                full = _safe_repo_abspath(fp)
                if not full:
                    payload = {"ok": False, "msg": "非法路径"}
                    if idem_key:
                        _idempotency_set(idem_key, payload, 400)
                    self.send_json(payload, 400)
                    return
                try:
                    if os.path.isdir(full):
                        payload = {"ok": False, "msg": "目标是目录"}
                        if idem_key:
                            _idempotency_set(idem_key, payload, 400)
                        self.send_json(payload, 400)
                        return
                    if os.path.exists(full):
                        os.remove(full)
                    invalidate_changed_files_cache()
                    notify_files_updated()
                    payload = {"ok": True, "msg": "删除成功"}
                    if idem_key:
                        _idempotency_set(idem_key, payload, 200)
                    self.send_json(payload)
                except Exception as e:
                    logger.error(f"删除文件失败: {fp} - {e}", exc_info=True)
                    payload = {"ok": False, "msg": str(e)}
                    if idem_key:
                        _idempotency_set(idem_key, payload, 500)
                    self.send_json(payload, 500)

            elif p == "/api/run_cmd":
                logger.info("处理 /api/run_cmd 请求")
                if not self._require_repo():
                    return
                cmd = (data.get("cmd") or "").strip()
                if not cmd:
                    self.send_json({"ok": False, "msg": "缺少 cmd"}, 400)
                    return

                undo_gid = (self.headers.get("X-Undo-Group") or "").strip()
                pre_map, pre_snaps = _undo_prepare_cmd_snapshots(undo_gid)
                try:
                    timeout = int(data.get("timeout") or 30)
                except Exception:
                    timeout = 30
                if timeout <= 0:
                    timeout = 30
                if timeout > 600:
                    timeout = 600

                cwd_raw = data.get("cwd")
                cwd = REPO_PATH
                if cwd_raw is not None:
                    rel = str(cwd_raw).strip()
                    if rel:
                        full = _safe_repo_abspath(rel)
                        if not full or (not os.path.isdir(full)):
                            self.send_json({"ok": False, "msg": "非法 cwd"}, 400)
                            return
                        cwd = full

                # Reject commands with cmd metacharacters to prevent injection
                dangerous = re.search(r'[&|;`$(){}\[\]<>]', cmd)
                if dangerous:
                    self.send_json({"ok": False, "msg": "命令包含不允许的特殊字符"}, 400)
                    return
                try:
                    r = subprocess.run(
                        ["cmd.exe", "/c", cmd],
                        cwd=cwd,
                        capture_output=True,
                        text=True,
                        timeout=timeout,
                    )

                    try:
                        logger.warning(f"run_cmd 完成 (ok={1 if r.returncode == 0 else 0}, exit_code={int(r.returncode)}, cmd={cmd[:180]})")
                    except Exception as e:
                        logger.debug(f"Exception ignored: {e}")
                    try:
                        _undo_finalize_cmd_snapshots(undo_gid, pre_map, pre_snaps)
                    except Exception as e:
                        logger.debug(f"Exception ignored: {e}")

                    try:
                        invalidate_changed_files_cache()
                        notify_files_updated()
                    except Exception as e:
                        logger.debug(f"Exception ignored: {e}")

                    self.send_json({
                        "ok": (r.returncode == 0),
                        "exit_code": int(r.returncode),
                        "stdout": r.stdout or "",
                        "stderr": r.stderr or "",
                    })
                except subprocess.TimeoutExpired:
                    self.send_json({"ok": False, "msg": f"命令超时(超过{timeout}秒)", "exit_code": 124, "stdout": "", "stderr": ""}, 500)
                except Exception as e:
                    logger.error(f"执行命令失败: {cmd} - {e}", exc_info=True)
                    self.send_json({"ok": False, "msg": str(e), "exit_code": 1, "stdout": "", "stderr": ""}, 500)

            elif p == "/api/run_cmd_script":
                logger.info("处理 /api/run_cmd_script 请求")
                if not self._require_repo():
                    return
                cmds = data.get("cmds")
                if not isinstance(cmds, list) or len(cmds) == 0:
                    script_text = str(data.get("script") or "").strip()
                    if not script_text:
                        self.send_json({"ok": False, "msg": "缺少 cmds 或 script"}, 400)
                        return
                    cmds = [l.strip() for l in script_text.split('\n') if l.strip() and not l.strip().startswith('#')]
                if not cmds:
                    self.send_json({"ok": False, "msg": "没有有效命令"}, 400)
                    return

                timeout = int(data.get("timeout") or 120)
                timeout = max(10, min(600, timeout))

                cwd = REPO_PATH
                cwd_raw = data.get("cwd")
                if cwd_raw:
                    safe = _safe_repo_abspath(str(cwd_raw).strip())
                    if safe and os.path.isdir(safe):
                        cwd = safe

                undo_gid = (self.headers.get("X-Undo-Group") or "").strip()
                pre_map, pre_snaps = _undo_prepare_cmd_snapshots(undo_gid)

                fd, script_path = tempfile.mkstemp(suffix='.bat', prefix='gm_script_', dir=REPO_PATH)
                try:
                    with os.fdopen(fd, 'w', encoding='utf-8', newline='') as f:
                        f.write('@echo off\r\n')
                        f.write('setlocal enabledelayedexpansion\r\n')
                        for i, cmd in enumerate(cmds):
                            c = cmd.strip()
                            if not c:
                                continue
                            f.write(f'{c}\r\n')
                            f.write(f'echo [GM:EXIT:{i}:!ERRORLEVEL!]\r\n')
                        f.write('endlocal\r\n')

                    try:
                        r = subprocess.run(
                            ["cmd.exe", "/c", script_path],
                            cwd=cwd,
                            capture_output=True,
                            text=True,
                            timeout=timeout,
                            encoding='utf-8',
                            errors='replace',
                        )
                        try:
                            logger.warning(f"run_cmd_script 完成 (ok={1 if r.returncode == 0 else 0}, cmds={len(cmds)})")
                        except Exception as e:
                            logger.debug(f"Exception ignored: {e}")

                        try:
                            _undo_finalize_cmd_snapshots(undo_gid, pre_map, pre_snaps)
                            invalidate_changed_files_cache()
                            notify_files_updated()
                        except Exception as e:
                            logger.debug(f"Exception ignored: {e}")

                        full_output = (r.stdout or "")
                        outputs = []
                        exit_matches = re.findall(r'\[GM:EXIT:(\d+):(-?\d+)\]', full_output)
                        parts = re.split(r'\[GM:EXIT:\d+:-?\d+\]', full_output)

                        for idx, match in enumerate(exit_matches):
                            cmd_idx = int(match[0])
                            exit_code = int(match[1])
                            out_text = (parts[idx].strip() if idx < len(parts) else "")
                            outputs.append({
                                "idx": cmd_idx,
                                "cmd": cmds[cmd_idx] if cmd_idx < len(cmds) else "",
                                "ok": exit_code == 0,
                                "exit_code": exit_code,
                                "output": out_text,
                            })

                        self.send_json({
                            "ok": r.returncode == 0,
                            "exit_code": r.returncode,
                            "stdout": r.stdout or "",
                            "stderr": r.stderr or "",
                            "outputs": outputs,
                        })
                    except subprocess.TimeoutExpired:
                        self.send_json({"ok": False, "msg": f"脚本执行超时({timeout}秒)", "exit_code": 124, "stdout": "", "stderr": ""}, 500)
                    except Exception as e:
                        logger.error(f"执行脚本失败 - {e}", exc_info=True)
                        self.send_json({"ok": False, "msg": str(e), "exit_code": 1, "stdout": "", "stderr": ""}, 500)
                except Exception as e:
                    logger.error(f"创建脚本失败: {e}", exc_info=True)
                    self.send_json({"ok": False, "msg": f"创建脚本失败: {e}"}, 500)
                finally:
                    try:
                        os.unlink(script_path)
                    except Exception:
                        pass

            elif p == "/api/rename_file":
                logger.info("处理 /api/rename_file 请求")
                if not self._require_repo():
                    return
                idem_key = (self.headers.get("X-Idempotency-Key") or "").strip()
                if idem_key:
                    ent = _idempotency_get(idem_key)
                    if ent and isinstance(ent.get("payload"), dict):
                        self.send_json(ent.get("payload"), int(ent.get("code") or 200))
                        return
                oldp = (data.get("old_path") or data.get("path") or "").strip()
                newp = (data.get("new_path") or "").strip()
                undo_gid = (self.headers.get("X-Undo-Group") or "").strip()
                if undo_gid:
                    s_old = _undo_capture_file_snapshot(oldp)
                    s_new = _undo_capture_file_snapshot(newp)
                    snaps = [x for x in (s_old, s_new) if x is not None]
                    _undo_record(undo_gid, {"type": "rename", "old_path": oldp, "new_path": newp, "snapshots": snaps})
                ok, msg = rename_file(oldp, newp)
                if ok:
                    invalidate_changed_files_cache()
                    notify_files_updated()
                payload = {"ok": ok, "msg": msg}
                if idem_key:
                    _idempotency_set(idem_key, payload, 200)
                self.send_json(payload)

            elif p == "/api/revert_file":
                logger.info("处理 /api/revert_file 请求")
                if not self._require_repo():
                    return
                try:
                    key = _rl_make_key(self)
                    ok_rl, delay_s, retry_after = _rl_check_and_get_delay(key, "modify")
                    if not ok_rl:
                        self.send_json({"ok": False, "msg": f"请求过于频繁，请在 {retry_after} 秒后重试", "retry_after": retry_after}, 429)
                        return
                    if delay_s > 0:
                        time.sleep(delay_s)
                except Exception as e:
                    logger.debug(f"Exception ignored: {e}")
                idem_key = (self.headers.get("X-Idempotency-Key") or "").strip()
                if idem_key:
                    ent = _idempotency_get(idem_key)
                    if ent and isinstance(ent.get("payload"), dict):
                        self.send_json(ent.get("payload"), int(ent.get("code") or 200))
                        return
                fp = (data.get("path") or "").strip()
                st = (data.get("status") or "M")
                undo_gid = (self.headers.get("X-Undo-Group") or "").strip()
                if undo_gid:
                    snap = _undo_capture_file_snapshot(fp)
                    if snap is not None:
                        _undo_record(undo_gid, {"type": "file_snapshot", "snapshot": snap})
                ok, msg = revert_file(fp, st)
                if ok:
                    invalidate_changed_files_cache()
                    notify_files_updated()
                payload = {"ok": ok, "msg": msg}
                if idem_key:
                    _idempotency_set(idem_key, payload, 200)
                self.send_json(payload)

            elif p == "/api/pull_file":
                logger.info("处理 /api/pull_file 请求")
                if not self._require_repo():
                    return
                fp = (data.get("path") or "").strip()
                full = _safe_repo_abspath(fp)
                if not full:
                    self.send_json({"ok": False, "msg": "非法路径"}, 400)
                    return

                # Ensure upstream exists and fetch latest.
                out_u, err_u, code_u = run_git(["rev-parse", "--abbrev-ref", "--symbolic-full-name", "@{u}"], timeout=30)
                if code_u != 0:
                    self.send_json({"ok": False, "msg": (err_u or "当前分支未设置上游，无法单文件更新")}, 400)
                    return

                out_f, err_f, code_f = run_git(["fetch", "--all"], timeout=300)
                if code_f != 0:
                    self.send_json({"ok": False, "msg": (err_f or out_f or "fetch 失败")}, 500)
                    return

                # Update the file from upstream tip.
                out_c, err_c, code_c = run_git(["checkout", "@{u}", "--", fp], timeout=120)
                if code_c != 0:
                    self.send_json({"ok": False, "msg": (err_c or out_c or "更新文件失败")}, 500)
                    return

                invalidate_changed_files_cache()
                notify_files_updated()
                self.send_json({"ok": True, "msg": "更新成功"})

            elif p == "/api/ai_config":
                logger.info("处理 /api/ai_config 请求")
                cfg = data.get("config") if isinstance(data, dict) else None
                if cfg is None and isinstance(data, dict):
                    cfg = data
                undo_gid = (self.headers.get("X-Undo-Group") or "").strip()
                with ai_config_lock:
                    if undo_gid:
                        prev_cfg = load_hivo_ai_config()
                        _undo_record(undo_gid, {"type": "ai_config", "prev": prev_cfg})
                    ok, msg = save_hivo_ai_config(cfg)
                self.send_json({"ok": ok, "msg": msg})

            elif p == "/api/hivo_ai_config":
                logger.info("处理 /api/hivo_ai_config 请求")
                cfg = data.get("config") if isinstance(data, dict) else None
                if cfg is None and isinstance(data, dict):
                    cfg = data
                if not isinstance(cfg, dict):
                    self.send_json({"ok": False, "msg": "config 必须是对象"}, 400)
                    return
                undo_gid = (self.headers.get("X-Undo-Group") or "").strip()
                if undo_gid:
                    prev_cfg = _hivo_load_cfg()
                    _undo_record(undo_gid, {"type": "hivo_ai_config", "prev": prev_cfg})
                ok, msg = _hivo_save_cfg(cfg)
                self.send_json({"ok": ok, "msg": msg})

            elif p == "/api/ai_feature":
                logger.info("处理 /api/ai_feature 请求")
                try:
                    enabled_in = None
                    if isinstance(data, dict):
                        enabled_in = data.get("web_search_enabled")
                    enabled = bool(enabled_in)
                except Exception:
                    enabled = False
                undo_gid = (self.headers.get("X-Undo-Group") or "").strip()
                with ai_config_lock:
                    if undo_gid:
                        prev_cfg = load_hivo_ai_config()
                        _undo_record(undo_gid, {"type": "ai_config", "prev": prev_cfg})
                    cur = _hivo_load_cfg()
                    if not isinstance(cur, dict):
                        cur = {}
                    feat = cur.get("features") if isinstance(cur.get("features"), dict) else {}
                    feat = dict(feat)
                    feat["web_search_enabled"] = bool(enabled)
                    cur["features"] = feat
                    ok, msg = _hivo_save_cfg(cur)
                self.send_json({"ok": ok, "msg": msg})

            elif p == "/api/undo":
                logger.info("处理 /api/undo 请求")
                req_gid = ""
                req_sid = ""
                try:
                    if isinstance(data, dict):
                        req_gid = str(data.get("group_id") or "").strip()
                        req_sid = str(data.get("session_id") or "").strip()
                except Exception:
                    req_gid = ""
                if req_gid:
                    gid, actions = _undo_pop_group_by_id(req_gid)
                elif req_sid:
                    gid, actions = _undo_pop_latest_group_for_session(req_sid)
                else:
                    gid, actions = _undo_pop_latest_group()
                if not actions:
                    if gid:
                        self.send_json({"ok": True, "group_id": gid, "msg": "无可撤销操作", "noop": True, "applied": False})
                    else:
                        self.send_json({"ok": False, "msg": "无可撤销操作", "noop": True, "applied": False}, 400)
                    return
                with _apply_revert_lock:
                    ok, msg = _undo_apply_actions(actions)
                if not ok:
                    self.send_json({"ok": False, "msg": msg or "撤销失败", "group_id": gid}, 500)
                    return
                self.send_json({"ok": True, "group_id": gid, "noop": False, "applied": True})

            elif p == "/api/ai_trim_session_chat":
                # 撤回时清理后端 session memory 中指定数量的对话记录
                logger.info("处理 /api/ai_trim_session_chat 请求")
                sid = ""
                trim_count = 0
                try:
                    if isinstance(data, dict):
                        sid = str(data.get("session_id") or "").strip()
                        trim_count = int(data.get("trim_count") or 0)
                except Exception:
                    sid = ""
                if not sid:
                    self.send_json({"ok": False, "msg": "缺少 session_id"}, 400)
                    return
                try:
                    st = _hivo_get_session_state(sid)
                    chat = st.get("chat") if isinstance(st.get("chat"), list) else []
                    if trim_count > 0 and trim_count <= len(chat):
                        st["chat"] = chat[:len(chat) - trim_count]
                    elif trim_count > 0 and trim_count >= len(chat):
                        st["chat"] = []
                    self.send_json({"ok": True, "remaining": len(st.get("chat") or [])})
                except Exception as e:
                    self.send_json({"ok": False, "error": str(e)}, 500)


            elif p == "/api/undo_clear":
                logger.info("处理 /api/undo_clear 请求")
                sid = ""
                try:
                    if isinstance(data, dict):
                        sid = str(data.get("session_id") or "").strip()
                except Exception:
                    sid = ""
                if not sid:
                    self.send_json({"ok": False, "msg": "缺少 session_id"}, 400)
                    return
                removed = _undo_clear_for_session(sid)
                self.send_json({"ok": True, "session_id": sid, "removed": int(removed), "session_undo_steps": _undo_get_steps_for_session(sid)})

            elif p == "/api/undo_drop":
                logger.info("处理 /api/undo_drop 请求")
                gids = []
                try:
                    if isinstance(data, dict):
                        raw = data.get("group_ids")
                        if isinstance(raw, list):
                            gids = [str(x) for x in raw if x is not None and str(x).strip()]
                        else:
                            g0 = str(data.get("group_id") or "").strip()
                            if g0:
                                gids = [g0]
                except Exception:
                    gids = []
                if not gids:
                    self.send_json({"ok": False, "msg": "缺少 group_ids"}, 400)
                    return
                removed = _undo_drop_groups(gids)
                self.send_json({"ok": True, "removed": int(removed)})

            elif p == "/api/ai_chat":
                t0 = time.time()
                msgs = data.get("messages") if isinstance(data, dict) else None
                temp = data.get("temperature") if isinstance(data, dict) else None
                pid = data.get("profile_id") if isinstance(data, dict) else None
                sid = data.get("session_id") if isinstance(data, dict) else None
                env_obs = data.get("env_observation") if isinstance(data, dict) else None

                try:
                    env_obs_s = str(env_obs or "").strip()
                except Exception:
                    env_obs_s = ""
                # Backend owns system-level context. Client-provided system messages are dropped.
                # Keep only recent visible history so the system context is never squeezed out.
                try:
                    base_msgs = msgs if isinstance(msgs, list) else []
                    visible = []
                    for m in base_msgs:
                        if not isinstance(m, dict):
                            continue
                        role = str(m.get("role") or "").strip()
                        if role in ("user", "assistant"):
                            c = m.get("content")
                            # 支持 content 为字符串或数组（OpenAI 多模态格式）
                            if isinstance(c, str) and c.strip():
                                visible.append({"role": role, "content": c})
                            elif isinstance(c, list):
                                # 保留数组格式（可能含 image_url）
                                visible.append({"role": role, "content": c})
                    visible = visible[-80:]
                except Exception:
                    visible = []

                sys0 = {"role": "system", "content": _ai_build_system_context_text()}
                if env_obs_s:
                    obs_msg = {"role": "system", "content": "【环境观察(自动获取)】\n" + env_obs_s}
                    msgs = [sys0, obs_msg] + visible
                else:
                    msgs = [sys0] + visible

                n_msgs = 0
                try:
                    n_msgs = len(msgs) if isinstance(msgs, list) else 0
                except Exception:
                    n_msgs = 0
                logger.info(f"处理 /api/ai_chat 请求 (profile_id={pid or ''}, session_id={sid or ''}, messages={n_msgs})")

                try:
                    key = _rl_make_key(self, profile_id=pid, session_id=sid)
                    ok_rl, delay_s, retry_after = _rl_check_and_get_delay(key, "chat")
                    if not ok_rl:
                        self.send_json({"ok": False, "msg": f"请求过于频繁，请在 {retry_after} 秒后重试", "retry_after": retry_after}, 429)
                        return
                    if delay_s > 0:
                        time.sleep(delay_s)
                except Exception as e:
                    logger.debug(f"Exception ignored: {e}")

                user_q = ""
                try:
                    if isinstance(msgs, list):
                        for m in reversed(msgs):
                            if isinstance(m, dict) and str(m.get("role") or "") == "user":
                                user_q = str(m.get("content") or "")
                                break
                except Exception:
                    user_q = ""

                try:
                    ent = _ai_cache_get(user_q, profile_id=pid)
                except Exception:
                    ent = None
                if ent is not None:
                    try:
                        cost_ms = int((time.time() - t0) * 1000)
                    except Exception:
                        cost_ms = -1
                    logger.info(f"/api/ai_chat 命中缓存 (profile_id={pid or ''}, session_id={sid or ''}, cost_ms={cost_ms})")
                    self.send_json({"ok": True, "content": ent.get("resp", ""), "raw": {"cached": True, "hits": int(ent.get("hits") or 0)}})
                    return

                ok, msg, result = ai_chat(msgs, temperature=temp, profile_id=pid)
                if ok:
                    content = result.get("content", "")
                    try:
                        _ai_cache_put(user_q, str(content or ""), profile_id=pid)
                    except Exception as e:
                        logger.debug(f"Exception ignored: {e}")
                    try:
                        cost_ms = int((time.time() - t0) * 1000)
                    except Exception:
                        cost_ms = -1
                    logger.info(f"/api/ai_chat 完成 (ok=1, profile_id={pid or ''}, session_id={sid or ''}, cost_ms={cost_ms}, resp_chars={len(str(content or ''))})")
                    self.send_json({"ok": True, "content": content, "raw": result.get("raw")})
                else:
                    try:
                        cost_ms = int((time.time() - t0) * 1000)
                    except Exception:
                        cost_ms = -1
                    logger.info(f"/api/ai_chat 完成 (ok=0, profile_id={pid or ''}, session_id={sid or ''}, cost_ms={cost_ms}, msg={str(msg or '')[:200]})")
                    self.send_json({"ok": False, "msg": msg}, 400)

            elif p == "/api/hivo_agent_cancel":
                logger.info("处理 /api/hivo_agent_cancel 请求")
                rid = ""
                try:
                    if isinstance(data, dict):
                        rid = str(data.get("run_id") or "").strip()
                except Exception:
                    rid = ""
                if not rid:
                    self.send_json({"ok": False, "msg": "缺少 run_id"}, 400)
                    return
                okc = _hivo_agent_request_cancel(rid)
                self.send_json({"ok": bool(okc), "msg": "已请求取消" if okc else "run_id 不存在"})

            elif p == "/api/ai_cache_clear":
                logger.info("处理 /api/ai_cache_clear 请求")
                try:
                    key = _rl_make_key(self)
                    ok_rl, delay_s, retry_after = _rl_check_and_get_delay(key, "modify")
                    if not ok_rl:
                        self.send_json({"ok": False, "msg": f"请求过于频繁，请在 {retry_after} 秒后重试", "retry_after": retry_after}, 429)
                        return
                    if delay_s > 0:
                        time.sleep(delay_s)
                except Exception as e:
                    logger.debug(f"Exception ignored: {e}")
                try:
                    _ai_cache_clear()
                except Exception as e:
                    logger.debug(f"Exception ignored: {e}")
                self.send_json({"ok": True, "msg": "已清空缓存"})

            elif p == "/api/ai_chat_history":
                logger.info("处理 /api/ai_chat_history 请求")
                pid = data.get("profile_id") if isinstance(data, dict) else None
                msgs = data.get("messages") if isinstance(data, dict) else None
                sid = data.get("session_id") if isinstance(data, dict) else None
                with ai_history_lock:
                    ok, msg = save_ai_chat_history(pid, msgs, session_id=sid)
                    info = list_ai_sessions(pid) if pid else {"active_session_id": "", "sessions": []}
                self.send_json({"ok": ok, "msg": msg, "active_session_id": info.get("active_session_id"), "sessions": info.get("sessions")})

            elif p == "/api/ai_sessions":
                logger.info("处理 /api/ai_sessions 请求")
                pid = data.get("profile_id") if isinstance(data, dict) else None
                action = str(data.get("action") or "") if isinstance(data, dict) else ""
                title = data.get("title") if isinstance(data, dict) else None
                sid = data.get("session_id") if isinstance(data, dict) else None
                order = data.get("session_ids") if isinstance(data, dict) else None

                with ai_history_lock:
                    if action == "create":
                        ok, msg, new_id = create_ai_session(pid, title=title)
                        info = list_ai_sessions(pid) if pid else {"active_session_id": "", "sessions": []}
                        self.send_json({"ok": ok, "msg": msg, "session_id": new_id, "active_session_id": info.get("active_session_id"), "sessions": info.get("sessions")})
                        return
                    if action == "reorder":
                        ids = []
                        if isinstance(order, list):
                            ids = [str(x) for x in order if x is not None and str(x).strip()]
                        ok, msg = reorder_ai_sessions(pid, ids)
                        info = list_ai_sessions(pid) if pid else {"active_session_id": "", "sessions": []}
                        self.send_json({"ok": ok, "msg": msg, "active_session_id": info.get("active_session_id"), "sessions": info.get("sessions")})
                        return
                    if action == "delete":
                        ok, msg = delete_ai_session(pid, sid)
                        info = list_ai_sessions(pid) if pid else {"active_session_id": "", "sessions": []}
                        self.send_json({"ok": ok, "msg": msg, "active_session_id": info.get("active_session_id"), "sessions": info.get("sessions")})
                        return
                    if action == "set_active":
                        ok, msg = set_ai_active_session(pid, sid)
                        info = list_ai_sessions(pid) if pid else {"active_session_id": "", "sessions": []}
                        self.send_json({"ok": ok, "msg": msg, "active_session_id": info.get("active_session_id"), "sessions": info.get("sessions")})
                        return
                    info = list_ai_sessions(pid) if pid else {"active_session_id": "", "sessions": []}
                self.send_json({"ok": True, "active_session_id": info.get("active_session_id"), "sessions": info.get("sessions")})

            elif p == "/api/ai_models":
                logger.info("处理 /api/ai_models 请求")
                base_url = str(data.get("base_url") or "") if isinstance(data, dict) else ""
                api_key = str(data.get("api_key") or "") if isinstance(data, dict) else ""
                ok, msg, models = ai_list_models(base_url, api_key)
                self.send_json({"ok": ok, "msg": msg, "models": models})

            elif p == "/api/ai_store_attachments":
                logger.info("处理 /api/ai_store_attachments 请求")
                arr = data.get("items") if isinstance(data, dict) else None
                if not isinstance(arr, list) or not arr:
                    self.send_json({"ok": True, "items": []})
                    return
                # Storage root under backend data dir
                base_dir = _hivo_ai_data_dir()
                full_dir = os.path.join(str(base_dir), "attachments")
                try:
                    os.makedirs(full_dir, exist_ok=True)
                except Exception as e:
                    self.send_json({"ok": False, "error": f"创建目录失败: {e}"}, 500)
                    return

                saved = []
                total_bytes = 0
                # Pull upload limits from config if present
                try:
                    cfg0 = _hivo_load_cfg()
                    limits0 = cfg0.get("limits") if isinstance(cfg0, dict) else {}
                except Exception:
                    limits0 = {}
                max_each = int(limits0.get("max_attachment_bytes") or (20 * 1024 * 1024))
                max_total = int(limits0.get("max_upload_total_bytes") or (100 * 1024 * 1024))

                for it in arr:
                    try:
                        if not isinstance(it, dict):
                            continue
                        nm = str(it.get("name") or "file").strip()
                        mime = str(it.get("mime") or "").strip().lower()
                        du = str(it.get("data_url") or "").strip()
                        if not du.startswith("data:"):
                            continue
                        # data URL format: data:mime;base64,xxxx
                        if ";base64," not in du:
                            continue
                        head, b64 = du.split(",", 1)
                        raw = base64.b64decode(b64)
                        size0 = len(raw)
                        if max_each > 0 and size0 > max_each:
                            logger.warning(f"丢弃过大的附件: {nm} size={size0} > each_limit={max_each}")
                            continue
                        if max_total > 0 and (total_bytes + size0) > max_total:
                            logger.warning(f"总大小超限，跳过附件: {nm} size={size0}, total={total_bytes}, limit={max_total}")
                            continue
                        total_bytes += size0
                        # Build filename
                        ext = mimetypes.guess_extension(mime) or ""
                        if not ext:
                            try:
                                base_nm = nm.rsplit(".", 1)[-1].lower()
                                ext = ("." + base_nm) if base_nm and (len(base_nm) <= 8) else ""
                            except Exception:
                                ext = ""
                        fname = time.strftime("%Y%m%d-%H%M%S") + "-" + uuid.uuid4().hex[:8] + (ext or ".bin")
                        full_path = os.path.join(full_dir, fname)
                        with open(full_path, "wb") as f:
                            f.write(raw)
                        rel_path = f"attachments/{fname}".replace("\\", "/")
                        url = f"/api/ai_attachment?path={urllib.parse.quote(rel_path)}"
                        saved.append({"name": nm, "mime": mime, "size": size0, "url": url, "path": rel_path})
                    except Exception as e:
                        logger.debug(f"Exception ignored while saving attachment: {e}")
                        continue

                self.send_json({"ok": True, "items": saved})

            elif p == "/api/hivo_agent":
                logger.info("处理 /api/hivo_agent 请求")
                pid = data.get("profile_id") if isinstance(data, dict) else None
                sid = data.get("session_id") if isinstance(data, dict) else None
                user_text = data.get("user_text") if isinstance(data, dict) else None
                msgs0 = data.get("messages") if isinstance(data, dict) else None
                dyn_ctx = data.get("dyn_context") if isinstance(data, dict) else ""
                ctx_refs = data.get("context_refs") if isinstance(data, dict) else None
                req_undo_gid = data.get("undo_gid") if isinstance(data, dict) else None
                env_obs = data.get("env_observation") if isinstance(data, dict) else None
                attachments = data.get("attachments") if isinstance(data, dict) else None
                resume = data.get("resume") if isinstance(data, dict) else None
                try:
                    env_obs_s = str(env_obs or "").strip()
                except Exception:
                    env_obs_s = ""
                if env_obs_s:
                    dyn_ctx = "【环境观察(自动获取)】\n" + env_obs_s + ("\n\n" + str(dyn_ctx or "") if dyn_ctx else "")

                try:
                    cfg0 = _hivo_load_cfg()
                    feat = cfg0.get("features") if isinstance(cfg0, dict) else None
                    feat = feat if isinstance(feat, dict) else {}
                    context_refs_enabled = bool(feat.get("context_refs_enabled", True))
                except Exception:
                    context_refs_enabled = True

                if context_refs_enabled:
                    try:
                        refs_block, refs_meta = _hivo_parse_context_refs_structured(ctx_refs if isinstance(ctx_refs, list) else [])
                    except Exception:
                        refs_block, refs_meta = "", []
                    if refs_block:
                        dyn_ctx = (str(dyn_ctx or "") + "\n\n" + refs_block).strip()
                    try:
                        if refs_meta:
                            st = _hivo_get_session_state(str(sid or "").strip())
                            mc = _hivo_mem_conf(cfg0)
                            tl = st.get("tool_log") if isinstance(st.get("tool_log"), list) else []
                            tl.append({"ts": time.time(), "type": "context_ref", "ok": True, "msg": "@引用解析", "data": {"refs": refs_meta}})
                            while len(tl) > int(mc.get("tool_log_items") or 80):
                                tl.pop(0)
                            st["tool_log"] = tl
                    except Exception as e:
                        logger.debug(f"Exception ignored: {e}")

                sid_s0 = str(sid or "").strip()
                if sid_s0:
                    try:
                        with _hivo_agent_run_lock:
                            active = str(_hivo_agent_session_active.get(sid_s0) or "").strip()
                            if active and (active in _hivo_agent_run_state) and (not bool((_hivo_agent_run_state.get(active) or {}).get("done"))):
                                self.send_json({"ok": False, "msg": "已有任务执行中，请先取消或等待完成", "run_id": active}, 409)
                                return
                    except Exception as e:
                        logger.debug(f"Exception ignored: {e}")

                run_id = uuid.uuid4().hex
                try:
                    _hivo_agent_mark_started(run_id, sid_s0)
                except Exception as e:
                    logger.debug(f"Exception ignored: {e}")

                # ack immediately; stream progress via websocket
                self.send_json({"ok": True, "run_id": run_id})

                def _bg():
                    try:
                        ok2, final2, _rid2 = hivo_agent_run(
                            run_id,
                            str(pid or "").strip(),
                            sid_s0,
                            str(user_text or ""),
                            history_messages=msgs0,
                            dyn_context=str(dyn_ctx or ""),
                            undo_gid=str(req_undo_gid or "").strip(),
                            attachments=attachments if isinstance(attachments, list) else None,
                            resume=resume if isinstance(resume, dict) else None,
                        )
                        try:
                            pid_s = str(pid or "").strip()
                            sid_s = str(sid or "").strip()
                            if pid_s and sid_s:
                                visible = []
                                base = msgs0 if isinstance(msgs0, list) else []
                                for m in base:
                                    if not isinstance(m, dict):
                                        continue
                                    role = str(m.get("role") or "").strip()
                                    if role not in ("user", "assistant"):
                                        continue
                                    if bool(m.get("pending")) and role == "assistant":
                                        continue
                                    c0 = m.get("content")
                                    content0 = ""
                                    if isinstance(c0, list):
                                        # store only text parts to keep history compact
                                        for p in c0:
                                            if not isinstance(p, dict):
                                                continue
                                            if str(p.get("type") or "") == "text":
                                                content0 += str(p.get("text") or "")
                                        content0 = content0.strip()
                                    else:
                                        content0 = str(c0 or "")
                                    if not content0:
                                        continue
                                    item0 = {"role": role, "content": content0}
                                    if role == "user":
                                        ug0 = str(m.get("undo_gid") or "").strip()
                                        if ug0:
                                            item0["undo_gid"] = ug0
                                        try:
                                            meta_in0 = m.get("_attMeta")
                                            if isinstance(meta_in0, list) and meta_in0:
                                                meta_out0 = []
                                                for a in meta_in0:
                                                    if not isinstance(a, dict):
                                                        continue
                                                    meta_out0.append({
                                                        "name": str(a.get("name") or "file"),
                                                        "mime": str(a.get("mime") or ""),
                                                        "size": int(a.get("size") or 0),
                                                        **({"url": str(a.get("url") or "")} if str(a.get("url") or "").strip() else {})
                                                    })
                                                if meta_out0:
                                                    item0["_attMeta"] = meta_out0
                                        except Exception:
                                            pass
                                    visible.append(item0)
                                if final2:
                                    visible.append({"role": "assistant", "content": str(final2)})
                                with ai_history_lock:
                                    save_ai_chat_history(pid_s, visible, session_id=sid_s)
                        except Exception as e:
                            logger.debug(f"Exception ignored: {e}")
                    except Exception as e:
                        try:
                            _hivo_ws_emit(run_id, sid_s0, "error", str(e))
                            _hivo_ws_emit_final(run_id, sid_s0, str(e), ok=False)
                        except Exception as e:
                            logger.debug(f"Exception ignored: {e}")
                    finally:
                        try:
                            _hivo_agent_mark_done(run_id)
                        except Exception as e:
                            logger.debug(f"Exception ignored: {e}")

                try:
                    th = threading.Thread(target=_bg, daemon=True)
                    th.start()
                except Exception as e:
                    logger.debug(f"Exception ignored: {e}")

            elif p == "/api/revert_all":
                logger.info("处理 /api/revert_all 请求")
                if not self._require_repo():
                    return
                try:
                    key = _rl_make_key(self)
                    ok_rl, delay_s, retry_after = _rl_check_and_get_delay(key, "modify")
                    if not ok_rl:
                        self.send_json({"ok": False, "msg": f"请求过于频繁，请在 {retry_after} 秒后重试", "retry_after": retry_after}, 429)
                        return
                    if delay_s > 0:
                        time.sleep(delay_s)
                except Exception as e:
                    logger.debug(f"Exception ignored: {e}")
                errors = []
                raw_paths = data.get("paths")
                paths = None
                if isinstance(raw_paths, list):
                    paths = [str(x) for x in raw_paths if x is not None and str(x).strip()]
                is_partial = bool(paths)

                all_files = get_changed_files()
                if is_partial:
                    sel = set((p or "").replace("\\", "/").lstrip("/") for p in paths)
                    files = [f for f in all_files if (f.get("path") or "").replace("\\", "/").lstrip("/") in sel]
                    logger.info(f"开始还原所选文件，共 {len(files)} 个")
                else:
                    files = all_files
                    logger.info(f"开始还原所有文件，共 {len(files)} 个")

                for f in files:
                    ok, msg = revert_file(f["path"], f["status"])
                    if not ok:
                        errors.append(f"{f['path']}: {msg}")

                if is_partial:
                    for f in files:
                        if (f.get("status") or "").strip().upper() != "U":
                            continue
                        fp = (f.get("path") or "").replace("\\", "/").lstrip("/")
                        full = _safe_repo_abspath(fp)
                        if not full:
                            continue
                        try:
                            if os.path.isdir(full):
                                shutil.rmtree(full)
                            elif os.path.exists(full):
                                os.remove(full)
                        except Exception as e:
                            errors.append(f"{fp}: {e}")
                else:
                    run_git(["clean", "-fd"])

                logger.info(f"还原文件完成，错误数: {len(errors)}")
                if not errors:
                    try:
                        invalidate_changed_files_cache()
                        notify_files_updated()
                    except Exception as e:
                        logger.debug(f"Exception ignored: {e}")
                self.send_json({"ok": not errors, "errors": errors})

            elif p == "/api/restore_file":
                logger.info("处理 /api/restore_file 请求")
                if not self._require_repo():
                    return
                commit = (data.get("hash") or "").strip()
                fp = (data.get("path") or "").strip()
                ok, msg = restore_file_from_commit(commit, fp)
                self.send_json({"ok": ok, "msg": msg})

            elif p == "/api/restore_workspace":
                logger.info("处理 /api/restore_workspace 请求")
                if not self._require_repo():
                    return
                commit = (data.get("hash") or "").strip()
                if not commit:
                    self.send_json({"error":"缺少 hash"}, 400)
                    return
                ok, msg = restore_workspace_to_commit(commit, force=False)
                self.send_json({"ok": ok, "msg": msg})

            elif p == "/api/restore_workspace_force":
                logger.info("处理 /api/restore_workspace_force 请求 (丢弃本地修改后还原)")
                if not self._require_repo():
                    return
                commit = (data.get("hash") or "").strip()
                if not commit:
                    self.send_json({"error":"缺少 hash"}, 400)
                    return
                ok, msg = restore_workspace_to_commit(commit, force=True)
                self.send_json({"ok": ok, "msg": msg})

            elif p == "/api/diff_commit":
                logger.info("处理 /api/diff_commit 请求")
                if not self._require_repo():
                    return
                commit = qget("hash")
                if not commit:
                    self.send_json({"error":"缺少 hash"}, 400)
                    return
                try:
                    out, err, code = run_git(["diff", "--name-status", commit], timeout=120)
                    if code != 0:
                        self.send_json({"ok": False, "error": err or "获取差异失败"}, 500)
                        return
                    
                    files = []
                    added = 0
                    removed = 0
                    modified = 0
                    
                    if out:
                        for line in out.strip().split('\n'):
                            line = line.strip()
                            if not line:
                                continue
                            parts = line.split('\t', 2)
                            if len(parts) >= 2:
                                status = parts[0]
                                path = parts[-1]
                                files.append({"status": status, "path": path})
                                if status == 'A':
                                    added += 1
                                elif status == 'D':
                                    removed += 1
                                elif status == 'M':
                                    modified += 1
                    
                    self.send_json({
                        "ok": True,
                        "files": files,
                        "added": added,
                        "removed": removed,
                        "modified": modified,
                        "total": len(files)
                    })
                except Exception as e:
                    logger.error(f"获取差异失败: {e}")
                    self.send_json({"ok": False, "error": str(e)}, 500)

            elif p == "/api/revert_commit":
                logger.info("处理 /api/revert_commit 请求")
                if not self._require_repo():
                    return
                commit = (data.get("hash") or "").strip()
                if not commit:
                    self.send_json({"error":"缺少 hash"}, 400)
                    return
                ok, msg, full_hash = revert_commit(commit)
                self.send_json({
                    "ok": ok,
                    "msg": msg,
                    "full_hash": full_hash,
                    "hash": full_hash[:7] if full_hash else "",
                })

            elif p == "/api/revert_commit_force":
                logger.info("处理 /api/revert_commit_force 请求 (丢弃本地修改后 Revert)")
                if not self._require_repo():
                    return
                commit = (data.get("hash") or "").strip()
                if not commit:
                    self.send_json({"error":"缺少 hash"}, 400)
                    return
                # 丢弃本地未提交修改
                _, rh_err, rh_code = run_git(["reset", "--hard", "HEAD"], timeout=120)
                if rh_code != 0:
                    self.send_json({"ok": False, "error": rh_err or "reset --hard 失败"}, 400)
                    return
                run_git(["clean", "-fd"], timeout=120)
                # 执行 Revert
                ok, msg, full_hash = revert_commit(commit)
                try:
                    invalidate_changed_files_cache()
                    notify_files_updated()
                except Exception as e:
                    logger.debug(f"Exception ignored: {e}")
                self.send_json({
                    "ok": ok,
                    "msg": msg,
                    "full_hash": full_hash,
                    "hash": full_hash[:7] if full_hash else "",
                })

            elif p == "/api/soft_reset_commit":
                logger.info("处理 /api/soft_reset_commit 请求")
                if not self._require_repo():
                    return
                commit = (data.get("hash") or "").strip()
                if not commit:
                    self.send_json({"error":"缺少 hash"}, 400)
                    return

                head_out, head_err, head_code = run_git(["rev-parse", "HEAD"])
                if head_code != 0:
                    self.send_json({"error": head_err or "无法获取 HEAD"}, 400)
                    return
                head_full = (head_out or "").strip()

                if commit != head_full and commit != head_full[:7]:
                    self.send_json({"error": "仅允许撤销当前分支最新一次提交（HEAD）"}, 400)
                    return

                pushed, _, perr = is_commit_pushed(head_full)
                if perr:
                    self.send_json({"error": perr}, 400)
                    return
                if pushed:
                    self.send_json({"error": "该提交已推送到远端，无法使用软回退撤销；请使用 Revert"}, 400)
                    return

                _, err, code = run_git(["reset", "--soft", f"{head_full}^"], timeout=120)
                if code != 0:
                    self.send_json({"error": err or "软回退失败"}, 400)
                    return

                # Soft reset keeps index staged, which breaks partial commit flows.
                # Clear staging area while preserving working tree changes.
                _, unstage_err, unstage_code = run_git(["reset"], timeout=120)
                if unstage_code != 0:
                    self.send_json({"error": unstage_err or "清空暂存区失败"}, 400)
                    return

                new_head_out, _, new_head_code = run_git(["rev-parse", "HEAD"])
                new_head_full = (new_head_out or "").strip() if new_head_code == 0 else ""
                self.send_json({
                    "ok": True,
                    "full_hash": new_head_full,
                    "hash": new_head_full[:7] if new_head_full else "",
                })

            elif p == "/api/blame":
                logger.info("处理 /api/blame 请求")
                if not self._require_repo():
                    return
                raw_path = (data.get("path") or "").strip()
                if not raw_path:
                    self.send_json({"error": "缺少 path"}, 400)
                    return
                # _safe_repo_abspath handles both absolute and repo-relative paths
                full = _safe_repo_abspath(raw_path)
                if not full or (not os.path.exists(full)):
                    self.send_json({"error": "文件不存在"}, 404)
                    return
                # Derive the repo-relative path from the resolved absolute path
                rel = os.path.relpath(full, os.path.abspath(REPO_PATH)).replace("\\", "/")
                # Use line-porcelain for structured blame output
                out, err, code = run_git(["blame", "--line-porcelain", "--", rel], timeout=120)
                if code != 0:
                    self.send_json({"error": err or out or "blame 失败"}, 400)
                    return
                lines = []
                cur = {
                    "hash": "",
                    "author": "",
                    "author_mail": "",
                    "author_time": "",
                    "author_time_unix": None,
                    "author_tz": "",
                    "author_date": "",
                    "author_datetime": "",
                    "summary": "",
                }
                ln = 0
                for raw in (out or "").splitlines():
                    if not raw:
                        continue
                    if raw[0] != '\t':
                        # header line(s)
                        parts = raw.split(' ')
                        if len(parts) >= 1 and re.fullmatch(r"[0-9a-f]{8,40}", parts[0].strip() or ""):
                            # commit sha header line; when a new group starts, cur applies to next content line
                            cur["hash"] = parts[0].strip()
                        elif raw.startswith("author "):
                            cur["author"] = raw[7:].strip()
                        elif raw.startswith("author-mail "):
                            cur["author_mail"] = raw[12:].strip()
                        elif raw.startswith("author-time "):
                            try:
                                ts = int(raw[12:].strip())
                                cur["author_time_unix"] = ts
                                cur["author_time"] = datetime.utcfromtimestamp(ts).isoformat() + "Z"
                            except Exception:
                                cur["author_time"] = ""
                                cur["author_time_unix"] = None
                        elif raw.startswith("author-tz "):
                            # timezone like +0800
                            cur["author_tz"] = raw[10:].strip()
                            try:
                                ts = cur.get("author_time_unix")
                                tzs = cur.get("author_tz") or ""
                                if isinstance(ts, int) and tzs and len(tzs) >= 5 and (tzs[0] in "+-"):
                                    sign = 1 if tzs[0] == '+' else -1
                                    hh = int(tzs[1:3])
                                    mm = int(tzs[3:5])
                                    offset = timedelta(hours=hh, minutes=mm) * sign
                                    tzinfo = timezone(offset)
                                    dt = datetime.fromtimestamp(ts, tz=tzinfo)
                                    cur["author_date"] = dt.strftime("%Y-%m-%d")
                                    cur["author_datetime"] = dt.strftime("%Y-%m-%d %H:%M:%S ") + tzs
                            except Exception:
                                pass
                        elif raw.startswith("summary "):
                            cur["summary"] = raw[8:].strip()
                        continue
                    # content line: emit a record tied to current header
                    ln += 1
                    lines.append({
                        "no": ln,
                        "hash": (cur.get("hash") or "")[:7],
                        "author": cur.get("author") or "",
                        "author_mail": cur.get("author_mail") or "",
                        "author_time": cur.get("author_time") or "",
                        "author_time_unix": cur.get("author_time_unix"),
                        "author_tz": cur.get("author_tz") or "",
                        "author_date": cur.get("author_date") or "",
                        "author_datetime": cur.get("author_datetime") or "",
                        "summary": cur.get("summary") or "",
                    })
                self.send_json({"ok": True, "lines": lines})

            elif p == "/api/stage_file":
                logger.info("处理 /api/stage_file 请求")
                if not self._require_repo():
                    return
                filepath = (data.get("path") or "").strip()
                if not filepath:
                    self.send_json({"error":"缺少文件路径"}, 400)
                    return
                _, err, code = run_git(["add", "--", filepath])
                if code != 0:
                    logger.error(f"暂存文件失败: {err}")
                    self.send_json({"error": err or "暂存文件失败"}, 400)
                    return
                logger.info(f"暂存文件成功: {filepath}")
                try:
                    invalidate_changed_files_cache()
                    notify_files_updated()
                except Exception as e:
                    logger.debug(f"Exception ignored: {e}")
                self.send_json({"ok": True})

            elif p == "/api/commit":
                logger.info("处理 /api/commit 请求")
                if not self._require_repo():
                    return
                msg   = data.get("message", "").strip()
                paths = data.get("files", [])
                if not msg:
                    logger.warning("提交失败: 提交信息为空")
                    self.send_json({"error":"提交信息不能为空"}, 400)
                    return
                logger.info(f"开始提交，文件数: {len(paths)}, 消息: {msg}")
                for fp in paths:
                    run_git(["add", "--", fp])
                _, err, code = run_git(["commit", "-m", msg])
                if code != 0:
                    logger.error(f"提交失败: {err}")
                    self.send_json({"error": err}, 400)
                    return
                logger.info("提交成功")
                full_hash, _, code2 = run_git(["rev-parse", "HEAD"])
                full_hash = (full_hash or "").strip() if code2 == 0 else ""
                self.send_json({
                    "ok": True,
                    "full_hash": full_hash,
                    "hash": full_hash[:7] if full_hash else ""
                })

            elif p == "/api/commit_hunks":
                logger.info("处理 /api/commit_hunks 请求")
                if not self._require_repo():
                    return

                msg = (data.get("message") or "").strip()
                filepath = (data.get("path") or "").strip()
                status = (data.get("status") or "").strip() or "M"
                hunks = data.get("hunks") or []
                ctx_lines = data.get("ctx", 5)

                if not msg:
                    self.send_json({"error":"提交信息不能为空"}, 400)
                    return
                if not filepath:
                    self.send_json({"error":"缺少 path"}, 400)
                    return
                if not isinstance(hunks, list) or not hunks:
                    self.send_json({"error":"未选择任何变更块"}, 400)
                    return

                staged, err = has_any_staged_changes()
                if err:
                    self.send_json({"error": err}, 400)
                    return
                if staged:
                    self.send_json({"error": "检测到已有暂存区内容，请先提交/取消暂存后再进行按块提交"}, 400)
                    return

                # Untracked new file: create an intent-to-add entry so git diff can produce a patch.
                # IMPORTANT: Do this after staged-check; add -N itself would make index non-empty.
                st = (status or "").strip().upper() or "M"
                if st == "U":
                    _, addn_err, addn_code = run_git(["add", "-N", "--", filepath], timeout=60)
                    if addn_code != 0:
                        self.send_json({"error": addn_err or "初始化新文件暂存失败"}, 400)
                        return

                raw_patch, patch_err = get_raw_file_diff_patch(filepath, status, ctx_lines=ctx_lines)
                if patch_err:
                    if st == "U":
                        run_git(["reset", "--", filepath])
                    self.send_json({"error": patch_err}, 400)
                    return

                picked_patch = extract_selected_hunks_from_patch(raw_patch, hunks)
                ok, apply_err = apply_patch_to_index(picked_patch)
                if not ok:
                    if st == "U":
                        run_git(["reset", "--", filepath])
                    self.send_json({"error": apply_err}, 400)
                    return

                _, err, code = run_git(["commit", "-m", msg])
                if code != 0:
                    logger.error(f"按块提交失败: {err}")
                    if st == "U":
                        run_git(["reset", "--", filepath])
                    run_git(["reset"])
                    self.send_json({"error": err}, 400)
                    return

                logger.info("按块提交成功")
                full_hash, _, code2 = run_git(["rev-parse", "HEAD"])
                full_hash = (full_hash or "").strip() if code2 == 0 else ""
                self.send_json({
                    "ok": True,
                    "full_hash": full_hash,
                    "hash": full_hash[:7] if full_hash else ""
                })

            elif p == "/api/commit_lines":
                logger.info("处理 /api/commit_lines 请求")
                if not self._require_repo():
                    return

                msg = (data.get("message") or "").strip()
                filepath = (data.get("path") or "").strip()
                status = (data.get("status") or "").strip() or "M"
                lines = data.get("lines") or []
                ctx_lines = data.get("ctx", 5)

                if not msg:
                    self.send_json({"error":"提交信息不能为空"}, 400)
                    return
                if not filepath:
                    self.send_json({"error":"缺少 path"}, 400)
                    return
                if not isinstance(lines, list) or not lines:
                    self.send_json({"error":"未选择任何变更行"}, 400)
                    return

                st = (status or "").strip().upper() or "M"
                if st in ("A", "D"):
                    self.send_json({"error": "该文件状态不支持按行提交"}, 400)
                    return

                staged, err = has_any_staged_changes()
                if err:
                    self.send_json({"error": err}, 400)
                    return
                if staged:
                    self.send_json({"error": "检测到已有暂存区内容，请先提交/取消暂存后再进行按行提交"}, 400)
                    return

                # Untracked new file: create an intent-to-add entry so git diff can produce a patch.
                # IMPORTANT: Do this after staged-check; add -N itself would make index non-empty.
                if st == "U":
                    _, addn_err, addn_code = run_git(["add", "-N", "--", filepath], timeout=60)
                    if addn_code != 0:
                        self.send_json({"error": addn_err or "初始化新文件暂存失败"}, 400)
                        return

                raw_patch, patch_err = get_raw_file_diff_patch(filepath, st, ctx_lines=ctx_lines)
                if patch_err:
                    self.send_json({"error": patch_err}, 400)
                    return

                picked_patch = extract_selected_lines_from_patch(raw_patch, lines)
                if not (picked_patch or "").strip():
                    if st == "U":
                        run_git(["reset", "--", filepath])
                    self.send_json({"error": "构建 patch 失败（可能选择的行已变更，请刷新后重试）"}, 400)
                    return

                ok, apply_err = apply_patch_to_index(picked_patch)
                if not ok:
                    if st == "U":
                        run_git(["reset", "--", filepath])
                    self.send_json({"error": apply_err}, 400)
                    return

                _, err, code = run_git(["commit", "-m", msg])
                if code != 0:
                    logger.error(f"按行提交失败: {err}")
                    if st == "U":
                        run_git(["reset", "--", filepath])
                    run_git(["reset"])
                    self.send_json({"error": err}, 400)
                    return

                logger.info("按行提交成功")
                full_hash, _, code2 = run_git(["rev-parse", "HEAD"])
                full_hash = (full_hash or "").strip() if code2 == 0 else ""
                self.send_json({
                    "ok": True,
                    "full_hash": full_hash,
                    "hash": full_hash[:7] if full_hash else ""
                })

            elif p == "/api/push":
                logger.info("处理 /api/push 请求")
                if not self._require_repo():
                    return
                _, err, code = run_git(["push"], timeout=300)
                if code == 0:
                    logger.info("推送成功")
                    # Refresh remote-tracking refs so commit_push_status becomes accurate immediately
                    run_git(["fetch", "--prune"], timeout=300)
                else:
                    logger.error(f"推送失败: {err}")
                self.send_json({"ok": code == 0, "msg": err})

            elif p == "/api/pull":
                logger.info("处理 /api/pull 请求")
                if not self._require_repo():
                    return

                out, err, code = run_git(["pull", "--no-edit"], timeout=300)
                output = (out or "")
                error = (err or "")

                # 检测本地修改会被覆盖的情况
                local_changes_conflict = False
                affected_files = []
                
                # Git会在error中输出类似:
                # "error: Your local changes to the following files would be overwritten by merge:"
                # 或者中文版本可能是其他提示
                if code != 0 and ("would be overwritten" in error.lower() or 
                                 "will be overwritten" in error.lower() or
                                 "本地修改" in error or
                                 "覆盖" in error):
                    local_changes_conflict = True
                    # 尝试提取受影响的文件列表
                    lines = error.split('\n')
                    in_file_list = False
                    for line in lines:
                        line = line.strip()
                        if 'would be overwritten' in line.lower() or 'will be overwritten' in line.lower():
                            in_file_list = True
                            continue
                        if in_file_list:
                            if line.startswith('Please') or line.startswith('Aborting') or not line:
                                break
                            # 去除可能的制表符或空格
                            if line and not line.startswith('error:') and not line.startswith('hint:'):
                                affected_files.append(line.strip())
                
                conflict_files, _ = get_unmerged_files()
                has_conflicts = bool(conflict_files)

                ok = (code == 0) and (not has_conflicts)
                self.send_json({
                    "ok": ok,
                    "output": output.strip(),
                    "error": error.strip(),
                    "has_conflicts": has_conflicts,
                    "conflict_files": conflict_files,
                    "local_changes_conflict": local_changes_conflict,
                    "affected_files": affected_files,
                })

            elif p == "/api/stash_and_pull":
                logger.info("处理 /api/stash_and_pull 请求 (暂存修改并更新)")
                if not self._require_repo():
                    return

                # 1. git stash
                stash_out, stash_err, stash_code = run_git(["stash", "push", "-u", "-m", "Auto stash before pull"], timeout=60)
                if stash_code != 0:
                    logger.error(f"暂存失败: {stash_err}")
                    self.send_json({"error": f"暂存失败: {stash_err}"}, 400)
                    return
                
                # 检查是否真的有内容被暂存（如果工作区干净，stash不会创建新条目）
                stashed = "No local changes to save" not in (stash_out or "")
                
                # 2. git pull
                pull_out, pull_err, pull_code = run_git(["pull", "--no-edit"], timeout=300)
 
                # stash pop 由前端在用户确认后触发（/api/stash_pop）
                conflict_files, _ = get_unmerged_files()
                has_conflicts = bool(conflict_files)
                ok = (pull_code == 0) and (not has_conflicts)
                
                response = {
                    "ok": ok,
                    "stashed": stashed,
                    "pending_pop": bool(stashed),
                    "pull_output": (pull_out or "").strip(),
                    "pull_error": (pull_err or "").strip(),
                    "has_conflicts": has_conflicts,
                    "conflict_files": conflict_files,
                }
                
                if ok:
                    response["message"] = "成功暂存修改并更新"
                else:
                    if stashed:
                        response["message"] = "更新未完成：你的本地修改已暂存（可稍后手动恢复）"
                    else:
                        response["message"] = "更新未完成"
                
                self.send_json(response)

            elif p == "/api/commit_and_pull":
                logger.info("处理 /api/commit_and_pull 请求 (提交并更新)")
                if not self._require_repo():
                    return

                commit_msg = (data.get("message") or "").strip()
                if not commit_msg:
                    self.send_json({"error": "提交信息不能为空"}, 400)
                    return

                # 1. git add -A (暂存所有修改)
                add_out, add_err, add_code = run_git(["add", "-A"], timeout=60)
                if add_code != 0:
                    logger.error(f"暂存文件失败: {add_err}")
                    self.send_json({"error": f"暂存文件失败: {add_err}"}, 400)
                    return

                # 2. git commit
                commit_out, commit_err, commit_code = run_git(["commit", "-m", commit_msg], timeout=60)
                if commit_code != 0:
                    logger.error(f"提交失败: {commit_err}")
                    self.send_json({"error": f"提交失败: {commit_err}"}, 400)
                    return

                # 3. git pull
                pull_out, pull_err, pull_code = run_git(["pull", "--no-edit"], timeout=300)
                
                # 检查合并冲突
                conflict_files, _ = get_unmerged_files()
                has_conflicts = bool(conflict_files)
                
                ok = (pull_code == 0) and (not has_conflicts)
                
                self.send_json({
                    "ok": ok,
                    "commit_output": (commit_out or "").strip(),
                    "pull_output": (pull_out or "").strip(),
                    "pull_error": (pull_err or "").strip(),
                    "has_conflicts": has_conflicts,
                    "conflict_files": conflict_files,
                    "message": "成功提交并更新" if ok else "提交成功但更新时发生冲突"
                })

            elif p == "/api/switch_branch":
                logger.info("处理 /api/switch_branch 请求")
                if not self._require_repo():
                    return

                target_branch = (data.get("branch") or "").strip()
                if not target_branch:
                    self.send_json({"error": "未指定目标分支"}, 400)
                    return

                # 获取当前分支
                current_out, _, current_code = run_git(["branch", "--show-current"], timeout=30)
                current_branch = (current_out or "").strip()
                
                if current_code != 0:
                    logger.error("获取当前分支失败")
                    self.send_json({"error": "获取当前分支失败"}, 500)
                    return

                # 如果目标分支就是当前分支，直接返回成功
                if target_branch == current_branch:
                    self.send_json({
                        "ok": True,
                        "current": current_branch,
                        "message": "已在目标分支上"
                    })
                    return

                # 检查工作区是否有未提交的修改
                status_out, _, status_code = run_git(["status", "--porcelain"], timeout=30)
                has_changes = bool((status_out or "").strip())

                is_remote, remote_ref, _raw = _normalize_remote_ref(target_branch)
                want_detached = bool(data.get("detached"))

                if not has_changes:
                    if is_remote and want_detached:
                        # remote_ref is normalized (origin/foo)
                        remote_ref = remote_ref or target_branch.replace("remotes/", "", 1)
                        # switch --detach avoids creating local branch
                        out, err, code = run_git(["switch", "--detach", remote_ref], timeout=120)
                        if code != 0:
                            out, err, code = run_git(["checkout", remote_ref], timeout=120)
                        if code != 0:
                            self.send_json({
                                "ok": False,
                                "error": err or "切换到远端分支失败",
                                "output": out or ""
                            })
                            return
                        _, cur = get_branches()
                        self.send_json({
                            "ok": True,
                            "current": cur,
                            "message": f"成功切换到分支 {cur}"
                        })
                        return

                    # 工作区干净：统一走 switch_branch（远端分支会自动创建本地并建立跟踪关系）
                    logger.info(f"工作区干净，切换到分支: {target_branch}")
                    ok, cur, err_msg, out_msg, safe_err = switch_branch(target_branch)
                    if ok:
                        logger.info(f"成功切换到分支: {cur}")
                        self.send_json({
                            "ok": True,
                            "current": cur,
                            "message": f"成功切换到分支 {cur}"
                        })
                    else:
                        logger.error(f"切换分支失败: {err_msg}")
                        self.send_json({
                            "ok": False,
                            "error": err_msg or "切换分支失败",
                            "output": out_msg or ""
                        })
                    return

                # 工作区有未提交的修改，检查是否会被覆盖
                if is_remote and (not want_detached):
                    # 远端分支切换需要创建/设置跟踪关系，避免在有未提交修改时产生不可预期结果；
                    # 交互上引导用户先处理修改再切换。
                    self.send_json({
                        "ok": False,
                        "needs_handling": True,
                        "has_uncommitted_changes": True,
                        "affected_files": [],
                        "error": "工作区有未提交的修改，请先提交/暂存后再从远端分支创建/切换本地分支",
                        "message": "工作区有未提交的修改，请先处理后再切换远端分支"
                    })
                    return
                logger.info("工作区有未提交修改，检查是否会被覆盖...")
                
                # 使用 git checkout --dry-run 来检测是否会有冲突
                # 注意：git checkout 本身没有 --dry-run 选项，我们用其他方式检测
                # 方法：直接尝试切换，如果失败则说明会覆盖
                test_out, test_err, test_code = run_git(["checkout", target_branch], timeout=60)
                
                if test_code == 0:
                    # 切换成功，说明修改不会被覆盖
                    logger.info(f"修改不会被覆盖，成功切换到分支: {target_branch}")
                    self.send_json({
                        "ok": True,
                        "current": target_branch,
                        "has_uncommitted_changes": True,
                        "message": f"成功切换到分支 {target_branch}，未提交的修改已保留"
                    })
                    return

                # 切换失败，检查是否是因为会覆盖文件
                error_msg = (test_err or "").lower()
                if "would be overwritten" in error_msg or "overwritten by checkout" in error_msg:
                    # 修改会被覆盖，需要用户处理
                    logger.warning("切换分支会覆盖未提交的修改")
                    
                    # 提取受影响的文件列表
                    affected_files = []
                    lines = (test_err or "").split('\n')
                    in_file_list = False
                    for line in lines:
                        line = line.strip()
                        if 'would be overwritten' in line.lower() or 'overwritten by checkout' in line.lower():
                            in_file_list = True
                            continue
                        if in_file_list:
                            if line.startswith('Please') or line.startswith('Aborting') or not line:
                                break
                            if line and not line.startswith('error:') and not line.startswith('hint:'):
                                affected_files.append(line.strip())
                    
                    self.send_json({
                        "ok": False,
                        "needs_handling": True,
                        "has_uncommitted_changes": True,
                        "affected_files": affected_files,
                        "error": "工作区有未提交的修改会被覆盖",
                        "message": "切换分支会覆盖当前未提交的修改，请先处理这些修改"
                    })
                    return

                # 其他错误
                logger.error(f"切换分支失败: {test_err}")
                self.send_json({
                    "ok": False,
                    "error": test_err or "切换分支失败",
                    "output": test_out or ""
                })

            elif p == "/api/stash_and_switch":
                logger.info("处理 /api/stash_and_switch 请求 (暂存修改并切换分支)")
                if not self._require_repo():
                    return

                target_branch = (data.get("branch") or "").strip()
                if not target_branch:
                    self.send_json({"error": "未指定目标分支"}, 400)
                    return

                # 1. git stash (include untracked to truly clean the working tree)
                stash_out, stash_err, stash_code = run_git(
                    ["stash", "push", "-u", "-m", f"Auto stash before switching to {target_branch}"],
                    timeout=60
                )
                if stash_code != 0:
                    logger.error(f"暂存失败: {stash_err}")
                    self.send_json({"error": f"暂存失败: {stash_err}"}, 400)
                    return
                
                # 检查是否真的有内容被暂存
                stashed = "No local changes to save" not in (stash_out or "")
                
                # After stashing, ensure worktree is clean; otherwise switching will be refused by switch_branch.
                dirty2, detail2 = _has_worktree_changes()
                if dirty2:
                    logger.warning("暂存后工作区仍有未提交更改，拒绝切换分支")
                    if stashed:
                        try:
                            run_git(["stash", "pop"], timeout=60)
                        except Exception as e:
                            logger.debug(f"Exception ignored: {e}")
                    self.send_json({
                        "ok": False,
                        "error": "暂存后工作区仍存在未提交更改，无法安全切换分支（可能包含未被 stash 的变更）",
                        "output": detail2 or "",
                    })
                    return

                # 2. git checkout
                ok, cur, err_msg, out_msg, safe_err = switch_branch(target_branch)

                if not ok:
                    logger.error(f"切换分支失败: {err_msg}")
                    # 切换失败，尝试恢复暂存
                    if stashed:
                        run_git(["stash", "pop"], timeout=60)
                    self.send_json({
                        "ok": False,
                        "error": f"切换分支失败: {err_msg}",
                        "output": out_msg or ""
                    })
                    return

                logger.info(f"成功切换到分支: {cur}")

                response = {
                    "ok": True,
                    "current": cur,
                    "stashed": stashed,
                    "message": f"成功切换到分支 {cur}"
                }
                
                if stashed:
                    response["has_stash"] = True
                    response["message"] += "，修改已暂存"
                
                self.send_json(response)

            elif p == "/api/commit_and_switch":
                logger.info("处理 /api/commit_and_switch 请求 (提交并切换分支)")
                if not self._require_repo():
                    return

                target_branch = (data.get("branch") or "").strip()
                commit_msg = (data.get("message") or "").strip()
                
                if not target_branch:
                    self.send_json({"error": "未指定目标分支"}, 400)
                    return
                if not commit_msg:
                    self.send_json({"error": "提交信息不能为空"}, 400)
                    return

                # 1. git add -A
                add_out, add_err, add_code = run_git(["add", "-A"], timeout=60)
                if add_code != 0:
                    logger.error(f"暂存文件失败: {add_err}")
                    self.send_json({"error": f"暂存文件失败: {add_err}"}, 400)
                    return

                # 2. git commit
                commit_out, commit_err, commit_code = run_git(["commit", "-m", commit_msg], timeout=60)
                if commit_code != 0:
                    logger.error(f"提交失败: {commit_err}")
                    self.send_json({"error": f"提交失败: {commit_err}"}, 400)
                    return

                # 3. switch (remote refs will create/switch to local tracking branch)
                ok, cur, err_msg, out_msg, safe_err = switch_branch(target_branch)

                if not ok:
                    logger.error(f"切换分支失败: {err_msg}")
                    self.send_json({
                        "ok": False,
                        "committed": True,
                        "error": f"提交成功但切换分支失败: {err_msg}",
                        "output": out_msg or ""
                    })
                    return

                logger.info(f"成功提交并切换到分支: {cur}")
                self.send_json({
                    "ok": True,
                    "current": cur,
                    "committed": True,
                    "message": f"成功提交并切换到分支 {cur}"
                })

            elif p == "/api/stash_pop":
                logger.info("处理 /api/stash_pop 请求 (恢复暂存的修改)")
                if not self._require_repo():
                    return

                ref = str((data or {}).get("ref") or "").strip()
                cmd = ["stash", "pop"]
                if ref:
                    cmd.append(ref)

                pop_out, pop_err, pop_code = run_git(cmd, timeout=60)

                conflict_files, _ = get_unmerged_files()
                has_conflict = (pop_code != 0) or bool(conflict_files) or ("CONFLICT" in (pop_out or "")) or ("CONFLICT" in (pop_err or ""))

                if has_conflict:
                    logger.warning("恢复暂存的修改时发生冲突")
                    list_out, list_err, list_code = run_git(["stash", "list"], timeout=30)
                    stash_kept = False
                    top = ""
                    try:
                        raw = (list_out or "").strip()
                        if raw:
                            stash_kept = True
                            top = raw.splitlines()[0].strip()
                    except Exception:
                        stash_kept = False
                        top = ""
                    self.send_json({
                        "ok": False,
                        "has_conflict": True,
                        "conflict_files": conflict_files,
                        "stash_kept": stash_kept,
                        "stash_top": top,
                        "error": "恢复暂存的修改时发生冲突：修改未完全恢复，不代表丢失。请先解决冲突后再继续。",
                        "output": (pop_out or "").strip(),
                        "error_detail": (pop_err or "").strip()
                    })
                    return

                logger.info("成功恢复暂存的修改")
                try:
                    invalidate_changed_files_cache()
                    notify_files_updated()
                except Exception as e:
                    logger.debug(f"Exception ignored: {e}")
                self.send_json({
                    "ok": True,
                    "message": "成功恢复暂存的修改",
                    "output": (pop_out or "").strip()
                })

            elif p == "/api/stash_push":
                logger.info("处理 /api/stash_push 请求 (仅暂存修改)")
                if not self._require_repo():
                    return

                status_out, status_err, status_code = run_git(["status", "--porcelain=v1"], timeout=30)
                if status_code != 0:
                    self.send_json({"ok": False, "error": status_err or "无法检测工作区状态"}, 400)
                    return
                if not (status_out or "").strip():
                    self.send_json({"ok": True, "stashed": False, "message": "无本地修改，无需暂存"})
                    return

                out, err, code = run_git(["stash", "push", "-u", "-m", "Auto stash"], timeout=60)
                if code != 0:
                    self.send_json({"ok": False, "error": err or out or "git stash 失败", "output": (out or "").strip()}, 400)
                    return
                try:
                    invalidate_changed_files_cache()
                    notify_files_updated()
                except Exception as e:
                    logger.debug(f"Exception ignored: {e}")
                self.send_json({"ok": True, "stashed": True, "output": (out or "").strip()})

            elif p == "/api/stash_drop":
                logger.info("处理 /api/stash_drop 请求 (丢弃 stash 条目)")
                if not self._require_repo():
                    return

                ref = str((data or {}).get("ref") or "").strip()
                if not ref:
                    self.send_json({"ok": False, "error": "未指定 ref"}, 400)
                    return

                out, err, code = run_git(["stash", "drop", ref], timeout=60)
                if code != 0:
                    self.send_json({"ok": False, "error": (err or out or "丢弃 stash 失败"), "output": (out or "").strip()}, 400)
                    return

                try:
                    invalidate_changed_files_cache()
                    notify_files_updated()
                except Exception as e:
                    logger.debug(f"Exception ignored: {e}")

                self.send_json({
                    "ok": True,
                    "message": "已丢弃 stash",
                    "output": (out or "").strip(),
                })

            elif p == "/api/delete_branch":
                logger.info("处理 /api/delete_branch 请求 (删除本地分支)")
                if not self._require_repo():
                    return

                branch = str((data or {}).get("branch") or "").strip()
                force = bool((data or {}).get("force"))
                if not branch:
                    self.send_json({"ok": False, "error": "未指定分支"}, 400)
                    return
                if branch.startswith("remotes/"):
                    self.send_json({"ok": False, "error": "不支持删除远端分支引用（remotes/*）"}, 400)
                    return

                cur_out, cur_err, cur_code = run_git(["branch", "--show-current"], timeout=30)
                if cur_code != 0:
                    self.send_json({"ok": False, "error": (cur_err or "获取当前分支失败"), "output": (cur_out or "").strip()}, 400)
                    return
                current_branch = (cur_out or "").strip()
                if branch == current_branch:
                    self.send_json({"ok": False, "error": "不能删除当前所在分支"}, 400)
                    return

                args = ["branch", "-D" if force else "-d", branch]
                out, err, code = run_git(args, timeout=60)
                if code != 0:
                    self.send_json({"ok": False, "error": (err or out or "删除分支失败"), "output": (out or "").strip()}, 400)
                    return
                try:
                    invalidate_changed_files_cache()
                    notify_files_updated()
                except Exception as e:
                    logger.debug(f"Exception ignored: {e}")
                self.send_json({"ok": True, "message": f"已删除分支 {branch}", "output": (out or "").strip()})

            elif p == "/api/fetch_remotes":
                logger.info("处理 /api/fetch_remotes 请求 (刷新远端引用)")
                if not self._require_repo():
                    return
                ok, msg, out = fetch_remotes()
                if not ok:
                    self.send_json({"ok": False, "error": msg or "fetch 失败", "output": (out or "").strip()}, 400)
                    return
                try:
                    invalidate_changed_files_cache()
                    notify_files_updated()
                except Exception as e:
                    logger.debug(f"Exception ignored: {e}")
                self.send_json({"ok": True, "output": (out or "").strip()})

            elif p == "/api/tag/create":
                # 创建 tag
                if not REPO_PATH:
                    self.send_json({"ok": False, "error": "未打开仓库"}, 400)
                    return
                name = str((data or {}).get("name") or "").strip()
                message = str((data or {}).get("message") or "").strip()
                target = str((data or {}).get("target") or "").strip() or "HEAD"
                annotated = bool((data or {}).get("annotated", True))
                if not name:
                    self.send_json({"ok": False, "error": "tag 名称不能为空"}, 400)
                    return
                # 验证 tag 名称
                if not re.match(r"^[A-Za-z0-9._/-]+$", name) or name.startswith("-"):
                    self.send_json({"ok": False, "error": "tag 名称不合法"}, 400)
                    return
                try:
                    args = ["tag"]
                    if annotated:
                        args += ["-a", name, "-m", message or f"Tag {name}", target]
                    else:
                        args += [name, target]
                    out, err, code = run_git(args, timeout=60)
                    if code != 0:
                        self.send_json({"ok": False, "error": (err or "创建 tag 失败").strip()}, 400)
                        return
                    self.send_json({"ok": True, "message": "tag 创建成功"})
                except Exception as e:
                    logger.error(f"处理 /api/tag/create 异常: {e}", exc_info=True)
                    self.send_json({"ok": False, "error": str(e)}, 500)

            elif p == "/api/tag/delete":
                # 删除 tag
                if not REPO_PATH:
                    self.send_json({"ok": False, "error": "未打开仓库"}, 400)
                    return
                name = str((data or {}).get("name") or "").strip()
                remote = bool((data or {}).get("remote", False))
                if not name:
                    self.send_json({"ok": False, "error": "tag 名称不能为空"}, 400)
                    return
                if not re.match(r"^[A-Za-z0-9._/-]+$", name) or name.startswith("-"):
                    self.send_json({"ok": False, "error": "tag 名称不合法"}, 400)
                    return
                try:
                    args = ["tag", "-d", name]
                    out, err, code = run_git(args, timeout=60)
                    if code != 0:
                        self.send_json({"ok": False, "error": (err or "删除 tag 失败").strip()}, 400)
                        return
                    # 可选：删除远端 tag
                    if remote:
                        try:
                            _, err2, code2 = run_git(["push", "origin", "--delete", name], timeout=120)
                            # 远端失败不阻断本地删除结果
                        except Exception as e2:
                            logger.debug(f"删除远端 tag 失败: {e2}")
                    self.send_json({"ok": True, "message": "tag 删除成功"})
                except Exception as e:
                    logger.error(f"处理 /api/tag/delete 异常: {e}", exc_info=True)
                    self.send_json({"ok": False, "error": str(e)}, 500)

            elif p == "/api/conflict/resolve":
                # 解决冲突：把指定文件中的冲突块替换为选定方案，然后 git add
                if not REPO_PATH:
                    self.send_json({"ok": False, "error": "未打开仓库"}, 400)
                    return
                path = str((data or {}).get("path") or "").strip()
                resolution = str((data or {}).get("resolution") or "").strip()  # ours / theirs / both / custom
                custom_content = (data or {}).get("custom_content")  # 完整自定义文件内容
                if not path:
                    self.send_json({"ok": False, "error": "path 不能为空"}, 400)
                    return
                try:
                    full = os.path.join(REPO_PATH, path)
                    if not os.path.exists(full):
                        self.send_json({"ok": False, "error": f"文件不存在: {path}"}, 400)
                        return
                    if custom_content is not None:
                        # 客户端提供完整新内容
                        new_content = str(custom_content)
                    else:
                        with open(full, "r", encoding="utf-8", errors="replace") as fp:
                            content = fp.read()
                        hunks = _extract_conflict_hunks(content)
                        if not hunks:
                            # 没有冲突标记，直接 git add
                            _, _, _ = run_git(["add", "--", path], timeout=60)
                            self.send_json({"ok": True, "message": "无冲突标记，已标记为已解决"})
                            return
                        # 按 resolution 替换每个冲突块
                        def resolve_block(h):
                            if resolution == "ours":
                                return h["ours"]
                            elif resolution == "theirs":
                                return h["theirs"]
                            elif resolution == "both":
                                return h["ours"] + "\n" + h["theirs"] if h["ours"] and h["theirs"] else (h["ours"] or h["theirs"])
                            elif resolution == "base":
                                return h["base"]
                            else:
                                # 默认 ours
                                return h["ours"]
                        # 重建内容
                        lines = content.split("\n")
                        result = []
                        i = 0
                        while i < len(lines):
                            if lines[i].startswith("<<<<<<<"):
                                # 找到对应的 >>>>>>> 结束行
                                j = i + 1
                                sep_seen = False
                                while j < len(lines):
                                    if lines[j].startswith("|||||||"):
                                        j += 1
                                        while j < len(lines) and not lines[j].startswith("======="):
                                            j += 1
                                        if j < len(lines) and lines[j].startswith("======="):
                                            sep_seen = True
                                            j += 1
                                        continue
                                    if lines[j].startswith("======="):
                                        sep_seen = True
                                        j += 1
                                        continue
                                    if lines[j].startswith(">>>>>>>"):
                                        break
                                    j += 1
                                # 找到 hunk 对象
                                hunk = None
                                for h in hunks:
                                    if h["start"] == i + 1:
                                        hunk = h
                                        break
                                if hunk:
                                    resolved = resolve_block(hunk)
                                    if resolved:
                                        result.append(resolved)
                                i = j + 1
                            else:
                                result.append(lines[i])
                                i += 1
                        new_content = "\n".join(result)
                    # 写回
                    with open(full, "w", encoding="utf-8") as fp:
                        fp.write(new_content)
                    # git add 标记已解决
                    _, _, _ = run_git(["add", "--", path], timeout=60)
                    self.send_json({"ok": True, "message": "冲突已解决"})
                except Exception as e:
                    logger.error(f"处理 /api/conflict/resolve 异常: {e}", exc_info=True)
                    self.send_json({"ok": False, "error": str(e)}, 500)

            else:
                logger.warning(f"未知的 POST 请求路径: {p}")
                self.send_json({"error":"Not found"}, 404)
        except Exception as e:
            logger.error(f"处理 POST 请求异常: {p} - {str(e)}", exc_info=True)
            self.send_json({"error": f"服务器错误: {str(e)}"}, 500)


def main():
    """主函数"""
    port = 7842
    max_attempts = 10

    srv = None

    try:
        # 在单独的线程中启动WebSocket服务器
        if WEBSOCKET_AVAILABLE:
            ws_thread = threading.Thread(target=start_websocket_server, daemon=True)
            ws_thread.start()
            logger.info("WebSocket服务器线程已启动")
        
        # 尝试启动HTTP服务器，如果端口被占用则尝试下一个端口
        for attempt in range(max_attempts):
            try:
                srv = ThreadingHTTPServer(("127.0.0.1", port), Handler)
                logger.info("=" * 60)
                logger.info("Git Manager 后端启动成功!")
                logger.info(f"HTTP 服务器: http://localhost:{port}")
                if WEBSOCKET_AVAILABLE:
                    logger.info(f"WebSocket 服务器: ws://localhost:{WS_PORT}")
                logger.info(f"日志目录: {LOG_DIR}")
                logger.info("=" * 60)
                print(f"\n{'='*60}")
                print(f"  Git Manager 已启动!")
                print(f"  浏览器打开 → http://localhost:{port}")
                if WEBSOCKET_AVAILABLE:
                    print(f"  WebSocket → ws://localhost:{WS_PORT} (实时通信)")
                print(f"  日志文件 → {log_file}")
                print(f"{'='*60}\n")
                
                srv.serve_forever()
            except OSError as e:
                if e.errno == 10048:  # 端口被占用
                    logger.warning(f"端口 {port} 被占用，尝试下一个端口...")
                    port += 1
                    continue
                else:
                    raise
    except KeyboardInterrupt:
        logger.info("收到键盘中断信号，正在关闭服务器...")
        print("\n正在关闭服务器...")
    except Exception as e:
        logger.critical(f"服务器启动失败: {e}", exc_info=True)
        print(f"\n服务器启动失败: {e}")
        sys.exit(1)
    finally:
        try:
            _pw_cleanup()
        except Exception as e:
            logger.debug(f"Exception ignored: {e}")
        try:
            if srv is not None:
                srv.server_close()
        except Exception as e:
            logger.debug(f"Exception ignored: {e}")
        logger.info("Git Manager 后端已停止")
        print("已停止")

if __name__ == "__main__":
    main()
